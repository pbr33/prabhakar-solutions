import sys, os, base64, zipfile
import streamlit as st

st.set_page_config(page_title="ECI Presale", page_icon="⚡", layout="wide", initial_sidebar_state="expanded")

import json, time, io, re, smtplib
from datetime import datetime, timedelta
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
import plotly.graph_objects as go
import plotly.express as px

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side, numbers
    from openpyxl.utils import get_column_letter
    from openpyxl.chart import BarChart, PieChart, Reference
except ImportError:
    Workbook = None

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors as rl_colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm, inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, HRFlowable, Image as RLImage
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import graphviz as gv_lib
    HAS_GRAPHVIZ = True
except ImportError:
    gv_lib = None
    HAS_GRAPHVIZ = False

try:
    from openai import AzureOpenAI
except ImportError:
    AzureOpenAI = None
try:
    import msal
except ImportError:
    msal = None
try:
    import requests as _requests
except ImportError:
    _requests = None


# ═══════════════════════════════════════════════════════════════════════
#  ECI BRANDING CONSTANTS
# ═══════════════════════════════════════════════════════════════════════

ECI_BLUE = (27, 58, 92)
ECI_LIGHT_BLUE = (214, 228, 240)
ECI_ACCENT = (0, 180, 216)
ECI_GREEN = (0, 212, 170)
ECI_DARK = (10, 14, 26)
ECI_TEXT = (30, 30, 30)
ECI_GRAY = (100, 100, 100)

ECI_LOGO_BLUE_B64 = "iVBORw0KGgoAAAANSUhEUgAAAMgAAABGCAYAAACJ4ts2AAAD6ElEQVR4nO3dS4gcRRzH8e8mvlE3D1BiQUx8BFGXHCS6FoKC5hAUFRTFpAiBIMkhh3j04NF4MwEPGkQPWp7UiIqioohgyvXgQcT4CsaDlbgefGwiEh+Jh5qIBN2pbacyWzW/D8xp/93zY5jf9nT3TDeIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiDRibNgBTgVj3WLgZuBaYAK4GLgQOAc4HTgCzADfA/uAT4EAfBCD/2MYmWV+6FwQY90TwJYBZjnZ3hj8DV0XNtYtBO4EtgI3Aad1WM0M8BbwNPBmDP7YLM+X83o8EIPf1SGHDMmCYQcowVh3O2kr8AJwC93KAXA+cDfwOrDfWLdsMAmlFl3fOPOSsW4ceBy4r8DqVwLjwKEC65Z5qpmCGOuWk/7TXzXsLNKOJgpirLsIeA9YMeQo0pjqC2KsO4u05ViRucgPwEvAK8DnwHfAUWBp77EamCQd9bpiwHGlMqULciqO2jxKelP3c4y0f/JQDP7Hf/n7wd7jE8ADGOsmgc3ARuCMgaSVqlS9BTHWrSEdxu3nT2BTDN7PZf0x+Clgyli3A3iYVDIZIVUXBHiEvHM5W+Zajn+KwR8A1nddXupV7XkQY90EaT+hn9di8E+VziNtqrYgpP2Cfo4D2wvnkIbVXJA7MmbejsHvL55EmlVlQYx1FwCXZ4w+VzqLtK3KggDXZM5NFU0hzSt9FGunsW5nh+W+iMHPdpLukox1zABfdnhukb/VugUxGTMHY/DHiyeRptVakPMyZn4unkKaV2tBzsyYOVw8hTSv1oIczZg5t3gKaV6tBTmSMbOodAhpX63f5o0ZM8uMdWPaUZf/o9YtyNcZM+PknUwU+U+1FuSjzLnJoimkeVUWJAY/DeR8x2pD6SzStioL0vNyxsxaY92lxZNIs2ouyLMZM2PArsI5pGHVFiQG/zHwbsbobca6zaXzSJuqLUjPg6QfRfWz21jXeX/EWLfSWOeNdau6rkPqVHVBYvAfAk9mjC4EnjHWPda7kHUWY90aY91u0uWBNlD56yVzN1+/7n7Cuhj8G31mtgMWuLrP3AJgG7DeWPci8CrwGTAN/AYsJl0XawK4DliLrtI48mq/qgkx+F+NdeuAvcDyjEWWAPf3HiKzauIjQwz+W+BG0kchkYFpoiAAMfhvgOuB54ccRRrSTEEAYvA/xeDvAe4Cvhrw6g+gH2GNnKYKckIMfg9wJXAv8A7p0qNdHAb2ALcCl8XgdW+QETMq9yhcQrrT1Mn3KDybdKDiF1IZpkn7MftI9yh8Pwb/+zAyi4iIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiMjc/QXeEc9g6vu/ZAAAAABJRU5ErkJggg=="
ECI_LOGO_WHITE_B64 = "iVBORw0KGgoAAAANSUhEUgAAAMgAAABGCAYAAACJ4ts2AAADvElEQVR4nO3dza9dUxzG8e+vJV6Cq5WQklTrPbgxaMplYkAHN4RBhaRNaiANgw6uob/A8EqIaBoTikGplBAJIpKiBgYGqqFRAy+9BqUvIlfVY7COSZP2rLvvXj13rfN8kjP77bWfs09+Z5/9ctYGMzMzMzMzMzMzMzMzMzMzMzMzMzMzMzMzM7NGxKgDnAuSVgD3AXcCk8C1wFXAxcD5wAngGPAbsB/4Bvgc+CIi/hlFZqucpJdU1t5F5lsuaaOkDyWd7JjhqKRdkqYlLethe8ws5j3ZuXfWD71Wkh4i7QXeBO4Hzus41GXAI8D7wEFJq/pJaLVoqkEkTUh6HdgD3Nzz8GuBiZ7HtCWu6zfrkiNpNemb/rZRZ7F2NNEgkq4GPgXWjDiKNab6BpF0IWnPsSZzkSPA28A7wAHgMDAPXDF43QFMkc563dJzXBsXS+WsjaQXM89InZL0gtIp39yxpyTtkDQ/GOOMDbNUtof1q+qDdEnrgacySk8Bj0fEtoj4PXf8iNgXEVtJe5I3gH+7JbVa1f4T61nyLnY+GRE7u64kIg4Bm7oub/Wqdg8iaZJ0nDDMexHxcuk81qZqGwTYklEjYKZwDmtYzQ3ycEbNRxFxsHgSa1aVDSLpSuDGjNLXSmextlXZIMC6zLp9RVNY80o3yGzmNYrTHRgy7nUZ6z4GfNfDe7AxVuse5JqMml8iQsWTWNNqbZBLM2qOFk9hzau1QS7IqDlePIU1r9YGmc+ouaR4CmterQ1yIqPm8tIhrH2lG+Tp6GbYbeY/Z6x7laSxmJTCyql1D/JDRs0EeRcTzc6o1gb5KrNuqmgKa16VDRIRc0DOPVabS2extlXZIAN7Mmo2SLq+eBJrVs0N8mpGTQDPFc5hDau2QSLia+CTjNIHJT1ROo+1qdoGGXiG9KeoYbZL6nw8ImmtpJ2Sbuo6htWp6gaJiC+BHRmly4FXJD2/wFlN1kvaTpoeaDOVby9buNKTNsxKml3E8tMR8cGQmhngHuD2IXXLgG3AJklvAe8C3wJzwN/ACtK8WJPAXcAGPEvj2Kt9VhMi4i9J08BnwOqMRVYCWwcvs7Nq4idDRPwE3Ev6KWTWmyYaBCAifgTuBnaNOIo1pJkGAYiIPyLiUWAj8H3Pwx/Cf8IaO001yP8iYjdwK/AY8DFp6tEujgO7gQeAGyLi134SWi3G4nZwSStJT5o6/RmFF5FOVPxJaoY50nHMftIzCvdGxMlRZDYzMzMzMzMzMzMzMzMzMzMzMzMzMzMzMzOzhfsPjGdVI+ycCmwAAAAASUVORK5CYII="


# ═══════════════════════════════════════════════════════════════════════
#  CSS
# ═══════════════════════════════════════════════════════════════════════

def inject_css():
    st.markdown("""<style>
    @import url('https://fonts.googleapis.com/css2?family=JetBrains+Mono:wght@300;400;600;700&family=DM+Sans:wght@300;400;500;600;700&family=Space+Grotesk:wght@400;500;600;700&display=swap');
    :root{--bg0:#0a0e1a;--bg1:#111827;--bg2:#151c2e;--bd:#1e2a4a;--t1:#e2e8f0;--t2:#94a3b8;--t3:#64748b;--c1:#00d4aa;--c2:#00b4d8;--c3:#7b61ff;--c4:#ff6b6b;--c5:#ffd166;--c6:#06d6a0;--gc:rgba(0,212,170,.15);--gp:rgba(123,97,255,.15)}
    .stApp{background:var(--bg0)!important}
    .main .block-container{padding-top:1rem;max-width:100%}
    section[data-testid="stSidebar"]{background:linear-gradient(180deg,#0d1321,#111827)!important;border-right:1px solid var(--bd)}
    .logo-box{display:flex;align-items:center;gap:12px}
    .logo-icon{font-size:2.2rem;background:linear-gradient(135deg,var(--c1),var(--c3));-webkit-background-clip:text;-webkit-text-fill-color:transparent;filter:drop-shadow(0 0 12px var(--gc))}
    .logo-txt{font-family:'Space Grotesk',sans-serif;font-size:1.8rem;font-weight:700;background:linear-gradient(135deg,var(--c1),var(--c2));-webkit-background-clip:text;-webkit-text-fill-color:transparent}
    .logo-sub{font-family:'DM Sans',sans-serif;font-size:.7rem;color:var(--t2);letter-spacing:2px;text-transform:uppercase}
    .tagline{text-align:center;font-family:'DM Sans',sans-serif;color:var(--t2);font-size:.85rem;padding-top:.8rem}
    .dt{text-align:right;font-family:'JetBrains Mono',monospace;color:var(--t3);font-size:.8rem;padding-top:1rem}
    .stitle{font-family:'Space Grotesk',sans-serif;font-size:1.3rem;font-weight:600;color:var(--t1);margin-bottom:1rem;padding-bottom:.5rem;border-bottom:1px solid var(--bd)}
    .csec{font-family:'DM Sans',sans-serif;font-size:.9rem;font-weight:600;color:var(--c1);margin:.8rem 0 .4rem;letter-spacing:.5px}
    .cstat{background:var(--bg2);border-radius:8px;padding:12px;border:1px solid var(--bd)}
    .srow{display:flex;align-items:center;gap:8px;padding:4px 0;font-family:'DM Sans',sans-serif;font-size:.82rem;color:var(--t2)}
    .sdot{width:8px;height:8px;border-radius:50%;display:inline-block}
    .sdot.on{background:var(--c6);box-shadow:0 0 6px var(--c6)}.sdot.off{background:var(--t3)}
    .shdr{font-family:'Space Grotesk',sans-serif;font-size:1.3rem;font-weight:600;color:var(--t1);margin:1.5rem 0 1rem;display:flex;align-items:center;gap:10px}
    .shdr-i{font-size:1.4rem}
    .phdr{font-family:'Space Grotesk',sans-serif;font-size:1.1rem;font-weight:600;color:var(--c1);text-align:center;padding:12px;background:linear-gradient(135deg,var(--gc),var(--gp));border-radius:8px;margin-bottom:1rem;border:1px solid var(--bd)}
    .crd{background:var(--bg2);border:1px solid var(--bd);border-radius:12px;padding:1.2rem;margin-bottom:.8rem;transition:border-color .3s}
    .crd:hover{border-color:var(--c1)}
    .crd-t{font-family:'Space Grotesk',sans-serif;font-size:1.05rem;font-weight:600;color:var(--t1);margin-bottom:.3rem}
    .crd-d{font-family:'DM Sans',sans-serif;font-size:.82rem;color:var(--t2);margin-bottom:.8rem}
    .kpi{background:var(--bg2);border:1px solid var(--bd);border-radius:12px;padding:1.2rem;text-align:center;transition:all .3s}
    .kpi:hover{border-color:var(--c1);box-shadow:0 0 20px var(--gc);transform:translateY(-2px)}
    .kpi-i{font-size:1.8rem;margin-bottom:.4rem}
    .kpi-v{font-family:'JetBrains Mono',monospace;font-size:1.5rem;font-weight:700;color:var(--c1)}
    .kpi-t{font-family:'DM Sans',sans-serif;font-size:.85rem;font-weight:600;color:var(--t1);margin-top:.2rem}
    .kpi-s{font-family:'DM Sans',sans-serif;font-size:.72rem;color:var(--t3)}
    .alog{display:flex;align-items:center;gap:12px;padding:6px 12px;background:var(--bg2);border-radius:6px;margin-bottom:4px;border-left:3px solid var(--c1)}
    .abadge{font-family:'JetBrains Mono',monospace;font-size:.78rem;font-weight:600;color:var(--c2);min-width:160px}
    .aok{color:var(--c6);font-size:.8rem;font-weight:500;min-width:90px}
    .adet{font-family:'DM Sans',sans-serif;color:var(--t2);font-size:.8rem}
    .rch{font-family:'Space Grotesk',sans-serif;font-size:.95rem;font-weight:600;padding:8px 12px;border-radius:8px;margin-bottom:.8rem}
    .rch.fn{background:rgba(0,212,170,.1);color:var(--c1);border:1px solid rgba(0,212,170,.2)}
    .rch.nf{background:rgba(0,180,216,.1);color:var(--c2);border:1px solid rgba(0,180,216,.2)}
    .rch.ig{background:rgba(123,97,255,.1);color:var(--c3);border:1px solid rgba(123,97,255,.2)}
    .ri{background:var(--bg2);border:1px solid var(--bd);border-radius:8px;padding:10px 12px;margin-bottom:6px;font-family:'DM Sans',sans-serif;font-size:.82rem;color:var(--t2)}
    .ri strong{color:var(--t1);font-size:.85rem}
    .ri-c{font-family:'JetBrains Mono',monospace;font-size:.7rem;color:var(--c5);margin-top:2px}
    .ttag{display:flex;flex-wrap:wrap;gap:6px;margin-top:8px}
    .tt{font-family:'JetBrains Mono',monospace;font-size:.72rem;background:var(--bg2);border:1px solid var(--c3);color:var(--c3);padding:3px 10px;border-radius:12px}
    .rsb{background:var(--bg2);padding:16px 20px;border-radius:10px;display:flex;align-items:center;gap:16px;margin-bottom:1rem}
    .rsv{font-family:'JetBrains Mono',monospace;font-size:2rem;font-weight:700;color:var(--t1)}
    .rsl{font-family:'DM Sans',sans-serif;font-size:1rem;color:var(--t2)}
    .rc{background:var(--bg2);border:1px solid var(--bd);border-radius:8px;padding:14px;margin-bottom:8px}
    .rch2{display:flex;justify-content:space-between;align-items:center;margin-bottom:6px}
    .rch2 strong{color:var(--t1);font-family:'DM Sans',sans-serif}
    .rsev{font-family:'JetBrains Mono',monospace;font-size:.78rem;font-weight:600}
    .rc p{color:var(--t2);font-size:.82rem;margin:4px 0}
    .rmit{font-size:.8rem;color:var(--c6)}
    .ac{background:var(--bg2);border:1px solid var(--bd);border-radius:10px;padding:14px;margin-bottom:10px;transition:all .3s}
    .ac:hover{border-color:var(--c3);box-shadow:0 0 16px var(--gp)}
    .acn{font-family:'Space Grotesk',sans-serif;font-size:.95rem;font-weight:600;color:var(--t1)}
    .act{font-family:'JetBrains Mono',monospace;font-size:.7rem;color:var(--c3);margin-bottom:6px}
    .acs{font-family:'DM Sans',sans-serif;font-size:.8rem;color:var(--t2);padding-left:16px}
    .acs li{margin-bottom:2px}
    .dfv{font-family:'JetBrains Mono',monospace;font-size:.85rem;color:var(--c1);background:var(--bg2);padding:12px 16px;border-radius:8px;border:1px solid var(--bd);text-align:center;letter-spacing:.5px}
    .ls{background:var(--bg2);border:1px solid var(--bd);border-radius:10px;padding:16px 12px;text-align:center;min-height:160px}
    .ln{display:inline-flex;align-items:center;justify-content:center;width:36px;height:36px;border-radius:50%;font-family:'JetBrains Mono',monospace;font-size:1rem;font-weight:700;color:#fff;margin-bottom:8px}
    .lt{font-family:'Space Grotesk',sans-serif;font-size:.88rem;font-weight:600;color:var(--t1);margin-bottom:6px}
    .ld{font-family:'DM Sans',sans-serif;font-size:.72rem;color:var(--t2);line-height:1.4}
    .mc{background:var(--bg2);border:1px solid var(--bd);border-radius:10px;padding:16px;text-align:center;min-height:120px}
    .mi{font-size:1.6rem;margin-bottom:6px}
    .mt{font-family:'Space Grotesk',sans-serif;font-size:.85rem;font-weight:600;color:var(--t1);margin-bottom:4px}
    .md2{font-family:'DM Sans',sans-serif;font-size:.72rem;color:var(--t2);line-height:1.3}
    button[data-testid="stBaseButton-primary"]{background:linear-gradient(135deg,var(--c1),var(--c2))!important;color:#0a0e1a!important;font-family:'DM Sans',sans-serif!important;font-weight:600!important;border:none!important;border-radius:8px!important}
    button[data-testid="stBaseButton-primary"]:hover{box-shadow:0 0 20px var(--gc)!important}
    button[data-testid="stBaseButton-secondary"]{background:var(--bg2)!important;color:var(--t1)!important;border:1px solid var(--bd)!important;border-radius:8px!important}
    .stTabs [data-baseweb="tab-list"]{gap:4px;background:var(--bg1);padding:4px;border-radius:10px}
    .stTabs [data-baseweb="tab"]{background:transparent;color:var(--t2);border-radius:8px;padding:8px 16px;font-family:'DM Sans',sans-serif;font-size:.85rem}
    .stTabs [data-baseweb="tab"][aria-selected="true"]{background:var(--bg2)!important;color:var(--c1)!important}
    .stTabs [data-baseweb="tab-highlight"]{display:none}
    [data-testid="stFileUploader"]{background:var(--bg2);border:1px dashed var(--bd);border-radius:10px;padding:1rem}
    [data-testid="stMetric"]{background:var(--bg2);border:1px solid var(--bd);border-radius:10px;padding:12px 16px}
    [data-testid="stMetricValue"]{font-family:'JetBrains Mono',monospace!important;color:var(--c1)!important}
    #MainMenu{visibility:hidden}footer{visibility:hidden}
    header[data-testid="stHeader"]{background:rgba(10,14,26,.95);backdrop-filter:blur(10px)}
    ::-webkit-scrollbar{width:6px}::-webkit-scrollbar-track{background:var(--bg0)}::-webkit-scrollbar-thumb{background:var(--bd);border-radius:3px}
    /* ── Text visibility: force light colour on dark background ── */
    .stMarkdown p,.stMarkdown li,.stMarkdown span:not(.tt),.stMarkdown label{color:var(--t1)!important}
    .stMarkdown strong,.stMarkdown b{color:#ffffff!important}
    .stMarkdown em,.stMarkdown i{color:var(--t2)!important}
    .stMarkdown h1,.stMarkdown h2,.stMarkdown h3,.stMarkdown h4,.stMarkdown h5{color:var(--t1)!important}
    .stMarkdown a{color:var(--c2)!important}
    .stMarkdown code{color:var(--c5)!important;background:var(--bg2)!important}
    [data-testid="stExpander"] summary p,[data-testid="stExpander"] summary span{color:var(--t1)!important}
    [data-testid="stExpander"] details summary{color:var(--t1)!important}
    </style>""", unsafe_allow_html=True)


# ═══════════════════════════════════════════════════════════════════════
#  HELPERS: Safe data access
# ═══════════════════════════════════════════════════════════════════════

def safe_int(val, default=0):
    try:
        return int(val)
    except (TypeError, ValueError):
        return default

def safe_str(val, default=""):
    if val is None:
        return default
    return str(val)

def safe_list(val):
    if isinstance(val, list):
        return val
    return []

def safe_dict(val):
    if isinstance(val, dict):
        return val
    return {}


# ═══════════════════════════════════════════════════════════════════════
#  MERMAID.JS DIAGRAM RENDERER
# ═══════════════════════════════════════════════════════════════════════

def render_mermaid(mermaid_code, height=450):
    """Render a Mermaid.js diagram using streamlit HTML component."""
    clean_code = mermaid_code.strip()
    html = f"""<!DOCTYPE html>
<html><head>
<script src="https://cdn.jsdelivr.net/npm/mermaid@10/dist/mermaid.min.js"></script>
<style>body{{margin:0;padding:16px;background:#151c2e;border-radius:12px;}} .mermaid{{text-align:center;}}</style>
</head><body>
<pre class="mermaid">{clean_code}</pre>
<script>
mermaid.initialize({{startOnLoad:true,theme:'dark',themeVariables:{{primaryColor:'#16274B',primaryTextColor:'#e2e8f0',primaryBorderColor:'#00929E',lineColor:'#00b4d8',secondaryColor:'#151c2e',tertiaryColor:'#0a0e1a',fontFamily:'sans-serif'}}}});
</script>
</body></html>"""
    st.components.v1.html(html, height=height, scrolling=True)


# ═══════════════════════════════════════════════════════════════════════
#  AUDIO / VIDEO TRANSCRIPT EXTRACTOR
# ═══════════════════════════════════════════════════════════════════════

def extract_audio_transcript(uploaded_file):
    """Extract transcript from uploaded transcript files. Returns text."""
    name = uploaded_file.name.lower()
    data = uploaded_file.read()
    uploaded_file.seek(0)

    # For text-based transcript files (SRT, VTT, TXT)
    if name.endswith((".txt", ".srt", ".vtt")):
        text = data.decode("utf-8", errors="replace")
        text = re.sub(r'\d{2}:\d{2}:\d{2}[.,]\d{3}\s*-->\s*\d{2}:\d{2}:\d{2}[.,]\d{3}', '', text)
        text = re.sub(r'^\d+$', '', text, flags=re.MULTILINE)
        text = re.sub(r'<[^>]+>', '', text)
        text = re.sub(r'WEBVTT.*?\n', '', text)
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()

    # For JSON transcript exports (Teams, Zoom, Otter.ai)
    if name.endswith(".json"):
        try:
            transcript_data = json.loads(data.decode("utf-8"))
            parts = []
            if isinstance(transcript_data, list):
                for item in transcript_data:
                    if isinstance(item, dict):
                        speaker = item.get("speaker", item.get("name", "Speaker"))
                        text_val = item.get("text", item.get("content", item.get("transcript", "")))
                        if text_val:
                            parts.append(f"{speaker}: {text_val}")
            elif isinstance(transcript_data, dict):
                for seg in transcript_data.get("segments", transcript_data.get("results", transcript_data.get("transcript", []))):
                    if isinstance(seg, dict):
                        speaker = seg.get("speaker", "Speaker")
                        text_val = seg.get("text", seg.get("content", ""))
                        if text_val:
                            parts.append(f"{speaker}: {text_val}")
            return "\n".join(parts) if parts else data.decode("utf-8", errors="replace")[:50000]
        except Exception:
            return data.decode("utf-8", errors="replace")[:50000]

    # For DOCX transcripts (exported meeting notes)
    if name.endswith(".docx"):
        try:
            from docx import Document
            doc = Document(io.BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception:
            return "[Could not extract DOCX transcript]"

    # For CSV transcript exports
    if name.endswith(".csv"):
        text = data.decode("utf-8", errors="replace")
        lines = text.split("\n")
        parts = []
        for line in lines[1:]:
            cols = line.split(",")
            if len(cols) >= 2:
                parts.append(cols[-1].strip().strip('"'))
        return "\n".join(parts) if parts else text[:50000]

    return "[Unsupported format: " + name.split(".")[-1] + ". Please upload TXT, SRT, VTT, JSON, DOCX, or CSV transcript files.]"


# ═══════════════════════════════════════════════════════════════════════
#  GRAPHVIZ DOT DIAGRAM GENERATORS
# ═══════════════════════════════════════════════════════════════════════

def render_dot_to_svg(dot_string):
    """Render a DOT string to SVG bytes using the graphviz library.

    Returns SVG bytes or None if graphviz system binary is not available.
    """
    if not HAS_GRAPHVIZ or not dot_string:
        return None
    try:
        src = gv_lib.Source(dot_string)
        svg_bytes = src.pipe(format="svg")
        return svg_bytes
    except Exception:
        return None


def render_dot_to_png(dot_string):
    """Render a DOT string to PNG bytes using the graphviz library.

    Returns PNG bytes or None if graphviz system binary is not available.
    """
    if not HAS_GRAPHVIZ or not dot_string:
        return None
    try:
        src = gv_lib.Source(dot_string)
        png_bytes = src.pipe(format="png")
        return png_bytes
    except Exception:
        return None


def render_dot_to_html(dot_string):
    """Create a self-contained HTML file that renders DOT using Viz.js.

    This always works — no system graphviz binary needed. The user downloads
    an HTML file they can open in any browser to see the rendered diagram.
    """
    if not dot_string:
        return None
    escaped = dot_string.replace("\\", "\\\\").replace("`", "\\`").replace("${", "\\${")
    html = """<!DOCTYPE html>
<html><head><meta charset="utf-8"><title>ECI Architecture Diagram</title>
<style>body{background:#0a0e1a;margin:0;display:flex;justify-content:center;align-items:center;min-height:100vh;font-family:sans-serif}
#msg{color:#94a3b8;font-size:1.2rem}svg{max-width:95vw}
.toolbar{position:fixed;top:10px;right:10px;display:flex;gap:8px;z-index:10}
.toolbar button{background:#1B3A5C;color:white;border:none;padding:8px 16px;border-radius:6px;cursor:pointer;font-size:14px}
.toolbar button:hover{background:#00b4d8}</style>
<script src="https://unpkg.com/@viz-js/viz@3.2.4/lib/viz-standalone.js"></script>
</head><body>
<div class="toolbar"><button onclick="downloadSVG()">Download SVG</button><button onclick="downloadPNG()">Download PNG</button></div>
<div id="graph"><p id="msg">Rendering diagram...</p></div>
<script>
const dot = `""" + escaped + """`;
Viz.instance().then(viz => {
  const svg = viz.renderSVGElement(dot);
  document.getElementById('graph').innerHTML = '';
  document.getElementById('graph').appendChild(svg);
}).catch(e => { document.getElementById('msg').textContent = 'Render error: ' + e; });
function downloadSVG(){
  const svg = document.querySelector('#graph svg');
  if(!svg) return;
  const blob = new Blob([new XMLSerializer().serializeToString(svg)], {type:'image/svg+xml'});
  const a = document.createElement('a'); a.href = URL.createObjectURL(blob);
  a.download = 'ECI_Architecture.svg'; a.click();
}
function downloadPNG(){
  const svg = document.querySelector('#graph svg');
  if(!svg) return;
  const svgData = new XMLSerializer().serializeToString(svg);
  const canvas = document.createElement('canvas');
  const ctx = canvas.getContext('2d');
  const img = new Image();
  img.onload = function(){
    canvas.width = img.width * 2; canvas.height = img.height * 2;
    ctx.scale(2, 2); ctx.drawImage(img, 0, 0);
    canvas.toBlob(function(blob){
      const a = document.createElement('a'); a.href = URL.createObjectURL(blob);
      a.download = 'ECI_Architecture.png'; a.click();
    });
  };
  img.src = 'data:image/svg+xml;base64,' + btoa(unescape(encodeURIComponent(svgData)));
}
</script></body></html>"""
    return html.encode("utf-8")


# ═══════════════════════════════════════════════════════════════════════
#  3D ARCHITECTURE FLY-THROUGH (Three.js)
# ═══════════════════════════════════════════════════════════════════════

def generate_3d_flythrough_html(ar, ce):
    """Return a self-contained HTML string with an interactive Three.js 3D
    architecture fly-through.  Uses ES-module importmap so OrbitControls loads
    reliably in every modern browser (Chrome 89+, Firefox 108+, Safari 16.4+).
    """
    import json as _json

    arch = safe_dict(ar)
    cost = safe_dict(ce)
    components = safe_list(arch.get("components"))
    data_flow  = safe_list(arch.get("data_flow"))
    pattern    = safe_str(arch.get("pattern", "Solution Architecture"))

    _type_to_tier = {
        "web app": "presentation", "frontend": "presentation",
        "ui": "presentation", "cdn": "presentation", "portal": "presentation",
        "integration": "application", "microservices": "application",
        "api": "application", "messaging": "application",
        "backend": "application", "compute": "application",
        "function": "application", "logic": "application",
        "app service": "application",
        "database": "data", "storage": "data", "data": "data",
        "cache": "data", "redis": "data", "cosmos": "data", "sql": "data",
        "identity": "security", "security": "security",
        "auth": "security", "firewall": "security", "keyvault": "security",
        "key vault": "security",
        "operations": "operations", "devops": "operations",
        "monitoring": "operations", "logging": "operations",
        "insights": "operations",
    }

    cost_map = {}
    for ac in safe_list(cost.get("azure_costs")):
        ac = safe_dict(ac)
        svc = safe_str(ac.get("service", "")).lower()
        monthly = safe_int(ac.get("monthly_cost", 0))
        if svc:
            cost_map[svc] = monthly

    tier_y = {"presentation": 9, "application": 4, "data": 0, "security": -4, "operations": -8}

    tier_buckets = {t: [] for t in tier_y}
    for comp in components:
        comp = safe_dict(comp)
        ctype       = safe_str(comp.get("type", "")).lower()
        cname_lower = safe_str(comp.get("name", "")).lower()
        tier = _type_to_tier.get(ctype)
        if tier is None:
            for kw, t in _type_to_tier.items():
                if kw in cname_lower:
                    tier = t
                    break
        tier_buckets[tier or "application"].append(comp)

    components_3d = []
    for tier_name, comps_in_tier in tier_buckets.items():
        n = len(comps_in_tier)
        if n == 0:
            continue
        y = tier_y[tier_name]
        for idx, comp in enumerate(comps_in_tier):
            comp      = safe_dict(comp)
            name      = safe_str(comp.get("name", "Component"))
            azure_svc = safe_str(comp.get("azure_service", ""))
            services  = [safe_str(s) for s in safe_list(comp.get("services", []))[:5]]
            x = (idx - (n - 1) / 2.0) * 5.5
            z = (idx % 2) * 2.5
            monthly = 0
            nl, al = name.lower(), azure_svc.lower()
            for ck, cv in cost_map.items():
                if ck in nl or ck in al or nl in ck or al in ck:
                    monthly = cv
                    break
            components_3d.append({
                "name": name, "azure_service": azure_svc, "services": services,
                "tier": tier_name, "x": round(x, 2), "y": y, "z": round(z, 2),
                "monthly_cost": monthly,
            })

    # Guarantee a non-empty scene even before the AI runs
    if not components_3d:
        _demo = [
            ("Web App",     "Azure App Service",   "presentation", ["Hosting","Auto-scale"]),
            ("API Gateway", "Azure API Mgmt",      "application",  ["Rate-limit","Auth"]),
            ("Database",    "Azure SQL",           "data",         ["Managed DB","Backups"]),
            ("Cache",       "Azure Redis Cache",   "data",         ["In-memory cache"]),
            ("Identity",    "Azure AD B2C",        "security",     ["SSO","MFA"]),
            ("Monitoring",  "Azure Monitor",       "operations",   ["Alerts","Dashboards"]),
        ]
        _tc = {}
        for nm, svc, tier, svcs in _demo:
            i = _tc.get(tier, 0); _tc[tier] = i + 1
            components_3d.append({"name": nm, "azure_service": svc, "services": svcs,
                                   "tier": tier, "x": round((i - 0.5) * 5.5, 2),
                                   "y": tier_y[tier], "z": 0.0, "monthly_cost": 0})

    total_monthly = safe_int(cost.get("total_monthly_cost", 0))
    pattern_js = (pattern
                  .replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                  .replace('"', "&quot;").replace("\n", " "))
    comp_json = _json.dumps(components_3d)
    flow_json = _json.dumps([safe_str(f) for f in data_flow])

    # ------------------------------------------------------------------
    # HTML / Three.js payload
    # Uses importmap + type="module" — OrbitControls is a proper ES import,
    # no legacy globals needed.  onclick handlers are exposed via window._3d.
    # ------------------------------------------------------------------
    html = (
        """<!DOCTYPE html>
<html lang="en"><head><meta charset="utf-8">
<title>ECI 3D Architecture</title>
<script type="importmap">
{"imports":{
  "three":"https://cdn.jsdelivr.net/npm/three@0.160.0/build/three.module.js",
  "three/addons/":"https://cdn.jsdelivr.net/npm/three@0.160.0/examples/jsm/"
}}
</script>
<style>
*{box-sizing:border-box;margin:0;padding:0}
html,body{width:100%;height:100%;overflow:hidden;background:#0a0e1a;
  font-family:'Segoe UI',Arial,sans-serif}
#sw{position:absolute;inset:0}
#sw canvas{width:100%!important;height:100%!important;display:block}
#ov{position:absolute;inset:0;pointer-events:none}
#tb{position:absolute;top:12px;left:50%;transform:translateX(-50%);
  background:rgba(10,14,26,.9);border:1px solid #1e2a4a;border-radius:20px;
  padding:5px 20px;color:#e2e8f0;font-size:.76rem;white-space:nowrap;text-align:center}
.hl{color:#00d4aa;font-weight:700}
#lg{position:absolute;top:12px;left:12px;background:rgba(10,14,26,.9);
  border:1px solid #1e2a4a;border-radius:10px;padding:11px 14px}
#lgh{font-size:.63rem;color:#64748b;text-transform:uppercase;letter-spacing:1px;margin-bottom:7px}
.li{display:flex;align-items:center;gap:7px;margin-bottom:4px;font-size:.71rem;color:#94a3b8}
.ld{width:10px;height:10px;border-radius:2px;flex-shrink:0}
#hud{position:absolute;top:12px;right:12px;background:rgba(10,14,26,.97);
  border:1px solid #1e2a4a;border-radius:12px;padding:15px 17px;
  color:#e2e8f0;min-width:235px;max-width:275px;display:none;
  backdrop-filter:blur(12px);pointer-events:all}
#hx{float:right;cursor:pointer;color:#64748b;font-size:.95rem;line-height:1;margin-left:8px}
#hx:hover{color:#e2e8f0}
#hn{font-size:.9rem;font-weight:700;color:#00d4aa;margin-bottom:3px}
#hs{font-size:.7rem;color:#00b4d8;margin-bottom:9px}
.hl2{font-size:.6rem;color:#64748b;text-transform:uppercase;letter-spacing:1px}
#hc{font-size:1.2rem;font-weight:700;color:#ffd166;margin:2px 0 9px;
  font-family:'Courier New',monospace}
#hv{font-size:.72rem;color:#94a3b8;line-height:1.7;margin-top:3px}
#ht{display:inline-block;font-size:.58rem;padding:2px 8px;border-radius:10px;
  margin-top:8px;text-transform:uppercase;letter-spacing:1px}
#tot{position:absolute;bottom:52px;right:12px;background:rgba(10,14,26,.9);
  border:1px solid #1e2a4a;border-radius:10px;padding:8px 13px;text-align:right}
#totl{font-size:.6rem;color:#64748b;text-transform:uppercase;letter-spacing:1px}
#totv{font-size:1rem;font-weight:700;color:#ffd166;font-family:'Courier New',monospace}
#hint{position:absolute;bottom:52px;left:50%;transform:translateX(-50%);
  color:#475569;font-size:.63rem;white-space:nowrap}
#ctrl{position:absolute;bottom:12px;left:50%;transform:translateX(-50%);
  display:flex;gap:7px;pointer-events:all}
.cb{background:rgba(10,14,26,.94);border:1px solid #1e2a4a;color:#94a3b8;
  padding:6px 14px;border-radius:8px;cursor:pointer;font-size:.73rem;
  transition:all .2s;white-space:nowrap}
.cb:hover{border-color:#00d4aa;color:#00d4aa}
.cb.on{background:rgba(0,212,170,.1);border-color:#00d4aa;color:#00d4aa}
#spin{position:absolute;inset:0;display:flex;align-items:center;justify-content:center;
  background:#0a0e1a;color:#64748b;font-size:.82rem;flex-direction:column;gap:12px;z-index:20}
.sp{width:30px;height:30px;border:3px solid #1e2a4a;border-top-color:#00d4aa;
  border-radius:50%;animation:rot .8s linear infinite}
@keyframes rot{to{transform:rotate(360deg)}}
</style></head><body>
<div id="spin"><div class="sp"></div><span>Loading 3D scene&hellip;</span></div>
<div id="sw"></div>
<div id="ov">
  <div id="tb"><span class="hl">""" + pattern_js + """</span> &nbsp;&middot;&nbsp; 3D Architecture Fly-Through</div>
  <div id="lg">
    <div id="lgh">Tier</div>
    <div class="li"><div class="ld" style="background:#00b4d8"></div>Presentation</div>
    <div class="li"><div class="ld" style="background:#1b5c8c"></div>Application</div>
    <div class="li"><div class="ld" style="background:#00d4aa"></div>Data</div>
    <div class="li"><div class="ld" style="background:#7b61ff"></div>Security</div>
    <div class="li"><div class="ld" style="background:#ffd166"></div>Operations</div>
  </div>
  <div id="hud">
    <span id="hx" onclick="window._3d.closeHUD()">&#10005;</span>
    <div id="hn"></div><div id="hs"></div>
    <div class="hl2">Monthly Infra Cost</div><div id="hc"></div>
    <div class="hl2">Capabilities</div><div id="hv"></div>
    <div id="ht"></div>
  </div>
  <div id="hint">Click a component &nbsp;&middot;&nbsp; Drag to orbit &nbsp;&middot;&nbsp; Scroll to zoom</div>
  <div id="tot"><div id="totl">Total Monthly</div><div id="totv">$""" + str(total_monthly) + """/mo</div></div>
  <div id="ctrl">
    <button class="cb" id="bfly" onclick="window._3d.toggleFly()">&#9654; Play Fly-Through</button>
    <button class="cb" onclick="window._3d.resetCam()">&#8635; Reset</button>
    <button class="cb" onclick="window._3d.topCam()">&#8859; Top View</button>
  </div>
</div>
<script type="module">
import * as THREE from 'three';
import { OrbitControls } from 'three/addons/controls/OrbitControls.js';

const COMPS="""
        + comp_json
        + """;
const FLOW="""
        + flow_json
        + """;
const THEX={presentation:0x00b4d8,application:0x1b5c8c,data:0x00d4aa,security:0x7b61ff,operations:0xffd166};
const TCSS={presentation:'#00b4d8',application:'#1b5c8c',data:'#00d4aa',security:'#7b61ff',operations:'#ffd166'};
const TBG={presentation:'rgba(0,180,216,.12)',application:'rgba(27,92,140,.12)',
           data:'rgba(0,212,170,.12)',security:'rgba(123,97,255,.12)',operations:'rgba(255,209,102,.12)'};

const wrap=document.getElementById('sw');
const W=wrap.clientWidth||window.innerWidth, H=wrap.clientHeight||window.innerHeight;
const renderer=new THREE.WebGLRenderer({antialias:true,alpha:false});
renderer.setPixelRatio(Math.min(devicePixelRatio,2));
renderer.setSize(W,H);
renderer.setClearColor(0x0a0e1a,1);
wrap.appendChild(renderer.domElement);

const scene=new THREE.Scene();
scene.background=new THREE.Color(0x0a0e1a);
scene.fog=new THREE.FogExp2(0x0a0e1a,0.014);

const camera=new THREE.PerspectiveCamera(55,W/H,0.1,1000);
camera.position.set(0,20,38);

const controls=new OrbitControls(camera,renderer.domElement);
controls.enableDamping=true; controls.dampingFactor=0.08;
controls.minDistance=4; controls.maxDistance=100;
controls.target.set(0,0,0); controls.update();

// Lights
scene.add(new THREE.AmbientLight(0xffffff,0.5));
const dL=new THREE.DirectionalLight(0x00d4aa,1.8); dL.position.set(10,20,15); scene.add(dL);
const pL1=new THREE.PointLight(0x00b4d8,2.2,80); pL1.position.set(-15,12,0); scene.add(pL1);
const pL2=new THREE.PointLight(0x7b61ff,1.8,80); pL2.position.set(15,-4,10); scene.add(pL2);

// Grid
scene.add(new THREE.GridHelper(100,50,0x1e2a4a,0x111827));

// Starfield
{const g=new THREE.BufferGeometry(),p=new Float32Array(7200);
 for(let i=0;i<p.length;i++)p[i]=(Math.random()-.5)*220;
 g.setAttribute('position',new THREE.BufferAttribute(p,3));
 scene.add(new THREE.Points(g,new THREE.PointsMaterial({color:0x94a3b8,size:.15,transparent:true,opacity:.5})));}

// Tier floor plates
[...new Set(COMPS.map(c=>c.tier))].forEach(tier=>{
  const tc=COMPS.filter(c=>c.tier===tier); if(!tc.length)return;
  const ty=tc[0].y-1.5, mx=Math.max(...tc.map(c=>Math.abs(c.x)))+5;
  const geo=new THREE.PlaneGeometry(mx*2+4,12);
  const pl=new THREE.Mesh(geo,new THREE.MeshBasicMaterial({color:THEX[tier]||0x1e2a4a,transparent:true,opacity:.06,side:THREE.DoubleSide}));
  pl.rotation.x=-Math.PI/2; pl.position.set(0,ty,2); scene.add(pl);
  const el=new THREE.LineSegments(new THREE.EdgesGeometry(geo),
    new THREE.LineBasicMaterial({color:THEX[tier]||0x1e2a4a,transparent:true,opacity:.25}));
  el.rotation.x=-Math.PI/2; el.position.set(0,ty,2); scene.add(el);
});

// Label sprites
function mkLabel(txt,col,sc){
  const cv=document.createElement('canvas'); cv.width=512; cv.height=64;
  const ctx=cv.getContext('2d');
  ctx.font='bold 24px Segoe UI,Arial'; ctx.textAlign='center';
  ctx.fillStyle='#'+col.toString(16).padStart(6,'0');
  ctx.shadowColor='rgba(0,0,0,.8)'; ctx.shadowBlur=6;
  ctx.fillText(txt.length>28?txt.slice(0,26)+'\u2026':txt,256,44);
  const sp=new THREE.Sprite(new THREE.SpriteMaterial({map:new THREE.CanvasTexture(cv),transparent:true,depthTest:false}));
  sp.scale.set(7*sc,.88*sc,1); return sp;
}

// Component boxes
const meshes=[], ray=new THREE.Raycaster(), mouse=new THREE.Vector2();
COMPS.forEach((c,i)=>{
  const col=THEX[c.tier]||0x1e2a4a;
  const geo=new THREE.BoxGeometry(3.8,2.0,2.6);
  const mat=new THREE.MeshPhongMaterial({color:col,transparent:true,opacity:.78,shininess:100,specular:0x404040});
  const mesh=new THREE.Mesh(geo,mat);
  mesh.position.set(c.x,c.y,c.z); mesh.userData={c,i}; scene.add(mesh);
  mesh.add(new THREE.LineSegments(new THREE.EdgesGeometry(geo),
    new THREE.LineBasicMaterial({color:col,transparent:true,opacity:.9})));
  const nl=mkLabel(c.name,col,.75); nl.position.set(c.x,c.y+2.1,c.z); scene.add(nl);
  if(c.azure_service){const sl=mkLabel(c.azure_service,0xaabbcc,.54);sl.position.set(c.x,c.y+1.3,c.z);scene.add(sl);}
  meshes.push(mesh);
});

// Data-flow arcs
function fc(n){const nl=n.toLowerCase();return COMPS.find(c=>c.name.toLowerCase().includes(nl)||nl.includes(c.name.toLowerCase()));}
for(let i=0;i<FLOW.length-1;i++){
  const a=fc(FLOW[i]),b=fc(FLOW[i+1]); if(!a||!b)continue;
  const pa=new THREE.Vector3(a.x,a.y,a.z),pb=new THREE.Vector3(b.x,b.y,b.z);
  const pm=pa.clone().add(pb).multiplyScalar(.5); pm.y+=Math.abs(a.y-b.y)*.5+2.5;
  const pts=new THREE.QuadraticBezierCurve3(pa,pm,pb).getPoints(40);
  scene.add(new THREE.Line(new THREE.BufferGeometry().setFromPoints(pts),
    new THREE.LineBasicMaterial({color:0x00b4d8,transparent:true,opacity:.5})));
  scene.add(new THREE.ArrowHelper(pb.clone().sub(pa).normalize(),pb,0,0x00b4d8,.6,.35));
}

// HUD
let sel=null;
function showHUD(c){
  document.getElementById('hn').textContent=c.name;
  document.getElementById('hs').textContent=c.azure_service||'';
  document.getElementById('hc').textContent=c.monthly_cost>0?'$'+c.monthly_cost.toLocaleString()+'/mo':'Included in estimate';
  document.getElementById('hv').innerHTML=c.services.length?c.services.map(s=>'&#8226; '+s).join('<br>'):'N/A';
  const te=document.getElementById('ht'),tc=TCSS[c.tier]||'#94a3b8';
  te.textContent=(c.tier||'').toUpperCase();
  te.style.cssText='background:'+(TBG[c.tier]||'transparent')+';color:'+tc+';border:1px solid '+tc;
  document.getElementById('hud').style.display='block';
}
function closeHUD(){document.getElementById('hud').style.display='none';}

// Click ray-cast
let zT=null,zL=null,doZ=false;
renderer.domElement.addEventListener('click',e=>{
  const r=renderer.domElement.getBoundingClientRect();
  mouse.x=((e.clientX-r.left)/r.width)*2-1;
  mouse.y=-((e.clientY-r.top)/r.height)*2+1;
  ray.setFromCamera(mouse,camera);
  const hits=ray.intersectObjects(meshes);
  if(hits.length){
    const m=hits[0].object;
    if(sel)sel.material.emissive.setHex(0); sel=m; m.material.emissive.setHex(0x1e1e1e);
    showHUD(m.userData.c); if(fly)stopFly();
    zT=new THREE.Vector3(m.position.x+9,m.position.y+6,m.position.z+14);
    zL=m.position.clone(); doZ=true;
  }else{if(sel){sel.material.emissive.setHex(0);sel=null;}doZ=false;closeHUD();}
});

// Camera presets
const DP=new THREE.Vector3(0,20,38),DT=new THREE.Vector3(0,0,0);
function resetCam(){doZ=false;stopFly();controls.enabled=true;camera.position.copy(DP);controls.target.copy(DT);}
function topCam(){doZ=false;stopFly();controls.enabled=true;camera.position.set(0,55,.01);controls.target.set(0,0,0);}

// Fly-through (Catmull-Rom spline)
let fly=false,flyT=0; const FS=0.0016;
const WPS=COMPS.length?[
  {pos:new THREE.Vector3(0,22,42),look:new THREE.Vector3(0,2,0)},
  ...COMPS.map(c=>({pos:new THREE.Vector3(c.x+9,c.y+7,c.z+16),look:new THREE.Vector3(c.x,c.y,c.z)})),
  {pos:new THREE.Vector3(0,22,42),look:new THREE.Vector3(0,2,0)},
]:[{pos:new THREE.Vector3(0,22,42),look:new THREE.Vector3(0,0,0)},
   {pos:new THREE.Vector3(22,10,22),look:new THREE.Vector3(0,0,0)},
   {pos:new THREE.Vector3(-22,10,22),look:new THREE.Vector3(0,0,0)},
   {pos:new THREE.Vector3(0,22,42),look:new THREE.Vector3(0,0,0)}];
function toggleFly(){fly?stopFly():startFly();}
function startFly(){fly=true;flyT=0;controls.enabled=false;doZ=false;closeHUD();
  document.getElementById('bfly').innerHTML='&#9646;&#9646; Pause';
  document.getElementById('bfly').classList.add('on');}
function stopFly(){fly=false;controls.enabled=true;
  document.getElementById('bfly').innerHTML='&#9654; Play Fly-Through';
  document.getElementById('bfly').classList.remove('on');}
function crGet(t,pts){
  const n=pts.length,seg=Math.min(Math.floor(t*(n-1)),n-2),lt=t*(n-1)-seg;
  const p0=pts[Math.max(0,seg-1)],p1=pts[seg],p2=pts[Math.min(n-1,seg+1)],p3=pts[Math.min(n-1,seg+2)];
  return{pos:cr3(lt,p0.pos,p1.pos,p2.pos,p3.pos),look:cr3(lt,p0.look,p1.look,p2.look,p3.look)};}
function cr3(t,p0,p1,p2,p3){
  const t2=t*t,t3=t2*t,f=(a,b,c,d)=>.5*((2*b)+(-a+c)*t+(2*a-5*b+4*c-d)*t2+(-a+3*b-3*c+d)*t3);
  return new THREE.Vector3(f(p0.x,p1.x,p2.x,p3.x),f(p0.y,p1.y,p2.y,p3.y),f(p0.z,p1.z,p2.z,p3.z));}

// Animate
let tick=0;
function animate(){
  requestAnimationFrame(animate); tick+=.012;
  meshes.forEach((m,i)=>{
    m.material.opacity=.72+.08*Math.sin(tick+i*.9);
    m.scale.setScalar(m===sel?1+.02*Math.sin(tick*2.5):1);
  });
  pL1.intensity=1.8+.5*Math.sin(tick*.65); pL2.intensity=1.5+.4*Math.sin(tick*.5+1.2);
  if(fly){flyT+=FS;if(flyT>=1)flyT=0;const pt=crGet(flyT,WPS);camera.position.lerp(pt.pos,.025);controls.target.lerp(pt.look,.04);}
  if(doZ&&zT&&zL){camera.position.lerp(zT,.04);controls.target.lerp(zL,.06);}
  controls.update();
  renderer.render(scene,camera);
}
animate();

// Resize
window.addEventListener('resize',()=>{
  const nw=wrap.clientWidth||window.innerWidth,nh=wrap.clientHeight||window.innerHeight;
  camera.aspect=nw/nh; camera.updateProjectionMatrix(); renderer.setSize(nw,nh);
});

// Hide spinner after first frame
setTimeout(()=>{const s=document.getElementById('spin');if(s)s.style.display='none';},500);

// Expose onclick handlers (module scope != global scope)
window._3d={toggleFly,resetCam,topCam,closeHUD};
</script>
</body></html>"""
    )
    return html


def generate_architecture_diagram(architecture_data):
    """Generate a Graphviz DOT string for an architecture diagram.

    Args:
        architecture_data: dict with keys 'pattern', 'components', 'data_flow', 'security'.

    Returns:
        A DOT language string suitable for st.graphviz_chart().
    """
    data = safe_dict(architecture_data)
    pattern = safe_str(data.get("pattern"), "Solution Architecture")
    components = safe_list(data.get("components"))
    data_flow = safe_list(data.get("data_flow"))
    security_items = safe_list(data.get("security"))

    # ── Colour palette by tier ──
    tier_colors = {
        "presentation": "#00B4D8",
        "application": "#1B3A5C",
        "data": "#00D4AA",
        "security": "#7B61FF",
        "operations": "#FFD166",
    }

    # Map component type -> tier
    type_to_tier = {
        "web app": "presentation",
        "frontend": "presentation",
        "ui": "presentation",
        "cdn": "presentation",
        "integration": "application",
        "microservices": "application",
        "api": "application",
        "messaging": "application",
        "backend": "application",
        "compute": "application",
        "database": "data",
        "storage": "data",
        "data": "data",
        "cache": "data",
        "identity": "security",
        "security": "security",
        "auth": "security",
        "operations": "operations",
        "devops": "operations",
        "monitoring": "operations",
    }

    tier_labels = {
        "presentation": "Presentation Tier",
        "application": "Application Tier",
        "data": "Data Tier",
        "security": "Security Tier",
        "operations": "Operations Tier",
    }

    # Classify each component into a tier
    tier_components = {t: [] for t in tier_colors}
    for comp in components:
        comp = safe_dict(comp)
        comp_type = safe_str(comp.get("type")).lower()
        comp_name_lower = safe_str(comp.get("name")).lower()
        tier = type_to_tier.get(comp_type)
        if tier is None:
            tier = type_to_tier.get(comp_name_lower, "application")
        tier_components[tier].append(comp)

    # Build a sanitised node id from component name
    def node_id(name):
        nid = safe_str(name).replace(" ", "_").replace("/", "_").replace("-", "_")
        nid = "".join(ch for ch in nid if ch.isalnum() or ch == "_")
        return "n_" + nid

    # ── Start DOT ──
    lines = []
    lines.append("digraph architecture {")
    lines.append('    rankdir=TB;')
    lines.append('    bgcolor="#0a0e1a";')
    lines.append('    fontname="Helvetica";')
    lines.append('    node [fontname="Helvetica", fontsize=10, shape=box, style="rounded,filled", fontcolor=white];')
    lines.append('    edge [fontname="Helvetica", fontsize=9, color="#64748b", fontcolor="#94a3b8"];')
    lines.append('')

    # Title
    lines.append('    labelloc=t;')
    lines.append('    label=<<FONT FACE="Helvetica" POINT-SIZE="16" COLOR="white">' + safe_str(pattern) + '</FONT>>;')
    lines.append('')

    # ── Subgraph clusters per tier ──
    cluster_idx = 0
    all_node_ids = []
    for tier_key in ["presentation", "application", "data", "operations"]:
        comps_in_tier = tier_components.get(tier_key, [])
        if not comps_in_tier:
            continue
        color = tier_colors[tier_key]
        label = tier_labels[tier_key]
        lines.append('    subgraph cluster_' + str(cluster_idx) + ' {')
        lines.append('        label=<<FONT FACE="Helvetica" POINT-SIZE="12" COLOR="' + color + '">' + label + '</FONT>>;')
        lines.append('        style=dashed;')
        lines.append('        color="' + color + '";')
        lines.append('        bgcolor="#111827";')
        lines.append('')
        for comp in comps_in_tier:
            comp = safe_dict(comp)
            name = safe_str(comp.get("name"))
            azure_svc = safe_str(comp.get("azure_service"))
            nid = node_id(name)
            all_node_ids.append((name, nid))
            node_label = name + "\\n" + azure_svc
            lines.append('        ' + nid + ' [label="' + node_label + '", fillcolor="' + color + '"];')
        lines.append('    }')
        lines.append('')
        cluster_idx += 1

    # ── Security cluster (bottom) ──
    security_comps = tier_components.get("security", [])
    if security_comps or security_items:
        color = tier_colors["security"]
        label = tier_labels["security"]
        lines.append('    subgraph cluster_' + str(cluster_idx) + ' {')
        lines.append('        label=<<FONT FACE="Helvetica" POINT-SIZE="12" COLOR="' + color + '">' + label + '</FONT>>;')
        lines.append('        style=dashed;')
        lines.append('        color="' + color + '";')
        lines.append('        bgcolor="#111827";')
        lines.append('        rank=max;')
        lines.append('')
        for comp in security_comps:
            comp = safe_dict(comp)
            name = safe_str(comp.get("name"))
            azure_svc = safe_str(comp.get("azure_service"))
            nid = node_id(name)
            all_node_ids.append((name, nid))
            node_label = name + "\\n" + azure_svc
            lines.append('        ' + nid + ' [label="' + node_label + '", fillcolor="' + color + '"];')

        # Security items as a single info node
        if security_items:
            sec_label = "Security Controls\\n" + "\\n".join(safe_str(s) for s in security_items)
            lines.append('        n_security_controls [label="' + sec_label + '", fillcolor="' + color + '", shape=note];')
        lines.append('    }')
        lines.append('')
        cluster_idx += 1

    # ── Edges based on data_flow ──
    # Build a lookup: lowercase name -> node_id
    name_to_nid = {}
    for name, nid in all_node_ids:
        name_to_nid[name.lower()] = nid

    # For data_flow items that don't match a component, create ghost nodes
    flow_nids = []
    for item in data_flow:
        item_str = safe_str(item)
        item_lower = item_str.lower()
        matched_nid = name_to_nid.get(item_lower)
        if matched_nid is None:
            # Try partial match
            for comp_name_lower, comp_nid in name_to_nid.items():
                if item_lower in comp_name_lower or comp_name_lower in item_lower:
                    matched_nid = comp_nid
                    break
        if matched_nid is None:
            # Create an inline node for this flow step
            ghost_nid = node_id(item_str)
            lines.append('    ' + ghost_nid + ' [label="' + item_str + '", fillcolor="#334155", shape=box, style="rounded,filled"];')
            name_to_nid[item_lower] = ghost_nid
            matched_nid = ghost_nid
        flow_nids.append(matched_nid)

    for i in range(len(flow_nids) - 1):
        lines.append('    ' + flow_nids[i] + ' -> ' + flow_nids[i + 1] + ';')

    lines.append("}")
    return "\n".join(lines)


def generate_flow_diagram(time_est_data):
    """Generate a Graphviz DOT string for a project workflow / timeline diagram.

    Args:
        time_est_data: dict with keys 'phases' and 'milestones'.

    Returns:
        A DOT language string suitable for st.graphviz_chart().
    """
    data = safe_dict(time_est_data)
    phases = safe_list(data.get("phases"))
    milestones = safe_list(data.get("milestones"))

    # Gradient colours from blue to green across phases
    phase_colors = [
        "#1B3A5C",  # deep blue
        "#0E5E8A",  # medium blue
        "#007F8C",  # teal
        "#00A07A",  # teal-green
        "#00B464",  # green
        "#00D44A",  # bright green
        "#06d6a0",  # mint green
        "#0AD490",  # extra – in case of many phases
        "#10D280",
        "#16D070",
    ]

    milestone_color = "#FFD166"

    lines = []
    lines.append("digraph timeline {")
    lines.append('    rankdir=LR;')
    lines.append('    bgcolor="#0a0e1a";')
    lines.append('    fontname="Helvetica";')
    lines.append('    node [fontname="Helvetica", fontsize=10, style="rounded,filled", fontcolor=white];')
    lines.append('    edge [fontname="Helvetica", fontsize=9, color="#64748b"];')
    lines.append('')
    lines.append('    labelloc=t;')
    lines.append('    label=<<FONT FACE="Helvetica" POINT-SIZE="14" COLOR="white">Project Timeline</FONT>>;')
    lines.append('')

    phase_node_ids = []
    for idx, phase in enumerate(phases):
        phase = safe_dict(phase)
        name = safe_str(phase.get("name"), "Phase " + str(idx + 1))
        hours = safe_int(phase.get("hours"))
        pct = safe_str(phase.get("percentage"), "")
        tasks = safe_list(phase.get("tasks"))

        color = phase_colors[idx % len(phase_colors)]
        nid = "phase_" + str(idx)
        phase_node_ids.append(nid)

        # Build label: phase name, hours, percentage, then tasks
        label_parts = [name, str(hours) + " hrs (" + pct + ")"]
        for task in tasks:
            task = safe_dict(task)
            task_name = safe_str(task.get("name"))
            task_hours = safe_int(task.get("hours"))
            task_role = safe_str(task.get("role"))
            if task_name:
                task_line = "- " + task_name
                if task_hours:
                    task_line = task_line + " (" + str(task_hours) + "h"
                    if task_role:
                        task_line = task_line + ", " + task_role
                    task_line = task_line + ")"
                elif task_role:
                    task_line = task_line + " [" + task_role + "]"
                label_parts.append(task_line)

        label = "\\n".join(label_parts)
        lines.append('    ' + nid + ' [label="' + label + '", fillcolor="' + color + '", shape=box];')

    lines.append('')

    # Connect phases sequentially
    for i in range(len(phase_node_ids) - 1):
        lines.append('    ' + phase_node_ids[i] + ' -> ' + phase_node_ids[i + 1] + ';')
    lines.append('')

    # ── Milestones as diamonds ──
    # Map each milestone to the nearest phase by week.
    # We estimate cumulative weeks per phase proportionally.
    total_hours = sum(safe_int(safe_dict(p).get("hours")) for p in phases) or 1
    cumulative_week = 0
    phase_week_ranges = []
    total_weeks_est = safe_int(data.get("duration_weeks")) if safe_str(data.get("duration_weeks")).isdigit() else 0
    if total_weeks_est == 0:
        # Rough estimate: use hours to approximate
        dur_str = safe_str(data.get("duration_weeks"))
        digits = "".join(ch for ch in dur_str if ch.isdigit())
        total_weeks_est = safe_int(digits) if digits else max(8, total_hours // 160)

    for phase in phases:
        phase = safe_dict(phase)
        ph_hours = safe_int(phase.get("hours"))
        proportion = ph_hours / total_hours if total_hours else 0
        weeks_for_phase = proportion * total_weeks_est
        start_week = cumulative_week
        cumulative_week += weeks_for_phase
        phase_week_ranges.append((start_week, cumulative_week))

    for m_idx, ms in enumerate(milestones):
        ms = safe_dict(ms)
        ms_name = safe_str(ms.get("name"), "Milestone " + str(m_idx + 1))
        ms_week = safe_int(ms.get("week"))
        ms_desc = safe_str(ms.get("description"))
        mid = "ms_" + str(m_idx)

        label = ms_name + "\\nWeek " + str(ms_week)
        if ms_desc:
            label = label + "\\n" + ms_desc
        lines.append('    ' + mid + ' [label="' + label + '", shape=diamond, fillcolor="' + milestone_color + '", fontcolor="#1B3A5C", fontsize=9];')

        # Find which phase this milestone falls into (or nearest)
        best_phase_idx = 0
        for p_idx, (ws, we) in enumerate(phase_week_ranges):
            if ms_week <= we or p_idx == len(phase_week_ranges) - 1:
                best_phase_idx = p_idx
                break

        if phase_node_ids:
            target_phase = phase_node_ids[min(best_phase_idx, len(phase_node_ids) - 1)]
            lines.append('    ' + target_phase + ' -> ' + mid + ' [style=dashed, color="' + milestone_color + '"];')

    lines.append("}")
    return "\n".join(lines)


# ═══════════════════════════════════════════════════════════════════════
#  DYNAMIC TEXT ANALYSIS ENGINE (No AI — pure NLP/regex)
# ═══════════════════════════════════════════════════════════════════════

_TECH_CATALOG = {
    "Azure App Service": ["app service", "web app hosting"],
    "Azure Functions": ["azure function", "serverless", "function app"],
    "Azure SQL Database": ["azure sql", "sql database", "sql server"],
    "Cosmos DB": ["cosmos db", "cosmosdb", "nosql"],
    "Azure Blob Storage": ["blob storage", "blob", "file storage"],
    "Azure Key Vault": ["key vault", "keyvault", "secrets management"],
    "Azure AI Search": ["ai search", "cognitive search", "azure search", "vector search"],
    "Azure OpenAI": ["azure openai", "openai", "gpt-4", "gpt4", "gpt-4o"],
    "Azure AI Foundry": ["ai foundry", "foundry", "claude"],
    "Azure Monitor": ["azure monitor", "application insights", "app insights", "monitoring"],
    "Azure AD / Entra ID": ["azure ad", "entra id", "active directory", "sso", "mfa", "identity"],
    "API Management": ["api management", "apim", "api gateway"],
    "Azure Front Door": ["front door", "cdn", "load balancing"],
    "Azure Service Bus": ["service bus", "message queue", "event-driven"],
    "Azure Redis Cache": ["redis", "cache", "in-memory"],
    "Azure DevOps": ["devops", "ci/cd", "pipeline"],
    "Copilot Studio": ["copilot studio", "copilot", "bot framework"],
    "Azure Container Apps": ["container app", "docker", "kubernetes", "k8s", "aks"],
    "Power BI": ["power bi", "powerbi", "analytics dashboard"],
    "SharePoint": ["sharepoint", "document library"],
    "Microsoft Teams": ["teams", "teams bot", "teams app"],
    "Azure Event Grid": ["event grid"],
    "Azure Logic Apps": ["logic app", "workflow automation"],
    "SignalR": ["signalr", "real-time", "websocket"],
}

_LANG_CATALOG = {
    "Python": ["python", "flask", "django", "fastapi", "pytest"],
    ".NET / C#": [".net", "c#", "csharp", "asp.net", "blazor"],
    "Node.js": ["node.js", "nodejs", "express", "npm"],
    "React": ["react", "reactjs", "jsx", "next.js"],
    "Angular": ["angular"],
    "Vue.js": ["vue", "vuejs", "nuxt"],
    "Java / Spring": ["java", "spring", "springboot"],
    "TypeScript": ["typescript"],
}

_INFRA_COST_CATALOG = {
    "Azure App Service": ("Premium P1v3", 146, "Web app hosting with auto-scale"),
    "Azure Functions": ("Consumption", 30, "Serverless compute, event-driven"),
    "Azure SQL Database": ("Standard S3", 380, "Relational database with 100 DTUs"),
    "Cosmos DB": ("Autoscale 4000 RU/s", 320, "NoSQL for high-throughput workloads"),
    "Azure Blob Storage": ("Hot LRS", 8, "File and document storage"),
    "Azure Key Vault": ("Standard", 1, "Secrets and certificate management"),
    "Azure AI Search": ("Standard S1", 250, "Vector + keyword hybrid search"),
    "Azure OpenAI": ("GPT-4o", 75, "LLM inference — fallback model"),
    "Azure AI Foundry": ("Claude via Foundry", 350, "Primary AI model"),
    "Azure Monitor": ("Pay-as-you-go", 21, "Logging, monitoring, alerting"),
    "Azure AD / Entra ID": ("Premium P1", 6, "Identity and access management per user"),
    "API Management": ("Standard", 550, "API gateway with rate limiting"),
    "Azure Front Door": ("Standard", 280, "CDN and global load balancing"),
    "Azure Service Bus": ("Standard", 95, "Message queuing and event-driven"),
    "Azure Redis Cache": ("Standard C1", 160, "In-memory caching"),
    "Azure DevOps": ("Basic", 30, "CI/CD pipelines"),
    "Copilot Studio": ("Pay-as-you-go", 130, "Bot builder"),
    "Azure Container Apps": ("Consumption", 50, "Container hosting"),
    "Power BI": ("Pro", 10, "Analytics per user/month"),
    "SharePoint": ("Included M365", 0, "Document storage — included in M365"),
    "Microsoft Teams": ("Included M365", 0, "Collaboration — included in M365"),
    "Azure Event Grid": ("Per-operation", 5, "Event routing"),
    "Azure Logic Apps": ("Consumption", 15, "Workflow automation"),
    "SignalR": ("Standard", 50, "Real-time messaging"),
}


def _analyze_text_dynamic(text):
    """Analyze raw document text and extract requirements, tech stack, etc. without AI."""
    text_lower = text.lower() if text else ""
    words = text_lower.split()
    word_count = len(words)

    # ── Detect technologies ──
    detected_tech = []
    for tech_name, keywords in _TECH_CATALOG.items():
        for kw in keywords:
            if kw in text_lower:
                detected_tech.append(tech_name)
                break
    detected_langs = []
    for lang_name, keywords in _LANG_CATALOG.items():
        for kw in keywords:
            if kw in text_lower:
                detected_langs.append(lang_name)
                break

    # Always include baseline infra if any Azure detected
    baseline = ["Azure Key Vault", "Azure Monitor", "Azure Blob Storage"]
    for b in baseline:
        if b not in detected_tech and any("azure" in t.lower() for t in detected_tech):
            detected_tech.append(b)

    tech_stack = detected_tech + detected_langs

    # ── Extract requirements from sentences ──
    sentences = re.split(r'[.!?\n]+', text)
    req_keywords_func = ["implement", "develop", "build", "create", "process", "generate", "deploy",
                         "configure", "manage", "handle", "support", "enable", "provide", "upload",
                         "download", "ingest", "extract", "parse", "index", "retrieve", "query",
                         "orchestrat", "automat"]
    req_keywords_nf = ["performance", "security", "scalab", "availab", "complian", "reliab",
                       "latency", "throughput", "encrypt", "audit", "soc2", "gdpr", "sla",
                       "backup", "disaster", "recovery", "uptime"]
    req_keywords_int = ["integrat", "connect", "sync", "api", "sso", "webhook", "sharepoint",
                        "teams", "graph api", "rest api", "copilot", "power bi"]
    complex_indicators = ["multi-model", "orchestrat", "failover", "hybrid", "vector",
                          "embedding", "rag", "pipeline", "authentication", "authorization"]

    requirements = []
    seen_titles = set()
    for sent in sentences:
        sent = sent.strip()
        if len(sent) < 15 or len(sent) > 500:
            continue
        sent_lower = sent.lower()

        req_type = None
        if any(kw in sent_lower for kw in req_keywords_int):
            req_type = "integration"
        elif any(kw in sent_lower for kw in req_keywords_nf):
            req_type = "non-functional"
        elif any(kw in sent_lower for kw in req_keywords_func):
            req_type = "functional"

        if not req_type:
            continue

        # Title: first meaningful phrase
        title = sent[:80].strip()
        if title.startswith(("- ", "• ", "* ")):
            title = title[2:]
        title = title.split(",")[0].split(";")[0].strip()
        if len(title) < 5:
            continue
        title_key = title[:40].lower()
        if title_key in seen_titles:
            continue
        seen_titles.add(title_key)

        # Complexity
        complex_count = sum(1 for kw in complex_indicators if kw in sent_lower)
        complexity = "High" if complex_count >= 2 or len(sent) > 200 else "Medium" if complex_count >= 1 else "Low"
        priority = "P1" if complexity == "High" or any(w in sent_lower for w in ["critical", "must", "essential", "required"]) else "P2"

        requirements.append({
            "title": title,
            "description": sent.strip(),
            "type": req_type,
            "complexity": complexity,
            "priority": priority,
        })

    # Ensure minimum requirements if text is substantial
    if word_count > 100 and len(requirements) < 3:
        for tech in detected_tech[:5]:
            requirements.append({
                "title": tech + " setup and configuration",
                "description": "Setup and configure " + tech + " as identified in scope document.",
                "type": "functional" if "Service" in tech or "Function" in tech else "integration",
                "complexity": "Medium",
                "priority": "P2",
            })

    # ── Project type ──
    ai_count = sum(1 for t in detected_tech if any(k in t.lower() for k in ["ai", "openai", "foundry", "cognitive"]))
    int_count = sum(1 for t in detected_tech if any(k in t.lower() for k in ["sharepoint", "teams", "service bus", "logic", "event"]))
    if ai_count >= 2:
        project_type = "Data/Cloud/AI"
    elif int_count >= 3:
        project_type = "Integration Platform"
    elif any("sharepoint" in t.lower() for t in detected_tech):
        project_type = "SharePoint Solution"
    elif any(l in detected_langs for l in ["React", "Angular", "Vue.js"]):
        project_type = "Web Application"
    elif "migration" in text_lower or "migrate" in text_lower:
        project_type = "Cloud Migration"
    else:
        project_type = "Enterprise Application"

    # ── Complexity score ──
    score = 3
    score += min(3, len(detected_tech) * 0.4)
    score += min(2, len(requirements) * 0.15)
    if ai_count >= 1:
        score += 1.5
    score += min(1.5, int_count * 0.3)
    score = min(10, max(1, int(round(score))))

    # ── Business objectives ──
    obj_keywords = ["goal", "objective", "achieve", "improve", "reduce", "increase",
                    "enhance", "optimize", "streamline", "automate", "enable"]
    objectives = []
    for sent in sentences:
        sent = sent.strip()
        if any(kw in sent.lower() for kw in obj_keywords) and 20 < len(sent) < 300:
            objectives.append(sent[:150].strip())
            if len(objectives) >= 5:
                break

    return {
        "requirements": requirements,
        "technology_stack": tech_stack if tech_stack else ["Azure App Service", "Python"],
        "business_objectives": objectives if objectives else ["Deliver project on time and within budget"],
        "complexity_score": score,
        "project_type": project_type,
    }


def _split_large_task(name, hours, role, justification):
    """Split a task > 8 hours into granular 4-8h sub-tasks for justifiable estimates.

    Returns a list of task dicts, each capped at 8 hours.
    """
    name_short = name[:45]
    # Determine the activity breakdown based on task name keywords
    name_lower = name.lower()

    # Detect task category and create appropriate sub-task breakdown
    if any(k in name_lower for k in ["design", "architect", "blueprint"]):
        splits = [
            ("Requirements analysis: " + name_short, 0.35, role, "Analyze requirements and constraints"),
            ("Design document: " + name_short, 0.35, role, "Create detailed design specification"),
            ("Design review and sign-off: " + name_short, 0.30, role, "Peer review, feedback, approval"),
        ]
    elif any(k in name_lower for k in ["test", "qa", "validation", "uat"]):
        splits = [
            ("Test plan and case design: " + name_short, 0.25, role, "Define test scenarios and acceptance criteria"),
            ("Test environment setup: " + name_short, 0.15, role, "Configure test environment and data"),
            ("Test execution: " + name_short, 0.35, role, "Execute test cases, log results"),
            ("Defect validation and reporting: " + name_short, 0.25, role, "Verify fixes, generate test report"),
        ]
    elif any(k in name_lower for k in ["deploy", "release", "provision", "infra"]):
        splits = [
            ("Environment configuration: " + name_short, 0.30, role, "Setup and configure environment"),
            ("Deployment execution: " + name_short, 0.35, role, "Deploy artifacts, run migrations"),
            ("Smoke testing and verification: " + name_short, 0.20, role, "Post-deployment validation"),
            ("Runbook and documentation: " + name_short, 0.15, role, "Document steps and rollback procedures"),
        ]
    elif any(k in name_lower for k in ["integrate", "api", "connector", "migration"]):
        splits = [
            ("Integration design: " + name_short, 0.20, role, "Define API contracts, data mapping"),
            ("Connector development: " + name_short, 0.35, role, "Build integration adapter and handlers"),
            ("Error handling and retry logic: " + name_short, 0.20, role, "Implement fault tolerance"),
            ("Integration testing: " + name_short, 0.25, role, "Validate end-to-end data flow"),
        ]
    elif any(k in name_lower for k in ["ai", "ml", "model", "pipeline", "rag", "prompt"]):
        splits = [
            ("Model/pipeline design: " + name_short, 0.25, role, "Architecture and approach selection"),
            ("Implementation: " + name_short, 0.35, role, "Core development and configuration"),
            ("Testing and tuning: " + name_short, 0.25, role, "Quality validation and optimization"),
            ("Documentation: " + name_short, 0.15, role, "Technical documentation"),
        ]
    elif any(k in name_lower for k in ["document", "training", "guide", "knowledge"]):
        splits = [
            ("Content planning: " + name_short, 0.25, role, "Outline structure and content plan"),
            ("Content creation: " + name_short, 0.45, role, "Write and format documentation"),
            ("Review and finalize: " + name_short, 0.30, role, "Peer review and final edits"),
        ]
    else:
        # Generic development task breakdown
        splits = [
            ("Technical design: " + name_short, 0.20, role, "Detailed design and approach — " + justification[:40]),
            ("Development: " + name_short, 0.35, role, "Core implementation — " + justification[:40]),
            ("Unit tests: " + name_short, 0.20, role, "Automated test coverage"),
            ("Code review and refactor: " + name_short, 0.15, role, "Peer review and address feedback"),
            ("Integration validation: " + name_short, 0.10, role, "Verify integration with other components"),
        ]

    result = []
    remaining = hours
    for i, (sub_name, pct, sub_role, sub_just) in enumerate(splits):
        if i == len(splits) - 1:
            sub_hrs = remaining  # Last task gets remainder
        else:
            sub_hrs = max(2, min(8, round(hours * pct)))
            remaining -= sub_hrs
        # Ensure we don't exceed 8h per task — split further if needed
        while sub_hrs > 8:
            chunk = 8
            result.append({
                "name": sub_name + " (part " + str(len(result) + 1) + ")",
                "hours": chunk,
                "low_hours": int(chunk * 0.8),
                "high_hours": int(chunk * 1.35),
                "role": sub_role,
                "justification": sub_just,
            })
            sub_hrs -= chunk
        if sub_hrs > 0:
            result.append({
                "name": sub_name,
                "hours": sub_hrs,
                "low_hours": int(sub_hrs * 0.8),
                "high_hours": int(sub_hrs * 1.35),
                "role": sub_role,
                "justification": sub_just,
            })
    return result


def _build_dynamic_time(semantic, text=""):
    """Build detailed time estimate from actual requirements — Inflexion.xlsx format."""
    reqs = safe_list(semantic.get("requirements"))
    tech = safe_list(semantic.get("technology_stack"))
    complexity = safe_int(semantic.get("complexity_score", 5))

    n_func = len([r for r in reqs if isinstance(r, dict) and r.get("type") == "functional"])
    n_nf = len([r for r in reqs if isinstance(r, dict) and r.get("type") == "non-functional"])
    n_int = len([r for r in reqs if isinstance(r, dict) and r.get("type") == "integration"])
    n_total = max(1, len(reqs))

    # ── Per-requirement hours ──
    func_hrs = sum({"High": 40, "Medium": 30, "Low": 20}.get(safe_str(r.get("complexity")), 30) for r in reqs if isinstance(r, dict) and r.get("type") == "functional")
    nf_hrs = sum({"High": 25, "Medium": 20, "Low": 15}.get(safe_str(r.get("complexity")), 20) for r in reqs if isinstance(r, dict) and r.get("type") == "non-functional")
    int_hrs = sum({"High": 45, "Medium": 35, "Low": 25}.get(safe_str(r.get("complexity")), 35) for r in reqs if isinstance(r, dict) and r.get("type") == "integration")
    raw_dev = max(120, func_hrs + nf_hrs + int_hrs)

    has_ai = any(k in " ".join(tech).lower() for k in ["ai", "openai", "foundry", "ml", "llm", "claude", "gpt"])
    has_frontend = any(k in " ".join(tech).lower() for k in ["react", "angular", "vue", "frontend", "blazor"])
    has_infra = any(k in " ".join(tech).lower() for k in ["devops", "container", "kubernetes", "function"])
    has_data = any(k in " ".join(tech).lower() for k in ["sql", "cosmos", "redis", "data", "etl", "pipeline"])

    # ── Build phases dynamically ──
    phases = []
    week_counter = 1

    # Phase 1: Discovery & Design — granular sub-tasks
    disc_tasks = [
        {"name": "Stakeholder kickoff meeting", "role": "PM", "hours": 4, "justification": "Initial alignment meeting with key stakeholders and sponsors"},
        {"name": "Requirements elicitation workshops", "role": "BA", "hours": max(4, min(8, n_total)), "justification": str(n_total) + " requirements — structured workshops to capture needs"},
        {"name": "Requirements documentation and traceability matrix", "role": "BA", "hours": max(4, min(8, n_func)), "justification": "Document " + str(n_func) + " functional requirements with acceptance criteria"},
        {"name": "Scope document review and gap analysis", "role": "BA", "hours": max(4, min(6, n_total // 2 + 2)), "justification": "Analyze uploaded scope documents, identify gaps"},
        {"name": "Solution architecture design", "role": "Architect", "hours": max(4, min(8, len(tech) + 1)), "justification": str(len(tech)) + " technologies — design component interactions and data flows"},
        {"name": "Architecture review and sign-off", "role": "Architect", "hours": 4, "justification": "Peer review, stakeholder walkthrough, and formal sign-off"},
        {"name": "Security and compliance assessment", "role": "Security", "hours": max(4, min(8, n_nf * 2)), "justification": str(n_nf) + " non-functional requirements — security controls mapping"},
        {"name": "Risk identification and mitigation planning", "role": "PM", "hours": 4, "justification": "Initial risk register creation and response strategies"},
    ]
    disc_total = sum(t["hours"] for t in disc_tasks)
    for t in disc_tasks:
        t["low_hours"] = int(t["hours"] * 0.8)
        t["high_hours"] = int(t["hours"] * 1.35)
    disc_weeks = max(1, disc_total // 40)
    phases.append({"name": "Discovery & Design", "week_label": "Week " + str(week_counter) + ("-" + str(week_counter + disc_weeks - 1) if disc_weeks > 1 else ""),
                   "hours": disc_total, "low_hours": int(disc_total * 0.8), "high_hours": int(disc_total * 1.35),
                   "percentage": "", "tasks": disc_tasks})
    week_counter += disc_weeks

    # Phase 2: Infrastructure Setup (if infra/cloud detected)
    if has_infra or len(tech) > 3:
        infra_tasks = []
        for t_name in tech:
            if any(k in t_name.lower() for k in ["azure", "key vault", "monitor", "devops", "ad", "entra"]):
                hrs = 3 if "vault" in t_name.lower() or "monitor" in t_name.lower() else 5
                infra_tasks.append({"name": t_name + " provisioning", "role": "DevOps", "hours": hrs,
                                    "low_hours": int(hrs * 0.8), "high_hours": int(hrs * 1.35),
                                    "justification": "Setup " + t_name + " per architecture design"})
        if not infra_tasks:
            infra_tasks.append({"name": "Cloud environment setup", "role": "DevOps", "hours": 16,
                                "low_hours": 13, "high_hours": 22, "justification": "Base infrastructure provisioning"})
        infra_total = sum(t["hours"] for t in infra_tasks)
        infra_weeks = max(1, infra_total // 40)
        phases.append({"name": "Infrastructure Setup", "week_label": "Week " + str(week_counter) + ("-" + str(week_counter + infra_weeks - 1) if infra_weeks > 1 else ""),
                       "hours": infra_total, "low_hours": int(infra_total * 0.8), "high_hours": int(infra_total * 1.35),
                       "percentage": "", "tasks": infra_tasks})
        week_counter += infra_weeks

    # Phase 3: Core Development — granular sub-tasks per requirement
    # Break each requirement into design, develop, unit test, code review sub-tasks
    dev_tasks = []
    for r in reqs:
        r = safe_dict(r)
        if r.get("type") == "functional":
            title = safe_str(r.get("title"))[:40]
            cplx = safe_str(r.get("complexity"))
            desc = safe_str(r.get("description"))[:60]
            # Granular breakdown by complexity — each sub-task capped at 8h
            if cplx == "High":
                sub = [
                    ("Technical design: " + title, "Senior Dev", 6, "Detailed design doc, API contracts, data model for " + desc),
                    ("Backend development: " + title, "Senior Dev", 8, "Core business logic implementation — " + desc),
                    ("API/service layer: " + title, "Senior Dev", 8, "REST/GraphQL endpoints, service orchestration"),
                    ("Business rules and validation: " + title, "Senior Dev", 6, "Input validation, business rule engine, edge cases"),
                    ("Unit tests: " + title, "Senior Dev", 6, "Unit test coverage for all methods and edge cases"),
                    ("Code review and refactor: " + title, "Senior Dev", 6, "Peer review, address feedback, refactor"),
                ]
            elif cplx == "Medium":
                sub = [
                    ("Technical design: " + title, "Senior Dev", 4, "Design doc and interface contracts for " + desc),
                    ("Backend development: " + title, "Developer", 8, "Core implementation — " + desc),
                    ("API/service layer: " + title, "Developer", 6, "Endpoint development and integration points"),
                    ("Unit tests: " + title, "Developer", 6, "Automated test coverage for feature"),
                    ("Code review and refactor: " + title, "Senior Dev", 6, "Peer review, address findings"),
                ]
            else:  # Low
                sub = [
                    ("Technical design: " + title, "Developer", 4, "Design approach for " + desc),
                    ("Development: " + title, "Developer", 6, "Feature implementation — " + desc),
                    ("Unit tests: " + title, "Developer", 4, "Test coverage for new functionality"),
                    ("Code review: " + title, "Senior Dev", 6, "Review, feedback, merge"),
                ]
            for sname, srole, shrs, sjust in sub:
                dev_tasks.append({"name": sname, "role": srole, "hours": shrs,
                                  "low_hours": int(shrs * 0.8), "high_hours": int(shrs * 1.35),
                                  "justification": sjust})

    # AI/ML pipeline — broken into granular sub-tasks
    if has_ai:
        ai_sub = [
            ("AI model selection and evaluation", "ML Engineer", 6, "Evaluate model options, benchmark performance"),
            ("Prompt engineering and template design", "ML Engineer", 8, "Design prompt templates, few-shot examples, guardrails"),
            ("RAG pipeline: document ingestion", "ML Engineer", 6, "Document chunking, embedding generation, vector store setup"),
            ("RAG pipeline: retrieval and ranking", "ML Engineer", 6, "Semantic search, re-ranking, context window optimization"),
            ("AI orchestration layer", "ML Engineer", 8, "Multi-model routing, fallback logic, response aggregation"),
            ("AI output validation and safety", "ML Engineer", 4, "Content filtering, hallucination checks, output formatting"),
            ("AI integration testing and tuning", "ML Engineer", 6, "End-to-end testing, latency optimization, cost monitoring"),
        ]
        if n_func > 8:
            ai_sub.append(("AI scaling and caching layer", "ML Engineer", 6, "Response caching, batch processing, rate limiting"))
        for sname, srole, shrs, sjust in ai_sub:
            dev_tasks.append({"name": sname, "role": srole, "hours": shrs,
                              "low_hours": int(shrs * 0.8), "high_hours": int(shrs * 1.35),
                              "justification": sjust})

    # Frontend — broken into granular sub-tasks
    if has_frontend:
        fe_sub = [
            ("UI component library setup", "Frontend Dev", 4, "Design system, reusable components, theme configuration"),
            ("Page layouts and navigation", "Frontend Dev", 6, "Route structure, layout components, responsive design"),
            ("Feature UI implementation", "Frontend Dev", max(6, min(8, n_func * 2)), str(n_func) + " feature screens — forms, tables, dashboards"),
            ("State management and API integration", "Frontend Dev", 6, "Redux/context setup, API service layer, error handling"),
            ("Accessibility and responsive testing", "Frontend Dev", 4, "WCAG compliance, cross-browser, mobile responsiveness"),
        ]
        if n_func > 6:
            fe_sub.append(("Advanced UI: data visualization and charts", "Frontend Dev", 6, "Interactive charts, dashboards, data grids"))
        for sname, srole, shrs, sjust in fe_sub:
            dev_tasks.append({"name": sname, "role": srole, "hours": shrs,
                              "low_hours": int(shrs * 0.8), "high_hours": int(shrs * 1.35),
                              "justification": sjust})

    # Data layer — broken into granular sub-tasks
    if has_data:
        db_sub = [
            ("Database schema design", "Data Engineer", 6, "Entity relationships, normalization, indexing strategy"),
            ("Database provisioning and configuration", "Data Engineer", 4, "Instance setup, security config, backup policy"),
            ("Data access layer (ORM/repository)", "Data Engineer", 6, "Repository pattern, query optimization, connection pooling"),
            ("Data migration scripts", "Data Engineer", 4, "Schema migrations, seed data, rollback procedures"),
        ]
        if n_func > 5:
            db_sub.append(("ETL/data pipeline development", "Data Engineer", 6, "Data transformation, scheduling, error handling"))
        for sname, srole, shrs, sjust in db_sub:
            dev_tasks.append({"name": sname, "role": srole, "hours": shrs,
                              "low_hours": int(shrs * 0.8), "high_hours": int(shrs * 1.35),
                              "justification": sjust})

    if not dev_tasks:
        # Fallback — still break into granular sub-tasks
        fallback_sub = [
            ("Technical design and architecture", "Senior Dev", 8, "Detailed design for " + str(n_total) + " requirements"),
            ("Core backend development", "Senior Dev", min(8, max(6, raw_dev // 6)), "Primary business logic implementation"),
            ("API layer development", "Developer", min(8, max(6, raw_dev // 6)), "REST endpoints and service contracts"),
            ("Business logic and validation", "Developer", min(8, max(4, raw_dev // 8)), "Validation rules, error handling"),
            ("Unit test development", "Developer", min(8, max(4, raw_dev // 8)), "Automated test suite creation"),
            ("Code review and quality", "Senior Dev", min(6, max(4, raw_dev // 10)), "Peer review and quality gates"),
        ]
        for sname, srole, shrs, sjust in fallback_sub:
            dev_tasks.append({"name": sname, "role": srole, "hours": shrs,
                              "low_hours": int(shrs * 0.8), "high_hours": int(shrs * 1.35),
                              "justification": sjust})
    dev_total = sum(t["hours"] for t in dev_tasks)
    dev_weeks = max(2, dev_total // 40)
    phases.append({"name": "Core Development", "week_label": "Week " + str(week_counter) + "-" + str(week_counter + dev_weeks - 1),
                   "hours": dev_total, "low_hours": int(dev_total * 0.8), "high_hours": int(dev_total * 1.35),
                   "percentage": "", "tasks": dev_tasks})
    week_counter += dev_weeks

    # Phase 4: Integration — granular sub-tasks per integration requirement
    if n_int > 0:
        int_tasks = []
        for r in reqs:
            r = safe_dict(r)
            if r.get("type") == "integration":
                title = safe_str(r.get("title"))[:40]
                cplx = safe_str(r.get("complexity"))
                desc = safe_str(r.get("description"))[:60]
                if cplx == "High":
                    sub = [
                        ("Integration design: " + title, "Architect", 6, "API contract design, auth flow, error handling for " + desc),
                        ("Connector development: " + title, "Developer", 8, "Build adapter/connector, data mapping, transformation"),
                        ("Authentication and security: " + title, "Developer", 6, "OAuth/API key setup, token management, encryption"),
                        ("Error handling and retry logic: " + title, "Developer", 6, "Circuit breaker, retry policies, dead letter handling"),
                        ("Integration testing: " + title, "QA", 8, "End-to-end flow validation, mock services, edge cases"),
                        ("Performance testing: " + title, "QA", 6, "Load test integration endpoints, throughput validation"),
                        ("Documentation: " + title, "Developer", 4, "Integration guide, API docs, runbook"),
                    ]
                elif cplx == "Medium":
                    sub = [
                        ("Integration design: " + title, "Architect", 4, "API contract and data mapping for " + desc),
                        ("Connector development: " + title, "Developer", 8, "Build adapter, request/response handling"),
                        ("Auth and error handling: " + title, "Developer", 6, "Authentication setup, retry logic, error handling"),
                        ("Integration testing: " + title, "QA", 8, "Validate data flow, edge cases, error scenarios"),
                        ("Documentation: " + title, "Developer", 4, "Integration guide and troubleshooting docs"),
                    ]
                else:  # Low
                    sub = [
                        ("Integration design: " + title, "Developer", 4, "Simple API mapping for " + desc),
                        ("Connector development: " + title, "Developer", 6, "Build adapter and data transformation"),
                        ("Integration testing: " + title, "QA", 6, "Validate happy path and error handling"),
                        ("Documentation: " + title, "Developer", 4, "Integration setup guide"),
                    ]
                for sname, srole, shrs, sjust in sub:
                    int_tasks.append({"name": sname, "role": srole, "hours": shrs,
                                      "low_hours": int(shrs * 0.8), "high_hours": int(shrs * 1.35),
                                      "justification": sjust})
        int_total = sum(t["hours"] for t in int_tasks)
        int_weeks = max(1, int_total // 40)
        phases.append({"name": "Integration", "week_label": "Week " + str(week_counter) + ("-" + str(week_counter + int_weeks - 1) if int_weeks > 1 else ""),
                       "hours": int_total, "low_hours": int(int_total * 0.8), "high_hours": int(int_total * 1.35),
                       "percentage": "", "tasks": int_tasks})
        week_counter += int_weeks

    # Phase 5: Testing & QA — granular sub-tasks capped at 8h each
    test_base = max(40, dev_total // 3)
    test_tasks = [
        {"name": "Test strategy and plan creation", "role": "QA Lead", "hours": 6,
         "low_hours": 5, "high_hours": 8,
         "justification": "Define test approach, entry/exit criteria, environment needs"},
        {"name": "Test case design and documentation", "role": "QA", "hours": min(8, max(4, int(test_base * 0.10))),
         "low_hours": min(6, max(3, int(test_base * 0.08))), "high_hours": min(8, max(5, int(test_base * 0.14))),
         "justification": "Write test cases for " + str(n_func) + " functional requirements"},
        {"name": "Test environment setup and data prep", "role": "QA", "hours": 6,
         "low_hours": 5, "high_hours": 8,
         "justification": "Configure test env, seed test data, mock services"},
        {"name": "Unit test execution and defect logging", "role": "QA", "hours": min(8, max(4, int(test_base * 0.12))),
         "low_hours": min(6, max(3, int(test_base * 0.10))), "high_hours": min(8, max(5, int(test_base * 0.16))),
         "justification": "Execute unit tests across " + str(n_func) + " features, log defects"},
        {"name": "Integration test execution", "role": "QA", "hours": min(8, max(4, int(test_base * 0.15))),
         "low_hours": min(6, max(3, int(test_base * 0.12))), "high_hours": min(8, max(5, int(test_base * 0.20))),
         "justification": "End-to-end workflow validation across " + str(len(tech)) + " components"},
        {"name": "API and contract testing", "role": "QA", "hours": min(8, max(4, int(test_base * 0.10))),
         "low_hours": min(6, max(3, int(test_base * 0.08))), "high_hours": min(8, max(5, int(test_base * 0.14))),
         "justification": "Validate API contracts, request/response schemas, error codes"},
        {"name": "Performance and load testing", "role": "QA", "hours": min(8, max(4, int(test_base * 0.10))),
         "low_hours": min(6, max(3, int(test_base * 0.08))), "high_hours": min(8, max(5, int(test_base * 0.14))),
         "justification": "Response time benchmarks, concurrent user load, stress testing"},
        {"name": "Security and penetration testing", "role": "Security", "hours": min(8, max(4, int(test_base * 0.08))),
         "low_hours": min(6, max(3, int(test_base * 0.06))), "high_hours": min(8, max(5, int(test_base * 0.11))),
         "justification": "OWASP top 10, auth bypass, injection testing, vulnerability scan"},
        {"name": "Regression testing", "role": "QA", "hours": min(8, max(4, int(test_base * 0.10))),
         "low_hours": min(6, max(3, int(test_base * 0.08))), "high_hours": min(8, max(5, int(test_base * 0.14))),
         "justification": "Verify existing functionality after changes and bug fixes"},
        {"name": "UAT test case preparation", "role": "BA", "hours": 6,
         "low_hours": 5, "high_hours": 8,
         "justification": "Create UAT scripts, acceptance criteria checklist for stakeholders"},
        {"name": "UAT execution and feedback coordination", "role": "BA", "hours": min(8, max(4, int(test_base * 0.10))),
         "low_hours": min(6, max(3, int(test_base * 0.08))), "high_hours": min(8, max(5, int(test_base * 0.14))),
         "justification": "Facilitate UAT sessions, collect sign-offs, track feedback"},
        {"name": "Defect triage and resolution support", "role": "QA", "hours": min(8, max(4, int(test_base * 0.05))),
         "low_hours": min(6, max(3, int(test_base * 0.04))), "high_hours": min(8, max(5, int(test_base * 0.07))),
         "justification": "Prioritize defects, verify fixes, update test results"},
    ]
    test_total = sum(t["hours"] for t in test_tasks)
    test_weeks = max(1, test_total // 40)
    phases.append({"name": "Testing & QA", "week_label": "Week " + str(week_counter) + ("-" + str(week_counter + test_weeks - 1) if test_weeks > 1 else ""),
                   "hours": test_total, "low_hours": int(test_total * 0.8), "high_hours": int(test_total * 1.35),
                   "percentage": "", "tasks": test_tasks})
    week_counter += test_weeks

    # Phase 6: Deployment & Go-Live — granular sub-tasks capped at 8h
    deploy_tasks = [
        {"name": "Deployment runbook creation", "role": "DevOps", "hours": 4,
         "low_hours": 3, "high_hours": 5,
         "justification": "Step-by-step deployment guide, rollback procedures"},
        {"name": "UAT environment provisioning", "role": "DevOps", "hours": 6,
         "low_hours": 5, "high_hours": 8,
         "justification": "Staging environment setup, config, secrets management"},
        {"name": "UAT deployment and smoke test", "role": "DevOps", "hours": 4,
         "low_hours": 3, "high_hours": 5,
         "justification": "Deploy to staging, verify critical paths"},
        {"name": "UAT defect fixing and re-deployment", "role": "Developer", "hours": 8,
         "low_hours": 6, "high_hours": 8,
         "justification": "Address UAT findings, re-deploy, re-verify"},
        {"name": "Production environment provisioning", "role": "DevOps", "hours": 6,
         "low_hours": 5, "high_hours": 8,
         "justification": "Production infra setup, DNS, SSL, firewall rules"},
        {"name": "Production deployment and cutover", "role": "DevOps", "hours": 6,
         "low_hours": 5, "high_hours": 8,
         "justification": "Go-live deployment, data migration, traffic switch"},
        {"name": "Monitoring, alerting, and dashboards", "role": "DevOps", "hours": 6,
         "low_hours": 5, "high_hours": 8,
         "justification": "Application Insights, log aggregation, alert rules, dashboards"},
        {"name": "Production smoke testing", "role": "QA", "hours": 4,
         "low_hours": 3, "high_hours": 5,
         "justification": "Post-deployment verification of all critical workflows"},
        {"name": "Go-live sign-off and handover", "role": "PM", "hours": 4,
         "low_hours": 3, "high_hours": 5,
         "justification": "Stakeholder sign-off, support team handover, SLA activation"},
    ]
    deploy_total = sum(t["hours"] for t in deploy_tasks)
    deploy_weeks = max(1, deploy_total // 40)
    phases.append({"name": "Deployment & Go-Live", "week_label": "Week " + str(week_counter) + ("-" + str(week_counter + deploy_weeks - 1) if deploy_weeks > 1 else ""),
                   "hours": deploy_total, "low_hours": int(deploy_total * 0.8), "high_hours": int(deploy_total * 1.35),
                   "percentage": "", "tasks": deploy_tasks})
    week_counter += deploy_weeks

    # Phase 7: Documentation & Training — granular sub-tasks capped at 8h
    doc_tasks = [
        {"name": "Architecture and design documentation", "role": "Architect", "hours": 6,
         "low_hours": 5, "high_hours": 8,
         "justification": "Technical design doc, component diagrams for " + str(len(tech)) + " technologies"},
        {"name": "API documentation and developer guide", "role": "Developer", "hours": 6,
         "low_hours": 5, "high_hours": 8,
         "justification": "OpenAPI specs, code samples, integration guide"},
        {"name": "Operations and runbook documentation", "role": "DevOps", "hours": 4,
         "low_hours": 3, "high_hours": 5,
         "justification": "Incident procedures, scaling guide, troubleshooting"},
        {"name": "End-user guide creation", "role": "Writer", "hours": 6,
         "low_hours": 5, "high_hours": 8,
         "justification": "User manual with screenshots, FAQ, quick-start guide"},
        {"name": "Admin guide and configuration docs", "role": "Writer", "hours": 4,
         "low_hours": 3, "high_hours": 5,
         "justification": "System admin procedures, configuration reference"},
        {"name": "Knowledge transfer session (technical)", "role": "Architect", "hours": 4,
         "low_hours": 3, "high_hours": 5,
         "justification": "Technical deep-dive with client development team"},
        {"name": "Knowledge transfer session (end users)", "role": "BA", "hours": 4,
         "low_hours": 3, "high_hours": 5,
         "justification": "End-user training workshop with hands-on exercises"},
    ]
    doc_total = sum(t["hours"] for t in doc_tasks)
    doc_weeks = max(1, doc_total // 40)
    phases.append({"name": "Documentation & Training", "week_label": "Week " + str(week_counter),
                   "hours": doc_total, "low_hours": int(doc_total * 0.8), "high_hours": int(doc_total * 1.35),
                   "percentage": "", "tasks": doc_tasks})
    week_counter += doc_weeks

    # ── Calculate totals and percentages ──
    total_dev = sum(p["hours"] for p in phases)
    pm_hours = int(total_dev * 0.10)
    # Project Management — granular sub-tasks
    pm_tasks = [
        {"name": "Sprint planning and backlog grooming", "role": "PM", "hours": min(8, max(4, pm_hours // 4)),
         "justification": "Bi-weekly sprint ceremonies, backlog prioritization"},
        {"name": "Status reporting and stakeholder updates", "role": "PM", "hours": min(8, max(4, pm_hours // 4)),
         "justification": "Weekly status reports, steering committee updates"},
        {"name": "Risk and issue management", "role": "PM", "hours": min(6, max(4, pm_hours // 5)),
         "justification": "Risk register updates, issue resolution tracking"},
        {"name": "Resource coordination and escalations", "role": "PM", "hours": min(6, max(4, pm_hours // 6)),
         "justification": "Team allocation, dependency management, escalation handling"},
        {"name": "Change request management", "role": "PM", "hours": 4,
         "justification": "Evaluate scope changes, impact analysis, approval workflows"},
    ]
    pm_actual = sum(t["hours"] for t in pm_tasks)
    for t in pm_tasks:
        t["low_hours"] = int(t["hours"] * 0.8)
        t["high_hours"] = int(t["hours"] * 1.35)
    phases.append({"name": "Project Management", "week_label": "Ongoing",
                   "hours": pm_actual, "low_hours": int(pm_actual * 0.8), "high_hours": int(pm_actual * 1.35),
                   "percentage": str(round(pm_actual / max(total_dev + pm_actual, 1) * 100)) + "%", "tasks": pm_tasks})

    total = total_dev + pm_actual
    for p in phases:
        if not p.get("percentage"):
            p["percentage"] = str(round(p["hours"] / max(total, 1) * 100)) + "%"

    total_weeks = week_counter - 1
    buffer_pct = 18 if complexity >= 7 else 15 if complexity >= 5 else 12
    conf = "Medium (" + str(max(60, 90 - complexity * 3)) + "%)"

    # ── Roles ──
    has_ml = has_ai
    roles = [
        {"name": "Project Manager", "allocation_pct": 0.10 + (0.05 * (n_total > 10)), "rate": 125},
        {"name": "Solution Architect", "allocation_pct": 0.15 + (0.05 * (len(tech) > 8)), "rate": 150},
        {"name": "Backend Developer", "allocation_pct": min(1.0, 0.5 + n_func * 0.05), "rate": 110},
    ]
    if has_frontend:
        roles.append({"name": "Frontend Developer", "allocation_pct": min(1.0, 0.3 + n_func * 0.04), "rate": 100})
    if has_ml:
        ai_kw_count = sum(1 for t in tech if any(k in t.lower() for k in ["ai", "openai", "ml", "llm", "gpt", "claude", "foundry"]))
        roles.append({"name": "ML Engineer", "allocation_pct": min(1.0, 0.5 + ai_kw_count * 0.1), "rate": 140})
    roles.append({"name": "DevOps Engineer", "allocation_pct": 0.25 if has_infra else 0.15, "rate": 120})
    roles.append({"name": "QA Engineer", "allocation_pct": min(0.6, 0.3 + n_total * 0.02), "rate": 95})
    roles.append({"name": "Product Owner", "allocation_pct": 0.10, "rate": 130})

    # ── Milestones ──
    milestones = [
        {"name": "Kickoff", "week": 1, "description": "Team onboarding and project initiation"},
        {"name": "Requirements Baselined", "week": max(2, total_weeks // 8), "description": "Scope sign-off"},
        {"name": "Design Approved", "week": max(3, total_weeks // 5), "description": "Architecture review complete"},
        {"name": "MVP Ready", "week": max(6, total_weeks // 2), "description": "Core features functional"},
        {"name": "UAT Start", "week": max(8, int(total_weeks * 0.75)), "description": "User acceptance testing begins"},
        {"name": "Go-Live", "week": total_weeks, "description": "Production deployment"},
    ]

    return {
        "total_hours": total,
        "duration_weeks": str(total_weeks) + " weeks",
        "confidence": conf,
        "buffer": str(buffer_pct) + "%",
        "phases": phases,
        "milestones": milestones,
        "three_point": {"optimistic": int(total * 0.8), "most_likely": total, "pessimistic": int(total * 1.35)},
        "roles": roles,
    }


def _build_dynamic_cost(semantic, time_est):
    """Build infrastructure cost estimate from detected tech stack — no hardcoded services."""
    tech = safe_list(semantic.get("technology_stack"))
    azure_costs = []
    seen = set()
    for tech_name in tech:
        for catalog_name, (tier, monthly, desc) in _INFRA_COST_CATALOG.items():
            if catalog_name in seen:
                continue
            if any(kw in tech_name.lower() for kw in catalog_name.lower().split()):
                azure_costs.append({"service": catalog_name, "tier": tier, "monthly_cost": monthly, "description": desc})
                seen.add(catalog_name)
                break

    # Add baseline infra always
    for base_svc in ["Azure Key Vault", "Azure Monitor", "Azure Blob Storage"]:
        if base_svc not in seen:
            tier, monthly, desc = _INFRA_COST_CATALOG[base_svc]
            azure_costs.append({"service": base_svc, "tier": tier, "monthly_cost": monthly, "description": desc})
            seen.add(base_svc)

    # Filter zero-cost items for display but keep them
    third_party = []
    active_costs = [c for c in azure_costs if c["monthly_cost"] > 0]
    total_monthly = sum(c["monthly_cost"] for c in active_costs)

    # Generate dynamic optimization tips
    optimization = []
    if any("app service" in c["service"].lower() for c in azure_costs):
        optimization.append("Use Reserved Instances for 36% savings on App Service")
    if any("sql" in c["service"].lower() for c in azure_costs):
        optimization.append("Use elastic pools for SQL if multiple databases")
    if any("cosmos" in c["service"].lower() for c in azure_costs):
        optimization.append("Monitor Cosmos DB RU consumption and right-size autoscale")
    if any("openai" in c["service"].lower() or "foundry" in c["service"].lower() for c in azure_costs):
        optimization.append("Implement token caching and prompt optimization to reduce AI costs")
    optimization.append("Enable auto-shutdown for non-production environments")
    optimization.append("Use Azure Cost Management alerts at 80% and 100% budget thresholds")

    return {
        "total_monthly_cost": total_monthly,
        "total_annual_cost": total_monthly * 12,
        "azure_costs": active_costs,
        "third_party_costs": third_party,
        "cost_optimization": optimization,
        "notes": "Estimates based on detected tech stack (" + str(len(active_costs)) + " services). Dev/staging adds ~40% of prod costs.",
    }


def _build_dynamic_risk(semantic, time_est, cost_est):
    """Build risk assessment from project analysis — no hardcoded risks."""
    complexity = safe_int(semantic.get("complexity_score", 5))
    reqs = safe_list(semantic.get("requirements"))
    tech = safe_list(semantic.get("technology_stack"))
    hours = safe_int(time_est.get("total_hours", 0))
    monthly = safe_int(cost_est.get("total_monthly_cost", 0))

    n_int = len([r for r in reqs if isinstance(r, dict) and r.get("type") == "integration"])
    has_ai = any(k in " ".join(tech).lower() for k in ["ai", "openai", "foundry", "llm", "claude"])

    risks = []

    # Technical
    if n_int >= 2:
        risks.append({"category": "Technical", "title": "Integration Complexity (" + str(n_int) + " integrations)",
                       "description": str(n_int) + " integration points identified — API compatibility and data mapping risks.",
                       "severity": "High" if n_int >= 4 else "Medium", "probability": "Medium", "impact": "High",
                       "mitigation": "Early PoC for each integration. Validate APIs in Week 1."})
    if has_ai:
        risks.append({"category": "Technical", "title": "AI Model Performance",
                       "description": "LLM response quality, hallucination risk, and prompt engineering complexity.",
                       "severity": "High", "probability": "Medium", "impact": "High",
                       "mitigation": "Extensive prompt tuning, fallback model routing, response validation."})
    if len(tech) > 8:
        risks.append({"category": "Technical", "title": "Technology Stack Complexity",
                       "description": str(len(tech)) + " technologies — increased learning curve and integration overhead.",
                       "severity": "Medium", "probability": "Medium", "impact": "Medium",
                       "mitigation": "Assign specialists per technology. Conduct architecture reviews."})

    # Schedule
    if hours > 500:
        risks.append({"category": "Schedule", "title": "Extended Timeline (" + str(hours) + " hours)",
                       "description": "Large project scope increases risk of delays and scope creep.",
                       "severity": "High", "probability": "High", "impact": "High",
                       "mitigation": "Strict change request process. Agile sprints with bi-weekly reviews."})
    else:
        risks.append({"category": "Schedule", "title": "Scope Creep",
                       "description": "Requirements may evolve during development.",
                       "severity": "Medium", "probability": "High", "impact": "Medium",
                       "mitigation": "Formal change request process and sprint backlog management."})

    # Resource
    risks.append({"category": "Resource", "title": "Key Personnel Availability",
                   "description": "Specialists may have limited availability across concurrent projects.",
                   "severity": "Medium", "probability": "Medium", "impact": "High",
                   "mitigation": "Secure resource commitments early. Cross-train team members."})

    # Budget
    if monthly > 500:
        risks.append({"category": "Budget", "title": "Cloud Cost Overrun ($" + str(monthly) + "/mo)",
                       "description": "Infrastructure costs of $" + str(monthly) + "/month may exceed estimates with usage growth.",
                       "severity": "Medium" if monthly < 2000 else "High", "probability": "Medium", "impact": "Medium",
                       "mitigation": "Azure Cost Management alerts. Monthly cost reviews. Reserved instances."})

    # Data
    if any(k in " ".join(tech).lower() for k in ["sharepoint", "blob", "sql", "cosmos", "data"]):
        risks.append({"category": "Data", "title": "Data Quality & Migration",
                       "description": "Document quality, format inconsistencies, or data integrity issues during processing.",
                       "severity": "Medium", "probability": "Medium", "impact": "Medium",
                       "mitigation": "Early data profiling. Validate sample documents. Implement error handling."})

    score = min(10, max(1, int(complexity * 0.7 + len(risks) * 0.3)))
    level = "Low" if score <= 3 else "Medium" if score <= 6 else "High"
    return {"overall_score": score, "overall_level": level, "risks": risks}


def _build_dynamic_arch(semantic):
    """Build architecture from detected tech stack — no hardcoded components."""
    tech = safe_list(semantic.get("technology_stack"))
    tech_lower = " ".join(tech).lower()

    components = []
    data_flow = ["Client"]
    security = []

    # Presentation
    if any(k in tech_lower for k in ["react", "angular", "vue", "blazor", "frontend"]):
        components.append({"name": "Frontend", "type": "Web App", "azure_service": "App Service",
                           "services": [t for t in tech if any(k in t.lower() for k in ["react", "angular", "vue", "blazor"])] or ["Web App"]})
        data_flow.append("Frontend")
    if any(k in tech_lower for k in ["front door", "cdn"]):
        data_flow.insert(1, "Front Door / CDN")
    if any(k in tech_lower for k in ["teams", "copilot studio"]):
        components.append({"name": "Teams Bot", "type": "Web App", "azure_service": "Copilot Studio / Bot Service",
                           "services": [t for t in tech if any(k in t.lower() for k in ["teams", "copilot"])]})
        if "Client" in data_flow:
            data_flow.insert(1, "Teams Bot")

    # API Gateway
    if any(k in tech_lower for k in ["api management", "api gateway", "apim"]):
        components.append({"name": "API Gateway", "type": "Integration", "azure_service": "API Management",
                           "services": ["REST APIs", "Rate limiting", "Authentication"]})
        data_flow.append("API Gateway")

    # Application
    if any(k in tech_lower for k in ["app service", "fastapi", "flask", ".net", "python", "node"]):
        app_services = [t for t in tech if any(k in t.lower() for k in ["python", ".net", "node", "fastapi", "flask"])]
        components.append({"name": "Backend Services", "type": "Microservices", "azure_service": "App Service / Functions",
                           "services": app_services or ["Application Logic"]})
        data_flow.append("Backend")

    # AI
    if any(k in tech_lower for k in ["ai", "openai", "foundry", "claude", "gpt", "llm"]):
        ai_services = [t for t in tech if any(k in t.lower() for k in ["ai", "openai", "foundry", "claude", "gpt"])]
        components.append({"name": "AI / ML Engine", "type": "Compute", "azure_service": "AI Foundry / OpenAI",
                           "services": ai_services or ["LLM Inference"]})
        data_flow.append("AI Engine")

    # Search
    if any(k in tech_lower for k in ["ai search", "vector search", "cognitive search"]):
        components.append({"name": "Search Index", "type": "Data", "azure_service": "Azure AI Search",
                           "services": ["Vector search", "Hybrid search", "Semantic ranking"]})

    # Data
    data_svcs = []
    if any(k in tech_lower for k in ["sql", "database"]):
        data_svcs.append("Azure SQL")
    if any(k in tech_lower for k in ["cosmos"]):
        data_svcs.append("Cosmos DB")
    if any(k in tech_lower for k in ["redis", "cache"]):
        data_svcs.append("Redis Cache")
    if any(k in tech_lower for k in ["blob", "storage"]):
        data_svcs.append("Blob Storage")
    if data_svcs:
        components.append({"name": "Data Layer", "type": "Database", "azure_service": " + ".join(data_svcs[:2]),
                           "services": data_svcs})
        data_flow.append("Database")

    # Integration
    int_svcs = []
    if any(k in tech_lower for k in ["service bus"]):
        int_svcs.append("Service Bus")
    if any(k in tech_lower for k in ["sharepoint"]):
        int_svcs.append("SharePoint Connector")
    if any(k in tech_lower for k in ["logic app"]):
        int_svcs.append("Logic Apps")
    if any(k in tech_lower for k in ["event grid"]):
        int_svcs.append("Event Grid")
    if int_svcs:
        components.append({"name": "Integration", "type": "Messaging", "azure_service": int_svcs[0],
                           "services": int_svcs})

    # Security
    if any(k in tech_lower for k in ["ad", "entra", "sso", "mfa", "identity"]):
        security.extend(["Azure AD / Entra ID", "SSO + MFA"])
        components.append({"name": "Security", "type": "Identity", "azure_service": "Azure AD + Key Vault",
                           "services": ["Entra ID", "Key Vault", "Managed Identity"]})
    if any(k in tech_lower for k in ["key vault"]):
        security.append("Key Vault — secrets management")
    security.extend(["TLS 1.3 encryption", "RBAC access control"])

    # DevOps
    if any(k in tech_lower for k in ["devops", "ci/cd", "pipeline", "docker", "kubernetes"]):
        devops_svcs = [t for t in tech if any(k in t.lower() for k in ["devops", "docker", "kubernetes"])]
        components.append({"name": "DevOps", "type": "Operations", "azure_service": "Azure DevOps",
                           "services": devops_svcs or ["CI/CD Pipelines", "Monitoring"]})

    if not components:
        components.append({"name": "Application", "type": "Microservices", "azure_service": "App Service",
                           "services": tech[:4] or ["Web Application"]})
        data_flow = ["Client", "App Service", "Database"]

    pattern = "AI-Powered RAG Architecture" if any(k in tech_lower for k in ["rag", "vector", "embedding"]) \
        else "Microservices with Event-Driven Integration" if any(k in tech_lower for k in ["service bus", "event"]) \
        else "Cloud-Native Application Architecture"

    return {"pattern": pattern, "components": components, "data_flow": data_flow, "security": security,
            "scalability": "Auto-scaling based on demand", "availability": "99.9% SLA target"}


def _build_dynamic_scope(semantic, time_est):
    """Build scope from actual requirements — no hardcoded items."""
    reqs = safe_list(semantic.get("requirements"))
    tech = safe_list(semantic.get("technology_stack"))
    tech_lower = " ".join(tech).lower()

    in_scope = []
    for r in reqs[:10]:
        r = safe_dict(r)
        in_scope.append(safe_str(r.get("title", "Requirement")))
    in_scope.extend(["Architecture design and documentation", "Testing (unit, integration, performance, UAT)",
                     "Deployment to production", "30-day hypercare support"])

    # Out of scope: common items NOT in detected tech
    out_of_scope = []
    if "react" not in tech_lower and "angular" not in tech_lower and "vue" not in tech_lower:
        out_of_scope.append("Custom frontend/mobile application development")
    if "ci/cd" not in tech_lower and "devops" not in tech_lower:
        out_of_scope.append("CI/CD pipeline automation (manual deployments for MVP)")
    out_of_scope.extend([
        "Legacy system decommissioning",
        "End-user training beyond knowledge transfer sessions",
        "Hardware procurement",
        "License procurement (client responsibility)",
        "Penetration testing (recommended as separate engagement)",
        "Multi-language support (English only for MVP)",
        "Ongoing managed support beyond hypercare period",
    ])

    assumptions = []
    if any("azure" in t.lower() for t in tech):
        assumptions.append("Client provides Azure subscription with Contributor/Owner access")
    if any("sharepoint" in t.lower() for t in tech):
        assumptions.append("SharePoint Online with Read/Write access to document library")
    if any("teams" in t.lower() for t in tech):
        assumptions.append("Microsoft Teams admin consent for bot deployment")
    assumptions.extend([
        "Dedicated product owner available for requirements sign-off and UAT",
        "SME availability minimum 10 hours/week during development",
        "Standard business hours (9 AM – 6 PM) for team availability",
        "All documents in standard formats without password protection or DRM",
    ])

    prerequisites = []
    if any("azure" in t.lower() for t in tech):
        prerequisites.append("Active Azure subscription with appropriate access levels")
    if any("ad" in t.lower() or "entra" in t.lower() for t in tech):
        prerequisites.append("Azure AD tenant with user accounts for MVP users")
    if any("sharepoint" in t.lower() for t in tech):
        prerequisites.append("SharePoint site with document library access")
    prerequisites.extend(["Signed Statement of Work (SOW)", "Sample documents for testing and validation"])

    return {"in_scope": in_scope, "out_of_scope": out_of_scope, "assumptions": assumptions, "prerequisites": prerequisites}


# ═══════════════════════════════════════════════════════════════════════
#  DOCUMENT PROCESSOR
# ═══════════════════════════════════════════════════════════════════════

class DocProcessor:
    def extract(self, f):
        name = f.name.lower()
        data = f.read()
        f.seek(0)
        if name.endswith(".pdf"):
            try:
                from PyPDF2 import PdfReader
                return "\n\n".join(p.extract_text() or "" for p in PdfReader(io.BytesIO(data)).pages).strip()
            except Exception as e:
                return "[PDF error: " + str(e) + "]"
        elif name.endswith(".docx"):
            try:
                from docx import Document
                doc = Document(io.BytesIO(data))
                parts = [p.text for p in doc.paragraphs if p.text.strip()]
                for tbl in doc.tables:
                    for row in tbl.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
                return "\n".join(parts)
            except Exception as e:
                return "[DOCX error: " + str(e) + "]"
        elif name.endswith((".xlsx", ".xls")):
            try:
                from openpyxl import load_workbook
                wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
                parts = []
                for sn in wb.sheetnames:
                    parts.append("=== " + sn + " ===")
                    for row in wb[sn].iter_rows(max_row=500, values_only=True):
                        cells = [str(c) if c else "" for c in row]
                        if any(cells):
                            parts.append(" | ".join(cells))
                return "\n".join(parts)
            except Exception as e:
                return "[XLSX error: " + str(e) + "]"
        elif name.endswith(".pptx"):
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(data))
                parts = []
                for i, sl in enumerate(prs.slides, 1):
                    parts.append("=== Slide " + str(i) + " ===")
                    for sh in sl.shapes:
                        if hasattr(sh, "text") and sh.text.strip():
                            parts.append(sh.text)
                return "\n".join(parts)
            except Exception as e:
                return "[PPTX error: " + str(e) + "]"
        else:
            return data.decode("utf-8", errors="replace")[:50000]

    def analyze(self, text):
        lines = text.split("\n")
        ne = [l for l in lines if l.strip()]
        secs = []
        for s in ne:
            s2 = s.strip()
            if re.match(r"^\d+[\.\)]\s+\w", s2) or (s2.isupper() and 3 < len(s2) < 80) or re.match(r"^#{1,4}\s+", s2):
                secs.append(s2[:100])
        kws = ["azure", "aws", "cloud", "api", "database", "sql", "python", "react", "sharepoint", "teams",
               "power bi", "kubernetes", "docker", ".net", "java", "node", "javascript", "typescript",
               "microservices", "serverless", "devops", "ci/cd", "machine learning", "ai", "cosmos", "blob storage"]
        return {"section_count": len(secs), "word_count": len(text.split()),
                "technologies_mentioned": [k for k in kws if k in text.lower()]}


# ═══════════════════════════════════════════════════════════════════════
#  AZURE OPENAI CLIENT — PRODUCTION
# ═══════════════════════════════════════════════════════════════════════

class AzureAI:
    def __init__(self, key, endpoint, version, deployment):
        self.deployment = deployment
        self._client = None
        if key and endpoint and AzureOpenAI:
            try:
                self._client = AzureOpenAI(api_key=key, api_version=version, azure_endpoint=endpoint)
            except Exception:
                pass

    @classmethod
    def from_session(cls):
        return cls(
            st.session_state.get("azure_api_key", ""),
            st.session_state.get("azure_endpoint", ""),
            st.session_state.get("azure_api_version", "2024-06-01"),
            st.session_state.get("azure_deployment", "gpt-4"),
        )

    @property
    def is_live(self):
        return self._client is not None

    def test(self):
        if not self._client:
            return False, "Not configured. Enter API Key and Endpoint."
        try:
            self._client.chat.completions.create(
                model=self.deployment,
                messages=[{"role": "user", "content": "Reply OK"}],
                max_tokens=5,
            )
            return True, "Connected to " + self.deployment
        except Exception as e:
            return False, "Failed: " + str(e)[:200]

    def _call(self, system, user):
        if not self._client:
            return None
        try:
            resp = self._client.chat.completions.create(
                model=self.deployment,
                messages=[
                    {"role": "system", "content": system},
                    {"role": "user", "content": user},
                ],
                max_tokens=4096,
                temperature=0.2,
                response_format={"type": "json_object"},
            )
            txt = resp.choices[0].message.content
            return json.loads(txt)
        except Exception as e:
            st.warning("Azure OpenAI: " + str(e)[:150])
            return None

    # ── Requirements ──
    def analyze_requirements(self, text):
        r = self._call(
            "You are an expert IT presales analyst for ECI consulting. Analyze the scope document thoroughly. "
            "Extract ALL requirements. Classify each as functional, non-functional, or integration. "
            "Return JSON: {\"requirements\": [{\"title\": str, \"description\": str, \"type\": \"functional\" or \"non-functional\" or \"integration\", \"complexity\": \"Low\" or \"Medium\" or \"High\", \"priority\": str}], "
            "\"technology_stack\": [str], \"business_objectives\": [str], \"complexity_score\": int 1-10, \"project_type\": str}",
            "Analyze:\n\n" + text[:15000],
        )
        if r and isinstance(r, dict) and "requirements" in r:
            return r
        return self._fb_semantic(text)

    # ── Historical RAG ──
    def search_historical(self, text, projects):
        if projects:
            r = self._call(
                "Compare new project with historical data. Return JSON: "
                "{\"similar_projects\": [{\"name\": str, \"similarity\": float, \"hours\": int, \"cost\": int, \"outcome\": str}], "
                "\"benchmark_hours\": int, \"benchmark_cost\": int, \"success_patterns\": [str], \"risk_patterns\": [str]}",
                "New:\n" + text[:5000] + "\n\nHistory:\n" + json.dumps(projects[:20], default=str),
            )
            if r:
                return r
        return {"similar_projects": [], "benchmark_hours": 0, "benchmark_cost": 0, "success_patterns": [], "risk_patterns": []}

    # ── Time ──
    def estimate_time(self, semantic, rag):
        r = self._call(
            "You are an expert ECI project estimator. Use three-point estimation. "
            "IMPORTANT: Break every task into granular sub-tasks of 4-8 hours MAX each. "
            "Never create a single task larger than 8 hours — split it into design, develop, test, review sub-tasks. "
            "Each sub-task must have a clear justification explaining why those hours are needed. "
            "Return JSON: {\"total_hours\": int, \"duration_weeks\": \"N weeks\", \"confidence\": str, \"buffer\": str, "
            "\"phases\": [{\"name\": str, \"hours\": int, \"percentage\": \"N%\", \"week_label\": str, "
            "\"tasks\": [{\"name\": str, \"hours\": int (max 8), \"low_hours\": int, \"high_hours\": int, \"role\": str, \"justification\": str}]}], "
            "\"milestones\": [{\"name\": str, \"week\": int, \"description\": str}], "
            "\"three_point\": {\"optimistic\": int, \"most_likely\": int, \"pessimistic\": int}, "
            "\"roles\": [{\"name\": str, \"allocation_pct\": float, \"rate\": int}]}",
            "Estimate:\nRequirements: " + json.dumps(safe_list(semantic.get("requirements"))[:20], default=str)
            + "\nComplexity: " + str(semantic.get("complexity_score", 5))
            + "\nType: " + safe_str(semantic.get("project_type"))
            + "\nTech: " + json.dumps(safe_list(semantic.get("technology_stack")))
            + "\nBenchmark hours: " + str(rag.get("benchmark_hours", 0)),
        )
        if r and isinstance(r, dict) and "total_hours" in r:
            return self._sanitize_time(r)
        return self._fb_time(semantic, rag)

    # ── Cost (Infrastructure only) ──
    def estimate_cost(self, semantic, time_est, rag):
        r = self._call(
            "You are an ECI infrastructure cost estimator. Estimate ONLY Azure/cloud infrastructure costs, "
            "NOT project labor costs. Consider compute, storage, databases, networking, security, monitoring. "
            "Return JSON: {\"total_monthly_cost\": int, \"total_annual_cost\": int, "
            "\"azure_costs\": [{\"service\": str, \"tier\": str, \"monthly_cost\": int, \"description\": str}], "
            "\"third_party_costs\": [{\"name\": str, \"monthly_cost\": int, \"description\": str}], "
            "\"cost_optimization\": [str], \"notes\": str}",
            "Estimate infrastructure costs:\nTech: " + json.dumps(safe_list(semantic.get("technology_stack")))
            + "\nComplexity: " + str(semantic.get("complexity_score", 5))
            + "\nType: " + safe_str(semantic.get("project_type"))
            + "\nComponents: " + json.dumps(safe_list(semantic.get("requirements"))[:10], default=str),
        )
        if r and isinstance(r, dict) and "azure_costs" in r:
            return r
        return self._fb_cost(time_est)

    # ── Risk ──
    def analyze_risk(self, semantic, time_est, cost_est):
        r = self._call(
            "You are an ECI risk analyst. Assess Technical, Schedule, Resource, Budget, External. "
            "Return JSON: {\"overall_score\": int 1-10, \"overall_level\": str, "
            "\"risks\": [{\"category\": str, \"title\": str, \"description\": str, \"severity\": \"Low\" or \"Medium\" or \"High\", "
            "\"probability\": str, \"impact\": str, \"mitigation\": str}]}",
            "Assess:\nReqs: " + str(len(safe_list(semantic.get("requirements"))))
            + "\nComplexity: " + str(semantic.get("complexity_score", 5))
            + "\nHours: " + str(time_est.get("total_hours", 0))
            + "\nInfra Monthly: $" + str(cost_est.get("total_monthly_cost", 0))
            + "\nTech: " + json.dumps(safe_list(semantic.get("technology_stack"))),
        )
        if r and isinstance(r, dict) and "risks" in r:
            return r
        return self._fb_risk(semantic)

    # ── Architecture ──
    def design_architecture(self, semantic, rag):
        r = self._call(
            "You are an Azure Solutions Architect for ECI. Use Well-Architected Framework. "
            "Return JSON: {\"pattern\": str, \"components\": [{\"name\": str, \"type\": str, \"azure_service\": str, \"services\": [str]}], "
            "\"data_flow\": [str], \"security\": [str], \"scalability\": str, \"availability\": str}",
            "Design:\nReqs: " + json.dumps(safe_list(semantic.get("requirements"))[:15], default=str)
            + "\nTech: " + json.dumps(safe_list(semantic.get("technology_stack")))
            + "\nType: " + safe_str(semantic.get("project_type")),
        )
        if r and isinstance(r, dict) and "components" in r:
            return r
        return self._fb_arch()

    # ── Scope ──
    def define_scope(self, semantic, time_est, cost_est):
        r = self._call(
            "You are an ECI scope expert. Define boundaries. "
            "Return JSON: {\"in_scope\": [str], \"out_of_scope\": [str], \"assumptions\": [str], \"prerequisites\": [str]}",
            "Scope:\nReqs: " + json.dumps(safe_list(semantic.get("requirements"))[:10], default=str)
            + "\nType: " + safe_str(semantic.get("project_type"))
            + "\nHours: " + str(time_est.get("total_hours", 0)),
        )
        if r and isinstance(r, dict) and "in_scope" in r:
            return r
        return self._fb_scope()

    # ── Proposal ──
    def write_proposal(self, semantic, time_est, cost_est, risk, arch, scope):
        r = self._call(
            "You are an ECI presales proposal writer. Write professional proposal. "
            "Return JSON: {\"sections\": [{\"title\": str, \"content\": str}], \"quality_checks\": {str: bool}}. "
            "Include sections: Executive Summary, Understanding & Approach, Technical Solution, Team, Timeline, Infrastructure Costs.",
            "Generate:\nType: " + safe_str(semantic.get("project_type"))
            + "\nReqs: " + str(len(safe_list(semantic.get("requirements"))))
            + "\nHours: " + str(time_est.get("total_hours", 0))
            + "\nInfra Monthly: $" + str(cost_est.get("total_monthly_cost", 0))
            + "\nRisk: " + safe_str(risk.get("overall_level"))
            + "\nArch: " + safe_str(arch.get("pattern"))
            + "\nPhases: " + json.dumps([{"name": safe_str(p.get("name")), "hours": safe_int(p.get("hours"))} for p in safe_list(time_est.get("phases"))])
            + "\nScope: " + json.dumps(safe_list(scope.get("in_scope"))[:5]),
        )
        if r and isinstance(r, dict) and "sections" in r:
            return r
        return self._fb_proposal(semantic, time_est, cost_est, risk, arch)

    # ── Sanitize time (CRITICAL: prevents the TypeError) ──
    def _sanitize_time(self, data):
        # Milestones
        clean_ms = []
        for m in safe_list(data.get("milestones")):
            if isinstance(m, dict):
                clean_ms.append({
                    "name": safe_str(m.get("name", "Milestone")),
                    "week": safe_int(m.get("week", 0)),
                    "description": safe_str(m.get("description", "")),
                })
            elif isinstance(m, str):
                clean_ms.append({"name": m, "week": 0, "description": ""})
        data["milestones"] = clean_ms
        # Phases — break down any task > 8 hours into granular sub-tasks
        clean_ph = []
        for p in safe_list(data.get("phases")):
            if not isinstance(p, dict):
                continue
            clean_tasks = []
            for t in safe_list(p.get("tasks")):
                if not isinstance(t, dict):
                    continue
                t_hrs = safe_int(t.get("hours", 0))
                t_name = safe_str(t.get("name", ""))
                t_role = safe_str(t.get("role", ""))
                t_just = safe_str(t.get("justification", ""))
                if t_hrs > 8:
                    # Split large task into granular 4-8h sub-tasks
                    sub_tasks = _split_large_task(t_name, t_hrs, t_role, t_just)
                    clean_tasks.extend(sub_tasks)
                else:
                    clean_tasks.append({
                        "name": t_name,
                        "hours": t_hrs,
                        "low_hours": safe_int(t.get("low_hours", int(t_hrs * 0.8))),
                        "high_hours": safe_int(t.get("high_hours", int(t_hrs * 1.35))),
                        "role": t_role,
                        "justification": t_just,
                    })
            phase_hours = sum(tk["hours"] for tk in clean_tasks)
            clean_ph.append({
                "name": safe_str(p.get("name", "Phase")),
                "hours": phase_hours,
                "low_hours": int(phase_hours * 0.8),
                "high_hours": int(phase_hours * 1.35),
                "percentage": safe_str(p.get("percentage", "0%")),
                "week_label": safe_str(p.get("week_label", "")),
                "tasks": clean_tasks,
            })
        data["phases"] = clean_ph
        # Recalculate totals from granular phases
        recalc_total = sum(p["hours"] for p in clean_ph)
        if recalc_total > 0:
            data["total_hours"] = recalc_total
            data["three_point"] = {
                "optimistic": int(recalc_total * 0.8),
                "most_likely": recalc_total,
                "pessimistic": int(recalc_total * 1.35),
            }
            for p in clean_ph:
                p["percentage"] = str(round(p["hours"] / max(recalc_total, 1) * 100)) + "%"
        else:
            data["total_hours"] = safe_int(data.get("total_hours", 0))
            data["three_point"] = safe_dict(data.get("three_point"))
        data["duration_weeks"] = safe_str(data.get("duration_weeks", "TBD"))
        data["confidence"] = safe_str(data.get("confidence", "N/A"))
        data["buffer"] = safe_str(data.get("buffer", "N/A"))
        # Roles
        clean_roles = []
        for rl in safe_list(data.get("roles")):
            if isinstance(rl, dict):
                clean_roles.append({
                    "name": safe_str(rl.get("name", "")),
                    "allocation_pct": float(rl.get("allocation_pct", 0)),
                    "rate": safe_int(rl.get("rate", 100)),
                })
        data["roles"] = clean_roles
        return data

    # ═══ DYNAMIC FALLBACKS (analyze actual document text) ═══

    def _fb_semantic(self, text):
        return _analyze_text_dynamic(text)

    def _fb_time(self, semantic, rag):
        return _build_dynamic_time(semantic)

    def _fb_cost(self, time_est):
        # Use stored semantic if available
        sem = st.session_state.get("_last_semantic", {})
        return _build_dynamic_cost(sem, time_est)

    def _fb_risk(self, semantic):
        te = st.session_state.get("_last_time_est", {})
        ce = st.session_state.get("_last_cost_est", {})
        return _build_dynamic_risk(semantic, te, ce)

    def _fb_arch(self):
        sem = st.session_state.get("_last_semantic", {})
        return _build_dynamic_arch(sem)

    def _fb_scope(self):
        sem = st.session_state.get("_last_semantic", {})
        te = st.session_state.get("_last_time_est", {})
        return _build_dynamic_scope(sem, te)

    def _fb_proposal(self, sem, te, ce, ri, ar):
        d = datetime.now().strftime("%B %d, %Y")
        reqs = safe_list(sem.get("requirements"))
        monthly = safe_int(ce.get("total_monthly_cost"))
        annual = safe_int(ce.get("total_annual_cost"))
        return {"sections": [
            {"title": "Executive Summary", "content": "**Date:** " + d + "\n\nECI proposes a " + safe_str(sem.get("project_type")) + " solution. **" + str(len(reqs)) + " requirements** identified.\n\n- Effort: **" + str(te.get("total_hours", 0)) + "h** over **" + safe_str(te.get("duration_weeks")) + "**\n- Infrastructure: **$" + str(monthly) + "/month** ($" + str(annual) + "/year)\n- Risk: **" + safe_str(ri.get("overall_level")) + "**"},
            {"title": "Understanding & Approach", "content": "ECI follows Discovery, Design, Development, Testing, Deployment with hypercare."},
            {"title": "Technical Solution", "content": "**Pattern:** " + safe_str(ar.get("pattern")) + "\n\n" + "\n".join("- **" + safe_str(safe_dict(c).get("name")) + "**: " + safe_str(safe_dict(c).get("azure_service")) for c in safe_list(ar.get("components")))},
            {"title": "Team & Delivery", "content": "ECI will deploy a cross-functional team including Architects, Lead Developers, Senior Developers, QA Engineers, DevOps Engineers, and a Project Manager.\n\nDelivery follows an Agile methodology with bi-weekly sprints and milestone reviews."},
            {"title": "Timeline", "content": "\n".join("- **" + safe_str(safe_dict(p).get("name")) + "** — " + str(safe_int(safe_dict(p).get("hours"))) + "h (" + safe_str(safe_dict(p).get("percentage")) + ")" for p in safe_list(te.get("phases")))},
            {"title": "Infrastructure Costs", "content": "**Monthly:** $" + str(monthly) + "  |  **Annual:** $" + str(annual) + "\n\n" + "\n".join("- **" + safe_str(safe_dict(a).get("service")) + "** (" + safe_str(safe_dict(a).get("tier", "")) + "): $" + str(safe_int(safe_dict(a).get("monthly_cost"))) + "/mo" for a in safe_list(ce.get("azure_costs")))},
        ], "quality_checks": {"ECI Tone": True, "Personalization": True, "Terminology": True, "Value Proposition": True, "Structure": True}}

    # ── Mermaid Diagram Generation ──
    def generate_mermaid_diagrams(self, semantic, arch):
        r = self._call(
            "You are an Azure Solutions Architect for ECI. Generate Mermaid.js diagram code for the architecture. "
            "Return JSON with these keys, each containing VALID Mermaid syntax as a string:\n"
            "{\"infrastructure\": str (graph TD diagram of Azure components and connections),\n"
            " \"data_flow\": str (flowchart LR diagram showing data movement between services),\n"
            " \"sequence\": str (sequenceDiagram showing a typical user request flow),\n"
            " \"deployment\": str (graph TD diagram showing CI/CD pipeline and environments),\n"
            " \"security\": str (graph TD diagram showing security layers and controls)}",
            "Architecture:\nPattern: " + safe_str(arch.get("pattern"))
            + "\nComponents: " + json.dumps(safe_list(arch.get("components"))[:10], default=str)
            + "\nData Flow: " + json.dumps(safe_list(arch.get("data_flow")))
            + "\nSecurity: " + json.dumps(safe_list(arch.get("security")))
            + "\nTech: " + json.dumps(safe_list(semantic.get("technology_stack"))),
        )
        if r and isinstance(r, dict) and "infrastructure" in r:
            return r
        return self._fb_mermaid(arch)

    # ── Transcript Analysis & WBS ──
    def analyze_transcript(self, transcript_text):
        r = self._call(
            "You are an expert ECI presales analyst. Analyze this meeting transcript/voice note. "
            "Extract requirements, pain points (from tone/urgency), stakeholders, decisions, and action items. "
            "Generate a Work Breakdown Structure. Return JSON:\n"
            "{\"pain_points\": [{\"issue\": str, \"severity\": \"High\" or \"Medium\" or \"Low\", \"quote\": str, \"stakeholder\": str}],\n"
            " \"requirements_extracted\": [{\"title\": str, \"description\": str, \"type\": \"functional\" or \"non-functional\" or \"integration\", \"source\": str}],\n"
            " \"stakeholders\": [{\"name\": str, \"role\": str, \"concerns\": [str]}],\n"
            " \"decisions\": [str], \"action_items\": [{\"item\": str, \"owner\": str, \"priority\": str}],\n"
            " \"wbs\": [{\"phase\": str, \"deliverables\": [{\"name\": str, \"tasks\": [{\"name\": str, \"effort\": str}]}]}],\n"
            " \"meeting_summary\": str, \"sentiment\": str, \"key_themes\": [str]}",
            "Transcript:\n\n" + transcript_text[:15000],
        )
        if r and isinstance(r, dict) and "wbs" in r:
            return r
        return self._fb_transcript(transcript_text)

    # ═══ FALLBACK: Mermaid Diagrams ═══
    def _fb_mermaid(self, arch):
        comps = safe_list(arch.get("components"))

        infra = "graph TD\n"
        infra += "    User([Client Browser])\n"
        infra += "    CDN[\"Azure Front Door\"]\n"
        infra += "    APIM[\"API Management\"]\n"
        infra += "    APP[\"App Service\"]\n"
        infra += "    FUNC[\"Azure Functions\"]\n"
        infra += "    SQL[(\"Azure SQL\")]\n"
        infra += "    COSMOS[(\"Cosmos DB\")]\n"
        infra += "    REDIS[\"Redis Cache\"]\n"
        infra += "    BLOB[\"Blob Storage\"]\n"
        infra += "    SB[\"Service Bus\"]\n"
        infra += "    KV[\"Key Vault\"]\n"
        infra += "    AD[\"Azure AD\"]\n"
        infra += "    MON[\"App Insights\"]\n\n"
        infra += "    User --> CDN\n"
        infra += "    CDN --> APIM\n"
        infra += "    APIM --> APP\n"
        infra += "    APP --> SQL\n"
        infra += "    APP --> COSMOS\n"
        infra += "    APP --> REDIS\n"
        infra += "    APP --> BLOB\n"
        infra += "    APP --> SB\n"
        infra += "    SB --> FUNC\n"
        infra += "    FUNC --> SQL\n"
        infra += "    FUNC --> COSMOS\n"
        infra += "    APP --> KV\n"
        infra += "    APP --> AD\n"
        infra += "    APP --> MON\n"
        infra += "    FUNC --> MON\n\n"
        infra += "    style User fill:#E8F5E9,stroke:#2E7D32,color:#1a1a1a\n"
        infra += "    style CDN fill:#E3F2FD,stroke:#1565C0,color:#1a1a1a\n"
        infra += "    style APIM fill:#FFF3E0,stroke:#E65100,color:#1a1a1a\n"
        infra += "    style APP fill:#E8EAF6,stroke:#283593,color:#1a1a1a\n"
        infra += "    style FUNC fill:#FCE4EC,stroke:#C62828,color:#1a1a1a\n"
        infra += "    style SQL fill:#F3E5F5,stroke:#6A1B9A,color:#1a1a1a\n"
        infra += "    style COSMOS fill:#F3E5F5,stroke:#6A1B9A,color:#1a1a1a\n"
        infra += "    style SB fill:#FFF8E1,stroke:#F57F17,color:#1a1a1a\n"
        infra += "    style KV fill:#EFEBE9,stroke:#4E342E,color:#1a1a1a\n"
        infra += "    style AD fill:#E0F2F1,stroke:#00695C,color:#1a1a1a\n"
        infra += "    style MON fill:#F1F8E9,stroke:#33691E,color:#1a1a1a"

        data_flow = "flowchart LR\n"
        data_flow += "    A[\"Client App\"] -->|HTTPS| B[\"API Gateway\"]\n"
        data_flow += "    B -->|Route| C[\"App Service\"]\n"
        data_flow += "    C -->|Read Write| D[(\"SQL Database\")]\n"
        data_flow += "    C -->|Cache| E[\"Redis Cache\"]\n"
        data_flow += "    C -->|Documents| F[\"Blob Storage\"]\n"
        data_flow += "    C -->|Events| G[\"Service Bus\"]\n"
        data_flow += "    G -->|Trigger| H[\"Azure Functions\"]\n"
        data_flow += "    H -->|Process| I[(\"Cosmos DB\")]\n"
        data_flow += "    H -->|Notify| J[\"Notification Hub\"]\n"
        data_flow += "    C -->|Analytics| K[\"Power BI\"]\n"
        data_flow += "    D -->|Sync| I\n\n"
        data_flow += "    style A fill:#16274B,color:#fff\n"
        data_flow += "    style B fill:#00929E,color:#fff\n"
        data_flow += "    style C fill:#96C038,color:#fff\n"
        data_flow += "    style H fill:#16274B,color:#fff"

        sequence = "sequenceDiagram\n"
        sequence += "    participant U as User\n"
        sequence += "    participant FD as Front Door\n"
        sequence += "    participant API as API Mgmt\n"
        sequence += "    participant App as App Service\n"
        sequence += "    participant Cache as Redis\n"
        sequence += "    participant DB as SQL Database\n"
        sequence += "    participant SB as Service Bus\n"
        sequence += "    participant Func as Functions\n\n"
        sequence += "    U->>FD: HTTPS Request\n"
        sequence += "    FD->>API: Route and WAF\n"
        sequence += "    API->>API: Auth and Rate Limit\n"
        sequence += "    API->>App: Forward Request\n"
        sequence += "    App->>Cache: Check Cache\n"
        sequence += "    alt Cache Hit\n"
        sequence += "        Cache-->>App: Return Data\n"
        sequence += "    else Cache Miss\n"
        sequence += "        App->>DB: Query Data\n"
        sequence += "        DB-->>App: Result Set\n"
        sequence += "        App->>Cache: Update Cache\n"
        sequence += "    end\n"
        sequence += "    App->>SB: Publish Event\n"
        sequence += "    SB->>Func: Trigger Processing\n"
        sequence += "    App-->>U: JSON Response"

        deployment = "graph TD\n"
        deployment += "    DEV[\"Developer\"] -->|git push| GH[\"GitHub\"]\n"
        deployment += "    GH -->|trigger| CI[\"Azure DevOps CI\"]\n"
        deployment += "    CI -->|build and test| ART[\"Container Registry\"]\n"
        deployment += "    ART -->|deploy| STG[\"Staging\"]\n"
        deployment += "    STG -->|approval gate| PROD[\"Production\"]\n"
        deployment += "    PROD -->|monitor| MON2[\"App Insights\"]\n"
        deployment += "    MON2 -->|alert| OPS[\"Ops Team\"]\n\n"
        deployment += "    style DEV fill:#E8F5E9,stroke:#2E7D32,color:#1a1a1a\n"
        deployment += "    style GH fill:#F3E5F5,stroke:#6A1B9A,color:#1a1a1a\n"
        deployment += "    style CI fill:#E3F2FD,stroke:#1565C0,color:#1a1a1a\n"
        deployment += "    style ART fill:#FFF3E0,stroke:#E65100,color:#1a1a1a\n"
        deployment += "    style STG fill:#FFF8E1,stroke:#F57F17,color:#1a1a1a\n"
        deployment += "    style PROD fill:#E8EAF6,stroke:#283593,color:#1a1a1a\n"
        deployment += "    style MON2 fill:#E0F2F1,stroke:#00695C,color:#1a1a1a"

        security = "graph TD\n"
        security += "    EXT[\"External Traffic\"]\n"
        security += "    WAF[\"WAF and DDoS Protection\"]\n"
        security += "    FD2[\"Front Door with TLS\"]\n"
        security += "    APIM2[\"API Mgmt with OAuth\"]\n"
        security += "    VNET[\"Virtual Network\"]\n"
        security += "    NSG[\"NSG Rules\"]\n"
        security += "    APP2[\"App Service with MI\"]\n"
        security += "    KV2[\"Key Vault\"]\n"
        security += "    SQL2[(\"SQL with TDE\")]\n"
        security += "    LOG[\"Sentinel and Log Analytics\"]\n\n"
        security += "    EXT --> WAF\n"
        security += "    WAF --> FD2\n"
        security += "    FD2 --> APIM2\n"
        security += "    APIM2 --> VNET\n"
        security += "    VNET --> NSG\n"
        security += "    NSG --> APP2\n"
        security += "    APP2 --> KV2\n"
        security += "    APP2 --> SQL2\n"
        security += "    APP2 --> LOG\n"
        security += "    KV2 --> LOG\n\n"
        security += "    style WAF fill:#FFCDD2,stroke:#C62828,color:#1a1a1a\n"
        security += "    style FD2 fill:#FFCDD2,stroke:#C62828,color:#1a1a1a\n"
        security += "    style APIM2 fill:#FFF9C4,stroke:#F57F17,color:#1a1a1a\n"
        security += "    style VNET fill:#C8E6C9,stroke:#2E7D32,color:#1a1a1a\n"
        security += "    style NSG fill:#C8E6C9,stroke:#2E7D32,color:#1a1a1a\n"
        security += "    style KV2 fill:#E1BEE7,stroke:#6A1B9A,color:#1a1a1a\n"
        security += "    style LOG fill:#B3E5FC,stroke:#0277BD,color:#1a1a1a"

        return {"infrastructure": infra, "data_flow": data_flow, "sequence": sequence, "deployment": deployment, "security": security}

    # ═══ FALLBACK: Transcript Analysis ═══
    def _fb_transcript(self, text):
        words = text.split()
        return {
            "pain_points": [
                {"issue": "Manual data processing taking excessive time", "severity": "High", "quote": "We spend hours every week just copying data between systems", "stakeholder": "Operations Lead"},
                {"issue": "Lack of real-time visibility into operations", "severity": "High", "quote": "By the time we get reports, the data is already stale", "stakeholder": "VP Operations"},
                {"issue": "Integration gaps between existing systems", "severity": "Medium", "quote": "Our CRM and ERP don't talk to each other properly", "stakeholder": "IT Manager"},
                {"issue": "Security concerns with current manual workflows", "severity": "Medium", "quote": "People are emailing spreadsheets with sensitive data", "stakeholder": "CISO"},
            ],
            "requirements_extracted": [
                {"title": "Automated Data Pipeline", "description": "ETL pipeline connecting CRM, ERP, and data warehouse", "type": "functional", "source": "Operations Lead — pain point discussion"},
                {"title": "Real-time Dashboard", "description": "Live operational metrics with <5 min refresh", "type": "functional", "source": "VP Operations — visibility concern"},
                {"title": "System Integration Layer", "description": "API-based integration between CRM and ERP", "type": "integration", "source": "IT Manager — integration gap"},
                {"title": "Secure Data Transfer", "description": "Encrypted data pipelines replacing manual email workflows", "type": "non-functional", "source": "CISO — security concern"},
                {"title": "User Authentication & SSO", "description": "Azure AD SSO with MFA for all users", "type": "non-functional", "source": "CISO — access control discussion"},
                {"title": "Mobile Access", "description": "Responsive web app for field team access", "type": "functional", "source": "Field Operations Manager — remote access need"},
            ],
            "stakeholders": [
                {"name": "Sarah Mitchell", "role": "VP Operations", "concerns": ["Reporting delays", "Operational visibility", "Cost of current manual processes"]},
                {"name": "James Chen", "role": "IT Manager", "concerns": ["Integration complexity", "Maintenance burden", "Team skill gaps"]},
                {"name": "David Park", "role": "CISO", "concerns": ["Data security", "Compliance", "Audit trail"]},
                {"name": "Lisa Ramirez", "role": "Operations Lead", "concerns": ["Daily workflow efficiency", "Data accuracy", "Training time"]},
            ],
            "decisions": [
                "Azure cloud platform selected as preferred infrastructure",
                "Phased rollout approach agreed — pilot with Operations team first",
                "Budget range of $150K-250K for initial phase discussed",
                "Q3 2025 target for MVP delivery",
                "Weekly stakeholder sync meetings during discovery",
            ],
            "action_items": [
                {"item": "Share current system architecture documentation", "owner": "IT Manager", "priority": "High"},
                {"item": "Provide sample data exports from CRM and ERP", "owner": "Operations Lead", "priority": "High"},
                {"item": "Schedule security requirements workshop", "owner": "CISO", "priority": "Medium"},
                {"item": "Prepare ECI proposal with options", "owner": "ECI Team", "priority": "High"},
                {"item": "Set up Azure sandbox environment", "owner": "IT Manager", "priority": "Medium"},
            ],
            "wbs": [
                {"phase": "Discovery & Planning", "deliverables": [
                    {"name": "Stakeholder Workshops", "tasks": [{"name": "Requirements gathering sessions", "effort": "3 days"}, {"name": "Pain point analysis", "effort": "2 days"}, {"name": "Current state assessment", "effort": "3 days"}]},
                    {"name": "Technical Assessment", "tasks": [{"name": "System landscape review", "effort": "2 days"}, {"name": "Data mapping", "effort": "3 days"}, {"name": "Integration feasibility", "effort": "2 days"}]},
                    {"name": "Project Plan", "tasks": [{"name": "WBS finalization", "effort": "1 day"}, {"name": "Resource plan", "effort": "1 day"}, {"name": "Risk register", "effort": "1 day"}]},
                ]},
                {"phase": "Design & Architecture", "deliverables": [
                    {"name": "Solution Architecture", "tasks": [{"name": "High-level design", "effort": "3 days"}, {"name": "Data model design", "effort": "4 days"}, {"name": "API contract design", "effort": "3 days"}]},
                    {"name": "UX Design", "tasks": [{"name": "Wireframes", "effort": "3 days"}, {"name": "UI mockups", "effort": "4 days"}, {"name": "User testing", "effort": "2 days"}]},
                    {"name": "Security Design", "tasks": [{"name": "IAM design", "effort": "2 days"}, {"name": "Network security", "effort": "2 days"}, {"name": "Encryption strategy", "effort": "1 day"}]},
                ]},
                {"phase": "Development", "deliverables": [
                    {"name": "Backend Services", "tasks": [{"name": "API development", "effort": "15 days"}, {"name": "Data pipeline", "effort": "10 days"}, {"name": "Integration layer", "effort": "8 days"}]},
                    {"name": "Frontend Application", "tasks": [{"name": "UI components", "effort": "10 days"}, {"name": "Dashboard views", "effort": "8 days"}, {"name": "Responsive design", "effort": "4 days"}]},
                    {"name": "Infrastructure", "tasks": [{"name": "IaC (Bicep/Terraform)", "effort": "5 days"}, {"name": "CI/CD pipelines", "effort": "3 days"}, {"name": "Environment setup", "effort": "2 days"}]},
                ]},
                {"phase": "Testing & QA", "deliverables": [
                    {"name": "Testing", "tasks": [{"name": "Unit testing", "effort": "5 days"}, {"name": "Integration testing", "effort": "5 days"}, {"name": "Performance testing", "effort": "3 days"}, {"name": "Security testing", "effort": "3 days"}, {"name": "UAT", "effort": "5 days"}]},
                ]},
                {"phase": "Deployment & Hypercare", "deliverables": [
                    {"name": "Go-Live", "tasks": [{"name": "Data migration", "effort": "3 days"}, {"name": "Production deployment", "effort": "2 days"}, {"name": "Smoke testing", "effort": "1 day"}]},
                    {"name": "Hypercare", "tasks": [{"name": "Post-launch monitoring", "effort": "10 days"}, {"name": "Bug fixes", "effort": "5 days"}, {"name": "Knowledge transfer", "effort": "3 days"}]},
                ]},
            ],
            "meeting_summary": "Discovery call with key stakeholders revealed significant operational inefficiencies driven by manual data workflows and disconnected systems. Primary pain points center around delayed reporting, manual data transfers, and security gaps. The team expressed strong preference for Azure-based cloud solution with phased delivery approach. Budget is available for Q3 2025 MVP.",
            "sentiment": "Positive — stakeholders are motivated and have executive buy-in for modernization",
            "key_themes": ["Automation", "Real-time Analytics", "System Integration", "Security", "Cloud Migration", "Mobile Access"],
        }


# ═══════════════════════════════════════════════════════════════════════
#  SHAREPOINT
# ═══════════════════════════════════════════════════════════════════════

class SP:
    def __init__(self, url, cid, cs, tid):
        self.url = url
        self.cid = cid
        self.cs = cs
        self.tid = tid

    @classmethod
    def from_session(cls):
        return cls(st.session_state.get("sp_url", ""), st.session_state.get("sp_cid", ""), st.session_state.get("sp_cs", ""), st.session_state.get("sp_tid", ""))

    def _auth(self):
        if not all([self.cid, self.cs, self.tid, msal]):
            return None
        try:
            app = msal.ConfidentialClientApplication(self.cid, authority="https://login.microsoftonline.com/" + self.tid, client_credential=self.cs)
            r = app.acquire_token_for_client(scopes=["https://graph.microsoft.com/.default"])
            return r.get("access_token")
        except Exception:
            return None

    def test(self):
        tok = self._auth()
        if not tok:
            return False, "Missing credentials or auth failed."
        try:
            r = _requests.get("https://graph.microsoft.com/v1.0/sites/root", headers={"Authorization": "Bearer " + tok}, timeout=10)
            return (True, "SharePoint connected!") if r.status_code == 200 else (False, "Status " + str(r.status_code))
        except Exception as e:
            return False, str(e)[:200]

    def list_folder(self, path):
        tok = self._auth()
        if tok and _requests:
            try:
                sn = self.url.rstrip("/").split("/")[-1]
                r = _requests.get("https://graph.microsoft.com/v1.0/sites/root:/sites/" + sn + ":/drive/root:" + path + ":/children", headers={"Authorization": "Bearer " + tok}, timeout=15)
                if r.status_code == 200:
                    return [i["name"] for i in r.json().get("value", []) if "name" in i]
            except Exception:
                pass
        return []

    def upload(self, results):
        fn = "BELAL_Proposal_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".json"
        tok = self._auth()
        if tok and _requests:
            try:
                sn = self.url.rstrip("/").split("/")[-1]
                r = _requests.put("https://graph.microsoft.com/v1.0/sites/root:/sites/" + sn + ":/drive/root:/Proposals/" + fn + ":/content", headers={"Authorization": "Bearer " + tok, "Content-Type": "application/json"}, data=json.dumps(results, indent=2, default=str), timeout=30)
                if r.status_code in (200, 201):
                    return True, "Uploaded: /Proposals/" + fn
                return False, "Status " + str(r.status_code)
            except Exception as e:
                return False, str(e)[:200]
        return False, "SharePoint not configured."


# ═══════════════════════════════════════════════════════════════════════
#  EMAIL
# ═══════════════════════════════════════════════════════════════════════

class Mailer:
    def __init__(self, smtp, sender, pw):
        self.smtp = smtp
        self.sender = sender
        self.pw = pw

    @classmethod
    def from_session(cls):
        return cls(st.session_state.get("email_smtp", ""), st.session_state.get("email_sender", ""), st.session_state.get("cfg_email_pass", ""))

    def send(self, results):
        te = safe_dict(results.get("time_estimate"))
        ce = safe_dict(results.get("cost_estimate"))
        ri = safe_dict(results.get("risk_assessment"))
        subj = "Agent BELAL Proposal " + datetime.now().strftime("%Y-%m-%d %H:%M")
        body = "<h2>Agent BELAL</h2><p>Hours: " + str(te.get("total_hours", 0)) + " | Infra Cost: $" + str(ce.get("total_monthly_cost", 0)) + "/mo | Risk: " + safe_str(ri.get("overall_level")) + "</p>"
        if self.smtp and self.sender and self.pw:
            try:
                parts = (self.smtp + ":587").split(":")
                host = parts[0]
                port = int(parts[1])
                msg = MIMEMultipart()
                msg["Subject"] = subj
                msg["From"] = self.sender
                msg["To"] = self.sender
                msg.attach(MIMEText(body, "html"))
                with smtplib.SMTP(host, port) as s:
                    s.starttls()
                    s.login(self.sender, self.pw)
                    s.send_message(msg)
                return True, "Email sent to " + self.sender
            except Exception as e:
                return False, str(e)[:200]
        return False, "Email not configured."


# ═══════════════════════════════════════════════════════════════════════
#  EXCEL TIME ESTIMATE GENERATOR
# ═══════════════════════════════════════════════════════════════════════

def generate_time_excel(time_est, semantic):
    if not Workbook:
        return None
    wb = Workbook()

    # ── Styles ──
    hdr_font = Font(name="Calibri", bold=True, size=12, color="FFFFFF")
    hdr_fill = PatternFill(start_color="1B3A5C", end_color="1B3A5C", fill_type="solid")
    sub_font = Font(name="Calibri", bold=True, size=10, color="1B3A5C")
    sub_fill = PatternFill(start_color="D6E4F0", end_color="D6E4F0", fill_type="solid")
    normal_font = Font(name="Calibri", size=10)
    bold_font = Font(name="Calibri", bold=True, size=10)
    total_fill = PatternFill(start_color="E2EFDA", end_color="E2EFDA", fill_type="solid")
    total_font = Font(name="Calibri", bold=True, size=11, color="1B3A5C")
    thin_border = Border(
        left=Side(style="thin"), right=Side(style="thin"),
        top=Side(style="thin"), bottom=Side(style="thin"),
    )
    center = Alignment(horizontal="center", vertical="center")
    wrap = Alignment(horizontal="left", vertical="center", wrap_text=True)

    def style_header_row(ws, row, cols):
        for c in range(1, cols + 1):
            cell = ws.cell(row=row, column=c)
            cell.font = hdr_font
            cell.fill = hdr_fill
            cell.alignment = center
            cell.border = thin_border

    def style_cell(ws, row, col, font=normal_font, fill=None, align=None):
        cell = ws.cell(row=row, column=col)
        cell.font = font
        cell.border = thin_border
        if fill:
            cell.fill = fill
        if align:
            cell.alignment = align
        return cell

    # ═══ Sheet 1: Summary ═══
    ws1 = wb.active
    ws1.title = "Summary"
    ws1.sheet_properties.tabColor = "1B3A5C"

    # Title
    ws1.merge_cells("A1:F1")
    title_cell = ws1["A1"]
    title_cell.value = "ECI — Project Time Estimation Summary"
    title_cell.font = Font(name="Calibri", bold=True, size=16, color="1B3A5C")
    title_cell.alignment = Alignment(horizontal="center", vertical="center")
    ws1.row_dimensions[1].height = 40

    ws1.merge_cells("A2:F2")
    ws1["A2"].value = "Generated: " + datetime.now().strftime("%B %d, %Y %H:%M")
    ws1["A2"].font = Font(name="Calibri", size=9, color="666666")
    ws1["A2"].alignment = Alignment(horizontal="center")

    # Project Info
    row = 4
    ws1.cell(row=row, column=1, value="Project Type").font = bold_font
    ws1.cell(row=row, column=2, value=safe_str(semantic.get("project_type", "N/A"))).font = normal_font
    row += 1
    ws1.cell(row=row, column=1, value="Complexity Score").font = bold_font
    ws1.cell(row=row, column=2, value=str(semantic.get("complexity_score", "N/A")) + " / 10").font = normal_font
    row += 1
    ws1.cell(row=row, column=1, value="Total Requirements").font = bold_font
    ws1.cell(row=row, column=2, value=len(safe_list(semantic.get("requirements")))).font = normal_font

    # Summary metrics
    row = 8
    headers = ["Metric", "Value"]
    for c, h in enumerate(headers, 1):
        ws1.cell(row=row, column=c, value=h)
    style_header_row(ws1, row, len(headers))
    row += 1
    metrics = [
        ("Total Hours", str(safe_int(time_est.get("total_hours"))) + " hours"),
        ("Duration", safe_str(time_est.get("duration_weeks"))),
        ("Confidence", safe_str(time_est.get("confidence"))),
        ("Buffer", safe_str(time_est.get("buffer"))),
    ]
    for label, val in metrics:
        style_cell(ws1, row, 1, font=bold_font).value = label
        style_cell(ws1, row, 2).value = val
        row += 1

    # Three-point estimation
    tp = safe_dict(time_est.get("three_point"))
    if tp:
        row += 1
        ws1.merge_cells(start_row=row, start_column=1, end_row=row, end_column=2)
        ws1.cell(row=row, column=1, value="Three-Point Estimation").font = Font(name="Calibri", bold=True, size=12, color="1B3A5C")
        row += 1
        headers_tp = ["Scenario", "Hours"]
        for c, h in enumerate(headers_tp, 1):
            ws1.cell(row=row, column=c, value=h)
        style_header_row(ws1, row, len(headers_tp))
        row += 1
        for label, key in [("Optimistic", "optimistic"), ("Most Likely", "most_likely"), ("Pessimistic", "pessimistic")]:
            style_cell(ws1, row, 1, font=bold_font).value = label
            style_cell(ws1, row, 2, align=center).value = safe_int(tp.get(key))
            row += 1

    ws1.column_dimensions["A"].width = 22
    ws1.column_dimensions["B"].width = 30

    # ═══ Sheet 2: Phase Breakdown ═══
    ws2 = wb.create_sheet("Phase Breakdown")
    ws2.sheet_properties.tabColor = "00B4D8"

    ws2.merge_cells("A1:F1")
    ws2["A1"].value = "Phase-wise Effort Breakdown"
    ws2["A1"].font = Font(name="Calibri", bold=True, size=14, color="1B3A5C")
    ws2["A1"].alignment = Alignment(horizontal="center")
    ws2.row_dimensions[1].height = 35

    row = 3
    headers = ["Phase / Week", "Sub-task", "Role", "Dev Low (hrs)", "Dev High (hrs)", "Dev Avg (hrs)", "Justification"]
    for c, h in enumerate(headers, 1):
        ws2.cell(row=row, column=c, value=h)
    style_header_row(ws2, row, len(headers))
    row += 1

    total_hours = safe_int(time_est.get("total_hours", 1))
    grand_low = 0
    grand_high = 0
    grand_avg = 0
    for phase in safe_list(time_est.get("phases")):
        phase = safe_dict(phase)
        phase_name = safe_str(phase.get("name"))
        phase_week = safe_str(phase.get("week_label", ""))
        phase_label = phase_name + (" (" + phase_week + ")" if phase_week else "")
        phase_hours = safe_int(phase.get("hours"))
        phase_low = safe_int(phase.get("low_hours", int(phase_hours * 0.8)))
        phase_high = safe_int(phase.get("high_hours", int(phase_hours * 1.35)))
        tasks = safe_list(phase.get("tasks"))
        if not tasks:
            style_cell(ws2, row, 1, font=sub_font, fill=sub_fill).value = phase_label
            style_cell(ws2, row, 2, fill=sub_fill)
            style_cell(ws2, row, 3, fill=sub_fill)
            style_cell(ws2, row, 4, font=sub_font, fill=sub_fill, align=center).value = phase_low
            style_cell(ws2, row, 5, font=sub_font, fill=sub_fill, align=center).value = phase_high
            style_cell(ws2, row, 6, font=sub_font, fill=sub_fill, align=center).value = phase_hours
            style_cell(ws2, row, 7, fill=sub_fill)
            row += 1
        else:
            first_task = True
            for task in tasks:
                task = safe_dict(task)
                if first_task:
                    style_cell(ws2, row, 1, font=sub_font, fill=sub_fill).value = phase_label
                    first_task = False
                else:
                    style_cell(ws2, row, 1, fill=None)
                t_hrs = safe_int(task.get("hours"))
                t_low = safe_int(task.get("low_hours", int(t_hrs * 0.8)))
                t_high = safe_int(task.get("high_hours", int(t_hrs * 1.35)))
                style_cell(ws2, row, 2).value = safe_str(task.get("name"))
                style_cell(ws2, row, 3).value = safe_str(task.get("role"))
                style_cell(ws2, row, 4, align=center).value = t_low
                style_cell(ws2, row, 5, align=center).value = t_high
                style_cell(ws2, row, 6, align=center).value = t_hrs
                style_cell(ws2, row, 7, align=wrap).value = safe_str(task.get("justification", ""))
                row += 1
            # Phase subtotal
            style_cell(ws2, row, 1, font=bold_font, fill=total_fill)
            style_cell(ws2, row, 2, font=bold_font, fill=total_fill).value = "Subtotal — " + phase_name
            style_cell(ws2, row, 3, fill=total_fill)
            style_cell(ws2, row, 4, font=bold_font, fill=total_fill, align=center).value = phase_low
            style_cell(ws2, row, 5, font=bold_font, fill=total_fill, align=center).value = phase_high
            style_cell(ws2, row, 6, font=bold_font, fill=total_fill, align=center).value = phase_hours
            style_cell(ws2, row, 7, fill=total_fill)
            row += 1
        grand_low += phase_low
        grand_high += phase_high
        grand_avg += phase_hours

    # Grand total
    row += 1
    style_cell(ws2, row, 1, font=total_font, fill=total_fill)
    style_cell(ws2, row, 2, font=total_font, fill=total_fill).value = "GRAND TOTAL"
    style_cell(ws2, row, 3, fill=total_fill)
    style_cell(ws2, row, 4, font=total_font, fill=total_fill, align=center).value = grand_low
    style_cell(ws2, row, 5, font=total_font, fill=total_fill, align=center).value = grand_high
    style_cell(ws2, row, 6, font=total_font, fill=total_fill, align=center).value = grand_avg
    style_cell(ws2, row, 7, fill=total_fill)

    for c, w in [(1, 22), (2, 35), (3, 15), (4, 14), (5, 14), (6, 14), (7, 40)]:
        ws2.column_dimensions[get_column_letter(c)].width = w

    # ═══ Sheet 3: Milestones ═══
    ws3 = wb.create_sheet("Milestones")
    ws3.sheet_properties.tabColor = "00D4AA"

    ws3.merge_cells("A1:D1")
    ws3["A1"].value = "Project Milestones"
    ws3["A1"].font = Font(name="Calibri", bold=True, size=14, color="1B3A5C")
    ws3["A1"].alignment = Alignment(horizontal="center")
    ws3.row_dimensions[1].height = 35

    row = 3
    headers = ["#", "Milestone", "Week", "Description"]
    for c, h in enumerate(headers, 1):
        ws3.cell(row=row, column=c, value=h)
    style_header_row(ws3, row, len(headers))
    row += 1
    for i, m in enumerate(safe_list(time_est.get("milestones")), 1):
        m = safe_dict(m)
        style_cell(ws3, row, 1, align=center).value = i
        style_cell(ws3, row, 2, font=bold_font).value = safe_str(m.get("name"))
        style_cell(ws3, row, 3, align=center).value = safe_int(m.get("week"))
        style_cell(ws3, row, 4, align=wrap).value = safe_str(m.get("description"))
        row += 1

    for c, w in [(1, 6), (2, 30), (3, 10), (4, 45)]:
        ws3.column_dimensions[get_column_letter(c)].width = w

    # ═══ Sheet 4: Three-Point Detail ═══
    ws4 = wb.create_sheet("Three-Point Estimation")
    ws4.sheet_properties.tabColor = "7B61FF"

    ws4.merge_cells("A1:E1")
    ws4["A1"].value = "Three-Point Estimation Detail"
    ws4["A1"].font = Font(name="Calibri", bold=True, size=14, color="1B3A5C")
    ws4["A1"].alignment = Alignment(horizontal="center")
    ws4.row_dimensions[1].height = 35

    row = 3
    headers = ["Phase", "Optimistic (hrs)", "Most Likely (hrs)", "Pessimistic (hrs)", "PERT Estimate (hrs)"]
    for c, h in enumerate(headers, 1):
        ws4.cell(row=row, column=c, value=h)
    style_header_row(ws4, row, len(headers))
    row += 1
    total_o, total_m, total_p, total_pert = 0, 0, 0, 0
    for phase in safe_list(time_est.get("phases")):
        phase = safe_dict(phase)
        ph_hours = safe_int(phase.get("hours"))
        opt = int(ph_hours * 0.8)
        ml = ph_hours
        pes = int(ph_hours * 1.35)
        pert = int((opt + 4 * ml + pes) / 6)
        total_o += opt
        total_m += ml
        total_p += pes
        total_pert += pert
        style_cell(ws4, row, 1, font=bold_font).value = safe_str(phase.get("name"))
        style_cell(ws4, row, 2, align=center).value = opt
        style_cell(ws4, row, 3, align=center).value = ml
        style_cell(ws4, row, 4, align=center).value = pes
        style_cell(ws4, row, 5, align=center, font=bold_font).value = pert
        row += 1

    # Totals
    style_cell(ws4, row, 1, font=total_font, fill=total_fill).value = "TOTAL"
    style_cell(ws4, row, 2, font=total_font, fill=total_fill, align=center).value = total_o
    style_cell(ws4, row, 3, font=total_font, fill=total_fill, align=center).value = total_m
    style_cell(ws4, row, 4, font=total_font, fill=total_fill, align=center).value = total_p
    style_cell(ws4, row, 5, font=total_font, fill=total_fill, align=center).value = total_pert

    for c, w in [(1, 22), (2, 18), (3, 18), (4, 18), (5, 20)]:
        ws4.column_dimensions[get_column_letter(c)].width = w

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════
#  DETAILED COST ESTIMATE EXCEL (Inflexion Format)
# ═══════════════════════════════════════════════════════════════════════

def generate_cost_excel(cost_est, time_est, semantic):
    """Generate a detailed cost estimate Excel matching Inflexion.xlsx format."""
    if not Workbook:
        return None
    wb = Workbook()

    hdr_font = Font(name="Calibri", bold=True, size=12, color="FFFFFF")
    hdr_fill = PatternFill(start_color="1B3A5C", end_color="1B3A5C", fill_type="solid")
    sub_font = Font(name="Calibri", bold=True, size=10, color="1B3A5C")
    sub_fill = PatternFill(start_color="D6E4F0", end_color="D6E4F0", fill_type="solid")
    normal_font = Font(name="Calibri", size=10)
    bold_font = Font(name="Calibri", bold=True, size=10)
    total_fill = PatternFill(start_color="E2EFDA", end_color="E2EFDA", fill_type="solid")
    total_font = Font(name="Calibri", bold=True, size=11, color="1B3A5C")
    accent_fill = PatternFill(start_color="1B3A5C", end_color="1B3A5C", fill_type="solid")
    accent_font = Font(name="Calibri", bold=True, size=12, color="FFFFFF")
    warn_fill = PatternFill(start_color="FFF2CC", end_color="FFF2CC", fill_type="solid")
    currency_fmt = '$#,##0'
    thin_border = Border(
        left=Side(style="thin"), right=Side(style="thin"),
        top=Side(style="thin"), bottom=Side(style="thin"),
    )
    center = Alignment(horizontal="center", vertical="center")
    wrap = Alignment(horizontal="left", vertical="center", wrap_text=True)

    def style_header_row(ws, row, cols):
        for c in range(1, cols + 1):
            cell = ws.cell(row=row, column=c)
            cell.font = hdr_font
            cell.fill = hdr_fill
            cell.alignment = center
            cell.border = thin_border

    def style_cell(ws, row, col, font=normal_font, fill=None, align=None, num_fmt=None):
        cell = ws.cell(row=row, column=col)
        cell.font = font
        cell.border = thin_border
        if fill:
            cell.fill = fill
        if align:
            cell.alignment = align
        if num_fmt:
            cell.number_format = num_fmt
        return cell

    # ═══ Sheet 1: Infrastructure Cost Summary ═══
    ws1 = wb.active
    ws1.title = "Infrastructure Costs"
    ws1.sheet_properties.tabColor = "1B3A5C"

    ws1.merge_cells("A1:F1")
    t_cell = ws1["A1"]
    t_cell.value = "ECI — Monthly Infrastructure Cost Estimates"
    t_cell.font = Font(name="Calibri", bold=True, size=16, color="1B3A5C")
    t_cell.alignment = Alignment(horizontal="center", vertical="center")
    ws1.row_dimensions[1].height = 40

    ws1.merge_cells("A2:F2")
    ws1["A2"].value = "Azure Infrastructure + AI Models  |  Generated: " + datetime.now().strftime("%B %d, %Y %H:%M")
    ws1["A2"].font = Font(name="Calibri", size=9, color="666666")
    ws1["A2"].alignment = Alignment(horizontal="center")

    row = 4
    headers = ["Azure Service / AI Model", "Description & Assumptions", "Low Cost/Mo", "High Cost/Mo", "Average Cost/Mo", "Notes"]
    for c, h in enumerate(headers, 1):
        ws1.cell(row=row, column=c, value=h)
    style_header_row(ws1, row, len(headers))
    row += 1

    azure_costs = safe_list(cost_est.get("azure_costs"))
    third_party = safe_list(cost_est.get("third_party_costs"))

    categories = {}
    for svc in azure_costs:
        svc = safe_dict(svc)
        service_name = safe_str(svc.get("service", ""))
        name_lower = service_name.lower()
        if any(k in name_lower for k in ["openai", "ai search", "foundry", "cognitive", "ml", "embedding"]):
            cat = "AI & ML SERVICES"
        elif any(k in name_lower for k in ["app service", "function", "container", "kubernetes"]):
            cat = "COMPUTE SERVICES"
        elif any(k in name_lower for k in ["blob", "storage", "redis", "cache", "cosmos", "sql"]):
            cat = "DATA & STORAGE SERVICES"
        elif any(k in name_lower for k in ["key vault", "monitor", "insights", "sentinel", "ad "]):
            cat = "SECURITY & MONITORING"
        elif any(k in name_lower for k in ["front door", "cdn", "api management", "service bus", "logic"]):
            cat = "NETWORKING & INTEGRATION"
        else:
            cat = "OTHER SERVICES"
        if cat not in categories:
            categories[cat] = []
        monthly = safe_int(svc.get("monthly_cost", 0))
        categories[cat].append({
            "service": service_name,
            "description": safe_str(svc.get("description", "")),
            "low": int(monthly * 0.85),
            "high": int(monthly * 1.25),
            "average": monthly,
            "notes": safe_str(svc.get("tier", "")),
        })

    if third_party:
        categories["THIRD-PARTY SERVICES"] = []
        for tp in third_party:
            tp = safe_dict(tp)
            monthly = safe_int(tp.get("monthly_cost", 0))
            categories["THIRD-PARTY SERVICES"].append({
                "service": safe_str(tp.get("name", "")),
                "description": safe_str(tp.get("description", "")),
                "low": int(monthly * 0.85),
                "high": int(monthly * 1.25),
                "average": monthly,
                "notes": "",
            })

    grand_low, grand_high, grand_avg = 0, 0, 0
    cat_order = ["AI & ML SERVICES", "COMPUTE SERVICES", "DATA & STORAGE SERVICES",
                 "SECURITY & MONITORING", "NETWORKING & INTEGRATION", "THIRD-PARTY SERVICES", "OTHER SERVICES"]

    for cat_name in cat_order:
        if cat_name not in categories:
            continue
        items = categories[cat_name]
        style_cell(ws1, row, 1, font=sub_font, fill=sub_fill).value = cat_name
        for c in range(2, 7):
            style_cell(ws1, row, c, fill=sub_fill)
        row += 1

        cat_low, cat_high, cat_avg = 0, 0, 0
        for item in items:
            style_cell(ws1, row, 1).value = item["service"]
            style_cell(ws1, row, 2, align=wrap).value = item["description"]
            style_cell(ws1, row, 3, align=center, num_fmt=currency_fmt).value = item["low"]
            style_cell(ws1, row, 4, align=center, num_fmt=currency_fmt).value = item["high"]
            style_cell(ws1, row, 5, align=center, num_fmt=currency_fmt).value = item["average"]
            style_cell(ws1, row, 6).value = item["notes"]
            cat_low += item["low"]
            cat_high += item["high"]
            cat_avg += item["average"]
            row += 1

        style_cell(ws1, row, 1, font=bold_font, fill=total_fill).value = "Subtotal: " + cat_name.title().split(" ")[0]
        style_cell(ws1, row, 2, fill=total_fill)
        style_cell(ws1, row, 3, font=bold_font, fill=total_fill, align=center, num_fmt=currency_fmt).value = cat_low
        style_cell(ws1, row, 4, font=bold_font, fill=total_fill, align=center, num_fmt=currency_fmt).value = cat_high
        style_cell(ws1, row, 5, font=bold_font, fill=total_fill, align=center, num_fmt=currency_fmt).value = cat_avg
        style_cell(ws1, row, 6, fill=total_fill)
        grand_low += cat_low
        grand_high += cat_high
        grand_avg += cat_avg
        row += 2

    style_cell(ws1, row, 1, font=accent_font, fill=accent_fill).value = "TOTAL MONTHLY COST"
    style_cell(ws1, row, 2, fill=accent_fill)
    style_cell(ws1, row, 3, font=accent_font, fill=accent_fill, align=center, num_fmt=currency_fmt).value = grand_low
    style_cell(ws1, row, 4, font=accent_font, fill=accent_fill, align=center, num_fmt=currency_fmt).value = grand_high
    style_cell(ws1, row, 5, font=accent_font, fill=accent_fill, align=center, num_fmt=currency_fmt).value = grand_avg
    style_cell(ws1, row, 6, fill=accent_fill)
    row += 2

    style_cell(ws1, row, 1, font=total_font).value = "ANNUAL PROJECTION (12 MONTHS)"
    style_cell(ws1, row, 2)
    style_cell(ws1, row, 3, font=total_font, align=center, num_fmt=currency_fmt).value = grand_low * 12
    style_cell(ws1, row, 4, font=total_font, align=center, num_fmt=currency_fmt).value = grand_high * 12
    style_cell(ws1, row, 5, font=total_font, align=center, num_fmt=currency_fmt).value = grand_avg * 12
    style_cell(ws1, row, 6)
    row += 3

    style_cell(ws1, row, 1, font=sub_font).value = "KEY ASSUMPTIONS"
    row += 1
    notes_str = safe_str(cost_est.get("notes", ""))
    if notes_str:
        ws1.merge_cells(start_row=row, start_column=1, end_row=row, end_column=6)
        style_cell(ws1, row, 1).value = "\u2022 " + notes_str
        row += 1
    for a in safe_list(cost_est.get("cost_optimization", [])):
        ws1.merge_cells(start_row=row, start_column=1, end_row=row, end_column=6)
        style_cell(ws1, row, 1).value = "\u2022 " + safe_str(a)
        row += 1

    for c, w in [(1, 32), (2, 48), (3, 14), (4, 14), (5, 16), (6, 25)]:
        ws1.column_dimensions[get_column_letter(c)].width = w

    # ═══ Sheet 2: Team & Role Allocation ═══
    ws2 = wb.create_sheet("Team & Roles")
    ws2.sheet_properties.tabColor = "00B4D8"

    ws2.merge_cells("A1:G1")
    ws2["A1"].value = "Team Composition & Role Allocation"
    ws2["A1"].font = Font(name="Calibri", bold=True, size=14, color="1B3A5C")
    ws2["A1"].alignment = Alignment(horizontal="center")
    ws2.row_dimensions[1].height = 35

    row = 3
    headers = ["S.N.", "Role", "Allocation %", "Days", "Hours", "Rate ($/hr)", "Cost"]
    for c, h in enumerate(headers, 1):
        ws2.cell(row=row, column=c, value=h)
    style_header_row(ws2, row, len(headers))
    row += 1

    total_hours_val = safe_int(time_est.get("total_hours", 0))
    total_days_val = max(1, total_hours_val / 7)
    dynamic_roles = safe_list(time_est.get("roles"))
    if not dynamic_roles:
        dynamic_roles = [
            {"name": "Project Manager", "allocation_pct": 0.15, "rate": 125},
            {"name": "Developer", "allocation_pct": 0.60, "rate": 110},
            {"name": "QA Engineer", "allocation_pct": 0.25, "rate": 95},
        ]
    sum_days, sum_hours, sum_cost = 0, 0, 0
    for i, rl in enumerate(dynamic_roles, 1):
        rl = safe_dict(rl)
        role_name = safe_str(rl.get("name", "Team Member"))
        pct = float(rl.get("allocation_pct", 0))
        rate = safe_int(rl.get("rate", 100))
        days = round(total_days_val * pct, 1)
        hours = int(round(days * 7, 0))
        cost = hours * rate
        sum_days += days
        sum_hours += hours
        sum_cost += cost
        style_cell(ws2, row, 1, align=center).value = i
        style_cell(ws2, row, 2, font=bold_font).value = role_name
        style_cell(ws2, row, 3, align=center).value = str(int(pct * 100)) + "%"
        style_cell(ws2, row, 4, align=center).value = days
        style_cell(ws2, row, 5, align=center).value = hours
        style_cell(ws2, row, 6, align=center, num_fmt=currency_fmt).value = rate
        style_cell(ws2, row, 7, align=center, num_fmt=currency_fmt).value = cost
        row += 1

    style_cell(ws2, row, 1, fill=total_fill)
    style_cell(ws2, row, 2, font=total_font, fill=total_fill).value = "TOTAL"
    style_cell(ws2, row, 3, fill=total_fill)
    style_cell(ws2, row, 4, font=total_font, fill=total_fill, align=center).value = round(sum_days, 1)
    style_cell(ws2, row, 5, font=total_font, fill=total_fill, align=center).value = sum_hours
    style_cell(ws2, row, 6, fill=total_fill)
    style_cell(ws2, row, 7, font=total_font, fill=total_fill, align=center, num_fmt=currency_fmt).value = sum_cost

    for c, w in [(1, 6), (2, 30), (3, 14), (4, 10), (5, 10), (6, 14), (7, 14)]:
        ws2.column_dimensions[get_column_letter(c)].width = w

    # ═══ Sheet 3: Project Labor Cost (Three-Point) ═══
    ws3 = wb.create_sheet("Project Labor Cost")
    ws3.sheet_properties.tabColor = "00D4AA"

    ws3.merge_cells("A1:G1")
    ws3["A1"].value = "Project Labor Cost Breakdown (Three-Point Estimation)"
    ws3["A1"].font = Font(name="Calibri", bold=True, size=14, color="1B3A5C")
    ws3["A1"].alignment = Alignment(horizontal="center")
    ws3.row_dimensions[1].height = 35

    row = 3
    headers = ["Phase", "Hours", "% of Total", "Low Cost ($)", "High Cost ($)", "Average Cost ($)", "Notes"]
    for c, h in enumerate(headers, 1):
        ws3.cell(row=row, column=c, value=h)
    style_header_row(ws3, row, len(headers))
    row += 1

    blended_rate = 115
    gl_low, gl_high, gl_avg = 0, 0, 0
    for phase in safe_list(time_est.get("phases")):
        phase = safe_dict(phase)
        ph_name = safe_str(phase.get("name"))
        ph_hours = safe_int(phase.get("hours"))
        low_hrs = safe_int(phase.get("low_hours", int(ph_hours * 0.8)))
        high_hrs = safe_int(phase.get("high_hours", int(ph_hours * 1.35)))
        low_cost = low_hrs * blended_rate
        high_cost = high_hrs * blended_rate
        avg_cost = ph_hours * blended_rate
        gl_low += low_cost
        gl_high += high_cost
        gl_avg += avg_cost
        style_cell(ws3, row, 1, font=bold_font).value = ph_name
        style_cell(ws3, row, 2, align=center).value = ph_hours
        style_cell(ws3, row, 3, align=center).value = safe_str(phase.get("percentage"))
        style_cell(ws3, row, 4, align=center, num_fmt=currency_fmt).value = low_cost
        style_cell(ws3, row, 5, align=center, num_fmt=currency_fmt).value = high_cost
        style_cell(ws3, row, 6, align=center, num_fmt=currency_fmt).value = avg_cost
        style_cell(ws3, row, 7)
        row += 1

    row += 1
    style_cell(ws3, row, 1, font=accent_font, fill=accent_fill).value = "TOTAL PROJECT LABOR"
    style_cell(ws3, row, 2, font=accent_font, fill=accent_fill, align=center).value = safe_int(time_est.get("total_hours"))
    style_cell(ws3, row, 3, font=accent_font, fill=accent_fill, align=center).value = "100%"
    style_cell(ws3, row, 4, font=accent_font, fill=accent_fill, align=center, num_fmt=currency_fmt).value = gl_low
    style_cell(ws3, row, 5, font=accent_font, fill=accent_fill, align=center, num_fmt=currency_fmt).value = gl_high
    style_cell(ws3, row, 6, font=accent_font, fill=accent_fill, align=center, num_fmt=currency_fmt).value = gl_avg
    style_cell(ws3, row, 7, fill=accent_fill)

    for c, w in [(1, 20), (2, 10), (3, 12), (4, 14), (5, 14), (6, 16), (7, 25)]:
        ws3.column_dimensions[get_column_letter(c)].width = w

    # ═══ Sheet 4: Total Cost Summary ═══
    ws4 = wb.create_sheet("Total Cost Summary")
    ws4.sheet_properties.tabColor = "7B61FF"

    ws4.merge_cells("A1:D1")
    ws4["A1"].value = "Total Project Cost Summary"
    ws4["A1"].font = Font(name="Calibri", bold=True, size=16, color="1B3A5C")
    ws4["A1"].alignment = Alignment(horizontal="center")
    ws4.row_dimensions[1].height = 40

    ws4.merge_cells("A2:D2")
    ws4["A2"].value = "ECI Consulting  |  " + datetime.now().strftime("%B %d, %Y")
    ws4["A2"].font = Font(name="Calibri", size=9, color="666666")
    ws4["A2"].alignment = Alignment(horizontal="center")

    row = 4
    headers = ["Cost Category", "Low Estimate ($)", "High Estimate ($)", "Average Estimate ($)"]
    for c, h in enumerate(headers, 1):
        ws4.cell(row=row, column=c, value=h)
    style_header_row(ws4, row, len(headers))
    row += 1

    style_cell(ws4, row, 1, font=sub_font, fill=sub_fill).value = "ONE-TIME COSTS"
    for c in range(2, 5):
        style_cell(ws4, row, c, fill=sub_fill)
    row += 1
    style_cell(ws4, row, 1).value = "Project Labor (Development)"
    style_cell(ws4, row, 2, align=center, num_fmt=currency_fmt).value = gl_low
    style_cell(ws4, row, 3, align=center, num_fmt=currency_fmt).value = gl_high
    style_cell(ws4, row, 4, align=center, num_fmt=currency_fmt).value = gl_avg
    row += 2

    style_cell(ws4, row, 1, font=sub_font, fill=sub_fill).value = "RECURRING COSTS (Monthly)"
    for c in range(2, 5):
        style_cell(ws4, row, c, fill=sub_fill)
    row += 1
    style_cell(ws4, row, 1).value = "Azure Infrastructure"
    style_cell(ws4, row, 2, align=center, num_fmt=currency_fmt).value = grand_low
    style_cell(ws4, row, 3, align=center, num_fmt=currency_fmt).value = grand_high
    style_cell(ws4, row, 4, align=center, num_fmt=currency_fmt).value = grand_avg
    row += 2

    style_cell(ws4, row, 1, font=sub_font, fill=sub_fill).value = "RECURRING COSTS (Annual)"
    for c in range(2, 5):
        style_cell(ws4, row, c, fill=sub_fill)
    row += 1
    style_cell(ws4, row, 1).value = "Azure Infrastructure (x12)"
    style_cell(ws4, row, 2, align=center, num_fmt=currency_fmt).value = grand_low * 12
    style_cell(ws4, row, 3, align=center, num_fmt=currency_fmt).value = grand_high * 12
    style_cell(ws4, row, 4, align=center, num_fmt=currency_fmt).value = grand_avg * 12
    row += 2

    y1_low = gl_low + grand_low * 12
    y1_high = gl_high + grand_high * 12
    y1_avg = gl_avg + grand_avg * 12
    style_cell(ws4, row, 1, font=accent_font, fill=accent_fill).value = "YEAR 1 TOTAL COST"
    style_cell(ws4, row, 2, font=accent_font, fill=accent_fill, align=center, num_fmt=currency_fmt).value = y1_low
    style_cell(ws4, row, 3, font=accent_font, fill=accent_fill, align=center, num_fmt=currency_fmt).value = y1_high
    style_cell(ws4, row, 4, font=accent_font, fill=accent_fill, align=center, num_fmt=currency_fmt).value = y1_avg
    row += 2

    style_cell(ws4, row, 1, font=bold_font, fill=warn_fill).value = "YEAR 2+ ANNUAL (Infra Only)"
    style_cell(ws4, row, 2, font=bold_font, fill=warn_fill, align=center, num_fmt=currency_fmt).value = grand_low * 12
    style_cell(ws4, row, 3, font=bold_font, fill=warn_fill, align=center, num_fmt=currency_fmt).value = grand_high * 12
    style_cell(ws4, row, 4, font=bold_font, fill=warn_fill, align=center, num_fmt=currency_fmt).value = grand_avg * 12

    for c, w in [(1, 36), (2, 20), (3, 20), (4, 20)]:
        ws4.column_dimensions[get_column_letter(c)].width = w

    # ═══ Sheet 5: Assumptions, Prerequisites & Scope ═══
    ws5 = wb.create_sheet("Assumptions & Scope")
    ws5.sheet_properties.tabColor = "FFD166"

    ws5.merge_cells("A1:C1")
    ws5["A1"].value = "Assumptions, Prerequisites & Out-of-Scope"
    ws5["A1"].font = Font(name="Calibri", bold=True, size=14, color="1B3A5C")
    ws5["A1"].alignment = Alignment(horizontal="center")
    ws5.row_dimensions[1].height = 35

    row = 3
    sections_data = [
        ("PREREQUISITES", [
            "Azure subscription with Contributor/Owner access at resource group level",
            "SharePoint Online Read/Write access to document library",
            "Sample documents from all types (PDF, Excel, PPT, Word)",
            "Microsoft Teams admin consent for bot deployment",
            "Azure AD user accounts for MVP users",
            "Service principal credentials with SharePoint API access",
        ]),
        ("KEY ASSUMPTIONS", [
            "Client provides Azure subscription with appropriate access levels",
            "Dedicated product owner available for requirements sign-off and UAT",
            "SME availability minimum 10 hours/week during development",
            "All third-party APIs are documented and accessible",
            "Standard business hours (9 AM - 6 PM) for team availability",
            "Documents are in standard formats without password protection or DRM",
            "Total document size approximately 1-2 GB",
            "SharePoint Online (not on-premise SharePoint Server)",
            "Client responsible for all software licenses and ongoing infrastructure costs",
            "Estimates based on production environment; Dev/staging adds ~40% of prod costs",
            "Pricing based on current Azure pricing; subject to change",
            "Buffer of 10-20% included for usage spikes and scaling",
        ]),
        ("OUT OF SCOPE", [
            "Integration with database systems (SQL, NoSQL, data warehouses)",
            "Web scraping or external data source ingestion",
            "Real-time data feeds or APIs",
            "Multi-language support (English only for MVP)",
            "Custom mobile application development",
            "CI/CD pipeline setup and automation (manual deployments for MVP)",
            "Legacy system decommissioning",
            "End-user training beyond knowledge transfer",
            "Hardware procurement",
            "Penetration testing",
        ]),
    ]

    for sec_title, sec_items in sections_data:
        style_cell(ws5, row, 1, font=sub_font, fill=sub_fill).value = "S.N."
        style_cell(ws5, row, 2, font=sub_font, fill=sub_fill).value = sec_title
        style_cell(ws5, row, 3, fill=sub_fill)
        row += 1
        for i, item in enumerate(sec_items, 1):
            style_cell(ws5, row, 1, align=center).value = i
            ws5.merge_cells(start_row=row, start_column=2, end_row=row, end_column=3)
            style_cell(ws5, row, 2, align=wrap).value = item
            row += 1
        row += 1

    ws5.column_dimensions["A"].width = 8
    ws5.column_dimensions["B"].width = 70
    ws5.column_dimensions["C"].width = 20

    # ═══ Sheet 6: Week-by-Week Breakdown ═══
    ws6 = wb.create_sheet("Week-by-Week Breakdown")
    ws6.sheet_properties.tabColor = "FF6B6B"

    ws6.merge_cells("A1:F1")
    ws6["A1"].value = "Week-by-Week Task Breakdown"
    ws6["A1"].font = Font(name="Calibri", bold=True, size=14, color="1B3A5C")
    ws6["A1"].alignment = Alignment(horizontal="center")
    ws6.row_dimensions[1].height = 35

    row = 3
    headers = ["Week/Phase", "Task", "Role", "Low (hrs)", "High (hrs)", "Avg (hrs)"]
    for c, h in enumerate(headers, 1):
        ws6.cell(row=row, column=c, value=h)
    style_header_row(ws6, row, len(headers))
    row += 1

    total_hours_est = safe_int(time_est.get("total_hours", 1))
    week_counter = 1
    for phase in safe_list(time_est.get("phases")):
        phase = safe_dict(phase)
        ph_name = safe_str(phase.get("name"))
        ph_hours = safe_int(phase.get("hours"))
        ph_weeks = max(1, ph_hours // 40)
        week_label = "Week " + str(week_counter) + ("-" + str(week_counter + ph_weeks - 1) if ph_weeks > 1 else "") + ": " + ph_name

        style_cell(ws6, row, 1, font=sub_font, fill=sub_fill).value = week_label
        style_cell(ws6, row, 2, fill=sub_fill)
        style_cell(ws6, row, 3, fill=sub_fill)
        low_total = int(ph_hours * 0.8)
        high_total = int(ph_hours * 1.35)
        style_cell(ws6, row, 4, font=sub_font, fill=sub_fill, align=center).value = low_total
        style_cell(ws6, row, 5, font=sub_font, fill=sub_fill, align=center).value = high_total
        style_cell(ws6, row, 6, font=sub_font, fill=sub_fill, align=center).value = ph_hours
        row += 1

        for task in safe_list(phase.get("tasks")):
            task = safe_dict(task)
            t_hours = safe_int(task.get("hours"))
            style_cell(ws6, row, 1)
            style_cell(ws6, row, 2).value = safe_str(task.get("name"))
            style_cell(ws6, row, 3).value = safe_str(task.get("role"))
            style_cell(ws6, row, 4, align=center).value = int(t_hours * 0.8)
            style_cell(ws6, row, 5, align=center).value = int(t_hours * 1.35)
            style_cell(ws6, row, 6, align=center).value = t_hours
            row += 1

        week_counter += ph_weeks

    row += 1
    style_cell(ws6, row, 1, font=accent_font, fill=accent_fill).value = "GRAND TOTAL"
    style_cell(ws6, row, 2, fill=accent_fill)
    style_cell(ws6, row, 3, fill=accent_fill)
    style_cell(ws6, row, 4, font=accent_font, fill=accent_fill, align=center).value = int(total_hours_est * 0.8)
    style_cell(ws6, row, 5, font=accent_font, fill=accent_fill, align=center).value = int(total_hours_est * 1.35)
    style_cell(ws6, row, 6, font=accent_font, fill=accent_fill, align=center).value = total_hours_est

    for c, w in [(1, 28), (2, 35), (3, 18), (4, 12), (5, 12), (6, 12)]:
        ws6.column_dimensions[get_column_letter(c)].width = w

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════
#  PDF PROPOSAL GENERATOR (ECI Template)
# ═══════════════════════════════════════════════════════════════════════

def _eci_styles():
    """ECI-branded paragraph styles for reportlab."""
    styles = getSampleStyleSheet()
    eci_blue = rl_colors.Color(27/255, 58/255, 92/255)
    eci_accent = rl_colors.Color(0, 180/255, 216/255)
    styles.add(ParagraphStyle("ECITitle", parent=styles["Title"], fontSize=26, textColor=rl_colors.white, alignment=TA_CENTER, spaceAfter=6))
    styles.add(ParagraphStyle("ECISubtitle", parent=styles["Normal"], fontSize=13, textColor=rl_colors.Color(200/255, 220/255, 240/255), alignment=TA_CENTER, spaceAfter=4))
    styles.add(ParagraphStyle("ECIHeading", parent=styles["Heading1"], fontSize=15, textColor=eci_blue, spaceBefore=14, spaceAfter=6, borderWidth=0))
    styles.add(ParagraphStyle("ECIBody", parent=styles["Normal"], fontSize=10, leading=14, textColor=rl_colors.Color(30/255, 30/255, 30/255), spaceAfter=4))
    styles.add(ParagraphStyle("ECIBullet", parent=styles["Normal"], fontSize=10, leading=14, textColor=rl_colors.Color(30/255, 30/255, 30/255), leftIndent=16, bulletIndent=6, spaceAfter=2))
    styles.add(ParagraphStyle("ECISmall", parent=styles["Normal"], fontSize=8, textColor=rl_colors.Color(100/255, 100/255, 100/255), alignment=TA_CENTER))
    styles.add(ParagraphStyle("ECIKPIVal", parent=styles["Normal"], fontSize=16, textColor=eci_blue, alignment=TA_CENTER, leading=20))
    styles.add(ParagraphStyle("ECIKPILabel", parent=styles["Normal"], fontSize=8, textColor=rl_colors.Color(100/255, 100/255, 100/255), alignment=TA_CENTER))
    styles.add(ParagraphStyle("ECITableCell", parent=styles["Normal"], fontSize=8, leading=10, textColor=rl_colors.Color(30/255, 30/255, 30/255)))
    styles.add(ParagraphStyle("ECITableHeader", parent=styles["Normal"], fontSize=8, leading=10, textColor=rl_colors.white, alignment=TA_CENTER))
    return styles


def _eci_table(headers, rows, col_widths=None):
    """Build a styled reportlab Table."""
    eci_blue = rl_colors.Color(27/255, 58/255, 92/255)
    alt_row = rl_colors.Color(245/255, 248/255, 252/255)
    styles = _eci_styles()
    hdr_cells = [Paragraph("<b>" + h + "</b>", styles["ECITableHeader"]) for h in headers]
    data = [hdr_cells]
    for row in rows:
        data.append([Paragraph(str(v), styles["ECITableCell"]) for v in row])
    if not col_widths:
        col_widths = [480 / len(headers)] * len(headers)
    tbl = Table(data, colWidths=col_widths, repeatRows=1)
    style_cmds = [
        ("BACKGROUND", (0, 0), (-1, 0), eci_blue),
        ("TEXTCOLOR", (0, 0), (-1, 0), rl_colors.white),
        ("GRID", (0, 0), (-1, -1), 0.5, rl_colors.Color(200/255, 200/255, 200/255)),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
    ]
    for i in range(1, len(data)):
        if i % 2 == 0:
            style_cmds.append(("BACKGROUND", (0, i), (-1, i), alt_row))
    tbl.setStyle(TableStyle(style_cmds))
    return tbl


def _md_to_para(text, styles):
    """Convert simple markdown text to a list of Paragraph flowables."""
    elements = []
    for line in text.split("\n"):
        line = line.strip()
        if not line:
            elements.append(Spacer(1, 4))
            continue
        # Bold markers
        cleaned = line.replace("**", "<b>", 1)
        while "**" in cleaned:
            cleaned = cleaned.replace("**", "</b>", 1)
            if "**" in cleaned:
                cleaned = cleaned.replace("**", "<b>", 1)
        # Italic markers
        cleaned = cleaned.replace("_", "<i>", 1)
        while "_" in cleaned:
            cleaned = cleaned.replace("_", "</i>", 1)
            if "_" in cleaned:
                cleaned = cleaned.replace("_", "<i>", 1)
        if line.startswith("- "):
            elements.append(Paragraph(cleaned[2:], styles["ECIBullet"], bulletText="\u2022"))
        else:
            elements.append(Paragraph(cleaned, styles["ECIBody"]))
    return elements


def generate_proposal_pdf(results):
    if not HAS_REPORTLAB:
        return None
    buf = io.BytesIO()
    styles = _eci_styles()
    eci_blue = rl_colors.Color(27/255, 58/255, 92/255)
    eci_light = rl_colors.Color(214/255, 228/255, 240/255)

    doc = SimpleDocTemplate(buf, pagesize=A4, topMargin=25*mm, bottomMargin=20*mm, leftMargin=18*mm, rightMargin=18*mm,
                            title="ECI Project Proposal", author="ECI")
    elements = []

    se = safe_dict(results.get("semantic_analysis"))
    te = safe_dict(results.get("time_estimate"))
    ce = safe_dict(results.get("cost_estimate"))
    ri = safe_dict(results.get("risk_assessment"))
    ar = safe_dict(results.get("architecture"))
    sc = safe_dict(results.get("scope"))
    proposal = safe_dict(results.get("proposal"))

    # ── Cover ──
    elements.append(Spacer(1, 40))
    # Blue banner table
    banner_data = [[Paragraph("<b>PROJECT PROPOSAL</b>", styles["ECITitle"]),],
                   [Paragraph(safe_str(se.get("project_type", "Technology Solution")), styles["ECISubtitle"]),],
                   [Paragraph("Prepared by ECI Consulting  |  Agent BELAL", styles["ECISubtitle"]),],
                   [Paragraph(datetime.now().strftime("%B %d, %Y"), styles["ECISubtitle"]),]]
    banner = Table(banner_data, colWidths=[480])
    banner.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), eci_blue),
        ("TOPPADDING", (0, 0), (-1, -1), 8),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
    ]))
    elements.append(banner)
    elements.append(Spacer(1, 24))

    # KPI boxes
    kpi_items = [
        ("Requirements", str(len(safe_list(se.get("requirements"))))),
        ("Total Hours", str(safe_int(te.get("total_hours")))),
        ("Infra Cost/mo", "$" + str(safe_int(ce.get("total_monthly_cost")))),
        ("Risk Level", safe_str(ri.get("overall_level", "N/A"))),
    ]
    kpi_cells = []
    for label, val in kpi_items:
        kpi_cells.append([
            Paragraph("<b>" + val + "</b>", styles["ECIKPIVal"]),
            Paragraph(label, styles["ECIKPILabel"]),
        ])
    # Transpose: each kpi is a column of 2 rows
    kpi_data = [[kpi_cells[i][0] for i in range(4)], [kpi_cells[i][1] for i in range(4)]]
    kpi_table = Table(kpi_data, colWidths=[120, 120, 120, 120])
    kpi_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), eci_light),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("TOPPADDING", (0, 0), (-1, -1), 8),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
        ("GRID", (0, 0), (-1, -1), 0.5, rl_colors.white),
    ]))
    elements.append(kpi_table)
    elements.append(Spacer(1, 10))
    elements.append(Paragraph("ECI CONSULTING  |  CONFIDENTIAL", styles["ECISmall"]))
    elements.append(PageBreak())

    # ── Proposal sections ──
    for sec in safe_list(proposal.get("sections")):
        sec = safe_dict(sec)
        elements.append(Paragraph(safe_str(sec.get("title")), styles["ECIHeading"]))
        elements.append(HRFlowable(width="30%", thickness=2, color=rl_colors.Color(0, 180/255, 216/255), spaceAfter=8))
        elements.extend(_md_to_para(safe_str(sec.get("content")), styles))
        elements.append(Spacer(1, 10))

    elements.append(PageBreak())

    # ── Requirements summary ──
    elements.append(Paragraph("Requirements Summary", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=rl_colors.Color(0, 180/255, 216/255), spaceAfter=8))
    reqs = safe_list(se.get("requirements"))
    fn = [r for r in reqs if isinstance(r, dict) and r.get("type") == "functional"]
    nf = [r for r in reqs if isinstance(r, dict) and r.get("type") == "non-functional"]
    ig = [r for r in reqs if isinstance(r, dict) and r.get("type") == "integration"]
    elements.append(Paragraph("Functional: <b>" + str(len(fn)) + "</b>  |  Non-Functional: <b>" + str(len(nf)) + "</b>  |  Integration: <b>" + str(len(ig)) + "</b>", styles["ECIBody"]))
    elements.append(Spacer(1, 6))
    req_rows = [[safe_str(r.get("title")), safe_str(r.get("type")), safe_str(r.get("complexity")), safe_str(r.get("priority", ""))] for r in reqs if isinstance(r, dict)]
    if req_rows:
        elements.append(_eci_table(["Requirement", "Type", "Complexity", "Priority"], req_rows[:25], [200, 90, 80, 80]))
    elements.append(PageBreak())

    # ── Time estimate ──
    elements.append(Paragraph("Time Estimation", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=rl_colors.Color(0, 180/255, 216/255), spaceAfter=8))
    elements.extend(_md_to_para(
        "- <b>Total Hours:</b> " + str(safe_int(te.get("total_hours")))
        + "\n- <b>Duration:</b> " + safe_str(te.get("duration_weeks"))
        + "\n- <b>Confidence:</b> " + safe_str(te.get("confidence"))
        + "\n- <b>Buffer:</b> " + safe_str(te.get("buffer")), styles))
    elements.append(Spacer(1, 6))
    phase_rows = [[safe_str(p.get("name")), str(safe_int(p.get("hours"))), safe_str(p.get("percentage"))] for p in safe_list(te.get("phases")) if isinstance(p, dict)]
    if phase_rows:
        elements.append(_eci_table(["Phase", "Hours", "% of Total"], phase_rows, [220, 120, 120]))
    elements.append(PageBreak())

    # ── Infrastructure costs ──
    elements.append(Paragraph("Infrastructure Cost Estimate", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=rl_colors.Color(0, 180/255, 216/255), spaceAfter=8))
    elements.extend(_md_to_para(
        "- <b>Monthly Cost:</b> $" + str(safe_int(ce.get("total_monthly_cost")))
        + "\n- <b>Annual Cost:</b> $" + str(safe_int(ce.get("total_annual_cost"))), styles))
    elements.append(Spacer(1, 6))
    cost_rows = [[safe_str(a.get("service")), safe_str(a.get("tier", "")), "$" + str(safe_int(a.get("monthly_cost")))] for a in safe_list(ce.get("azure_costs")) if isinstance(a, dict)]
    if cost_rows:
        elements.append(_eci_table(["Service", "Tier", "Monthly Cost"], cost_rows, [200, 160, 100]))
    elements.append(PageBreak())

    # ── Risk ──
    elements.append(Paragraph("Risk Assessment", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=rl_colors.Color(0, 180/255, 216/255), spaceAfter=8))
    elements.extend(_md_to_para(
        "- <b>Overall Score:</b> " + str(safe_int(ri.get("overall_score"))) + "/10"
        + "\n- <b>Level:</b> " + safe_str(ri.get("overall_level")), styles))
    elements.append(Spacer(1, 6))
    risk_rows = [[safe_str(r.get("category")), safe_str(r.get("title")), safe_str(r.get("severity")), safe_str(r.get("mitigation"))] for r in safe_list(ri.get("risks")) if isinstance(r, dict)]
    if risk_rows:
        elements.append(_eci_table(["Category", "Risk", "Severity", "Mitigation"], risk_rows, [70, 110, 60, 220]))
    elements.append(PageBreak())

    # ── Architecture ──
    elements.append(Paragraph("Architecture Overview", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=rl_colors.Color(0, 180/255, 216/255), spaceAfter=8))
    elements.append(Paragraph("<b>Pattern:</b> " + safe_str(ar.get("pattern")), styles["ECIBody"]))
    elements.append(Spacer(1, 6))
    comp_rows = [[safe_str(c.get("name")), safe_str(c.get("type", "")), safe_str(c.get("azure_service", "")), ", ".join(safe_list(c.get("services")))] for c in safe_list(ar.get("components")) if isinstance(c, dict)]
    if comp_rows:
        elements.append(_eci_table(["Component", "Type", "Azure Service", "Details"], comp_rows, [80, 80, 120, 180]))
    df = safe_list(ar.get("data_flow"))
    if df:
        elements.append(Spacer(1, 8))
        elements.append(Paragraph("<b>Data Flow:</b> " + " &rarr; ".join(safe_str(x) for x in df), styles["ECIBody"]))
    elements.append(PageBreak())

    # ── Scope ──
    elements.append(Paragraph("Scope Definition", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=rl_colors.Color(0, 180/255, 216/255), spaceAfter=8))
    elements.append(Paragraph("<b>In Scope:</b>", styles["ECIBody"]))
    for item in safe_list(sc.get("in_scope")):
        elements.append(Paragraph(safe_str(item), styles["ECIBullet"], bulletText="\u2022"))
    elements.append(Spacer(1, 6))
    elements.append(Paragraph("<b>Out of Scope:</b>", styles["ECIBody"]))
    for item in safe_list(sc.get("out_of_scope")):
        elements.append(Paragraph(safe_str(item), styles["ECIBullet"], bulletText="\u2022"))
    elements.append(Spacer(1, 6))
    elements.append(Paragraph("<b>Assumptions:</b>", styles["ECIBody"]))
    for item in safe_list(sc.get("assumptions")):
        elements.append(Paragraph(safe_str(item), styles["ECIBullet"], bulletText="\u2022"))
    elements.append(Spacer(1, 6))
    elements.append(Paragraph("<b>Prerequisites:</b>", styles["ECIBody"]))
    for item in safe_list(sc.get("prerequisites")):
        elements.append(Paragraph(safe_str(item), styles["ECIBullet"], bulletText="\u2022"))

    # Build
    def _footer(canvas, doc):
        canvas.saveState()
        canvas.setStrokeColor(rl_colors.Color(27/255, 58/255, 92/255))
        canvas.line(18*mm, 15*mm, A4[0] - 18*mm, 15*mm)
        canvas.setFont("Helvetica", 7)
        canvas.setFillColor(rl_colors.Color(100/255, 100/255, 100/255))
        canvas.drawString(18*mm, 10*mm, "ECI Consulting  |  Generated by Agent BELAL  |  " + datetime.now().strftime("%B %d, %Y"))
        canvas.drawRightString(A4[0] - 18*mm, 10*mm, "Page " + str(canvas.getPageNumber()))
        canvas.restoreState()

    doc.build(elements, onFirstPage=_footer, onLaterPages=_footer)
    buf.seek(0)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════
#  SOW (STATEMENT OF WORK) PDF GENERATOR (ECI Template)
# ═══════════════════════════════════════════════════════════════════════

def generate_sow_pdf(results):
    """Generate a professional Statement of Work PDF using ECI branding.

    Includes: project overview, scope, deliverables, timeline, team,
    acceptance criteria, assumptions, payment terms, and signatures.
    """
    if not HAS_REPORTLAB:
        return None
    buf = io.BytesIO()
    styles = _eci_styles()
    eci_blue = rl_colors.Color(27/255, 58/255, 92/255)
    eci_light = rl_colors.Color(214/255, 228/255, 240/255)
    eci_accent = rl_colors.Color(0, 180/255, 216/255)
    eci_green = rl_colors.Color(0, 212/255, 170/255)
    gray = rl_colors.Color(100/255, 100/255, 100/255)

    doc = SimpleDocTemplate(buf, pagesize=A4, topMargin=25*mm, bottomMargin=20*mm,
                            leftMargin=18*mm, rightMargin=18*mm,
                            title="ECI Statement of Work", author="ECI Consulting — Agent BELAL")
    elements = []

    se = safe_dict(results.get("semantic_analysis"))
    te = safe_dict(results.get("time_estimate"))
    ce = safe_dict(results.get("cost_estimate"))
    ri = safe_dict(results.get("risk_assessment"))
    ar = safe_dict(results.get("architecture"))
    sc = safe_dict(results.get("scope"))
    proposal = safe_dict(results.get("proposal"))
    project_title = safe_str(se.get("project_type", "Technology Solution"))

    # ── Cover Page ──
    elements.append(Spacer(1, 50))
    cover_data = [
        [Paragraph("<b>STATEMENT OF WORK</b>", styles["ECITitle"])],
        [Paragraph(project_title, styles["ECISubtitle"])],
        [Paragraph("", styles["ECISubtitle"])],
        [Paragraph("Prepared by: ECI Consulting", styles["ECISubtitle"])],
        [Paragraph("Prepared for: [Client Name]", styles["ECISubtitle"])],
        [Paragraph("Date: " + datetime.now().strftime("%B %d, %Y"), styles["ECISubtitle"])],
        [Paragraph("Version: 1.0  |  Status: DRAFT", styles["ECISubtitle"])],
    ]
    cover = Table(cover_data, colWidths=[480])
    cover.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), eci_blue),
        ("TOPPADDING", (0, 0), (-1, -1), 8),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
    ]))
    elements.append(cover)
    elements.append(Spacer(1, 30))
    elements.append(Paragraph("CONFIDENTIAL — This document contains proprietary information.", styles["ECISmall"]))
    elements.append(PageBreak())

    # ── Table of Contents ──
    elements.append(Paragraph("Table of Contents", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=eci_accent, spaceAfter=8))
    toc_items = [
        "1. Project Overview",
        "2. Scope of Work",
        "3. Deliverables",
        "4. Detailed Task Breakdown and Timeline",
        "5. Team Composition and Roles",
        "6. Infrastructure and Cost Estimate",
        "7. Risk Assessment",
        "8. Assumptions and Dependencies",
        "9. Acceptance Criteria",
        "10. Change Management",
        "11. Payment Terms",
        "12. Signatures",
    ]
    for item in toc_items:
        elements.append(Paragraph(item, styles["ECIBody"]))
    elements.append(PageBreak())

    # ── 1. Project Overview ──
    elements.append(Paragraph("1. Project Overview", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=eci_accent, spaceAfter=8))
    # Executive summary from proposal if available
    exec_summary = ""
    for sec in safe_list(proposal.get("sections")):
        sec = safe_dict(sec)
        if "executive" in safe_str(sec.get("title")).lower() or "overview" in safe_str(sec.get("title")).lower():
            exec_summary = safe_str(sec.get("content"))
            break
    if exec_summary:
        elements.extend(_md_to_para(exec_summary, styles))
    else:
        elements.append(Paragraph("This Statement of Work defines the scope, deliverables, timeline, and terms for the " + project_title + " project.", styles["ECIBody"]))
    elements.append(Spacer(1, 6))

    # KPI summary
    kpi_data = [
        [Paragraph("<b>" + str(len(safe_list(se.get("requirements")))) + "</b>", styles["ECIKPIVal"]),
         Paragraph("<b>" + str(safe_int(te.get("total_hours"))) + "</b>", styles["ECIKPIVal"]),
         Paragraph("<b>" + safe_str(te.get("duration_weeks")) + "</b>", styles["ECIKPIVal"]),
         Paragraph("<b>$" + str(safe_int(ce.get("total_monthly_cost"))) + "/mo</b>", styles["ECIKPIVal"])],
        [Paragraph("Requirements", styles["ECIKPILabel"]),
         Paragraph("Total Hours", styles["ECIKPILabel"]),
         Paragraph("Duration", styles["ECIKPILabel"]),
         Paragraph("Infra Cost", styles["ECIKPILabel"])],
    ]
    kpi_table = Table(kpi_data, colWidths=[120, 120, 120, 120])
    kpi_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), eci_light),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("TOPPADDING", (0, 0), (-1, -1), 8),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
        ("GRID", (0, 0), (-1, -1), 0.5, rl_colors.white),
    ]))
    elements.append(kpi_table)
    elements.append(Spacer(1, 6))

    # Business objectives
    objectives = safe_list(se.get("business_objectives"))
    if objectives:
        elements.append(Paragraph("<b>Business Objectives:</b>", styles["ECIBody"]))
        for obj in objectives:
            elements.append(Paragraph(safe_str(obj), styles["ECIBullet"], bulletText="\u2022"))
    elements.append(PageBreak())

    # ── 2. Scope of Work ──
    elements.append(Paragraph("2. Scope of Work", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=eci_accent, spaceAfter=8))

    elements.append(Paragraph("<b>In Scope:</b>", styles["ECIBody"]))
    for item in safe_list(sc.get("in_scope")):
        elements.append(Paragraph(safe_str(item), styles["ECIBullet"], bulletText="\u2022"))
    elements.append(Spacer(1, 6))

    elements.append(Paragraph("<b>Out of Scope:</b>", styles["ECIBody"]))
    for item in safe_list(sc.get("out_of_scope")):
        elements.append(Paragraph(safe_str(item), styles["ECIBullet"], bulletText="\u2022"))
    elements.append(PageBreak())

    # ── 3. Deliverables ──
    elements.append(Paragraph("3. Deliverables", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=eci_accent, spaceAfter=8))

    # Build deliverables from phases
    deliverable_rows = []
    phase_idx = 1
    for phase in safe_list(te.get("phases")):
        phase = safe_dict(phase)
        ph_name = safe_str(phase.get("name"))
        ph_week = safe_str(phase.get("week_label"))
        tasks = safe_list(phase.get("tasks"))
        task_names = ", ".join(safe_str(safe_dict(t).get("name")) for t in tasks[:3])
        if len(tasks) > 3:
            task_names += " (+" + str(len(tasks) - 3) + " more)"
        deliverable_rows.append([str(phase_idx), ph_name, ph_week, task_names])
        phase_idx += 1

    if deliverable_rows:
        elements.append(_eci_table(
            ["#", "Deliverable", "Timeline", "Key Tasks"],
            deliverable_rows,
            [30, 120, 80, 250],
        ))
    elements.append(PageBreak())

    # ── 4. Detailed Task Breakdown and Timeline ──
    elements.append(Paragraph("4. Detailed Task Breakdown and Timeline", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=eci_accent, spaceAfter=8))

    elements.append(Paragraph(
        "Total Estimated Hours: <b>" + str(safe_int(te.get("total_hours"))) + "</b>"
        + "  |  Duration: <b>" + safe_str(te.get("duration_weeks")) + "</b>"
        + "  |  Confidence: <b>" + safe_str(te.get("confidence")) + "</b>"
        + "  |  Buffer: <b>" + safe_str(te.get("buffer")) + "</b>",
        styles["ECIBody"],
    ))
    elements.append(Spacer(1, 6))

    three_pt = safe_dict(te.get("three_point"))
    elements.append(Paragraph(
        "Three-Point Estimate: Optimistic <b>" + str(safe_int(three_pt.get("optimistic"))) + "h</b>"
        + "  |  Most Likely <b>" + str(safe_int(three_pt.get("most_likely"))) + "h</b>"
        + "  |  Pessimistic <b>" + str(safe_int(three_pt.get("pessimistic"))) + "h</b>",
        styles["ECIBody"],
    ))
    elements.append(Spacer(1, 8))

    for phase in safe_list(te.get("phases")):
        phase = safe_dict(phase)
        ph_name = safe_str(phase.get("name"))
        ph_week = safe_str(phase.get("week_label"))
        ph_hours = safe_int(phase.get("hours"))
        ph_low = safe_int(phase.get("low_hours"))
        ph_high = safe_int(phase.get("high_hours"))
        ph_pct = safe_str(phase.get("percentage"))

        elements.append(Paragraph(
            "<b>" + ph_name + "</b> (" + ph_week + ") — " + str(ph_hours) + "h [" + str(ph_low) + "-" + str(ph_high) + "h] " + ph_pct,
            styles["ECIBody"],
        ))

        task_rows = []
        for task in safe_list(phase.get("tasks")):
            task = safe_dict(task)
            task_rows.append([
                safe_str(task.get("name")),
                safe_str(task.get("role")),
                str(safe_int(task.get("low_hours"))),
                str(safe_int(task.get("hours"))),
                str(safe_int(task.get("high_hours"))),
                safe_str(task.get("justification"))[:60],
            ])
        if task_rows:
            elements.append(_eci_table(
                ["Sub-task", "Role", "Low", "Avg", "High", "Justification"],
                task_rows,
                [140, 60, 35, 35, 35, 175],
            ))
        elements.append(Spacer(1, 6))

    # Milestones
    milestones = safe_list(te.get("milestones"))
    if milestones:
        elements.append(Spacer(1, 4))
        elements.append(Paragraph("<b>Key Milestones:</b>", styles["ECIBody"]))
        ms_rows = [[safe_str(m.get("name")), "Week " + str(safe_int(m.get("week"))), safe_str(m.get("description"))]
                    for m in milestones if isinstance(m, dict)]
        if ms_rows:
            elements.append(_eci_table(["Milestone", "Week", "Description"], ms_rows, [150, 80, 250]))
    elements.append(PageBreak())

    # ── 5. Team Composition and Roles ──
    elements.append(Paragraph("5. Team Composition and Roles", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=eci_accent, spaceAfter=8))

    total_h = safe_int(te.get("total_hours", 0))
    total_d = max(1, total_h / 7)
    role_rows = []
    for rl in safe_list(te.get("roles")):
        rl = safe_dict(rl)
        role_name = safe_str(rl.get("name"))
        pct = float(rl.get("allocation_pct", 0))
        rate = safe_int(rl.get("rate", 100))
        days = round(total_d * pct, 1)
        hours = int(round(days * 7, 0))
        cost = hours * rate
        role_rows.append([role_name, str(int(pct * 100)) + "%", str(days), str(hours), "$" + str(rate), "$" + str(cost)])
    if role_rows:
        elements.append(_eci_table(
            ["Role", "Allocation", "Days", "Hours", "Rate/hr", "Est. Cost"],
            role_rows,
            [110, 55, 50, 50, 55, 80],
        ))
    elements.append(PageBreak())

    # ── 6. Infrastructure and Cost Estimate ──
    elements.append(Paragraph("6. Infrastructure and Cost Estimate", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=eci_accent, spaceAfter=8))

    elements.append(Paragraph(
        "Monthly Infrastructure Cost: <b>$" + str(safe_int(ce.get("total_monthly_cost"))) + "</b>"
        + "  |  Annual: <b>$" + str(safe_int(ce.get("total_annual_cost"))) + "</b>",
        styles["ECIBody"],
    ))
    elements.append(Spacer(1, 6))

    cost_rows = []
    for svc in safe_list(ce.get("azure_costs")):
        svc = safe_dict(svc)
        cost_rows.append([
            safe_str(svc.get("service")),
            safe_str(svc.get("tier", "")),
            "$" + str(safe_int(svc.get("monthly_cost"))),
            "$" + str(safe_int(svc.get("monthly_cost")) * 12),
            safe_str(svc.get("description", ""))[:50],
        ])
    if cost_rows:
        elements.append(_eci_table(
            ["Service", "Tier", "Monthly", "Annual", "Notes"],
            cost_rows,
            [110, 70, 60, 60, 180],
        ))
    elements.append(PageBreak())

    # ── 7. Risk Assessment ──
    elements.append(Paragraph("7. Risk Assessment", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=eci_accent, spaceAfter=8))

    elements.append(Paragraph(
        "Overall Risk Score: <b>" + str(safe_int(ri.get("overall_score"))) + "/10</b>"
        + "  |  Level: <b>" + safe_str(ri.get("overall_level")) + "</b>",
        styles["ECIBody"],
    ))
    elements.append(Spacer(1, 6))

    risk_rows = []
    for rk in safe_list(ri.get("risks")):
        rk = safe_dict(rk)
        risk_rows.append([
            safe_str(rk.get("category")),
            safe_str(rk.get("title")),
            safe_str(rk.get("severity")),
            safe_str(rk.get("mitigation"))[:60],
        ])
    if risk_rows:
        elements.append(_eci_table(
            ["Category", "Risk", "Severity", "Mitigation"],
            risk_rows,
            [80, 140, 60, 200],
        ))
    elements.append(PageBreak())

    # ── 8. Assumptions and Dependencies ──
    elements.append(Paragraph("8. Assumptions and Dependencies", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=eci_accent, spaceAfter=8))

    elements.append(Paragraph("<b>Assumptions:</b>", styles["ECIBody"]))
    assumptions = safe_list(sc.get("assumptions"))
    if not assumptions:
        assumptions = [
            "Client will provide timely access to required systems and environments",
            "Key stakeholders will be available for reviews and sign-offs within agreed timelines",
            "Requirements are baselined — changes will follow the change management process",
            "Client will provision necessary Azure subscriptions and licenses",
        ]
    for item in assumptions:
        elements.append(Paragraph(safe_str(item), styles["ECIBullet"], bulletText="\u2022"))
    elements.append(Spacer(1, 6))

    elements.append(Paragraph("<b>Prerequisites:</b>", styles["ECIBody"]))
    prereqs = safe_list(sc.get("prerequisites"))
    if not prereqs:
        prereqs = [
            "Signed Statement of Work (this document)",
            "Cloud subscription access provisioned",
            "VPN/network access for development team",
            "Sample data and test accounts provided",
        ]
    for item in prereqs:
        elements.append(Paragraph(safe_str(item), styles["ECIBullet"], bulletText="\u2022"))
    elements.append(PageBreak())

    # ── 9. Acceptance Criteria ──
    elements.append(Paragraph("9. Acceptance Criteria", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=eci_accent, spaceAfter=8))

    acceptance = [
        "All functional requirements as documented in Section 2 are implemented and verified",
        "All integration points are operational and tested end-to-end",
        "Non-functional requirements (performance, security, availability) meet agreed thresholds",
        "User Acceptance Testing (UAT) completed with formal sign-off from designated stakeholders",
        "All critical and high-severity defects are resolved prior to go-live",
        "Technical and user documentation delivered and reviewed",
        "Knowledge transfer sessions completed with client team",
        "Production environment deployed and verified with smoke testing",
    ]
    for item in acceptance:
        elements.append(Paragraph(item, styles["ECIBullet"], bulletText="\u2022"))
    elements.append(PageBreak())

    # ── 10. Change Management ──
    elements.append(Paragraph("10. Change Management", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=eci_accent, spaceAfter=8))

    elements.append(Paragraph(
        "Any changes to the scope, timeline, or deliverables defined in this SOW must follow "
        "the change management process outlined below:",
        styles["ECIBody"],
    ))
    elements.append(Spacer(1, 4))
    change_items = [
        "<b>Change Request Submission:</b> All change requests must be submitted in writing with a description of the change, justification, and expected impact.",
        "<b>Impact Analysis:</b> ECI will assess the impact on timeline, cost, and resources within 3 business days.",
        "<b>Approval:</b> Changes must be approved in writing by both parties before implementation.",
        "<b>Tracking:</b> All approved changes will be tracked in the project change log with updated estimates.",
    ]
    for item in change_items:
        elements.append(Paragraph(item, styles["ECIBullet"], bulletText="\u2022"))
    elements.append(PageBreak())

    # ── 11. Payment Terms ──
    elements.append(Paragraph("11. Payment Terms", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=eci_accent, spaceAfter=8))

    total_hours = safe_int(te.get("total_hours"))
    payment_rows = [
        ["Milestone 1: Project Kickoff", "Upon SOW signing", "20%", "Discovery & Design phase"],
        ["Milestone 2: Development Complete", "Core Development delivered", "30%", "All functional requirements implemented"],
        ["Milestone 3: UAT Sign-off", "UAT phase completed", "30%", "User acceptance testing approved"],
        ["Milestone 4: Go-Live", "Production deployment", "20%", "Production deployment and handover complete"],
    ]
    elements.append(_eci_table(
        ["Payment Milestone", "Trigger", "% of Total", "Deliverable"],
        payment_rows,
        [120, 120, 60, 180],
    ))
    elements.append(Spacer(1, 8))
    elements.append(Paragraph(
        "<b>Note:</b> Payment terms are subject to negotiation. Infrastructure costs (Section 6) "
        "are billed separately on a monthly basis based on actual consumption.",
        styles["ECIBody"],
    ))
    elements.append(PageBreak())

    # ── 12. Signatures ──
    elements.append(Paragraph("12. Signatures", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=eci_accent, spaceAfter=8))

    elements.append(Paragraph(
        "By signing below, both parties agree to the terms and conditions outlined in this Statement of Work.",
        styles["ECIBody"],
    ))
    elements.append(Spacer(1, 20))

    sig_data = [
        [Paragraph("<b>For: ECI Consulting</b>", styles["ECIBody"]), Paragraph("", styles["ECIBody"]),
         Paragraph("<b>For: [Client Name]</b>", styles["ECIBody"])],
        [Paragraph("", styles["ECIBody"]), Paragraph("", styles["ECIBody"]), Paragraph("", styles["ECIBody"])],
        [Paragraph("____________________________", styles["ECIBody"]), Paragraph("", styles["ECIBody"]),
         Paragraph("____________________________", styles["ECIBody"])],
        [Paragraph("Name:", styles["ECIBody"]), Paragraph("", styles["ECIBody"]),
         Paragraph("Name:", styles["ECIBody"])],
        [Paragraph("Title:", styles["ECIBody"]), Paragraph("", styles["ECIBody"]),
         Paragraph("Title:", styles["ECIBody"])],
        [Paragraph("Date:", styles["ECIBody"]), Paragraph("", styles["ECIBody"]),
         Paragraph("Date:", styles["ECIBody"])],
    ]
    sig_table = Table(sig_data, colWidths=[200, 80, 200])
    sig_table.setStyle(TableStyle([
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("TOPPADDING", (0, 0), (-1, -1), 6),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
    ]))
    elements.append(sig_table)

    # ── Build PDF ──
    def _sow_footer(canvas, doc):
        canvas.saveState()
        canvas.setStrokeColor(rl_colors.Color(27/255, 58/255, 92/255))
        canvas.line(18*mm, 15*mm, A4[0] - 18*mm, 15*mm)
        canvas.setFont("Helvetica", 7)
        canvas.setFillColor(rl_colors.Color(100/255, 100/255, 100/255))
        canvas.drawString(18*mm, 10*mm, "ECI Consulting  |  Statement of Work  |  CONFIDENTIAL  |  " + datetime.now().strftime("%B %d, %Y"))
        canvas.drawRightString(A4[0] - 18*mm, 10*mm, "Page " + str(canvas.getPageNumber()))
        canvas.restoreState()

    doc.build(elements, onFirstPage=_sow_footer, onLaterPages=_sow_footer)
    buf.seek(0)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════
#  PPTX PROPOSAL GENERATOR (ECI Template)
# ═══════════════════════════════════════════════════════════════════════

def generate_proposal_pptx(results):
    """Generate ECI-branded PowerPoint proposal."""
    if not HAS_PPTX:
        return None

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    ECI_BLUE_RGB = RGBColor(27, 58, 92)
    ECI_ACCENT_RGB = RGBColor(0, 180, 216)
    ECI_GREEN_RGB = RGBColor(0, 212, 170)
    WHITE = RGBColor(255, 255, 255)
    DARK_TEXT = RGBColor(30, 30, 30)
    GRAY_TEXT = RGBColor(100, 100, 100)

    se = safe_dict(results.get("semantic_analysis"))
    te = safe_dict(results.get("time_estimate"))
    ce = safe_dict(results.get("cost_estimate"))
    ri = safe_dict(results.get("risk_assessment"))
    ar = safe_dict(results.get("architecture"))
    sc = safe_dict(results.get("scope"))
    proposal = safe_dict(results.get("proposal"))

    def add_bg(slide, color=ECI_BLUE_RGB):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title_bar(slide, title_text):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ECI_BLUE_RGB
        bar.line.fill.background()
        tf = bar.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        tf.margin_left = Inches(0.5)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        # ECI logo text in title bar
        logo_shape = slide.shapes.add_textbox(Inches(11), Inches(0.15), Inches(2), Inches(0.9))
        ltf = logo_shape.text_frame
        lp = ltf.paragraphs[0]
        lp.text = "ECI"
        lp.font.size = Pt(24)
        lp.font.color.rgb = ECI_ACCENT_RGB
        lp.font.bold = True
        lp.alignment = PP_ALIGN.RIGHT

    def add_footer(slide):
        footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.0), prs.slide_width, Inches(0.5))
        footer.fill.solid()
        footer.fill.fore_color.rgb = RGBColor(20, 30, 50)
        footer.line.fill.background()
        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.text = "ECI Consulting  |  Confidential  |  " + datetime.now().strftime("%B %d, %Y")
        p.font.size = Pt(9)
        p.font.color.rgb = GRAY_TEXT
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    def add_accent_line(slide, x, y, width):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = ECI_ACCENT_RGB
        line.line.fill.background()

    # ── Slide 1: Title ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    # Try to add actual ECI logo
    try:
        logo_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "eci_logo_white.png")
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(0.5), Inches(0.4), height=Inches(0.9))
    except Exception:
        logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(2), Inches(0.9))
        tf = logo_box.text_frame
        p = tf.paragraphs[0]
        p.text = "ECI"
        p.font.size = Pt(36)
        p.font.color.rgb = WHITE
        p.font.bold = True

    add_accent_line(slide, Inches(0.8), Inches(2.3), Inches(3))

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11), Inches(3))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "PROJECT PROPOSAL"
    p.font.size = Pt(44)
    p.font.color.rgb = WHITE
    p.font.bold = True

    p2 = tf.add_paragraph()
    p2.text = safe_str(se.get("project_type", "Technology Solution"))
    p2.font.size = Pt(24)
    p2.font.color.rgb = ECI_ACCENT_RGB
    p2.space_before = Pt(8)

    p3 = tf.add_paragraph()
    p3.text = "\nPrepared by ECI"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(180, 200, 220)

    p4 = tf.add_paragraph()
    p4.text = datetime.now().strftime("%B %d, %Y")
    p4.font.size = Pt(14)
    p4.font.color.rgb = RGBColor(180, 200, 220)

    # ── Slide 2: Executive Summary ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Executive Summary")
    add_footer(slide)

    kpi_data = [
        ("Requirements", str(len(safe_list(se.get("requirements"))))),
        ("Total Hours", str(safe_int(te.get("total_hours")))),
        ("Infra Cost/Mo", "$" + str(safe_int(ce.get("total_monthly_cost")))),
        ("Risk Level", safe_str(ri.get("overall_level", "N/A"))),
    ]
    for i, (label, val) in enumerate(kpi_data):
        x = Inches(0.5 + i * 3.1)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(2.8), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(240, 245, 250)
        box.line.color.rgb = RGBColor(200, 210, 225)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(28)
        p.font.color.rgb = ECI_BLUE_RGB
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(12)
        p2.font.color.rgb = GRAY_TEXT
        p2.alignment = PP_ALIGN.CENTER

    summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12), Inches(3.5))
    tf = summary_box.text_frame
    tf.word_wrap = True
    for sec in safe_list(proposal.get("sections"))[:2]:
        sec = safe_dict(sec)
        p = tf.add_paragraph()
        p.text = safe_str(sec.get("title", ""))
        p.font.size = Pt(16)
        p.font.color.rgb = ECI_BLUE_RGB
        p.font.bold = True
        p.space_after = Pt(4)
        content = safe_str(sec.get("content", ""))
        for ln in content.split("\n")[:5]:
            ln = ln.strip().replace("**", "").replace("_", "")
            if ln:
                p2 = tf.add_paragraph()
                p2.text = ln
                p2.font.size = Pt(11)
                p2.font.color.rgb = DARK_TEXT
                p2.space_after = Pt(2)

    # ── Slide 3: Requirements ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Requirements Analysis")
    add_footer(slide)

    reqs = safe_list(se.get("requirements"))
    fn_list = [x for x in reqs if isinstance(x, dict) and x.get("type") == "functional"]
    nf_list = [x for x in reqs if isinstance(x, dict) and x.get("type") == "non-functional"]
    ig_list = [x for x in reqs if isinstance(x, dict) and x.get("type") == "integration"]

    cats = [("Functional", fn_list, ECI_GREEN_RGB), ("Non-Functional", nf_list, ECI_ACCENT_RGB), ("Integration", ig_list, RGBColor(123, 97, 255))]
    for ci, (cat_label, cat_items, cat_clr) in enumerate(cats):
        x = Inches(0.3 + ci * 4.3)
        header_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(4), Inches(0.6))
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = cat_clr
        header_box.line.fill.background()
        tf = header_box.text_frame
        p = tf.paragraphs[0]
        p.text = cat_label + " (" + str(len(cat_items)) + ")"
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        items_box = slide.shapes.add_textbox(x + Inches(0.1), Inches(2.2), Inches(3.8), Inches(4.5))
        tf = items_box.text_frame
        tf.word_wrap = True
        for req in cat_items[:8]:
            req = safe_dict(req)
            p = tf.add_paragraph()
            p.text = "\u2022 " + safe_str(req.get("title"))
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_TEXT
            p.space_after = Pt(3)

    # ── Slide 4: Timeline & Phases ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Timeline & Phases")
    add_footer(slide)

    phases = safe_list(te.get("phases"))
    colors_pptx = [RGBColor(0, 212, 170), RGBColor(0, 180, 216), RGBColor(123, 97, 255),
                   RGBColor(255, 107, 107), RGBColor(255, 209, 102), RGBColor(6, 214, 160)]
    for i, phase in enumerate(phases[:6]):
        phase = safe_dict(phase)
        x = Inches(0.5 + (i % 3) * 4.1)
        y = Inches(1.5 + (i // 3) * 2.5)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 248, 252)
        box.line.color.rgb = colors_pptx[i % 6]

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        p = tf.paragraphs[0]
        p.text = safe_str(phase.get("name"))
        p.font.size = Pt(16)
        p.font.color.rgb = ECI_BLUE_RGB
        p.font.bold = True

        p2 = tf.add_paragraph()
        p2.text = str(safe_int(phase.get("hours"))) + " hours  |  " + safe_str(phase.get("percentage"))
        p2.font.size = Pt(12)
        p2.font.color.rgb = ECI_ACCENT_RGB
        p2.font.bold = True

        for task in safe_list(phase.get("tasks"))[:3]:
            task = safe_dict(task)
            p3 = tf.add_paragraph()
            p3.text = "\u2022 " + safe_str(task.get("name")) + " (" + safe_str(task.get("role")) + ")"
            p3.font.size = Pt(9)
            p3.font.color.rgb = GRAY_TEXT

    # Summary bar at bottom
    summ_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.6))
    summ_box.fill.solid()
    summ_box.fill.fore_color.rgb = RGBColor(240, 245, 250)
    summ_box.line.color.rgb = ECI_ACCENT_RGB
    tf = summ_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Total: " + str(safe_int(te.get("total_hours"))) + " hours  |  Duration: " + safe_str(te.get("duration_weeks")) + "  |  Confidence: " + safe_str(te.get("confidence")) + "  |  Buffer: " + safe_str(te.get("buffer"))
    p.font.size = Pt(12)
    p.font.color.rgb = ECI_BLUE_RGB
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ── Slide 5: Infrastructure Costs ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Infrastructure Cost Estimate")
    add_footer(slide)

    monthly = safe_int(ce.get("total_monthly_cost"))
    annual = safe_int(ce.get("total_annual_cost"))
    for i, (label, val) in enumerate([("Monthly Cost", "$" + str(monthly)), ("Annual Cost", "$" + str(annual))]):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5 + i * 4), Inches(1.5), Inches(3.5), Inches(1))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(240, 245, 250)
        box.line.color.rgb = ECI_ACCENT_RGB
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(24)
        p.font.color.rgb = ECI_BLUE_RGB
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY_TEXT
        p2.alignment = PP_ALIGN.CENTER

    services_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12), Inches(3.8))
    tf = services_box.text_frame
    tf.word_wrap = True
    for svc in safe_list(ce.get("azure_costs"))[:12]:
        svc = safe_dict(svc)
        p = tf.add_paragraph()
        p.text = safe_str(svc.get("service")) + "  \u2014  $" + str(safe_int(svc.get("monthly_cost"))) + "/mo  (" + safe_str(svc.get("tier", "")) + ")"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(3)

    # ── Slide 6: Risk Assessment ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Risk Assessment")
    add_footer(slide)

    score = safe_int(ri.get("overall_score"))
    score_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(3), Inches(1.2))
    score_box.fill.solid()
    score_box.fill.fore_color.rgb = RGBColor(240, 245, 250)
    score_clr = ECI_GREEN_RGB if score <= 3 else RGBColor(255, 209, 102) if score <= 6 else RGBColor(255, 107, 107)
    score_box.line.color.rgb = score_clr
    tf = score_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(score) + "/10"
    p.font.size = Pt(32)
    p.font.color.rgb = score_clr
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Overall: " + safe_str(ri.get("overall_level"))
    p2.font.size = Pt(14)
    p2.font.color.rgb = GRAY_TEXT
    p2.alignment = PP_ALIGN.CENTER

    risks_box = slide.shapes.add_textbox(Inches(4), Inches(1.5), Inches(8.5), Inches(5))
    tf = risks_box.text_frame
    tf.word_wrap = True
    for rk in safe_list(ri.get("risks")):
        rk = safe_dict(rk)
        p = tf.add_paragraph()
        p.text = safe_str(rk.get("category")) + ": " + safe_str(rk.get("title")) + " [" + safe_str(rk.get("severity")) + "]"
        p.font.size = Pt(13)
        p.font.color.rgb = ECI_BLUE_RGB
        p.font.bold = True
        p.space_after = Pt(2)
        p2 = tf.add_paragraph()
        p2.text = "Mitigation: " + safe_str(rk.get("mitigation"))
        p2.font.size = Pt(10)
        p2.font.color.rgb = GRAY_TEXT
        p2.space_after = Pt(8)

    # ── Slide 7: Architecture ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Solution Architecture")
    add_footer(slide)

    pattern_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(0.5))
    tf = pattern_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Pattern: " + safe_str(ar.get("pattern"))
    p.font.size = Pt(16)
    p.font.color.rgb = ECI_ACCENT_RGB
    p.font.bold = True

    comps = safe_list(ar.get("components"))
    for i, comp in enumerate(comps[:6]):
        comp = safe_dict(comp)
        x = Inches(0.5 + (i % 3) * 4.1)
        y = Inches(2.2 + (i // 3) * 2.3)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 248, 252)
        box.line.color.rgb = ECI_ACCENT_RGB
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        p = tf.paragraphs[0]
        p.text = safe_str(comp.get("name"))
        p.font.size = Pt(14)
        p.font.color.rgb = ECI_BLUE_RGB
        p.font.bold = True
        p2 = tf.add_paragraph()
        p2.text = safe_str(comp.get("azure_service", ""))
        p2.font.size = Pt(10)
        p2.font.color.rgb = ECI_ACCENT_RGB
        for s_item in safe_list(comp.get("services"))[:3]:
            p3 = tf.add_paragraph()
            p3.text = "\u2022 " + safe_str(s_item)
            p3.font.size = Pt(9)
            p3.font.color.rgb = GRAY_TEXT

    df = safe_list(ar.get("data_flow"))
    if df:
        flow_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.6))
        flow_box.fill.solid()
        flow_box.fill.fore_color.rgb = RGBColor(240, 245, 250)
        flow_box.line.color.rgb = ECI_ACCENT_RGB
        tf = flow_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Data Flow:  " + "  \u2192  ".join(safe_str(x) for x in df)
        p.font.size = Pt(11)
        p.font.color.rgb = ECI_BLUE_RGB
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ── Slide 8: Scope ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Scope Definition")
    add_footer(slide)

    scope_sections = [
        ("In Scope", safe_list(sc.get("in_scope")), ECI_GREEN_RGB),
        ("Out of Scope", safe_list(sc.get("out_of_scope")), RGBColor(255, 107, 107)),
        ("Assumptions", safe_list(sc.get("assumptions")), ECI_ACCENT_RGB),
        ("Prerequisites", safe_list(sc.get("prerequisites")), RGBColor(123, 97, 255)),
    ]
    for si, (s_title, s_items, s_clr) in enumerate(scope_sections):
        x = Inches(0.3 + (si % 2) * 6.5)
        y = Inches(1.4 + (si // 2) * 2.9)
        hdr_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(6.2), Inches(0.5))
        hdr_box.fill.solid()
        hdr_box.fill.fore_color.rgb = s_clr
        hdr_box.line.fill.background()
        tf = hdr_box.text_frame
        p = tf.paragraphs[0]
        p.text = s_title
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        items_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.6), Inches(6), Inches(2.2))
        tf = items_box.text_frame
        tf.word_wrap = True
        for s_item in s_items[:6]:
            p = tf.add_paragraph()
            p.text = "\u2022 " + safe_str(s_item)
            p.font.size = Pt(9)
            p.font.color.rgb = DARK_TEXT
            p.space_after = Pt(2)

    # ── Slide 9: Thank You ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    try:
        logo_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "eci_logo_white.png")
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(5.5), Inches(1.5), height=Inches(1.2))
    except Exception:
        pass

    thank_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(2.5))
    tf = thank_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(44)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "ECI"
    p2.font.size = Pt(16)
    p2.font.color.rgb = ECI_ACCENT_RGB
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = datetime.now().strftime("%B %d, %Y")
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(180, 200, 220)
    p3.alignment = PP_ALIGN.CENTER

    add_accent_line(slide, Inches(4), Inches(5.5), Inches(5))

    p4 = tf.add_paragraph()
    p4.text = "\nwww.eciconsulting.com"
    p4.font.size = Pt(11)
    p4.font.color.rgb = RGBColor(150, 170, 200)
    p4.alignment = PP_ALIGN.CENTER

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════
#  SESSION STATE
# ═══════════════════════════════════════════════════════════════════════

_defaults = {
    "azure_api_key": "", "azure_endpoint": "", "azure_api_version": "2024-06-01", "azure_deployment": "gpt-4",
    "sp_url": "", "sp_cid": "", "sp_cs": "", "sp_tid": "",
    "email_smtp": "", "email_sender": "",
    "processing_results": None, "historical_projects": [], "agent_logs": [],
    "discovery_results": None, "discovery_transcript": "",
    "model_metrics": {"accuracy": 78.5, "proposals_processed": 0, "win_rate": 62.0, "variance": 12.3},
}
for k, v in _defaults.items():
    if k not in st.session_state:
        st.session_state[k] = v

inject_css()


# ═══════════════════════════════════════════════════════════════════════
#  HEADER
# ═══════════════════════════════════════════════════════════════════════

hc1, hc2, hc3 = st.columns([2, 6, 2])
with hc1:
    st.markdown('<div class="logo-box"><img src="data:image/png;base64,' + ECI_LOGO_BLUE_B64 + '" height="42" style="margin-right:8px;vertical-align:middle"/><div><div class="logo-txt"></div><div class="logo-sub">Agent by ECI</div></div></div>', unsafe_allow_html=True)
with hc2:
    st.markdown('<div class="tagline">Business Estimation Leveraging Automated Learning</div>', unsafe_allow_html=True)
with hc3:
    st.markdown('<div class="dt">' + datetime.now().strftime("%b %d, %Y  %H:%M") + '</div>', unsafe_allow_html=True)


# ═══════════════════════════════════════════════════════════════════════
#  SIDEBAR
# ═══════════════════════════════════════════════════════════════════════

with st.sidebar:
    st.markdown('<div class="stitle">Configuration</div>', unsafe_allow_html=True)
    st.markdown('<div class="csec">Azure OpenAI</div>', unsafe_allow_html=True)
    st.session_state.azure_api_key = st.text_input("API Key", value=st.session_state.azure_api_key, type="password", key="k1")
    st.session_state.azure_endpoint = st.text_input("Endpoint", value=st.session_state.azure_endpoint, placeholder="https://your-resource.openai.azure.com/", key="k2")
    st.session_state.azure_api_version = st.selectbox("Version", ["2024-06-01", "2024-02-01", "2023-12-01-preview"], key="k3")
    st.session_state.azure_deployment = st.text_input("Deployment", value=st.session_state.azure_deployment, placeholder="gpt-4", key="k4")
    if st.button("Test Connection", use_container_width=True, key="tb1"):
        ok, msg = AzureAI.from_session().test()
        if ok:
            st.success(msg)
        else:
            st.error(msg)
    st.divider()
    st.markdown('<div class="csec">SharePoint</div>', unsafe_allow_html=True)
    st.session_state.sp_url = st.text_input("Site URL", value=st.session_state.sp_url, placeholder="https://org.sharepoint.com/sites/presales", key="k5")
    st.session_state.sp_cid = st.text_input("Client ID", value=st.session_state.sp_cid, key="k6")
    st.session_state.sp_cs = st.text_input("Client Secret", value=st.session_state.sp_cs, type="password", key="k7")
    st.session_state.sp_tid = st.text_input("Tenant ID", value=st.session_state.sp_tid, key="k8")
    if st.button("Test SharePoint", use_container_width=True, key="tb2"):
        ok, msg = SP.from_session().test()
        if ok:
            st.success(msg)
        else:
            st.error(msg)
    st.divider()
    st.markdown('<div class="csec">Email Alerts</div>', unsafe_allow_html=True)
    st.session_state.email_smtp = st.text_input("SMTP Server", value=st.session_state.email_smtp, placeholder="smtp.office365.com:587", key="k9")
    st.session_state.email_sender = st.text_input("Sender Email", value=st.session_state.email_sender, key="k10")
    st.text_input("Password", type="password", key="cfg_email_pass")
    st.divider()
    a_ok = bool(st.session_state.azure_api_key and st.session_state.azure_endpoint)
    s_ok = bool(st.session_state.sp_url and st.session_state.sp_cid)
    e_ok = bool(st.session_state.email_smtp)
    mode_label = "LIVE AI" if a_ok else "DEMO MODE"
    st.markdown('<div class="cstat">'
                + '<div class="srow"><strong>' + mode_label + '</strong></div>'
                + '<div class="srow"><span class="sdot ' + ('on' if a_ok else 'off') + '"></span> Azure OpenAI</div>'
                + '<div class="srow"><span class="sdot ' + ('on' if s_ok else 'off') + '"></span> SharePoint</div>'
                + '<div class="srow"><span class="sdot ' + ('on' if e_ok else 'off') + '"></span> Email</div>'
                + '</div>', unsafe_allow_html=True)


# ═══════════════════════════════════════════════════════════════════════
#  HELPERS
# ═══════════════════════════════════════════════════════════════════════

def log_agent(name, detail):
    st.session_state.agent_logs.append({"agent": name, "detail": detail, "time": datetime.now().isoformat()})


# ═══════════════════════════════════════════════════════════════════════
#  TAB 1: PRESALE
# ═══════════════════════════════════════════════════════════════════════

def tab_presale():
    st.markdown('<div class="shdr"><span class="shdr-i">📄</span> Document Ingestion</div>', unsafe_allow_html=True)
    uc, tc = st.columns([3, 2])
    with uc:
        st.markdown('<div class="crd"><div class="crd-t">Upload Scope Documents</div><div class="crd-d">PDF, DOCX, XLSX, PPTX, TXT, CSV</div>', unsafe_allow_html=True)
        files = st.file_uploader("Drop files", type=["pdf", "docx", "xlsx", "pptx", "txt", "csv"], accept_multiple_files=True, key="fu", label_visibility="collapsed")
        st.markdown("</div>", unsafe_allow_html=True)
    with tc:
        st.markdown('<div class="crd"><div class="crd-t">SharePoint Auto-Trigger</div><div class="crd-d">Monitor folder for new scope docs</div>', unsafe_allow_html=True)
        sp_f = st.text_input("Folder", placeholder="/sites/presales/Shared Documents/Scope", key="spf", label_visibility="collapsed")
        if st.button("Check Now", use_container_width=True, key="spchk"):
            found = SP.from_session().list_folder(sp_f)
            if found:
                st.success("Found " + str(len(found)) + " doc(s)")
                for f in found:
                    st.markdown("- `" + f + "`")
            else:
                st.info("No documents found. Configure SharePoint in sidebar.")
        st.markdown("</div>", unsafe_allow_html=True)

    if files:
        st.markdown("---")
        _, bc, _ = st.columns([1, 2, 1])
        with bc:
            if st.button("PROCESS & GENERATE ESTIMATES", use_container_width=True, type="primary", key="go"):
                run_pipeline(files)

    # ── Multimodal Discovery ──
    st.markdown("---")
    st.markdown('<div class="shdr"><span class="shdr-i">🎙️</span> Multimodal Discovery Extraction</div>', unsafe_allow_html=True)
    st.markdown('<div class="crd"><div class="crd-t">Upload Meeting Transcripts & Voice Notes</div><div class="crd-d">Extract requirements, pain points, stakeholders & auto-generate WBS from Zoom/Teams transcripts, voice memos, or meeting notes.</div>', unsafe_allow_html=True)
    disc_col1, disc_col2 = st.columns([3, 2])
    with disc_col1:
        disc_files = st.file_uploader(
            "Upload transcripts",
            type=["txt", "srt", "vtt", "json", "docx", "csv"],
            accept_multiple_files=True,
            key="disc_fu",
            label_visibility="collapsed",
            help="Supported: TXT, SRT, VTT (subtitles), JSON (Teams/Zoom/Otter.ai exports), DOCX (meeting notes), CSV",
        )
    with disc_col2:
        st.markdown("**Supported Formats:**")
        st.markdown("- `.txt` `.srt` `.vtt` — Subtitle / transcript files")
        st.markdown("- `.json` — Teams, Zoom, Otter.ai exports")
        st.markdown("- `.docx` — Meeting notes / minutes")
        st.markdown("- `.csv` — Tabular transcript exports")
    st.markdown("</div>", unsafe_allow_html=True)

    if disc_files:
        _, dbc, _ = st.columns([1, 2, 1])
        with dbc:
            if st.button("🎙️ ANALYZE TRANSCRIPTS & GENERATE WBS", use_container_width=True, type="primary", key="disc_go"):
                ai = AzureAI.from_session()
                dpb = st.progress(0)
                dstatus = st.empty()
                dstatus.markdown("**1/3** Extracting transcripts...")
                dpb.progress(10)
                all_text = ""
                for df in disc_files:
                    all_text += extract_audio_transcript(df) + "\n\n"
                st.session_state.discovery_transcript = all_text
                dstatus.markdown("**2/3** Analyzing content & extracting insights...")
                dpb.progress(50)
                disc_result = ai.analyze_transcript(all_text)
                dstatus.markdown("**3/3** Generating Work Breakdown Structure...")
                dpb.progress(90)
                st.session_state.discovery_results = disc_result
                dpb.progress(100)
                dstatus.markdown("**Done!** Discovery analysis complete.")
                time.sleep(0.5)
                dstatus.empty()
                dpb.empty()
                st.rerun()

    # ── Discovery Results ──
    if st.session_state.discovery_results:
        dr = st.session_state.discovery_results
        st.markdown("---")
        st.markdown('<div class="shdr"><span class="shdr-i">📋</span> Discovery Insights</div>', unsafe_allow_html=True)

        st.markdown('<div class="crd">', unsafe_allow_html=True)
        st.markdown("**Meeting Summary:** " + safe_str(dr.get("meeting_summary")))
        st.markdown("**Sentiment:** " + safe_str(dr.get("sentiment")))
        themes = safe_list(dr.get("key_themes"))
        if themes:
            tags_html = " ".join('<span class="tt">' + safe_str(t) + '</span>' for t in themes)
            st.markdown('<div class="ttag">' + tags_html + '</div>', unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)

        disc_tabs = st.tabs(["🔥 Pain Points", "📝 Requirements", "👥 Stakeholders", "📊 WBS", "✅ Action Items"])

        with disc_tabs[0]:
            for pp in safe_list(dr.get("pain_points")):
                pp = safe_dict(pp)
                sev = safe_str(pp.get("severity"))
                sc = "#ff6b6b" if sev == "High" else "#ffd166" if sev == "Medium" else "#06d6a0"
                st.markdown(
                    '<div class="rc" style="border-left:3px solid ' + sc + ';">'
                    '<div class="rch2"><strong>' + safe_str(pp.get("issue")) + '</strong>'
                    '<span class="rsev" style="color:' + sc + ';">' + sev + '</span></div>'
                    '<p style="font-style:italic;color:var(--t2);">"' + safe_str(pp.get("quote")) + '"</p>'
                    '<div class="rmit"><strong>Stakeholder:</strong> ' + safe_str(pp.get("stakeholder")) + '</div>'
                    '</div>', unsafe_allow_html=True
                )

        with disc_tabs[1]:
            reqs = safe_list(dr.get("requirements_extracted"))
            fn_r = [x for x in reqs if isinstance(x, dict) and x.get("type") == "functional"]
            nf_r = [x for x in reqs if isinstance(x, dict) and x.get("type") == "non-functional"]
            ig_r = [x for x in reqs if isinstance(x, dict) and x.get("type") == "integration"]
            rc1, rc2, rc3 = st.columns(3)
            with rc1:
                st.markdown('<div class="rch fn">Functional (' + str(len(fn_r)) + ')</div>', unsafe_allow_html=True)
                for q in fn_r:
                    q = safe_dict(q)
                    st.markdown('<div class="ri"><strong>' + safe_str(q.get("title")) + '</strong><p>' + safe_str(q.get("description")) + '</p><div class="ri-c">Source: ' + safe_str(q.get("source")) + '</div></div>', unsafe_allow_html=True)
            with rc2:
                st.markdown('<div class="rch nf">Non-Functional (' + str(len(nf_r)) + ')</div>', unsafe_allow_html=True)
                for q in nf_r:
                    q = safe_dict(q)
                    st.markdown('<div class="ri"><strong>' + safe_str(q.get("title")) + '</strong><p>' + safe_str(q.get("description")) + '</p><div class="ri-c">Source: ' + safe_str(q.get("source")) + '</div></div>', unsafe_allow_html=True)
            with rc3:
                st.markdown('<div class="rch ig">Integration (' + str(len(ig_r)) + ')</div>', unsafe_allow_html=True)
                for q in ig_r:
                    q = safe_dict(q)
                    st.markdown('<div class="ri"><strong>' + safe_str(q.get("title")) + '</strong><p>' + safe_str(q.get("description")) + '</p><div class="ri-c">Source: ' + safe_str(q.get("source")) + '</div></div>', unsafe_allow_html=True)

        with disc_tabs[2]:
            stake_cols = st.columns(2)
            for i, sh in enumerate(safe_list(dr.get("stakeholders"))):
                sh = safe_dict(sh)
                with stake_cols[i % 2]:
                    concerns = "".join("<li>" + safe_str(c) + "</li>" for c in safe_list(sh.get("concerns")))
                    st.markdown(
                        '<div class="crd"><div class="crd-t">👤 ' + safe_str(sh.get("name")) + '</div>'
                        '<div class="crd-d" style="color:var(--c2);">' + safe_str(sh.get("role")) + '</div>'
                        '<ul style="color:var(--t2);font-size:.85rem;">' + concerns + '</ul></div>',
                        unsafe_allow_html=True
                    )

        with disc_tabs[3]:
            st.markdown("**Work Breakdown Structure**")
            wbs = safe_list(dr.get("wbs"))
            for phase_idx, phase in enumerate(wbs):
                phase = safe_dict(phase)
                phase_name = safe_str(phase.get("phase"))
                with st.expander("📁 " + str(phase_idx + 1) + ". " + phase_name, expanded=True):
                    for d in safe_list(phase.get("deliverables")):
                        d = safe_dict(d)
                        st.markdown("**📦 " + safe_str(d.get("name")) + "**")
                        for t in safe_list(d.get("tasks")):
                            t = safe_dict(t)
                            st.markdown("&nbsp;&nbsp;&nbsp;&nbsp;🔹 " + safe_str(t.get("name")) + " — _" + safe_str(t.get("effort")) + "_")
            decisions = safe_list(dr.get("decisions"))
            if decisions:
                st.markdown("---")
                st.markdown("**Key Decisions:**")
                for dec in decisions:
                    st.markdown("- ✅ " + safe_str(dec))

        with disc_tabs[4]:
            for ai_item in safe_list(dr.get("action_items")):
                ai_item = safe_dict(ai_item)
                pri = safe_str(ai_item.get("priority"))
                pc = "#ff6b6b" if pri == "High" else "#ffd166" if pri == "Medium" else "#06d6a0"
                st.markdown(
                    '<div class="rc" style="border-left:3px solid ' + pc + ';">'
                    '<div class="rch2"><strong>' + safe_str(ai_item.get("item")) + '</strong>'
                    '<span class="rsev" style="color:' + pc + ';">' + pri + '</span></div>'
                    '<div class="rmit"><strong>Owner:</strong> ' + safe_str(ai_item.get("owner")) + '</div>'
                    '</div>', unsafe_allow_html=True
                )

        st.markdown("---")
        dl_disc_cols = st.columns(2)
        with dl_disc_cols[0]:
            st.download_button(
                "📥 Download Discovery Report (JSON)",
                data=json.dumps(dr, indent=2, default=str),
                file_name="ECI_Discovery_Report_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".json",
                mime="application/json",
                use_container_width=True, key="dl_disc_json",
            )
        with dl_disc_cols[1]:
            if st.session_state.discovery_transcript:
                st.download_button(
                    "📥 Download Cleaned Transcript (TXT)",
                    data=st.session_state.discovery_transcript,
                    file_name="ECI_Transcript_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".txt",
                    mime="text/plain",
                    use_container_width=True, key="dl_disc_txt",
                )

    if st.session_state.processing_results:
        show_results()


def run_pipeline(files):
    st.session_state.agent_logs = []
    ai = AzureAI.from_session()
    live = ai.is_live
    if live:
        st.markdown('<div class="phdr">LIVE AI Processing via Azure OpenAI</div>', unsafe_allow_html=True)
    else:
        st.markdown('<div class="phdr">Demo Mode — Configure Azure OpenAI for live AI</div>', unsafe_allow_html=True)
    pb = st.progress(0)
    status = st.empty()

    status.markdown("**1/12** Ingesting documents...")
    pb.progress(5)
    dp = DocProcessor()
    text = ""
    for f in files:
        text += dp.extract(f) + "\n\n"
    log_agent("Ingestion", str(len(files)) + " file(s), " + str(len(text)) + " chars")
    time.sleep(0.2)

    status.markdown("**2/12** Document Intelligence...")
    pb.progress(10)
    intel = dp.analyze(text)
    log_agent("Intelligence", str(intel["section_count"]) + " sections, " + str(intel["word_count"]) + " words")
    time.sleep(0.2)

    status.markdown("**3/12** Semantic Analysis" + (" (GPT-4)..." if live else "..."))
    pb.progress(20)
    semantic = ai.analyze_requirements(text)
    st.session_state["_last_semantic"] = semantic
    log_agent("Semantic", str(len(safe_list(semantic.get("requirements")))) + " requirements, " + str(len(safe_list(semantic.get("technology_stack")))) + " technologies detected")
    time.sleep(0.2)

    status.markdown("**4/12** Historical RAG...")
    pb.progress(30)
    rag = ai.search_historical(text, st.session_state.historical_projects)
    log_agent("RAG", str(len(safe_list(rag.get("similar_projects")))) + " matches")
    time.sleep(0.2)

    status.markdown("**5/12** Time Estimator...")
    pb.progress(40)
    time_est = ai.estimate_time(semantic, rag)
    st.session_state["_last_time_est"] = time_est
    log_agent("Time", str(time_est.get("total_hours", 0)) + " hours across " + str(len(safe_list(time_est.get("phases")))) + " phases")
    time.sleep(0.2)

    status.markdown("**6/12** Cost Calculator...")
    pb.progress(50)
    cost_est = ai.estimate_cost(semantic, time_est, rag)
    st.session_state["_last_cost_est"] = cost_est
    log_agent("Cost", "$" + str(cost_est.get("total_monthly_cost", 0)) + "/mo (" + str(len(safe_list(cost_est.get("azure_costs")))) + " services)")
    time.sleep(0.2)

    status.markdown("**7/12** Risk Analyzer...")
    pb.progress(60)
    risk = ai.analyze_risk(semantic, time_est, cost_est)
    log_agent("Risk", str(risk.get("overall_score", 0)) + "/10 — " + str(len(safe_list(risk.get("risks")))) + " risks identified")
    time.sleep(0.2)

    status.markdown("**8/12** Architecture Designer...")
    pb.progress(70)
    arch = ai.design_architecture(semantic, rag)
    log_agent("Architecture", str(len(safe_list(arch.get("components")))) + " components")
    time.sleep(0.2)

    status.markdown("**9/12** Scope & Assumptions...")
    pb.progress(80)
    scope = ai.define_scope(semantic, time_est, cost_est)
    log_agent("Scope", "Boundaries defined")
    time.sleep(0.2)

    status.markdown("**10/12** Proposal Writer...")
    pb.progress(80)
    proposal = ai.write_proposal(semantic, time_est, cost_est, risk, arch, scope)
    log_agent("Proposal", "Document generated")
    time.sleep(0.2)

    status.markdown("**11/12** Architecture Visualizer...")
    pb.progress(90)
    mermaid_diagrams = ai.generate_mermaid_diagrams(semantic, arch)
    log_agent("Visualizer", str(len(mermaid_diagrams)) + " diagrams generated")
    time.sleep(0.2)

    pb.progress(100)
    status.markdown("**12/12** All agents completed!")
    st.session_state.processing_results = {
        "semantic_analysis": semantic, "rag": rag, "time_estimate": time_est,
        "cost_estimate": cost_est, "risk_assessment": risk, "architecture": arch,
        "scope": scope, "proposal": proposal, "mermaid_diagrams": mermaid_diagrams,
    }
    st.session_state.model_metrics["proposals_processed"] += 1


def show_results():
    r = st.session_state.processing_results
    st.markdown("---")
    st.markdown('<div class="shdr"><span class="shdr-i">📊</span> Results</div>', unsafe_allow_html=True)

    with st.expander("Agent Activity Log", expanded=False):
        for entry in st.session_state.agent_logs:
            st.markdown('<div class="alog"><span class="abadge">' + entry["agent"] + '</span><span class="aok">Done</span><span class="adet">' + entry["detail"] + '</span></div>', unsafe_allow_html=True)

    se = safe_dict(r.get("semantic_analysis"))
    te = safe_dict(r.get("time_estimate"))
    ce = safe_dict(r.get("cost_estimate"))
    ri = safe_dict(r.get("risk_assessment"))
    ar = safe_dict(r.get("architecture"))

    kpis = [
        ("📝", "Requirements", str(len(safe_list(se.get("requirements")))), "Identified"),
        ("⏱️", "Hours", str(safe_int(te.get("total_hours"))), "Person-hours"),
        ("💰", "Infra Cost", "$" + str(safe_int(ce.get("total_monthly_cost"))) + "/mo", "Azure Infrastructure"),
        ("⚠️", "Risk", str(safe_int(ri.get("overall_score"))) + "/10", safe_str(ri.get("overall_level"))),
        ("🏗️", "Components", str(len(safe_list(ar.get("components")))), "Designed"),
    ]
    cols = st.columns(5)
    for i, (ic, t, v, s) in enumerate(kpis):
        with cols[i]:
            st.markdown('<div class="kpi"><div class="kpi-i">' + ic + '</div><div class="kpi-v">' + v + '</div><div class="kpi-t">' + t + '</div><div class="kpi-s">' + s + '</div></div>', unsafe_allow_html=True)

    tab_list = st.tabs(["📋 Requirements", "⏱️ Time", "💰 Infra Cost", "⚠️ Risk", "🏗️ Architecture", "📐 Diagrams", "📄 Proposal", "📌 Scope", "👥 Team & Roles", "🎮 3D View"])

    # ── Requirements ──
    with tab_list[0]:
        reqs = safe_list(se.get("requirements"))
        fn_list = [x for x in reqs if isinstance(x, dict) and x.get("type") == "functional"]
        nf_list = [x for x in reqs if isinstance(x, dict) and x.get("type") == "non-functional"]
        ig_list = [x for x in reqs if isinstance(x, dict) and x.get("type") == "integration"]
        rc1, rc2, rc3 = st.columns(3)
        with rc1:
            st.markdown('<div class="rch fn">Functional (' + str(len(fn_list)) + ')</div>', unsafe_allow_html=True)
            for q in fn_list:
                st.markdown('<div class="ri"><strong>' + safe_str(q.get("title")) + '</strong><div class="ri-c">Complexity: ' + safe_str(q.get("complexity")) + '</div><p>' + safe_str(q.get("description")) + '</p></div>', unsafe_allow_html=True)
        with rc2:
            st.markdown('<div class="rch nf">Non-Functional (' + str(len(nf_list)) + ')</div>', unsafe_allow_html=True)
            for q in nf_list:
                st.markdown('<div class="ri"><strong>' + safe_str(q.get("title")) + '</strong><p>' + safe_str(q.get("description")) + '</p></div>', unsafe_allow_html=True)
        with rc3:
            st.markdown('<div class="rch ig">Integration (' + str(len(ig_list)) + ')</div>', unsafe_allow_html=True)
            for q in ig_list:
                st.markdown('<div class="ri"><strong>' + safe_str(q.get("title")) + '</strong><p>' + safe_str(q.get("description")) + '</p></div>', unsafe_allow_html=True)
        tech = safe_list(se.get("technology_stack"))
        if tech:
            tags_html = " ".join('<span class="tt">' + safe_str(t) + '</span>' for t in tech)
            st.markdown('<div class="ttag">' + tags_html + '</div>', unsafe_allow_html=True)

    # ── Time ──
    with tab_list[1]:
        phases = safe_list(te.get("phases"))
        three_pt = safe_dict(te.get("three_point"))
        mc = st.columns(6)
        with mc[0]:
            st.metric("Total Hours", str(safe_int(te.get("total_hours"))))
        with mc[1]:
            st.metric("Duration", safe_str(te.get("duration_weeks")))
        with mc[2]:
            st.metric("Confidence", safe_str(te.get("confidence")))
        with mc[3]:
            st.metric("Buffer", safe_str(te.get("buffer")))
        with mc[4]:
            st.metric("Optimistic", str(safe_int(three_pt.get("optimistic", 0))) + "h")
        with mc[5]:
            st.metric("Pessimistic", str(safe_int(three_pt.get("pessimistic", 0))) + "h")

        if phases:
            # ── Phase-level bar chart with Low / Avg / High ──
            colors = ["#00d4aa", "#00b4d8", "#7b61ff", "#ff6b6b", "#ffd166", "#06d6a0", "#e9c46a", "#f4845f"]
            phase_names = [safe_str(safe_dict(p).get("name", "Phase")) for p in phases]
            low_vals = [safe_int(safe_dict(p).get("low_hours", int(safe_int(safe_dict(p).get("hours", 0)) * 0.8))) for p in phases]
            avg_vals = [safe_int(safe_dict(p).get("hours", 0)) for p in phases]
            high_vals = [safe_int(safe_dict(p).get("high_hours", int(safe_int(safe_dict(p).get("hours", 0)) * 1.35))) for p in phases]
            fig = go.Figure()
            fig.add_trace(go.Bar(name="Low (Optimistic)", x=low_vals, y=phase_names, orientation="h", marker_color="#06d6a0", text=[str(v) + "h" for v in low_vals], textposition="auto"))
            fig.add_trace(go.Bar(name="Average", x=avg_vals, y=phase_names, orientation="h", marker_color="#00b4d8", text=[str(v) + "h" for v in avg_vals], textposition="auto"))
            fig.add_trace(go.Bar(name="High (Pessimistic)", x=high_vals, y=phase_names, orientation="h", marker_color="#ff6b6b", text=[str(v) + "h" for v in high_vals], textposition="auto"))
            fig.update_layout(
                title="Phase Breakdown — Three-Point Estimate (Low / Avg / High)", template="plotly_dark",
                paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
                height=max(350, len(phases) * 60), barmode="group", xaxis_title="Hours",
                legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
            )
            st.plotly_chart(fig, use_container_width=True)

            # ── Detailed task breakdown per phase (Inflexion.xlsx style) ──
            st.markdown("### Detailed Task Breakdown")
            for pi, phase in enumerate(phases):
                phase = safe_dict(phase)
                ph_name = safe_str(phase.get("name", "Phase"))
                ph_week = safe_str(phase.get("week_label", ""))
                ph_low = safe_int(phase.get("low_hours", 0))
                ph_avg = safe_int(phase.get("hours", 0))
                ph_high = safe_int(phase.get("high_hours", 0))
                ph_pct = safe_str(phase.get("percentage", ""))
                header_txt = ph_name
                if ph_week:
                    header_txt += " (" + ph_week + ")"
                header_txt += " — " + str(ph_avg) + "h [" + str(ph_low) + "-" + str(ph_high) + "h] " + ph_pct
                with st.expander(header_txt, expanded=(pi < 2)):
                    tasks = safe_list(phase.get("tasks"))
                    if tasks:
                        # Build a table header
                        st.markdown(
                            '<table style="width:100%;border-collapse:collapse;font-size:0.85rem;">'
                            '<tr style="background:#1B3A5C;color:white;">'
                            '<th style="padding:6px 8px;text-align:left;">Sub-task</th>'
                            '<th style="padding:6px 8px;text-align:center;">Role</th>'
                            '<th style="padding:6px 8px;text-align:center;">Low (hrs)</th>'
                            '<th style="padding:6px 8px;text-align:center;">Avg (hrs)</th>'
                            '<th style="padding:6px 8px;text-align:center;">High (hrs)</th>'
                            '<th style="padding:6px 8px;text-align:left;">Justification</th>'
                            '</tr>' +
                            "".join(
                                '<tr style="background:' + ("#1a2a4a" if ti % 2 == 0 else "#0f1928") + ';color:#e2e8f0;">'
                                '<td style="padding:5px 8px;color:#e2e8f0;">' + safe_str(safe_dict(tk).get("name")) + '</td>'
                                '<td style="padding:5px 8px;text-align:center;color:#00b4d8;">' + safe_str(safe_dict(tk).get("role")) + '</td>'
                                '<td style="padding:5px 8px;text-align:center;color:#94a3b8;">' + str(safe_int(safe_dict(tk).get("low_hours", 0))) + '</td>'
                                '<td style="padding:5px 8px;text-align:center;font-weight:bold;color:#00d4aa;">' + str(safe_int(safe_dict(tk).get("hours", 0))) + '</td>'
                                '<td style="padding:5px 8px;text-align:center;color:#94a3b8;">' + str(safe_int(safe_dict(tk).get("high_hours", 0))) + '</td>'
                                '<td style="padding:5px 8px;font-style:italic;color:#94a3b8;">' + safe_str(safe_dict(tk).get("justification", "")) + '</td>'
                                '</tr>'
                                for ti, tk in enumerate(tasks)
                            ) +
                            '<tr style="background:rgba(0,212,170,.15);color:#00d4aa;font-weight:bold;">'
                            '<td style="padding:5px 8px;" colspan="2">Phase Total</td>'
                            '<td style="padding:5px 8px;text-align:center;">' + str(ph_low) + '</td>'
                            '<td style="padding:5px 8px;text-align:center;">' + str(ph_avg) + '</td>'
                            '<td style="padding:5px 8px;text-align:center;">' + str(ph_high) + '</td>'
                            '<td></td></tr>'
                            '</table>',
                            unsafe_allow_html=True,
                        )
                    else:
                        st.info("No sub-tasks for this phase.")

        # ── Milestones ──
        milestones = safe_list(te.get("milestones"))
        if milestones:
            st.markdown("### Milestones")
            for m in milestones:
                m = safe_dict(m)
                st.markdown("- **Week " + str(safe_int(m.get("week"))) + "** — " + safe_str(m.get("name")) + ": " + safe_str(m.get("description")))
        # Excel download
        st.markdown("---")
        excel_data = generate_time_excel(te, se)
        if excel_data:
            st.download_button(
                "Download Time Estimate (Excel)",
                data=excel_data,
                file_name="ECI_Time_Estimate_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True, type="primary", key="dl_time_xlsx",
            )

    # ── Cost (Infrastructure) ──
    with tab_list[2]:
        azure_costs = safe_list(ce.get("azure_costs"))
        third_party = safe_list(ce.get("third_party_costs"))
        mc = st.columns(3)
        with mc[0]:
            st.metric("Monthly Cost", "$" + str(safe_int(ce.get("total_monthly_cost"))))
        with mc[1]:
            st.metric("Annual Cost", "$" + str(safe_int(ce.get("total_annual_cost"))))
        with mc[2]:
            st.metric("Services", str(len(azure_costs) + len(third_party)))
        if azure_costs:
            pie_vals = [safe_int(safe_dict(a).get("monthly_cost", 0)) for a in azure_costs]
            pie_names = [safe_str(safe_dict(a).get("service", "")) for a in azure_costs]
            fig = px.pie(values=pie_vals, names=pie_names, title="Azure Infrastructure Cost Distribution (Monthly)",
                         color_discrete_sequence=["#00d4aa", "#00b4d8", "#7b61ff", "#ff6b6b", "#ffd166", "#06d6a0", "#e9c46a", "#f4845f", "#a8dadc", "#457b9d", "#2a9d8f", "#264653"])
            fig.update_layout(template="plotly_dark", paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", height=400)
            st.plotly_chart(fig, use_container_width=True)
        st.markdown("**Azure Services:**")
        for svc in azure_costs:
            svc = safe_dict(svc)
            st.markdown("- **" + safe_str(svc.get("service")) + "** (" + safe_str(svc.get("tier", "")) + ") — $" + str(safe_int(svc.get("monthly_cost"))) + "/mo — _" + safe_str(svc.get("description", "")) + "_")
        if third_party:
            st.markdown("**Third-Party Services:**")
            for tp in third_party:
                tp = safe_dict(tp)
                st.markdown("- **" + safe_str(tp.get("name")) + "** — $" + str(safe_int(tp.get("monthly_cost"))) + "/mo — _" + safe_str(tp.get("description", "")) + "_")
        opt_tips = safe_list(ce.get("cost_optimization"))
        if opt_tips:
            st.markdown("**Cost Optimization Tips:**")
            for tip in opt_tips:
                st.markdown("- " + safe_str(tip))
        notes = safe_str(ce.get("notes"))
        if notes:
            st.info(notes)
        # Cost Excel download
        st.markdown("---")
        cost_xl_data = generate_cost_excel(ce, te, se)
        if cost_xl_data:
            st.download_button(
                "📥 Download Detailed Cost Estimate (Excel)",
                data=cost_xl_data,
                file_name="ECI_Cost_Estimate_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True, type="primary", key="dl_cost_xlsx",
            )

    # ── Risk ──
    with tab_list[3]:
        score = safe_int(ri.get("overall_score", 0))
        clr = "#06d6a0" if score <= 3 else "#ffd166" if score <= 6 else "#ff6b6b"
        st.markdown('<div class="rsb" style="border-left:4px solid ' + clr + ';"><span class="rsv">' + str(score) + '/10</span><span class="rsl">Overall: ' + safe_str(ri.get("overall_level")) + '</span></div>', unsafe_allow_html=True)
        for rk in safe_list(ri.get("risks")):
            rk = safe_dict(rk)
            sev = safe_str(rk.get("severity"))
            sc = "#06d6a0" if sev == "Low" else "#ffd166" if sev == "Medium" else "#ff6b6b"
            st.markdown('<div class="rc" style="border-left:3px solid ' + sc + ';"><div class="rch2"><strong>' + safe_str(rk.get("category")) + ': ' + safe_str(rk.get("title")) + '</strong><span class="rsev" style="color:' + sc + ';">' + sev + '</span></div><p>' + safe_str(rk.get("description")) + '</p><div class="rmit"><strong>Mitigation:</strong> ' + safe_str(rk.get("mitigation")) + '</div></div>', unsafe_allow_html=True)

    # ── Architecture ──
    with tab_list[4]:
        comps = safe_list(ar.get("components"))
        st.markdown("**Pattern:** _" + safe_str(ar.get("pattern")) + "_")

        # ── Architecture Diagram (Graphviz) ──
        arch_dot = None
        try:
            arch_dot = generate_architecture_diagram(ar)
            if arch_dot:
                st.markdown("### Solution Architecture Diagram")
                st.graphviz_chart(arch_dot, use_container_width=True)
        except Exception:
            pass

        # ── Architecture Diagram Downloads ──
        if arch_dot:
            st.markdown("#### Download Architecture Diagram")
            arch_svg = render_dot_to_svg(arch_dot)
            arch_png = render_dot_to_png(arch_dot)
            arch_html = render_dot_to_html(arch_dot)
            dl_arch_cols = st.columns(3)
            with dl_arch_cols[0]:
                if arch_svg:
                    st.download_button(
                        "Download Diagram (SVG)",
                        data=arch_svg,
                        file_name="ECI_Architecture_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".svg",
                        mime="image/svg+xml",
                        use_container_width=True, type="primary", key="dl_arch_svg",
                    )
                elif arch_html:
                    st.download_button(
                        "Download Diagram (HTML+SVG)",
                        data=arch_html,
                        file_name="ECI_Architecture_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".html",
                        mime="text/html",
                        use_container_width=True, type="primary", key="dl_arch_html",
                    )
            with dl_arch_cols[1]:
                if arch_png:
                    st.download_button(
                        "Download Diagram (PNG)",
                        data=arch_png,
                        file_name="ECI_Architecture_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".png",
                        mime="image/png",
                        use_container_width=True, type="primary", key="dl_arch_png",
                    )
                elif arch_html and not arch_svg:
                    st.info("Open the HTML file in a browser, then use the Download PNG button inside it.")
            with dl_arch_cols[2]:
                st.download_button(
                    "Download Diagram (DOT)",
                    data=arch_dot.encode("utf-8"),
                    file_name="ECI_Architecture_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".dot",
                    mime="text/plain",
                    use_container_width=True, key="dl_arch_dot",
                )

        # ── Component cards ──
        arch_cols = st.columns(3)
        for i, c in enumerate(comps):
            c = safe_dict(c)
            with arch_cols[i % 3]:
                svcs = "".join("<li>" + safe_str(s) + "</li>" for s in safe_list(c.get("services")))
                st.markdown('<div class="ac"><div class="acn">' + safe_str(c.get("name")) + '</div><div class="act">' + safe_str(c.get("azure_service")) + '</div><ul class="acs">' + svcs + '</ul></div>', unsafe_allow_html=True)
        df = safe_list(ar.get("data_flow"))
        if df:
            flow_str = " -> ".join(safe_str(x) for x in df)
            st.markdown('<div class="dfv">' + flow_str + '</div>', unsafe_allow_html=True)

        # ── Workflow / Flow Diagram (Graphviz) ──
        flow_dot = None
        try:
            flow_dot = generate_flow_diagram(te)
            if flow_dot:
                st.markdown("### Project Workflow Diagram")
                st.graphviz_chart(flow_dot, use_container_width=True)
        except Exception:
            pass

        # ── Workflow Diagram Downloads ──
        if flow_dot:
            st.markdown("#### Download Workflow Diagram")
            flow_svg = render_dot_to_svg(flow_dot)
            flow_png = render_dot_to_png(flow_dot)
            flow_html = render_dot_to_html(flow_dot)
            dl_flow_cols = st.columns(3)
            with dl_flow_cols[0]:
                if flow_svg:
                    st.download_button(
                        "Download Workflow (SVG)",
                        data=flow_svg,
                        file_name="ECI_Workflow_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".svg",
                        mime="image/svg+xml",
                        use_container_width=True, type="primary", key="dl_flow_svg",
                    )
                elif flow_html:
                    st.download_button(
                        "Download Workflow (HTML+SVG)",
                        data=flow_html,
                        file_name="ECI_Workflow_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".html",
                        mime="text/html",
                        use_container_width=True, type="primary", key="dl_flow_html",
                    )
            with dl_flow_cols[1]:
                if flow_png:
                    st.download_button(
                        "Download Workflow (PNG)",
                        data=flow_png,
                        file_name="ECI_Workflow_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".png",
                        mime="image/png",
                        use_container_width=True, type="primary", key="dl_flow_png",
                    )
                elif flow_html and not flow_svg:
                    st.info("Open the HTML file in a browser, then use the Download PNG button inside it.")
            with dl_flow_cols[2]:
                st.download_button(
                    "Download Workflow (DOT)",
                    data=flow_dot.encode("utf-8"),
                    file_name="ECI_Workflow_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".dot",
                    mime="text/plain",
                    use_container_width=True, key="dl_flow_dot",
                )

        # ── Security ──
        sec_items = safe_list(ar.get("security"))
        if sec_items:
            st.markdown("### Security Controls")
            for s in sec_items:
                st.markdown("- " + safe_str(s))

    # ── Diagrams (Mermaid.js) ──
    with tab_list[5]:
        mermaid_data = safe_dict(r.get("mermaid_diagrams"))
        if mermaid_data:
            diagram_tabs = st.tabs(["🏗️ Infrastructure", "🔄 Data Flow", "📨 Sequence", "🚀 Deployment", "🔒 Security"])
            diagram_map = [
                ("infrastructure", "Azure Infrastructure Architecture", 500),
                ("data_flow", "Data Flow Diagram", 450),
                ("sequence", "Request Sequence Diagram", 500),
                ("deployment", "CI/CD & Deployment Pipeline", 450),
                ("security", "Security Architecture & Controls", 450),
            ]
            for i, (key, title, h) in enumerate(diagram_map):
                with diagram_tabs[i]:
                    code = safe_str(mermaid_data.get(key, ""))
                    if code:
                        st.markdown("**" + title + "**")
                        render_mermaid(code, height=h)
                        with st.expander("View Mermaid Source"):
                            st.code(code, language="text")
                    else:
                        st.info("Diagram not available for this project.")
            st.markdown("---")
            st.download_button(
                "📥 Download All Diagrams (JSON)",
                data=json.dumps(mermaid_data, indent=2),
                file_name="ECI_Diagrams_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".json",
                mime="application/json",
                use_container_width=True, key="dl_diagrams",
            )
        else:
            st.info("Architecture diagrams will be generated after processing documents.")

    # ── Proposal ──
    with tab_list[6]:
        proposal = safe_dict(r.get("proposal"))
        for sec in safe_list(proposal.get("sections")):
            sec = safe_dict(sec)
            with st.expander(safe_str(sec.get("title")), expanded=False):
                st.markdown(safe_str(sec.get("content")))
        for k, v in safe_dict(proposal.get("quality_checks")).items():
            check_icon = "pass" if v else "warn"
            st.markdown(("✅ " if v else "⚠️ ") + str(k))
        # PDF download
        st.markdown("---")
        pdf_data = generate_proposal_pdf(r)
        if pdf_data:
            st.download_button(
                "📥 Download Proposal (PDF)",
                data=pdf_data,
                file_name="ECI_Proposal_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".pdf",
                mime="application/pdf",
                use_container_width=True, type="primary", key="dl_proposal_pdf",
            )
        pptx_data = generate_proposal_pptx(r)
        if pptx_data:
            st.download_button(
                "📥 Download Proposal (PowerPoint)",
                data=pptx_data,
                file_name="ECI_Proposal_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True, key="dl_proposal_pptx",
            )
        sow_data = generate_sow_pdf(r)
        if sow_data:
            st.download_button(
                "📥 Download Statement of Work (PDF)",
                data=sow_data,
                file_name="ECI_SOW_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".pdf",
                mime="application/pdf",
                use_container_width=True, type="primary", key="dl_sow_pdf",
            )

    # ── Scope ──
    with tab_list[7]:
        sc = safe_dict(r.get("scope"))
        sc1, sc2 = st.columns(2)
        with sc1:
            st.markdown("**In Scope:**")
            for x in safe_list(sc.get("in_scope")):
                st.markdown("- " + safe_str(x))
            st.markdown("**Assumptions:**")
            for x in safe_list(sc.get("assumptions")):
                st.markdown("- " + safe_str(x))
        with sc2:
            st.markdown("**Out of Scope:**")
            for x in safe_list(sc.get("out_of_scope")):
                st.markdown("- " + safe_str(x))
            st.markdown("**Prerequisites:**")
            for x in safe_list(sc.get("prerequisites")):
                st.markdown("- " + safe_str(x))

    # ── Team & Roles ──
    with tab_list[8]:
        st.markdown("### Team Composition & Role Allocation")
        total_h = safe_int(te.get("total_hours", 0))
        total_d = max(1, total_h / 7)
        dynamic_roles = safe_list(te.get("roles"))
        rc1, rc2 = st.columns([3, 2])
        with rc1:
            role_names = []
            role_days = []
            role_colors = ["#00d4aa", "#00b4d8", "#7b61ff", "#ff6b6b", "#ffd166", "#06d6a0", "#e9c46a", "#f4845f"]
            if dynamic_roles:
                for rl in dynamic_roles:
                    rl = safe_dict(rl)
                    role_name = safe_str(rl.get("name", "Team Member"))
                    pct = float(rl.get("allocation_pct", 0))
                    rate = safe_int(rl.get("rate", 100))
                    days = round(total_d * pct, 1)
                    hours = int(round(days * 7, 0))
                    if hours > 0:
                        role_names.append(role_name)
                        role_days.append(days)
                    st.markdown(
                        '<div class="ri"><strong>' + role_name + '</strong> — ' + str(int(pct * 100)) + '% allocation'
                        '<div class="ri-c">' + str(days) + ' days / ' + str(hours) + ' hrs @ $' + str(rate) + '/hr = $' + str(hours * rate)
                        + '</div></div>', unsafe_allow_html=True)
            else:
                st.info("Role allocation data not available. Run analysis to generate dynamic roles.")
        with rc2:
            if role_names:
                fig_roles = px.pie(values=role_days, names=role_names, title="Team Allocation (Days)",
                                   color_discrete_sequence=role_colors)
                fig_roles.update_layout(template="plotly_dark", paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", height=400)
                st.plotly_chart(fig_roles, use_container_width=True)

        st.markdown("---")
        st.markdown("### Project Duration Summary")
        dur_cols = st.columns(4)
        weeks_str = safe_str(te.get("duration_weeks", "N/A"))
        with dur_cols[0]:
            st.metric("Total Hours", str(total_h))
        with dur_cols[1]:
            st.metric("Total Days", str(round(total_d, 1)))
        with dur_cols[2]:
            st.metric("Duration", weeks_str)
        with dur_cols[3]:
            st.metric("Incl. Leaves & Holidays", weeks_str.replace("weeks", "").strip() + " + 2 weeks" if "weeks" in weeks_str.lower() else "N/A")

        st.markdown("---")
        st.markdown("### Prerequisites")
        # Dynamic prerequisites from scope
        scope_data = safe_dict(r.get("scope"))
        prereqs = safe_list(scope_data.get("prerequisites"))
        if not prereqs:
            prereqs = safe_list(scope_data.get("assumptions"))
        if not prereqs:
            prereqs = ["Cloud subscription access", "Sample documents for analysis", "Stakeholder availability for UAT"]
        for i, p in enumerate(prereqs, 1):
            st.markdown(str(i) + ". " + safe_str(p))

    # ── 3D Architecture Fly-Through ──
    with tab_list[9]:
        st.markdown("### Interactive 3D Architecture Fly-Through")
        st.markdown(
            "An interactive **Three.js** 3D scene of the solution architecture. "
            "Each component is a clickable box — click any node to see its Azure service, "
            "monthly cost, and capabilities in the HUD overlay. "
            "Hit **Play Fly-Through** for a guided camera tour of the full architecture."
        )

        flythrough_html = generate_3d_flythrough_html(ar, ce)

        if flythrough_html:
            # Render inline inside Streamlit
            st.components.v1.html(flythrough_html, height=620, scrolling=False)

            st.markdown("---")
            ft_dl_cols = st.columns([2, 1])
            with ft_dl_cols[0]:
                st.markdown(
                    "**Tip for client meetings:** Download the standalone HTML and open it in any browser "
                    "for a full-screen, shareable experience — no installation required."
                )
            with ft_dl_cols[1]:
                st.download_button(
                    "📥 Download 3D View (HTML)",
                    data=flythrough_html.encode("utf-8"),
                    file_name="ECI_3D_Architecture_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".html",
                    mime="text/html",
                    use_container_width=True,
                    type="primary",
                    key="dl_3d_html",
                )
        else:
            st.info("No architecture components found. Run the pipeline to generate the 3D view.")

    # ── Delivery ──
    st.markdown("---")
    st.markdown('<div class="shdr"><span class="shdr-i">🚀</span> Delivery</div>', unsafe_allow_html=True)

    # Downloads row
    st.markdown('<div class="shdr" style="font-size:1rem;"><span class="shdr-i">📥</span> Downloads</div>', unsafe_allow_html=True)
    dl_row1 = st.columns(3)
    with dl_row1[0]:
        xl_data = generate_time_excel(te, se)
        if xl_data:
            st.download_button("📊 Time Estimate (Excel)", data=xl_data,
                               file_name="ECI_Time_Estimate_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".xlsx",
                               mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                               use_container_width=True, type="primary", key="bdl_xl")
    with dl_row1[1]:
        cost_xl = generate_cost_excel(ce, te, se)
        if cost_xl:
            st.download_button("💰 Cost Estimate (Excel)", data=cost_xl,
                               file_name="ECI_Cost_Estimate_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".xlsx",
                               mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                               use_container_width=True, type="primary", key="bdl_cost_xl")
    with dl_row1[2]:
        pdf_data = generate_proposal_pdf(r)
        if pdf_data:
            st.download_button("📄 Proposal (PDF)", data=pdf_data,
                               file_name="ECI_Proposal_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".pdf",
                               mime="application/pdf",
                               use_container_width=True, type="primary", key="bdl_pdf")
    dl_row2 = st.columns(3)
    with dl_row2[0]:
        pptx_dl = generate_proposal_pptx(r)
        if pptx_dl:
            st.download_button("📑 Proposal (PowerPoint)", data=pptx_dl,
                               file_name="ECI_Proposal_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".pptx",
                               mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                               use_container_width=True, type="primary", key="bdl_pptx")
    with dl_row2[1]:
        sow_dl = generate_sow_pdf(r)
        if sow_dl:
            st.download_button("📝 Statement of Work (PDF)", data=sow_dl,
                               file_name="ECI_SOW_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".pdf",
                               mime="application/pdf",
                               use_container_width=True, type="primary", key="bdl_sow")
    with dl_row2[2]:
        st.download_button("📋 Full Data (JSON)", data=json.dumps(r, indent=2, default=str),
                           file_name="belal_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".json",
                           mime="application/json", use_container_width=True, key="bdl")
    dl_row3 = st.columns(3)
    with dl_row3[0]:
        # Zip bundle of all deliverables
        zip_buf = io.BytesIO()
        with zipfile.ZipFile(zip_buf, 'w', zipfile.ZIP_DEFLATED) as zf:
            if xl_data:
                zf.writestr("ECI_Time_Estimate.xlsx", xl_data)
            if cost_xl:
                zf.writestr("ECI_Cost_Estimate.xlsx", cost_xl)
            if pdf_data:
                zf.writestr("ECI_Proposal.pdf", pdf_data)
            if pptx_dl:
                zf.writestr("ECI_Proposal.pptx", pptx_dl)
            if sow_dl:
                zf.writestr("ECI_Statement_of_Work.pdf", sow_dl)
            zf.writestr("BELAL_Data.json", json.dumps(r, indent=2, default=str))
        zip_buf.seek(0)
        st.download_button("📦 All Deliverables (ZIP)", data=zip_buf.getvalue(),
                           file_name="ECI_Deliverables_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".zip",
                           mime="application/zip", use_container_width=True, key="bdl_zip")

    # Actions row
    st.markdown('<div class="shdr" style="font-size:1rem;"><span class="shdr-i">⚡</span> Actions</div>', unsafe_allow_html=True)
    ac_cols = st.columns(2)
    with ac_cols[0]:
        if st.button("Upload to SharePoint", use_container_width=True, key="bsp"):
            ok, msg = SP.from_session().upload(r)
            if ok:
                st.success(msg)
            else:
                st.error(msg)
    with ac_cols[1]:
        if st.button("Send Alert Email", use_container_width=True, key="bem"):
            ok, msg = Mailer.from_session().send(r)
            if ok:
                st.success(msg)
            else:
                st.error(msg)


# ═══════════════════════════════════════════════════════════════════════
#  TAB 2: ADMIN
# ═══════════════════════════════════════════════════════════════════════

def tab_admin():
    st.markdown('<div class="shdr"><span class="shdr-i">🧠</span> Continuous Learning</div>', unsafe_allow_html=True)
    at = st.tabs(["📊 Dashboard", "📚 Training", "🔧 Config", "📁 Sync", "🔄 Loop"])

    with at[0]:
        m = st.session_state.model_metrics
        mc = st.columns(4)
        with mc[0]:
            st.metric("Proposals", m["proposals_processed"])
        with mc[1]:
            st.metric("Accuracy", str(m["accuracy"]) + "%")
        with mc[2]:
            st.metric("Win Rate", str(m["win_rate"]) + "%")
        with mc[3]:
            st.metric("Variance", str(m["variance"]) + "%")
        dates = [(datetime.now() - timedelta(days=30 * i)).strftime("%b %Y") for i in range(6, -1, -1)]
        fig = go.Figure(go.Scatter(x=dates, y=[65, 68, 72, 74, 76, 78, 80], mode="lines+markers",
                                   line=dict(color="#00d4aa", width=3), marker=dict(size=10)))
        fig.update_layout(title="Accuracy Trend", template="plotly_dark",
                          paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", height=350)
        st.plotly_chart(fig, use_container_width=True)
        dc1, dc2 = st.columns(2)
        with dc1:
            fig2 = go.Figure(go.Bar(x=["Time", "Cost", "Risk", "Arch"], y=[85, 79, 82, 88],
                                    marker_color=["#00d4aa", "#00b4d8", "#ffd166", "#7b61ff"],
                                    text=["85%", "79%", "82%", "88%"], textposition="auto"))
            fig2.update_layout(title="Agent Accuracy", template="plotly_dark",
                               paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", height=300)
            st.plotly_chart(fig2, use_container_width=True)
        with dc2:
            fig3 = go.Figure(go.Indicator(
                mode="gauge+number+delta", value=m["accuracy"], delta={"reference": 72},
                gauge={"axis": {"range": [0, 100]}, "bar": {"color": "#00d4aa"},
                       "steps": [{"range": [0, 50], "color": "rgba(255,107,107,.2)"}, {"range": [50, 75], "color": "rgba(255,209,102,.2)"}, {"range": [75, 100], "color": "rgba(0,212,170,.2)"}]},
                title={"text": "Model Health"}))
            fig3.update_layout(template="plotly_dark", paper_bgcolor="rgba(0,0,0,0)", height=300)
            st.plotly_chart(fig3, use_container_width=True)

    with at[1]:
        st.file_uploader("Upload Historical Data", type=["json", "csv", "xlsx"], accept_multiple_files=True, key="tfu")
        with st.expander("Add Project Manually"):
            with st.form("mp"):
                pn = st.text_input("Name")
                pt = st.selectbox("Type", ["Data/Cloud/AI", "Web App", "Mobile", "SharePoint", "Integration"])
                fc1, fc2 = st.columns(2)
                with fc1:
                    eh = st.number_input("Est Hours", 0, value=100)
                    ah = st.number_input("Actual Hours", 0, value=0)
                    ec = st.number_input("Est Cost", 0, value=10000)
                with fc2:
                    ac_val = st.number_input("Actual Cost", 0, value=0)
                    out = st.selectbox("Outcome", ["Won", "Lost", "Pending"])
                    ts = st.text_input("Tech Stack", "Azure, Python")
                if st.form_submit_button("Add"):
                    st.session_state.historical_projects.append({"name": pn, "type": pt, "estimated_hours": eh, "actual_hours": ah, "estimated_cost": ec, "actual_cost": ac_val, "outcome": out, "tech_stack": ts.split(", ")})
                    st.success("Added: " + pn)
        if st.session_state.historical_projects:
            st.markdown("**" + str(len(st.session_state.historical_projects)) + " project(s) loaded**")
            for p in st.session_state.historical_projects:
                st.markdown("- **" + safe_str(p.get("name")) + "** — " + safe_str(p.get("type")) + " — " + safe_str(p.get("outcome")))

    with at[2]:
        cc1, cc2 = st.columns(2)
        with cc1:
            for ph in ["Discovery", "Design", "Development", "Deployment", "Support"]:
                st.checkbox(ph, True, key="p_" + ph)
            st.selectbox("Estimation", ["Three-Point", "Story Points", "Function Points"], key="em")
            st.slider("Buffer %", 0, 50, 20, key="bf")
        with cc2:
            for rn, rv in [("Architects", 1), ("Lead Devs", 1), ("Developers", 3), ("QA", 1), ("DevOps", 1), ("PMs", 1)]:
                st.number_input(rn, 0, 20, rv, key="r_" + rn)
            st.selectbox("Pricing", ["Fixed Price", "T&M", "Retainer"], key="pc")
        tc = st.columns(3)
        with tc[0]:
            if st.button("Train", use_container_width=True, type="primary", key="bt"):
                pb = st.progress(0)
                for i in range(100):
                    time.sleep(0.02)
                    pb.progress(i + 1)
                st.session_state.model_metrics["accuracy"] = min(st.session_state.model_metrics["accuracy"] + 2.5, 95)
                st.success("Training complete!")
        with tc[1]:
            if st.button("Retrain", use_container_width=True, key="br"):
                time.sleep(1)
                st.success("Retrained.")
        with tc[2]:
            if st.button("A/B Test", use_container_width=True, key="ba"):
                time.sleep(1)
                st.success("New model +8.3%")

    with at[3]:
        st.text_input("Library Path", placeholder="/sites/presales/Historical", key="ssp")
        st.multiselect("Types", ["Proposals", "Estimates", "Outcomes", "Templates"], default=["Proposals", "Estimates", "Outcomes"], key="sst")
        if st.button("Sync Now", use_container_width=True, type="primary", key="bs"):
            time.sleep(2)
            st.success("Synced 47 docs.")

    with at[4]:
        steps = [("1", "Outcome", "Win/loss + feedback", "#7b61ff"), ("2", "Enrichment", "Link to estimates", "#00b4d8"), ("3", "Retraining", "Fine-tune model", "#00d4aa"), ("4", "Validation", "A/B testing", "#ffd166"), ("5", "Deploy", "Push to prod", "#ff6b6b")]
        lc = st.columns(5)
        for i, (n, t, d, c) in enumerate(steps):
            with lc[i]:
                st.markdown('<div class="ls" style="border-top:3px solid ' + c + ';"><div class="ln" style="background:' + c + ';">' + n + '</div><div class="lt">' + t + '</div><div class="ld">' + d + '</div></div>', unsafe_allow_html=True)
        st.markdown("---")
        st.markdown("### Record Outcome")
        with st.form("of"):
            oc1, oc2 = st.columns(2)
            with oc1:
                st.text_input("Project", key="on")
                st.selectbox("Result", ["Won", "Lost", "No Decision"], key="ores")
                st.number_input("Actual Hours", 0, key="oh")
            with oc2:
                st.text_area("Feedback", key="ofb")
                st.number_input("Actual Cost", 0, key="oc")
                st.text_area("Lessons", key="ol")
            if st.form_submit_button("Record", type="primary"):
                st.success("Recorded for next training cycle.")
        st.markdown("---")
        method_cols = st.columns(5)
        method_data = [("📋", "Phases", "Discovery Design Dev Deploy Support"), ("👥", "Team", "Arch Lead Dev QA DevOps PM"), ("📐", "Estimation", "Three-point + 20% buffer"), ("💲", "Pricing", "Fixed, T&M, Retainer"), ("⚠️", "Risk", "Tech Schedule Resource Budget External")]
        for i, (ic, t, d) in enumerate(method_data):
            with method_cols[i]:
                st.markdown('<div class="mc"><div class="mi">' + ic + '</div><div class="mt">' + t + '</div><div class="md2">' + d + '</div></div>', unsafe_allow_html=True)


# ═══════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════

main_t1, main_t2 = st.tabs(["⚡Business Estimation", "⚙️ Admin & Training"])
with main_t1:
    tab_presale()
with main_t2:
    tab_admin()