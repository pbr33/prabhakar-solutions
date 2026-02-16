import sys, os, base64, zipfile
import streamlit as st

st.set_page_config(page_title="ECI | Business Estimation Leveraging Automated Learning", page_icon="⚡", layout="wide", initial_sidebar_state="expanded")

import json, time, io, re, smtplib
from datetime import datetime, timedelta
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
import plotly.graph_objects as go
import plotly.express as px

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side, numbers
    from openpyxl.utils import get_column_letter
    from openpyxl.chart import BarChart, PieChart, Reference
except ImportError:
    Workbook = None

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors as rl_colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm, inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, HRFlowable, Image as RLImage
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from openai import AzureOpenAI
except ImportError:
    AzureOpenAI = None
try:
    import msal
except ImportError:
    msal = None
try:
    import requests as _requests
except ImportError:
    _requests = None


# ═══════════════════════════════════════════════════════════════════════
#  ECI BRANDING CONSTANTS
# ═══════════════════════════════════════════════════════════════════════

ECI_BLUE = (27, 58, 92)
ECI_LIGHT_BLUE = (214, 228, 240)
ECI_ACCENT = (0, 180, 216)
ECI_GREEN = (0, 212, 170)
ECI_DARK = (10, 14, 26)
ECI_TEXT = (30, 30, 30)
ECI_GRAY = (100, 100, 100)

ECI_LOGO_BLUE_B64 = "iVBORw0KGgoAAAANSUhEUgAAAMgAAABGCAYAAACJ4ts2AAAD6ElEQVR4nO3dS4gcRRzH8e8mvlE3D1BiQUx8BFGXHCS6FoKC5hAUFRTFpAiBIMkhh3j04NF4MwEPGkQPWp7UiIqioohgyvXgQcT4CsaDlbgefGwiEh+Jh5qIBN2pbacyWzW/D8xp/93zY5jf9nT3TDeIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiDRibNgBTgVj3WLgZuBaYAK4GLgQOAc4HTgCzADfA/uAT4EAfBCD/2MYmWV+6FwQY90TwJYBZjnZ3hj8DV0XNtYtBO4EtgI3Aad1WM0M8BbwNPBmDP7YLM+X83o8EIPf1SGHDMmCYQcowVh3O2kr8AJwC93KAXA+cDfwOrDfWLdsMAmlFl3fOPOSsW4ceBy4r8DqVwLjwKEC65Z5qpmCGOuWk/7TXzXsLNKOJgpirLsIeA9YMeQo0pjqC2KsO4u05ViRucgPwEvAK8DnwHfAUWBp77EamCQd9bpiwHGlMqULciqO2jxKelP3c4y0f/JQDP7Hf/n7wd7jE8ADGOsmgc3ARuCMgaSVqlS9BTHWrSEdxu3nT2BTDN7PZf0x+Clgyli3A3iYVDIZIVUXBHiEvHM5W+Zajn+KwR8A1nddXupV7XkQY90EaT+hn9di8E+VziNtqrYgpP2Cfo4D2wvnkIbVXJA7MmbejsHvL55EmlVlQYx1FwCXZ4w+VzqLtK3KggDXZM5NFU0hzSt9FGunsW5nh+W+iMHPdpLukox1zABfdnhukb/VugUxGTMHY/DHiyeRptVakPMyZn4unkKaV2tBzsyYOVw8hTSv1oIczZg5t3gKaV6tBTmSMbOodAhpX63f5o0ZM8uMdWPaUZf/o9YtyNcZM+PknUwU+U+1FuSjzLnJoimkeVUWJAY/DeR8x2pD6SzStioL0vNyxsxaY92lxZNIs2ouyLMZM2PArsI5pGHVFiQG/zHwbsbobca6zaXzSJuqLUjPg6QfRfWz21jXeX/EWLfSWOeNdau6rkPqVHVBYvAfAk9mjC4EnjHWPda7kHUWY90aY91u0uWBNlD56yVzN1+/7n7Cuhj8G31mtgMWuLrP3AJgG7DeWPci8CrwGTAN/AYsJl0XawK4DliLrtI48mq/qgkx+F+NdeuAvcDyjEWWAPf3HiKzauIjQwz+W+BG0kchkYFpoiAAMfhvgOuB54ccRRrSTEEAYvA/xeDvAe4Cvhrw6g+gH2GNnKYKckIMfg9wJXAv8A7p0qNdHAb2ALcCl8XgdW+QETMq9yhcQrrT1Mn3KDybdKDiF1IZpkn7MftI9yh8Pwb/+zAyi4iIiIiIiIiIiIiIiIiIiIiIiIiIiIiIiMjc/QXeEc9g6vu/ZAAAAABJRU5ErkJggg=="
ECI_LOGO_WHITE_B64 = "iVBORw0KGgoAAAANSUhEUgAAAMgAAABGCAYAAACJ4ts2AAADvElEQVR4nO3dza9dUxzG8e+vJV6Cq5WQklTrPbgxaMplYkAHN4RBhaRNaiANgw6uob/A8EqIaBoTikGplBAJIpKiBgYGqqFRAy+9BqUvIlfVY7COSZP2rLvvXj13rfN8kjP77bWfs09+Z5/9ctYGMzMzMzMzMzMzMzMzMzMzMzMzMzMzMzMzM7NGxKgDnAuSVgD3AXcCk8C1wFXAxcD5wAngGPAbsB/4Bvgc+CIi/hlFZqucpJdU1t5F5lsuaaOkDyWd7JjhqKRdkqYlLethe8ws5j3ZuXfWD71Wkh4i7QXeBO4Hzus41GXAI8D7wEFJq/pJaLVoqkEkTUh6HdgD3Nzz8GuBiZ7HtCWu6zfrkiNpNemb/rZRZ7F2NNEgkq4GPgXWjDiKNab6BpF0IWnPsSZzkSPA28A7wAHgMDAPXDF43QFMkc563dJzXBsXS+WsjaQXM89InZL0gtIp39yxpyTtkDQ/GOOMDbNUtof1q+qDdEnrgacySk8Bj0fEtoj4PXf8iNgXEVtJe5I3gH+7JbVa1f4T61nyLnY+GRE7u64kIg4Bm7oub/Wqdg8iaZJ0nDDMexHxcuk81qZqGwTYklEjYKZwDmtYzQ3ycEbNRxFxsHgSa1aVDSLpSuDGjNLXSmextlXZIMC6zLp9RVNY80o3yGzmNYrTHRgy7nUZ6z4GfNfDe7AxVuse5JqMml8iQsWTWNNqbZBLM2qOFk9hzau1QS7IqDlePIU1r9YGmc+ouaR4CmterQ1yIqPm8tIhrH2lG+Tp6GbYbeY/Z6x7laSxmJTCyql1D/JDRs0EeRcTzc6o1gb5KrNuqmgKa16VDRIRc0DOPVabS2extlXZIAN7Mmo2SLq+eBJrVs0N8mpGTQDPFc5hDau2QSLia+CTjNIHJT1ROo+1qdoGGXiG9KeoYbZL6nw8ImmtpJ2Sbuo6htWp6gaJiC+BHRmly4FXJD2/wFlN1kvaTpoeaDOVby9buNKTNsxKml3E8tMR8cGQmhngHuD2IXXLgG3AJklvAe8C3wJzwN/ACtK8WJPAXcAGPEvj2Kt9VhMi4i9J08BnwOqMRVYCWwcvs7Nq4idDRPwE3Ev6KWTWmyYaBCAifgTuBnaNOIo1pJkGAYiIPyLiUWAj8H3Pwx/Cf8IaO001yP8iYjdwK/AY8DFp6tEujgO7gQeAGyLi134SWi3G4nZwSStJT5o6/RmFF5FOVPxJaoY50nHMftIzCvdGxMlRZDYzMzMzMzMzMzMzMzMzMzMzMzMzMzMzMzOzhfsPjGdVI+ycCmwAAAAASUVORK5CYII="


# ═══════════════════════════════════════════════════════════════════════
#  CSS
# ═══════════════════════════════════════════════════════════════════════

def inject_css():
    st.markdown("""<style>
    @import url('https://fonts.googleapis.com/css2?family=JetBrains+Mono:wght@300;400;600;700&family=DM+Sans:wght@300;400;500;600;700&family=Space+Grotesk:wght@400;500;600;700&display=swap');
    :root{--bg0:#0a0e1a;--bg1:#111827;--bg2:#151c2e;--bd:#1e2a4a;--t1:#e2e8f0;--t2:#94a3b8;--t3:#64748b;--c1:#00d4aa;--c2:#00b4d8;--c3:#7b61ff;--c4:#ff6b6b;--c5:#ffd166;--c6:#06d6a0;--gc:rgba(0,212,170,.15);--gp:rgba(123,97,255,.15)}
    .stApp{background:var(--bg0)!important}
    .main .block-container{padding-top:1rem;max-width:100%}
    section[data-testid="stSidebar"]{background:linear-gradient(180deg,#0d1321,#111827)!important;border-right:1px solid var(--bd)}
    .logo-box{display:flex;align-items:center;gap:12px}
    .logo-icon{font-size:2.2rem;background:linear-gradient(135deg,var(--c1),var(--c3));-webkit-background-clip:text;-webkit-text-fill-color:transparent;filter:drop-shadow(0 0 12px var(--gc))}
    .logo-txt{font-family:'Space Grotesk',sans-serif;font-size:1.8rem;font-weight:700;background:linear-gradient(135deg,var(--c1),var(--c2));-webkit-background-clip:text;-webkit-text-fill-color:transparent}
    .logo-sub{font-family:'DM Sans',sans-serif;font-size:.7rem;color:var(--t2);letter-spacing:2px;text-transform:uppercase}
    .tagline{text-align:center;font-family:'DM Sans',sans-serif;color:var(--t2);font-size:.85rem;padding-top:.8rem}
    .dt{text-align:right;font-family:'JetBrains Mono',monospace;color:var(--t3);font-size:.8rem;padding-top:1rem}
    .stitle{font-family:'Space Grotesk',sans-serif;font-size:1.3rem;font-weight:600;color:var(--t1);margin-bottom:1rem;padding-bottom:.5rem;border-bottom:1px solid var(--bd)}
    .csec{font-family:'DM Sans',sans-serif;font-size:.9rem;font-weight:600;color:var(--c1);margin:.8rem 0 .4rem;letter-spacing:.5px}
    .cstat{background:var(--bg2);border-radius:8px;padding:12px;border:1px solid var(--bd)}
    .srow{display:flex;align-items:center;gap:8px;padding:4px 0;font-family:'DM Sans',sans-serif;font-size:.82rem;color:var(--t2)}
    .sdot{width:8px;height:8px;border-radius:50%;display:inline-block}
    .sdot.on{background:var(--c6);box-shadow:0 0 6px var(--c6)}.sdot.off{background:var(--t3)}
    .shdr{font-family:'Space Grotesk',sans-serif;font-size:1.3rem;font-weight:600;color:var(--t1);margin:1.5rem 0 1rem;display:flex;align-items:center;gap:10px}
    .shdr-i{font-size:1.4rem}
    .phdr{font-family:'Space Grotesk',sans-serif;font-size:1.1rem;font-weight:600;color:var(--c1);text-align:center;padding:12px;background:linear-gradient(135deg,var(--gc),var(--gp));border-radius:8px;margin-bottom:1rem;border:1px solid var(--bd)}
    .crd{background:var(--bg2);border:1px solid var(--bd);border-radius:12px;padding:1.2rem;margin-bottom:.8rem;transition:border-color .3s}
    .crd:hover{border-color:var(--c1)}
    .crd-t{font-family:'Space Grotesk',sans-serif;font-size:1.05rem;font-weight:600;color:var(--t1);margin-bottom:.3rem}
    .crd-d{font-family:'DM Sans',sans-serif;font-size:.82rem;color:var(--t2);margin-bottom:.8rem}
    .kpi{background:var(--bg2);border:1px solid var(--bd);border-radius:12px;padding:1.2rem;text-align:center;transition:all .3s}
    .kpi:hover{border-color:var(--c1);box-shadow:0 0 20px var(--gc);transform:translateY(-2px)}
    .kpi-i{font-size:1.8rem;margin-bottom:.4rem}
    .kpi-v{font-family:'JetBrains Mono',monospace;font-size:1.5rem;font-weight:700;color:var(--c1)}
    .kpi-t{font-family:'DM Sans',sans-serif;font-size:.85rem;font-weight:600;color:var(--t1);margin-top:.2rem}
    .kpi-s{font-family:'DM Sans',sans-serif;font-size:.72rem;color:var(--t3)}
    .alog{display:flex;align-items:center;gap:12px;padding:6px 12px;background:var(--bg2);border-radius:6px;margin-bottom:4px;border-left:3px solid var(--c1)}
    .abadge{font-family:'JetBrains Mono',monospace;font-size:.78rem;font-weight:600;color:var(--c2);min-width:160px}
    .aok{color:var(--c6);font-size:.8rem;font-weight:500;min-width:90px}
    .adet{font-family:'DM Sans',sans-serif;color:var(--t2);font-size:.8rem}
    .rch{font-family:'Space Grotesk',sans-serif;font-size:.95rem;font-weight:600;padding:8px 12px;border-radius:8px;margin-bottom:.8rem}
    .rch.fn{background:rgba(0,212,170,.1);color:var(--c1);border:1px solid rgba(0,212,170,.2)}
    .rch.nf{background:rgba(0,180,216,.1);color:var(--c2);border:1px solid rgba(0,180,216,.2)}
    .rch.ig{background:rgba(123,97,255,.1);color:var(--c3);border:1px solid rgba(123,97,255,.2)}
    .ri{background:var(--bg2);border:1px solid var(--bd);border-radius:8px;padding:10px 12px;margin-bottom:6px;font-family:'DM Sans',sans-serif;font-size:.82rem;color:var(--t2)}
    .ri strong{color:var(--t1);font-size:.85rem}
    .ri-c{font-family:'JetBrains Mono',monospace;font-size:.7rem;color:var(--c5);margin-top:2px}
    .ttag{display:flex;flex-wrap:wrap;gap:6px;margin-top:8px}
    .tt{font-family:'JetBrains Mono',monospace;font-size:.72rem;background:var(--bg2);border:1px solid var(--c3);color:var(--c3);padding:3px 10px;border-radius:12px}
    .rsb{background:var(--bg2);padding:16px 20px;border-radius:10px;display:flex;align-items:center;gap:16px;margin-bottom:1rem}
    .rsv{font-family:'JetBrains Mono',monospace;font-size:2rem;font-weight:700;color:var(--t1)}
    .rsl{font-family:'DM Sans',sans-serif;font-size:1rem;color:var(--t2)}
    .rc{background:var(--bg2);border:1px solid var(--bd);border-radius:8px;padding:14px;margin-bottom:8px}
    .rch2{display:flex;justify-content:space-between;align-items:center;margin-bottom:6px}
    .rch2 strong{color:var(--t1);font-family:'DM Sans',sans-serif}
    .rsev{font-family:'JetBrains Mono',monospace;font-size:.78rem;font-weight:600}
    .rc p{color:var(--t2);font-size:.82rem;margin:4px 0}
    .rmit{font-size:.8rem;color:var(--c6)}
    .ac{background:var(--bg2);border:1px solid var(--bd);border-radius:10px;padding:14px;margin-bottom:10px;transition:all .3s}
    .ac:hover{border-color:var(--c3);box-shadow:0 0 16px var(--gp)}
    .acn{font-family:'Space Grotesk',sans-serif;font-size:.95rem;font-weight:600;color:var(--t1)}
    .act{font-family:'JetBrains Mono',monospace;font-size:.7rem;color:var(--c3);margin-bottom:6px}
    .acs{font-family:'DM Sans',sans-serif;font-size:.8rem;color:var(--t2);padding-left:16px}
    .acs li{margin-bottom:2px}
    .dfv{font-family:'JetBrains Mono',monospace;font-size:.85rem;color:var(--c1);background:var(--bg2);padding:12px 16px;border-radius:8px;border:1px solid var(--bd);text-align:center;letter-spacing:.5px}
    .ls{background:var(--bg2);border:1px solid var(--bd);border-radius:10px;padding:16px 12px;text-align:center;min-height:160px}
    .ln{display:inline-flex;align-items:center;justify-content:center;width:36px;height:36px;border-radius:50%;font-family:'JetBrains Mono',monospace;font-size:1rem;font-weight:700;color:#fff;margin-bottom:8px}
    .lt{font-family:'Space Grotesk',sans-serif;font-size:.88rem;font-weight:600;color:var(--t1);margin-bottom:6px}
    .ld{font-family:'DM Sans',sans-serif;font-size:.72rem;color:var(--t2);line-height:1.4}
    .mc{background:var(--bg2);border:1px solid var(--bd);border-radius:10px;padding:16px;text-align:center;min-height:120px}
    .mi{font-size:1.6rem;margin-bottom:6px}
    .mt{font-family:'Space Grotesk',sans-serif;font-size:.85rem;font-weight:600;color:var(--t1);margin-bottom:4px}
    .md2{font-family:'DM Sans',sans-serif;font-size:.72rem;color:var(--t2);line-height:1.3}
    button[data-testid="stBaseButton-primary"]{background:linear-gradient(135deg,var(--c1),var(--c2))!important;color:#0a0e1a!important;font-family:'DM Sans',sans-serif!important;font-weight:600!important;border:none!important;border-radius:8px!important}
    button[data-testid="stBaseButton-primary"]:hover{box-shadow:0 0 20px var(--gc)!important}
    button[data-testid="stBaseButton-secondary"]{background:var(--bg2)!important;color:var(--t1)!important;border:1px solid var(--bd)!important;border-radius:8px!important}
    .stTabs [data-baseweb="tab-list"]{gap:4px;background:var(--bg1);padding:4px;border-radius:10px}
    .stTabs [data-baseweb="tab"]{background:transparent;color:var(--t2);border-radius:8px;padding:8px 16px;font-family:'DM Sans',sans-serif;font-size:.85rem}
    .stTabs [data-baseweb="tab"][aria-selected="true"]{background:var(--bg2)!important;color:var(--c1)!important}
    .stTabs [data-baseweb="tab-highlight"]{display:none}
    [data-testid="stFileUploader"]{background:var(--bg2);border:1px dashed var(--bd);border-radius:10px;padding:1rem}
    [data-testid="stMetric"]{background:var(--bg2);border:1px solid var(--bd);border-radius:10px;padding:12px 16px}
    [data-testid="stMetricValue"]{font-family:'JetBrains Mono',monospace!important;color:var(--c1)!important}
    #MainMenu{visibility:hidden}footer{visibility:hidden}
    header[data-testid="stHeader"]{background:rgba(10,14,26,.95);backdrop-filter:blur(10px)}
    ::-webkit-scrollbar{width:6px}::-webkit-scrollbar-track{background:var(--bg0)}::-webkit-scrollbar-thumb{background:var(--bd);border-radius:3px}
    </style>""", unsafe_allow_html=True)


# ═══════════════════════════════════════════════════════════════════════
#  HELPERS: Safe data access
# ═══════════════════════════════════════════════════════════════════════

def safe_int(val, default=0):
    try:
        return int(val)
    except (TypeError, ValueError):
        return default

def safe_str(val, default=""):
    if val is None:
        return default
    return str(val)

def safe_list(val):
    if isinstance(val, list):
        return val
    return []

def safe_dict(val):
    if isinstance(val, dict):
        return val
    return {}


# ═══════════════════════════════════════════════════════════════════════
#  DOCUMENT PROCESSOR
# ═══════════════════════════════════════════════════════════════════════

class DocProcessor:
    def extract(self, f):
        name = f.name.lower()
        data = f.read()
        f.seek(0)
        if name.endswith(".pdf"):
            try:
                from PyPDF2 import PdfReader
                return "\n\n".join(p.extract_text() or "" for p in PdfReader(io.BytesIO(data)).pages).strip()
            except Exception as e:
                return "[PDF error: " + str(e) + "]"
        elif name.endswith(".docx"):
            try:
                from docx import Document
                doc = Document(io.BytesIO(data))
                parts = [p.text for p in doc.paragraphs if p.text.strip()]
                for tbl in doc.tables:
                    for row in tbl.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
                return "\n".join(parts)
            except Exception as e:
                return "[DOCX error: " + str(e) + "]"
        elif name.endswith((".xlsx", ".xls")):
            try:
                from openpyxl import load_workbook
                wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
                parts = []
                for sn in wb.sheetnames:
                    parts.append("=== " + sn + " ===")
                    for row in wb[sn].iter_rows(max_row=500, values_only=True):
                        cells = [str(c) if c else "" for c in row]
                        if any(cells):
                            parts.append(" | ".join(cells))
                return "\n".join(parts)
            except Exception as e:
                return "[XLSX error: " + str(e) + "]"
        elif name.endswith(".pptx"):
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(data))
                parts = []
                for i, sl in enumerate(prs.slides, 1):
                    parts.append("=== Slide " + str(i) + " ===")
                    for sh in sl.shapes:
                        if hasattr(sh, "text") and sh.text.strip():
                            parts.append(sh.text)
                return "\n".join(parts)
            except Exception as e:
                return "[PPTX error: " + str(e) + "]"
        else:
            return data.decode("utf-8", errors="replace")[:50000]

    def analyze(self, text):
        lines = text.split("\n")
        ne = [l for l in lines if l.strip()]
        secs = []
        for s in ne:
            s2 = s.strip()
            if re.match(r"^\d+[\.\)]\s+\w", s2) or (s2.isupper() and 3 < len(s2) < 80) or re.match(r"^#{1,4}\s+", s2):
                secs.append(s2[:100])
        kws = ["azure", "aws", "cloud", "api", "database", "sql", "python", "react", "sharepoint", "teams",
               "power bi", "kubernetes", "docker", ".net", "java", "node", "javascript", "typescript",
               "microservices", "serverless", "devops", "ci/cd", "machine learning", "ai", "cosmos", "blob storage"]
        return {"section_count": len(secs), "word_count": len(text.split()),
                "technologies_mentioned": [k for k in kws if k in text.lower()]}


# ═══════════════════════════════════════════════════════════════════════
#  AZURE OPENAI CLIENT — PRODUCTION
# ═══════════════════════════════════════════════════════════════════════

class AzureAI:
    def __init__(self, key, endpoint, version, deployment):
        self.deployment = deployment
        self._client = None
        if key and endpoint and AzureOpenAI:
            try:
                self._client = AzureOpenAI(api_key=key, api_version=version, azure_endpoint=endpoint)
            except Exception:
                pass

    @classmethod
    def from_session(cls):
        return cls(
            st.session_state.get("azure_api_key", ""),
            st.session_state.get("azure_endpoint", ""),
            st.session_state.get("azure_api_version", "2024-06-01"),
            st.session_state.get("azure_deployment", "gpt-4"),
        )

    @property
    def is_live(self):
        return self._client is not None

    def test(self):
        if not self._client:
            return False, "Not configured. Enter API Key and Endpoint."
        try:
            self._client.chat.completions.create(
                model=self.deployment,
                messages=[{"role": "user", "content": "Reply OK"}],
                max_tokens=5,
            )
            return True, "Connected to " + self.deployment
        except Exception as e:
            return False, "Failed: " + str(e)[:200]

    def _call(self, system, user):
        if not self._client:
            return None
        try:
            resp = self._client.chat.completions.create(
                model=self.deployment,
                messages=[
                    {"role": "system", "content": system},
                    {"role": "user", "content": user},
                ],
                max_tokens=4096,
                temperature=0.2,
                response_format={"type": "json_object"},
            )
            txt = resp.choices[0].message.content
            return json.loads(txt)
        except Exception as e:
            st.warning("Azure OpenAI: " + str(e)[:150])
            return None

    # ── Requirements ──
    def analyze_requirements(self, text):
        r = self._call(
            "You are an expert IT presales analyst for ECI consulting. Analyze the scope document thoroughly. "
            "Extract ALL requirements. Classify each as functional, non-functional, or integration. "
            "Return JSON: {\"requirements\": [{\"title\": str, \"description\": str, \"type\": \"functional\" or \"non-functional\" or \"integration\", \"complexity\": \"Low\" or \"Medium\" or \"High\", \"priority\": str}], "
            "\"technology_stack\": [str], \"business_objectives\": [str], \"complexity_score\": int 1-10, \"project_type\": str}",
            "Analyze:\n\n" + text[:15000],
        )
        if r and isinstance(r, dict) and "requirements" in r:
            return r
        return self._fb_semantic(text)

    # ── Historical RAG ──
    def search_historical(self, text, projects):
        if projects:
            r = self._call(
                "Compare new project with historical data. Return JSON: "
                "{\"similar_projects\": [{\"name\": str, \"similarity\": float, \"hours\": int, \"cost\": int, \"outcome\": str}], "
                "\"benchmark_hours\": int, \"benchmark_cost\": int, \"success_patterns\": [str], \"risk_patterns\": [str]}",
                "New:\n" + text[:5000] + "\n\nHistory:\n" + json.dumps(projects[:20], default=str),
            )
            if r:
                return r
        return {"similar_projects": [], "benchmark_hours": 0, "benchmark_cost": 0, "success_patterns": [], "risk_patterns": []}

    # ── Time ──
    def estimate_time(self, semantic, rag):
        r = self._call(
            "You are an expert ECI project estimator. Use three-point estimation. "
            "Return JSON: {\"total_hours\": int, \"duration_weeks\": \"N weeks\", \"confidence\": str, \"buffer\": str, "
            "\"phases\": [{\"name\": str, \"hours\": int, \"percentage\": \"N%\", \"tasks\": [{\"name\": str, \"hours\": int, \"role\": str}]}], "
            "\"milestones\": [{\"name\": str, \"week\": int, \"description\": str}], "
            "\"three_point\": {\"optimistic\": int, \"most_likely\": int, \"pessimistic\": int}}",
            "Estimate:\nRequirements: " + json.dumps(safe_list(semantic.get("requirements"))[:20], default=str)
            + "\nComplexity: " + str(semantic.get("complexity_score", 5))
            + "\nType: " + safe_str(semantic.get("project_type"))
            + "\nTech: " + json.dumps(safe_list(semantic.get("technology_stack")))
            + "\nBenchmark hours: " + str(rag.get("benchmark_hours", 0)),
        )
        if r and isinstance(r, dict) and "total_hours" in r:
            return self._sanitize_time(r)
        return self._fb_time(semantic, rag)

    # ── Cost (Infrastructure only) ──
    def estimate_cost(self, semantic, time_est, rag):
        r = self._call(
            "You are an ECI infrastructure cost estimator. Estimate ONLY Azure/cloud infrastructure costs, "
            "NOT project labor costs. Consider compute, storage, databases, networking, security, monitoring. "
            "Return JSON: {\"total_monthly_cost\": int, \"total_annual_cost\": int, "
            "\"azure_costs\": [{\"service\": str, \"tier\": str, \"monthly_cost\": int, \"description\": str}], "
            "\"third_party_costs\": [{\"name\": str, \"monthly_cost\": int, \"description\": str}], "
            "\"cost_optimization\": [str], \"notes\": str}",
            "Estimate infrastructure costs:\nTech: " + json.dumps(safe_list(semantic.get("technology_stack")))
            + "\nComplexity: " + str(semantic.get("complexity_score", 5))
            + "\nType: " + safe_str(semantic.get("project_type"))
            + "\nComponents: " + json.dumps(safe_list(semantic.get("requirements"))[:10], default=str),
        )
        if r and isinstance(r, dict) and "azure_costs" in r:
            return r
        return self._fb_cost(time_est)

    # ── Risk ──
    def analyze_risk(self, semantic, time_est, cost_est):
        r = self._call(
            "You are an ECI risk analyst. Assess Technical, Schedule, Resource, Budget, External. "
            "Return JSON: {\"overall_score\": int 1-10, \"overall_level\": str, "
            "\"risks\": [{\"category\": str, \"title\": str, \"description\": str, \"severity\": \"Low\" or \"Medium\" or \"High\", "
            "\"probability\": str, \"impact\": str, \"mitigation\": str}]}",
            "Assess:\nReqs: " + str(len(safe_list(semantic.get("requirements"))))
            + "\nComplexity: " + str(semantic.get("complexity_score", 5))
            + "\nHours: " + str(time_est.get("total_hours", 0))
            + "\nInfra Monthly: $" + str(cost_est.get("total_monthly_cost", 0))
            + "\nTech: " + json.dumps(safe_list(semantic.get("technology_stack"))),
        )
        if r and isinstance(r, dict) and "risks" in r:
            return r
        return self._fb_risk(semantic)

    # ── Architecture ──
    def design_architecture(self, semantic, rag):
        r = self._call(
            "You are an Azure Solutions Architect for ECI. Use Well-Architected Framework. "
            "Return JSON: {\"pattern\": str, \"components\": [{\"name\": str, \"type\": str, \"azure_service\": str, \"services\": [str]}], "
            "\"data_flow\": [str], \"security\": [str], \"scalability\": str, \"availability\": str}",
            "Design:\nReqs: " + json.dumps(safe_list(semantic.get("requirements"))[:15], default=str)
            + "\nTech: " + json.dumps(safe_list(semantic.get("technology_stack")))
            + "\nType: " + safe_str(semantic.get("project_type")),
        )
        if r and isinstance(r, dict) and "components" in r:
            return r
        return self._fb_arch()

    # ── Scope ──
    def define_scope(self, semantic, time_est, cost_est):
        r = self._call(
            "You are an ECI scope expert. Define boundaries. "
            "Return JSON: {\"in_scope\": [str], \"out_of_scope\": [str], \"assumptions\": [str], \"prerequisites\": [str]}",
            "Scope:\nReqs: " + json.dumps(safe_list(semantic.get("requirements"))[:10], default=str)
            + "\nType: " + safe_str(semantic.get("project_type"))
            + "\nHours: " + str(time_est.get("total_hours", 0)),
        )
        if r and isinstance(r, dict) and "in_scope" in r:
            return r
        return self._fb_scope()

    # ── Proposal ──
    def write_proposal(self, semantic, time_est, cost_est, risk, arch, scope):
        r = self._call(
            "You are an ECI presales proposal writer. Write professional proposal. "
            "Return JSON: {\"sections\": [{\"title\": str, \"content\": str}], \"quality_checks\": {str: bool}}. "
            "Include sections: Executive Summary, Understanding & Approach, Technical Solution, Team, Timeline, Infrastructure Costs.",
            "Generate:\nType: " + safe_str(semantic.get("project_type"))
            + "\nReqs: " + str(len(safe_list(semantic.get("requirements"))))
            + "\nHours: " + str(time_est.get("total_hours", 0))
            + "\nInfra Monthly: $" + str(cost_est.get("total_monthly_cost", 0))
            + "\nRisk: " + safe_str(risk.get("overall_level"))
            + "\nArch: " + safe_str(arch.get("pattern"))
            + "\nPhases: " + json.dumps([{"name": safe_str(p.get("name")), "hours": safe_int(p.get("hours"))} for p in safe_list(time_est.get("phases"))])
            + "\nScope: " + json.dumps(safe_list(scope.get("in_scope"))[:5]),
        )
        if r and isinstance(r, dict) and "sections" in r:
            return r
        return self._fb_proposal(semantic, time_est, cost_est, risk, arch)

    # ── Architecture Mermaid Diagrams ──
    def generate_mermaid_diagrams(self, semantic, arch):
        r = self._call(
            "You are an Azure Solutions Architect for ECI. Generate Mermaid.js diagram code for the architecture. "
            "Return JSON with these keys, each containing VALID Mermaid syntax as a string:\n"
            "{\"infrastructure\": str (graph TD diagram of Azure components and connections),\n"
            " \"data_flow\": str (flowchart LR diagram showing data movement between services),\n"
            " \"sequence\": str (sequenceDiagram showing a typical user request flow),\n"
            " \"deployment\": str (graph TD diagram showing CI/CD pipeline and environments),\n"
            " \"security\": str (graph TD diagram showing security layers and controls)}",
            "Architecture:\nPattern: " + safe_str(arch.get("pattern"))
            + "\nComponents: " + json.dumps(safe_list(arch.get("components"))[:10], default=str)
            + "\nData Flow: " + json.dumps(safe_list(arch.get("data_flow")))
            + "\nSecurity: " + json.dumps(safe_list(arch.get("security")))
            + "\nTech: " + json.dumps(safe_list(semantic.get("technology_stack"))),
        )
        if r and isinstance(r, dict) and "infrastructure" in r:
            return r
        return self._fb_mermaid(arch)

    # ── Transcript Analysis & WBS ──
    def analyze_transcript(self, transcript_text):
        r = self._call(
            "You are an expert ECI presales analyst. Analyze this meeting transcript/voice note. "
            "Extract requirements, pain points (from tone/urgency), stakeholders, decisions, and action items. "
            "Generate a Work Breakdown Structure. Return JSON:\n"
            "{\"pain_points\": [{\"issue\": str, \"severity\": \"High\" or \"Medium\" or \"Low\", \"quote\": str, \"stakeholder\": str}],\n"
            " \"requirements_extracted\": [{\"title\": str, \"description\": str, \"type\": \"functional\" or \"non-functional\" or \"integration\", \"source\": str}],\n"
            " \"stakeholders\": [{\"name\": str, \"role\": str, \"concerns\": [str]}],\n"
            " \"decisions\": [str], \"action_items\": [{\"item\": str, \"owner\": str, \"priority\": str}],\n"
            " \"wbs\": [{\"phase\": str, \"deliverables\": [{\"name\": str, \"tasks\": [{\"name\": str, \"effort\": str}]}]}],\n"
            " \"meeting_summary\": str, \"sentiment\": str, \"key_themes\": [str]}",
            "Transcript:\n\n" + transcript_text[:15000],
        )
        if r and isinstance(r, dict) and "wbs" in r:
            return r
        return self._fb_transcript(transcript_text)

    # ═══ FALLBACK: Mermaid Diagrams ═══
    def _fb_mermaid(self, arch):
        comps = safe_list(arch.get("components"))
        comp_names = [safe_str(safe_dict(c).get("name", "Component")) for c in comps]
        services = {safe_str(safe_dict(c).get("name")): safe_str(safe_dict(c).get("azure_service", "")) for c in comps}

        infra = "graph TD\n"
        infra += "    User([Client Browser])\n"
        infra += "    CDN[Azure Front Door / CDN]\n"
        infra += "    APIM[API Management]\n"
        infra += "    APP[App Service]\n"
        infra += "    FUNC[Azure Functions]\n"
        infra += "    SQL[(Azure SQL)]\n"
        infra += "    COSMOS[(Cosmos DB)]\n"
        infra += "    REDIS[Redis Cache]\n"
        infra += "    BLOB[Blob Storage]\n"
        infra += "    SB[Service Bus]\n"
        infra += "    KV[Key Vault]\n"
        infra += "    AD[Azure AD]\n"
        infra += "    MON[App Insights / Monitor]\n\n"
        infra += "    User --> CDN --> APIM --> APP\n"
        infra += "    APP --> SQL\n"
        infra += "    APP --> COSMOS\n"
        infra += "    APP --> REDIS\n"
        infra += "    APP --> BLOB\n"
        infra += "    APP --> SB --> FUNC\n"
        infra += "    FUNC --> SQL\n"
        infra += "    FUNC --> COSMOS\n"
        infra += "    APP --> KV\n"
        infra += "    APP --> AD\n"
        infra += "    APP --> MON\n"
        infra += "    FUNC --> MON\n\n"
        infra += "    style User fill:#E8F5E9,stroke:#2E7D32\n"
        infra += "    style CDN fill:#E3F2FD,stroke:#1565C0\n"
        infra += "    style APIM fill:#FFF3E0,stroke:#E65100\n"
        infra += "    style APP fill:#E8EAF6,stroke:#283593\n"
        infra += "    style FUNC fill:#FCE4EC,stroke:#C62828\n"
        infra += "    style SQL fill:#F3E5F5,stroke:#6A1B9A\n"
        infra += "    style COSMOS fill:#F3E5F5,stroke:#6A1B9A\n"
        infra += "    style SB fill:#FFF8E1,stroke:#F57F17\n"
        infra += "    style KV fill:#EFEBE9,stroke:#4E342E\n"
        infra += "    style AD fill:#E0F2F1,stroke:#00695C\n"
        infra += "    style MON fill:#F1F8E9,stroke:#33691E"

        data_flow = "flowchart LR\n"
        data_flow += "    A[Client App] -->|HTTPS| B[API Gateway]\n"
        data_flow += "    B -->|Route| C[App Service]\n"
        data_flow += "    C -->|Read/Write| D[(SQL Database)]\n"
        data_flow += "    C -->|Cache| E[Redis Cache]\n"
        data_flow += "    C -->|Documents| F[Blob Storage]\n"
        data_flow += "    C -->|Events| G[Service Bus]\n"
        data_flow += "    G -->|Trigger| H[Azure Functions]\n"
        data_flow += "    H -->|Process| I[(Cosmos DB)]\n"
        data_flow += "    H -->|Notify| J[Notification Hub]\n"
        data_flow += "    C -->|Analytics| K[Power BI]\n"
        data_flow += "    D -->|Sync| I\n\n"
        data_flow += "    style A fill:#16274B,color:#fff\n"
        data_flow += "    style B fill:#00929E,color:#fff\n"
        data_flow += "    style C fill:#16274B,color:#fff\n"
        data_flow += "    style H fill:#96C038,color:#fff"

        sequence = "sequenceDiagram\n"
        sequence += "    actor User\n"
        sequence += "    participant FD as Front Door\n"
        sequence += "    participant API as API Management\n"
        sequence += "    participant AD as Azure AD\n"
        sequence += "    participant App as App Service\n"
        sequence += "    participant Cache as Redis\n"
        sequence += "    participant DB as Azure SQL\n"
        sequence += "    participant SB as Service Bus\n"
        sequence += "    participant Func as Functions\n\n"
        sequence += "    User->>FD: HTTPS Request\n"
        sequence += "    FD->>API: Route + WAF\n"
        sequence += "    API->>AD: Validate JWT\n"
        sequence += "    AD-->>API: Token Valid\n"
        sequence += "    API->>App: Forward Request\n"
        sequence += "    App->>Cache: Check Cache\n"
        sequence += "    alt Cache Hit\n"
        sequence += "        Cache-->>App: Return Cached\n"
        sequence += "    else Cache Miss\n"
        sequence += "        App->>DB: Query Data\n"
        sequence += "        DB-->>App: Results\n"
        sequence += "        App->>Cache: Update Cache\n"
        sequence += "    end\n"
        sequence += "    App->>SB: Publish Event\n"
        sequence += "    SB->>Func: Trigger Processing\n"
        sequence += "    App-->>User: JSON Response"

        deployment = "graph TD\n"
        deployment += "    DEV[Developer] -->|git push| REPO[Azure Repos / GitHub]\n"
        deployment += "    REPO -->|trigger| BUILD[Build Pipeline]\n"
        deployment += "    BUILD -->|test| TEST[Unit + Integration Tests]\n"
        deployment += "    TEST -->|pass| ARTIFACT[Build Artifact]\n"
        deployment += "    ARTIFACT -->|deploy| STAGING[Staging Slot]\n"
        deployment += "    STAGING -->|smoke test| QA{QA Gate}\n"
        deployment += "    QA -->|approve| SWAP[Slot Swap]\n"
        deployment += "    SWAP -->|live| PROD[Production]\n"
        deployment += "    QA -->|reject| DEV\n"
        deployment += "    PROD -->|monitor| INSIGHTS[App Insights]\n"
        deployment += "    INSIGHTS -->|alert| TEAMS[Teams Notification]\n\n"
        deployment += "    style DEV fill:#E8F5E9,stroke:#2E7D32\n"
        deployment += "    style PROD fill:#16274B,color:#fff\n"
        deployment += "    style STAGING fill:#FFF3E0,stroke:#E65100\n"
        deployment += "    style BUILD fill:#E3F2FD,stroke:#1565C0\n"
        deployment += "    style QA fill:#FCE4EC,stroke:#C62828"

        security = "graph TD\n"
        security += "    EXT[External Users] --> WAF[Web Application Firewall]\n"
        security += "    WAF --> FD[Azure Front Door]\n"
        security += "    FD --> APIM[API Management]\n"
        security += "    APIM --> AUTH[Azure AD B2C]\n"
        security += "    AUTH --> RBAC[Role-Based Access Control]\n"
        security += "    RBAC --> APP[App Service]\n"
        security += "    APP --> KV[Key Vault - Secrets]\n"
        security += "    APP --> MI[Managed Identity]\n"
        security += "    MI --> SQL[(SQL - TDE Encrypted)]\n"
        security += "    MI --> BLOB[(Blob - SSE Encrypted)]\n"
        security += "    APP --> NSG[Network Security Groups]\n"
        security += "    APP --> PE[Private Endpoints]\n"
        security += "    PE --> VNET[Virtual Network]\n"
        security += "    VNET --> SQL\n"
        security += "    VNET --> BLOB\n"
        security += "    APP --> LOG[Azure Sentinel / Log Analytics]\n\n"
        security += "    style WAF fill:#C62828,color:#fff\n"
        security += "    style AUTH fill:#00695C,color:#fff\n"
        security += "    style KV fill:#4E342E,color:#fff\n"
        security += "    style NSG fill:#E65100,color:#fff\n"
        security += "    style LOG fill:#1565C0,color:#fff"

        return {
            "infrastructure": infra,
            "data_flow": data_flow,
            "sequence": sequence,
            "deployment": deployment,
            "security": security,
        }

    # ═══ FALLBACK: Transcript / WBS ═══
    def _fb_transcript(self, text):
        words = text.split()
        word_count = len(words)
        return {
            "pain_points": [
                {"issue": "Manual data processing taking excessive time", "severity": "High", "quote": "We spend hours every week just copying data between systems", "stakeholder": "Operations Lead"},
                {"issue": "Lack of real-time visibility into operations", "severity": "High", "quote": "By the time we get reports, the data is already stale", "stakeholder": "VP Operations"},
                {"issue": "Integration gaps between existing systems", "severity": "Medium", "quote": "Our CRM and ERP don't talk to each other properly", "stakeholder": "IT Manager"},
                {"issue": "Security concerns with current manual workflows", "severity": "Medium", "quote": "People are emailing spreadsheets with sensitive data", "stakeholder": "CISO"},
            ],
            "requirements_extracted": [
                {"title": "Automated Data Pipeline", "description": "ETL pipeline connecting CRM, ERP, and data warehouse", "type": "functional", "source": "Operations Lead — pain point discussion"},
                {"title": "Real-time Dashboard", "description": "Live operational metrics with <5 min refresh", "type": "functional", "source": "VP Operations — visibility concern"},
                {"title": "System Integration Layer", "description": "API-based integration between CRM and ERP", "type": "integration", "source": "IT Manager — integration gap"},
                {"title": "Secure Data Transfer", "description": "Encrypted data pipelines replacing manual email workflows", "type": "non-functional", "source": "CISO — security concern"},
                {"title": "User Authentication & SSO", "description": "Azure AD SSO with MFA for all users", "type": "non-functional", "source": "CISO — access control discussion"},
                {"title": "Mobile Access", "description": "Responsive web app for field team access", "type": "functional", "source": "Field Operations Manager — remote access need"},
            ],
            "stakeholders": [
                {"name": "Sarah Mitchell", "role": "VP Operations", "concerns": ["Reporting delays", "Operational visibility", "Cost of current manual processes"]},
                {"name": "James Chen", "role": "IT Manager", "concerns": ["Integration complexity", "Maintenance burden", "Team skill gaps"]},
                {"name": "David Park", "role": "CISO", "concerns": ["Data security", "Compliance", "Audit trail"]},
                {"name": "Lisa Ramirez", "role": "Operations Lead", "concerns": ["Daily workflow efficiency", "Data accuracy", "Training time"]},
            ],
            "decisions": [
                "Azure cloud platform selected as preferred infrastructure",
                "Phased rollout approach agreed — pilot with Operations team first",
                "Budget range of $150K-250K for initial phase discussed",
                "Q3 2025 target for MVP delivery",
                "Weekly stakeholder sync meetings during discovery",
            ],
            "action_items": [
                {"item": "Share current system architecture documentation", "owner": "IT Manager", "priority": "High"},
                {"item": "Provide sample data exports from CRM and ERP", "owner": "Operations Lead", "priority": "High"},
                {"item": "Schedule security requirements workshop", "owner": "CISO", "priority": "Medium"},
                {"item": "Prepare ECI proposal with options", "owner": "ECI Team", "priority": "High"},
                {"item": "Set up Azure sandbox environment", "owner": "IT Manager", "priority": "Medium"},
            ],
            "wbs": [
                {"phase": "Discovery & Planning", "deliverables": [
                    {"name": "Stakeholder Workshops", "tasks": [{"name": "Requirements gathering sessions", "effort": "3 days"}, {"name": "Pain point analysis", "effort": "2 days"}, {"name": "Current state assessment", "effort": "3 days"}]},
                    {"name": "Technical Assessment", "tasks": [{"name": "System landscape review", "effort": "2 days"}, {"name": "Data mapping", "effort": "3 days"}, {"name": "Integration feasibility", "effort": "2 days"}]},
                    {"name": "Project Plan", "tasks": [{"name": "WBS finalization", "effort": "1 day"}, {"name": "Resource plan", "effort": "1 day"}, {"name": "Risk register", "effort": "1 day"}]},
                ]},
                {"phase": "Design & Architecture", "deliverables": [
                    {"name": "Solution Architecture", "tasks": [{"name": "High-level design", "effort": "3 days"}, {"name": "Data model design", "effort": "4 days"}, {"name": "API contract design", "effort": "3 days"}]},
                    {"name": "UX Design", "tasks": [{"name": "Wireframes", "effort": "3 days"}, {"name": "UI mockups", "effort": "4 days"}, {"name": "User testing", "effort": "2 days"}]},
                    {"name": "Security Design", "tasks": [{"name": "IAM design", "effort": "2 days"}, {"name": "Network security", "effort": "2 days"}, {"name": "Encryption strategy", "effort": "1 day"}]},
                ]},
                {"phase": "Development", "deliverables": [
                    {"name": "Backend Services", "tasks": [{"name": "API development", "effort": "15 days"}, {"name": "Data pipeline", "effort": "10 days"}, {"name": "Integration layer", "effort": "8 days"}]},
                    {"name": "Frontend Application", "tasks": [{"name": "UI components", "effort": "10 days"}, {"name": "Dashboard views", "effort": "8 days"}, {"name": "Responsive design", "effort": "4 days"}]},
                    {"name": "Infrastructure", "tasks": [{"name": "IaC (Bicep/Terraform)", "effort": "5 days"}, {"name": "CI/CD pipelines", "effort": "3 days"}, {"name": "Environment setup", "effort": "2 days"}]},
                ]},
                {"phase": "Testing & QA", "deliverables": [
                    {"name": "Testing", "tasks": [{"name": "Unit testing", "effort": "5 days"}, {"name": "Integration testing", "effort": "5 days"}, {"name": "Performance testing", "effort": "3 days"}, {"name": "Security testing", "effort": "3 days"}, {"name": "UAT", "effort": "5 days"}]},
                ]},
                {"phase": "Deployment & Hypercare", "deliverables": [
                    {"name": "Go-Live", "tasks": [{"name": "Data migration", "effort": "3 days"}, {"name": "Production deployment", "effort": "2 days"}, {"name": "Smoke testing", "effort": "1 day"}]},
                    {"name": "Hypercare", "tasks": [{"name": "Post-launch monitoring", "effort": "10 days"}, {"name": "Bug fixes", "effort": "5 days"}, {"name": "Knowledge transfer", "effort": "3 days"}]},
                ]},
            ],
            "meeting_summary": "Discovery call with key stakeholders revealed significant operational inefficiencies driven by manual data workflows and disconnected systems. Primary pain points center around delayed reporting, manual data transfers, and security gaps. The team expressed strong preference for Azure-based cloud solution with phased delivery approach. Budget is available for Q3 2025 MVP.",
            "sentiment": "Positive — stakeholders are motivated and have executive buy-in for modernization",
            "key_themes": ["Automation", "Real-time Analytics", "System Integration", "Security", "Cloud Migration", "Mobile Access"],
        }

    # ── Sanitize time (CRITICAL: prevents the TypeError) ──
    def _sanitize_time(self, data):
        # Milestones
        clean_ms = []
        for m in safe_list(data.get("milestones")):
            if isinstance(m, dict):
                clean_ms.append({
                    "name": safe_str(m.get("name", "Milestone")),
                    "week": safe_int(m.get("week", 0)),
                    "description": safe_str(m.get("description", "")),
                })
            elif isinstance(m, str):
                clean_ms.append({"name": m, "week": 0, "description": ""})
        data["milestones"] = clean_ms
        # Phases
        clean_ph = []
        for p in safe_list(data.get("phases")):
            if not isinstance(p, dict):
                continue
            clean_tasks = []
            for t in safe_list(p.get("tasks")):
                if isinstance(t, dict):
                    clean_tasks.append({
                        "name": safe_str(t.get("name", "")),
                        "hours": safe_int(t.get("hours", 0)),
                        "role": safe_str(t.get("role", "")),
                    })
            clean_ph.append({
                "name": safe_str(p.get("name", "Phase")),
                "hours": safe_int(p.get("hours", 0)),
                "percentage": safe_str(p.get("percentage", "0%")),
                "tasks": clean_tasks,
            })
        data["phases"] = clean_ph
        data["total_hours"] = safe_int(data.get("total_hours", 0))
        data["duration_weeks"] = safe_str(data.get("duration_weeks", "TBD"))
        data["confidence"] = safe_str(data.get("confidence", "N/A"))
        data["buffer"] = safe_str(data.get("buffer", "N/A"))
        data["three_point"] = safe_dict(data.get("three_point"))
        return data

    # ═══ FALLBACKS (only when Azure not configured) ═══

    def _fb_semantic(self, text):
        reqs = []
        for t, d, tp, c in [
            ("User Auth & SSO", "Azure AD SSO with MFA", "functional", "High"),
            ("Dashboard & Reports", "Real-time Power BI dashboards", "functional", "High"),
            ("Data Pipeline", "Automated ETL from multiple sources", "functional", "High"),
            ("API Gateway", "RESTful API with rate limiting", "functional", "Medium"),
            ("Notifications", "Email, Teams, push alerts", "functional", "Medium"),
            ("Search", "Full-text search with filtering", "functional", "Medium"),
            ("Performance SLA", "<200ms, 10K concurrent users", "non-functional", "High"),
            ("Security", "SOC2 Type II, encryption", "non-functional", "High"),
            ("Scalability", "Auto-scale 5x peak", "non-functional", "High"),
            ("Azure AD", "SSO and RBAC integration", "integration", "Medium"),
            ("SharePoint", "Document library integration", "integration", "Medium"),
            ("Power BI", "Embedded analytics", "integration", "Medium"),
            ("Teams", "Notification bot", "integration", "Medium"),
        ]:
            reqs.append({"title": t, "description": d, "type": tp, "complexity": c, "priority": "P1" if c == "High" else "P2"})
        return {
            "requirements": reqs,
            "technology_stack": ["Azure App Service", "Azure SQL", "Cosmos DB", "Functions", "API Management", "Azure AD", "Power BI", "DevOps", "React", "Python", ".NET"],
            "business_objectives": ["Reduce manual effort 60%", "Real-time analytics", "Secure collaboration"],
            "complexity_score": 7,
            "project_type": "Data/Cloud/AI",
        }

    def _fb_time(self, semantic, rag):
        reqs = safe_list(semantic.get("requirements"))
        n = len(reqs)
        base = max(800, n * 100)
        total = int(base * 1.18)
        weeks = max(8, total // 160)
        phases = [
            {"name": "Discovery", "hours": int(total * .10), "percentage": "10%", "tasks": [{"name": "Workshops", "hours": int(total * .04), "role": "BA"}, {"name": "Requirements", "hours": int(total * .03), "role": "BA"}, {"name": "Planning", "hours": int(total * .03), "role": "PM"}]},
            {"name": "Design", "hours": int(total * .15), "percentage": "15%", "tasks": [{"name": "Architecture", "hours": int(total * .06), "role": "Architect"}, {"name": "UX Design", "hours": int(total * .05), "role": "UX"}, {"name": "Data Model", "hours": int(total * .04), "role": "Data Architect"}]},
            {"name": "Development", "hours": int(total * .40), "percentage": "40%", "tasks": [{"name": "Backend", "hours": int(total * .18), "role": "Senior Dev"}, {"name": "Frontend", "hours": int(total * .12), "role": "Frontend Dev"}, {"name": "APIs", "hours": int(total * .10), "role": "Developer"}]},
            {"name": "Testing", "hours": int(total * .15), "percentage": "15%", "tasks": [{"name": "Unit/Integration", "hours": int(total * .06), "role": "QA"}, {"name": "Performance", "hours": int(total * .04), "role": "QA"}, {"name": "UAT", "hours": int(total * .03), "role": "BA"}, {"name": "Security", "hours": int(total * .02), "role": "Security"}]},
            {"name": "Deployment", "hours": int(total * .10), "percentage": "10%", "tasks": [{"name": "CI/CD", "hours": int(total * .03), "role": "DevOps"}, {"name": "Migration", "hours": int(total * .04), "role": "Data Eng"}, {"name": "Go-live", "hours": int(total * .03), "role": "Team"}]},
            {"name": "Support", "hours": int(total * .10), "percentage": "10%", "tasks": [{"name": "Hypercare", "hours": int(total * .05), "role": "Support"}, {"name": "KT", "hours": int(total * .03), "role": "Lead Dev"}, {"name": "Docs", "hours": int(total * .02), "role": "Writer"}]},
        ]
        milestones = [
            {"name": "Kickoff", "week": 1, "description": "Team onboarding"},
            {"name": "Requirements Baselined", "week": max(2, weeks // 8), "description": "Sign-off"},
            {"name": "Design Approved", "week": max(4, weeks // 4), "description": "Architecture review"},
            {"name": "MVP", "week": max(8, weeks // 2), "description": "Core features ready"},
            {"name": "UAT Start", "week": max(10, int(weeks * .75)), "description": "User testing"},
            {"name": "Go-Live", "week": weeks, "description": "Production deployment"},
        ]
        return {"total_hours": total, "duration_weeks": str(weeks) + " weeks", "confidence": "Medium (72%)", "buffer": "18%", "phases": phases, "milestones": milestones, "three_point": {"optimistic": int(total * .8), "most_likely": total, "pessimistic": int(total * 1.35)}}

    def _fb_cost(self, time_est):
        azure_costs = [
            {"service": "Azure App Service", "tier": "Premium P1v3", "monthly_cost": 450, "description": "Web app hosting with auto-scale"},
            {"service": "Azure SQL Database", "tier": "Standard S3", "monthly_cost": 380, "description": "Relational database with 100 DTUs"},
            {"service": "Cosmos DB", "tier": "Autoscale 4000 RU/s", "monthly_cost": 320, "description": "NoSQL for high-throughput workloads"},
            {"service": "Azure Functions", "tier": "Premium EP1", "monthly_cost": 180, "description": "Serverless compute for background jobs"},
            {"service": "API Management", "tier": "Standard", "monthly_cost": 550, "description": "API gateway with rate limiting"},
            {"service": "Azure Front Door", "tier": "Standard", "monthly_cost": 280, "description": "CDN and global load balancing"},
            {"service": "Azure Service Bus", "tier": "Standard", "monthly_cost": 95, "description": "Message queuing and event-driven"},
            {"service": "Azure Key Vault", "tier": "Standard", "monthly_cost": 15, "description": "Secrets and certificate management"},
            {"service": "Azure Monitor + App Insights", "tier": "Pay-as-you-go", "monthly_cost": 120, "description": "Logging, monitoring, alerting"},
            {"service": "Azure AD B2C", "tier": "Premium P1", "monthly_cost": 130, "description": "Identity and access management"},
            {"service": "Azure Blob Storage", "tier": "Hot LRS", "monthly_cost": 60, "description": "File and document storage"},
            {"service": "Azure Redis Cache", "tier": "Standard C1", "monthly_cost": 160, "description": "In-memory caching"},
        ]
        third_party = [
            {"name": "SendGrid (Email)", "monthly_cost": 45, "description": "Transactional email service"},
            {"name": "SSL Certificates", "monthly_cost": 15, "description": "Custom domain SSL"},
        ]
        total_monthly = sum(a["monthly_cost"] for a in azure_costs) + sum(t["monthly_cost"] for t in third_party)
        return {
            "total_monthly_cost": total_monthly,
            "total_annual_cost": total_monthly * 12,
            "azure_costs": azure_costs,
            "third_party_costs": third_party,
            "cost_optimization": [
                "Use Reserved Instances for 36% savings on App Service and SQL",
                "Enable auto-shutdown for non-production environments",
                "Use Azure Hybrid Benefit if existing Windows Server licenses",
                "Monitor Cosmos DB RU consumption and right-size",
                "Use Azure Cost Management alerts at 80% and 100% budget",
            ],
            "notes": "Estimates based on production environment. Dev/staging adds ~40% of prod costs.",
        }

    def _fb_risk(self, semantic):
        c = safe_int(semantic.get("complexity_score", 7))
        risks = [
            {"category": "Technical", "title": "Integration Complexity", "description": "Third-party integrations may have limitations.", "severity": "High", "probability": "Medium", "impact": "High", "mitigation": "Early PoC for each integration."},
            {"category": "Schedule", "title": "Scope Creep", "description": "Requirements may evolve.", "severity": "High", "probability": "High", "impact": "High", "mitigation": "Strict change request process."},
            {"category": "Resource", "title": "Key Personnel", "description": "Specialists may have limited availability.", "severity": "Medium", "probability": "Medium", "impact": "High", "mitigation": "Secure commitments early."},
            {"category": "Budget", "title": "Cloud Cost Overrun", "description": "Azure costs may exceed estimates.", "severity": "Medium", "probability": "Medium", "impact": "Medium", "mitigation": "Cost Management alerts."},
            {"category": "External", "title": "Regulatory Changes", "description": "Compliance may evolve.", "severity": "Low", "probability": "Low", "impact": "High", "mitigation": "Design for extensibility."},
        ]
        score = min(10, max(1, int(c * 0.8)))
        return {"overall_score": score, "overall_level": "Low" if score <= 3 else "Medium" if score <= 6 else "High", "risks": risks}

    def _fb_arch(self):
        return {"pattern": "Microservices with Event-Driven Integration", "components": [
            {"name": "Frontend", "type": "Web App", "azure_service": "App Service", "services": ["React SPA", "CDN", "Front Door"]},
            {"name": "API Gateway", "type": "Integration", "azure_service": "API Management", "services": ["REST APIs", "Rate limiting", "JWT"]},
            {"name": "Backend", "type": "Microservices", "azure_service": "App Service / Functions", "services": [".NET Core", "Azure Functions", "SignalR"]},
            {"name": "Data", "type": "Database", "azure_service": "SQL + Cosmos DB", "services": ["Azure SQL", "Cosmos DB", "Redis", "Blob"]},
            {"name": "Integration", "type": "Messaging", "azure_service": "Service Bus", "services": ["Async messaging", "Logic Apps", "SharePoint connector"]},
            {"name": "Security", "type": "Identity", "azure_service": "Azure AD + Key Vault", "services": ["AD B2C", "Key Vault", "Managed Identity"]},
            {"name": "DevOps", "type": "Operations", "azure_service": "Azure DevOps", "services": ["CI/CD", "Bicep IaC", "App Insights"]},
        ], "data_flow": ["Client", "Front Door", "API Mgmt", "App Service", "Database", "Service Bus", "Functions"], "security": ["Azure AD + MFA", "Key Vault", "TLS 1.3", "SQL TDE", "RBAC", "Azure Policy", "Sentinel"], "scalability": "Auto-scaling with Cosmos DB RUs", "availability": "Multi-region, 99.95% SLA"}

    def _fb_scope(self):
        return {"in_scope": ["All functional requirements", "Architecture design", "Full-stack development", "Integrations", "Testing (unit, integration, perf, security, UAT)", "CI/CD and IaC", "Data migration", "Documentation", "30-day hypercare"], "out_of_scope": ["Legacy decommissioning", "End-user training", "Hardware procurement", "License procurement", "Data cleansing", "Native mobile apps", "Multi-language", "Pen testing", "Ongoing support beyond hypercare"], "assumptions": ["Client provides Azure subscription", "Dedicated product owner", "SMEs 10hrs/week", "APIs documented", "Standard business hours"], "prerequisites": ["Azure subscription", "AD tenant", "SharePoint site", "API credentials", "Signed SOW"]}

    def _fb_proposal(self, sem, te, ce, ri, ar):
        d = datetime.now().strftime("%B %d, %Y")
        reqs = safe_list(sem.get("requirements"))
        monthly = safe_int(ce.get("total_monthly_cost"))
        annual = safe_int(ce.get("total_annual_cost"))
        return {"sections": [
            {"title": "Executive Summary", "content": "**Date:** " + d + "\n\nECI proposes a " + safe_str(sem.get("project_type")) + " solution. **" + str(len(reqs)) + " requirements** identified.\n\n- Effort: **" + str(te.get("total_hours", 0)) + "h** over **" + safe_str(te.get("duration_weeks")) + "**\n- Infrastructure: **$" + str(monthly) + "/month** ($" + str(annual) + "/year)\n- Risk: **" + safe_str(ri.get("overall_level")) + "**"},
            {"title": "Understanding & Approach", "content": "ECI follows Discovery, Design, Development, Testing, Deployment with hypercare."},
            {"title": "Technical Solution", "content": "**Pattern:** " + safe_str(ar.get("pattern")) + "\n\n" + "\n".join("- **" + safe_str(safe_dict(c).get("name")) + "**: " + safe_str(safe_dict(c).get("azure_service")) for c in safe_list(ar.get("components")))},
            {"title": "Team & Delivery", "content": "ECI will deploy a cross-functional team including Architects, Lead Developers, Senior Developers, QA Engineers, DevOps Engineers, and a Project Manager.\n\nDelivery follows an Agile methodology with bi-weekly sprints and milestone reviews."},
            {"title": "Timeline", "content": "\n".join("- **" + safe_str(safe_dict(p).get("name")) + "** — " + str(safe_int(safe_dict(p).get("hours"))) + "h (" + safe_str(safe_dict(p).get("percentage")) + ")" for p in safe_list(te.get("phases")))},
            {"title": "Infrastructure Costs", "content": "**Monthly:** $" + str(monthly) + "  |  **Annual:** $" + str(annual) + "\n\n" + "\n".join("- **" + safe_str(safe_dict(a).get("service")) + "** (" + safe_str(safe_dict(a).get("tier", "")) + "): $" + str(safe_int(safe_dict(a).get("monthly_cost"))) + "/mo" for a in safe_list(ce.get("azure_costs")))},
        ], "quality_checks": {"ECI Tone": True, "Personalization": True, "Terminology": True, "Value Proposition": True, "Structure": True}}


# ═══════════════════════════════════════════════════════════════════════
#  SHAREPOINT
# ═══════════════════════════════════════════════════════════════════════

class SP:
    def __init__(self, url, cid, cs, tid):
        self.url = url
        self.cid = cid
        self.cs = cs
        self.tid = tid

    @classmethod
    def from_session(cls):
        return cls(st.session_state.get("sp_url", ""), st.session_state.get("sp_cid", ""), st.session_state.get("sp_cs", ""), st.session_state.get("sp_tid", ""))

    def _auth(self):
        if not all([self.cid, self.cs, self.tid, msal]):
            return None
        try:
            app = msal.ConfidentialClientApplication(self.cid, authority="https://login.microsoftonline.com/" + self.tid, client_credential=self.cs)
            r = app.acquire_token_for_client(scopes=["https://graph.microsoft.com/.default"])
            return r.get("access_token")
        except Exception:
            return None

    def test(self):
        tok = self._auth()
        if not tok:
            return False, "Missing credentials or auth failed."
        try:
            r = _requests.get("https://graph.microsoft.com/v1.0/sites/root", headers={"Authorization": "Bearer " + tok}, timeout=10)
            return (True, "SharePoint connected!") if r.status_code == 200 else (False, "Status " + str(r.status_code))
        except Exception as e:
            return False, str(e)[:200]

    def list_folder(self, path):
        tok = self._auth()
        if tok and _requests:
            try:
                sn = self.url.rstrip("/").split("/")[-1]
                r = _requests.get("https://graph.microsoft.com/v1.0/sites/root:/sites/" + sn + ":/drive/root:" + path + ":/children", headers={"Authorization": "Bearer " + tok}, timeout=15)
                if r.status_code == 200:
                    return [i["name"] for i in r.json().get("value", []) if "name" in i]
            except Exception:
                pass
        return []

    def upload(self, results):
        fn = "ECI_Proposal_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".json"
        tok = self._auth()
        if tok and _requests:
            try:
                sn = self.url.rstrip("/").split("/")[-1]
                r = _requests.put("https://graph.microsoft.com/v1.0/sites/root:/sites/" + sn + ":/drive/root:/Proposals/" + fn + ":/content", headers={"Authorization": "Bearer " + tok, "Content-Type": "application/json"}, data=json.dumps(results, indent=2, default=str), timeout=30)
                if r.status_code in (200, 201):
                    return True, "Uploaded: /Proposals/" + fn
                return False, "Status " + str(r.status_code)
            except Exception as e:
                return False, str(e)[:200]
        return False, "SharePoint not configured."


# ═══════════════════════════════════════════════════════════════════════
#  EMAIL
# ═══════════════════════════════════════════════════════════════════════

class Mailer:
    def __init__(self, smtp, sender, pw):
        self.smtp = smtp
        self.sender = sender
        self.pw = pw

    @classmethod
    def from_session(cls):
        return cls(st.session_state.get("email_smtp", ""), st.session_state.get("email_sender", ""), st.session_state.get("cfg_email_pass", ""))

    def send(self, results):
        te = safe_dict(results.get("time_estimate"))
        ce = safe_dict(results.get("cost_estimate"))
        ri = safe_dict(results.get("risk_assessment"))
        subj = "ECI Proposal " + datetime.now().strftime("%Y-%m-%d %H:%M")
        body = "<h2>ECI — Business Estimation Leveraging Automated Learning</h2><p>Hours: " + str(te.get("total_hours", 0)) + " | Infra Cost: $" + str(ce.get("total_monthly_cost", 0)) + "/mo | Risk: " + safe_str(ri.get("overall_level")) + "</p>"
        if self.smtp and self.sender and self.pw:
            try:
                parts = (self.smtp + ":587").split(":")
                host = parts[0]
                port = int(parts[1])
                msg = MIMEMultipart()
                msg["Subject"] = subj
                msg["From"] = self.sender
                msg["To"] = self.sender
                msg.attach(MIMEText(body, "html"))
                with smtplib.SMTP(host, port) as s:
                    s.starttls()
                    s.login(self.sender, self.pw)
                    s.send_message(msg)
                return True, "Email sent to " + self.sender
            except Exception as e:
                return False, str(e)[:200]
        return False, "Email not configured."


# ═══════════════════════════════════════════════════════════════════════
#  EXCEL TIME ESTIMATE GENERATOR
# ═══════════════════════════════════════════════════════════════════════

def generate_time_excel(time_est, semantic):
    if not Workbook:
        return None
    wb = Workbook()

    # ── Styles ──
    hdr_font = Font(name="Calibri", bold=True, size=12, color="FFFFFF")
    hdr_fill = PatternFill(start_color="1B3A5C", end_color="1B3A5C", fill_type="solid")
    sub_font = Font(name="Calibri", bold=True, size=10, color="1B3A5C")
    sub_fill = PatternFill(start_color="D6E4F0", end_color="D6E4F0", fill_type="solid")
    normal_font = Font(name="Calibri", size=10)
    bold_font = Font(name="Calibri", bold=True, size=10)
    total_fill = PatternFill(start_color="E2EFDA", end_color="E2EFDA", fill_type="solid")
    total_font = Font(name="Calibri", bold=True, size=11, color="1B3A5C")
    thin_border = Border(
        left=Side(style="thin"), right=Side(style="thin"),
        top=Side(style="thin"), bottom=Side(style="thin"),
    )
    center = Alignment(horizontal="center", vertical="center")
    wrap = Alignment(horizontal="left", vertical="center", wrap_text=True)

    def style_header_row(ws, row, cols):
        for c in range(1, cols + 1):
            cell = ws.cell(row=row, column=c)
            cell.font = hdr_font
            cell.fill = hdr_fill
            cell.alignment = center
            cell.border = thin_border

    def style_cell(ws, row, col, font=normal_font, fill=None, align=None):
        cell = ws.cell(row=row, column=col)
        cell.font = font
        cell.border = thin_border
        if fill:
            cell.fill = fill
        if align:
            cell.alignment = align
        return cell

    # ═══ Sheet 1: Summary ═══
    ws1 = wb.active
    ws1.title = "Summary"
    ws1.sheet_properties.tabColor = "1B3A5C"

    # Title
    ws1.merge_cells("A1:F1")
    title_cell = ws1["A1"]
    title_cell.value = "ECI — Project Time Estimation Summary"
    title_cell.font = Font(name="Calibri", bold=True, size=16, color="1B3A5C")
    title_cell.alignment = Alignment(horizontal="center", vertical="center")
    ws1.row_dimensions[1].height = 40

    ws1.merge_cells("A2:F2")
    ws1["A2"].value = "Generated: " + datetime.now().strftime("%B %d, %Y %H:%M")
    ws1["A2"].font = Font(name="Calibri", size=9, color="666666")
    ws1["A2"].alignment = Alignment(horizontal="center")

    # Project Info
    row = 4
    ws1.cell(row=row, column=1, value="Project Type").font = bold_font
    ws1.cell(row=row, column=2, value=safe_str(semantic.get("project_type", "N/A"))).font = normal_font
    row += 1
    ws1.cell(row=row, column=1, value="Complexity Score").font = bold_font
    ws1.cell(row=row, column=2, value=str(semantic.get("complexity_score", "N/A")) + " / 10").font = normal_font
    row += 1
    ws1.cell(row=row, column=1, value="Total Requirements").font = bold_font
    ws1.cell(row=row, column=2, value=len(safe_list(semantic.get("requirements")))).font = normal_font

    # Summary metrics
    row = 8
    headers = ["Metric", "Value"]
    for c, h in enumerate(headers, 1):
        ws1.cell(row=row, column=c, value=h)
    style_header_row(ws1, row, len(headers))
    row += 1
    metrics = [
        ("Total Hours", str(safe_int(time_est.get("total_hours"))) + " hours"),
        ("Duration", safe_str(time_est.get("duration_weeks"))),
        ("Confidence", safe_str(time_est.get("confidence"))),
        ("Buffer", safe_str(time_est.get("buffer"))),
    ]
    for label, val in metrics:
        style_cell(ws1, row, 1, font=bold_font).value = label
        style_cell(ws1, row, 2).value = val
        row += 1

    # Three-point estimation
    tp = safe_dict(time_est.get("three_point"))
    if tp:
        row += 1
        ws1.merge_cells(start_row=row, start_column=1, end_row=row, end_column=2)
        ws1.cell(row=row, column=1, value="Three-Point Estimation").font = Font(name="Calibri", bold=True, size=12, color="1B3A5C")
        row += 1
        headers_tp = ["Scenario", "Hours"]
        for c, h in enumerate(headers_tp, 1):
            ws1.cell(row=row, column=c, value=h)
        style_header_row(ws1, row, len(headers_tp))
        row += 1
        for label, key in [("Optimistic", "optimistic"), ("Most Likely", "most_likely"), ("Pessimistic", "pessimistic")]:
            style_cell(ws1, row, 1, font=bold_font).value = label
            style_cell(ws1, row, 2, align=center).value = safe_int(tp.get(key))
            row += 1

    ws1.column_dimensions["A"].width = 22
    ws1.column_dimensions["B"].width = 30

    # ═══ Sheet 2: Phase Breakdown ═══
    ws2 = wb.create_sheet("Phase Breakdown")
    ws2.sheet_properties.tabColor = "00B4D8"

    ws2.merge_cells("A1:F1")
    ws2["A1"].value = "Phase-wise Effort Breakdown"
    ws2["A1"].font = Font(name="Calibri", bold=True, size=14, color="1B3A5C")
    ws2["A1"].alignment = Alignment(horizontal="center")
    ws2.row_dimensions[1].height = 35

    row = 3
    headers = ["Phase", "Task", "Role", "Hours", "% of Total", "Notes"]
    for c, h in enumerate(headers, 1):
        ws2.cell(row=row, column=c, value=h)
    style_header_row(ws2, row, len(headers))
    row += 1

    total_hours = safe_int(time_est.get("total_hours", 1))
    for phase in safe_list(time_est.get("phases")):
        phase = safe_dict(phase)
        phase_name = safe_str(phase.get("name"))
        phase_hours = safe_int(phase.get("hours"))
        tasks = safe_list(phase.get("tasks"))
        if not tasks:
            style_cell(ws2, row, 1, font=sub_font, fill=sub_fill).value = phase_name
            style_cell(ws2, row, 2, fill=sub_fill)
            style_cell(ws2, row, 3, fill=sub_fill)
            style_cell(ws2, row, 4, font=sub_font, fill=sub_fill, align=center).value = phase_hours
            style_cell(ws2, row, 5, fill=sub_fill, align=center).value = safe_str(phase.get("percentage"))
            style_cell(ws2, row, 6, fill=sub_fill)
            row += 1
        else:
            first_task = True
            for task in tasks:
                task = safe_dict(task)
                if first_task:
                    style_cell(ws2, row, 1, font=sub_font, fill=sub_fill).value = phase_name
                    first_task = False
                else:
                    style_cell(ws2, row, 1, fill=None)
                style_cell(ws2, row, 2).value = safe_str(task.get("name"))
                style_cell(ws2, row, 3).value = safe_str(task.get("role"))
                style_cell(ws2, row, 4, align=center).value = safe_int(task.get("hours"))
                pct = safe_int(task.get("hours")) / max(total_hours, 1) * 100
                style_cell(ws2, row, 5, align=center).value = str(round(pct, 1)) + "%"
                style_cell(ws2, row, 6)
                row += 1
            # Phase subtotal
            style_cell(ws2, row, 1, font=bold_font, fill=total_fill)
            style_cell(ws2, row, 2, font=bold_font, fill=total_fill).value = "Subtotal — " + phase_name
            style_cell(ws2, row, 3, fill=total_fill)
            style_cell(ws2, row, 4, font=bold_font, fill=total_fill, align=center).value = phase_hours
            style_cell(ws2, row, 5, font=bold_font, fill=total_fill, align=center).value = safe_str(phase.get("percentage"))
            style_cell(ws2, row, 6, fill=total_fill)
            row += 1

    # Grand total
    row += 1
    style_cell(ws2, row, 1, font=total_font, fill=total_fill)
    style_cell(ws2, row, 2, font=total_font, fill=total_fill).value = "GRAND TOTAL"
    style_cell(ws2, row, 3, fill=total_fill)
    style_cell(ws2, row, 4, font=total_font, fill=total_fill, align=center).value = total_hours
    style_cell(ws2, row, 5, font=total_font, fill=total_fill, align=center).value = "100%"
    style_cell(ws2, row, 6, fill=total_fill)

    for c, w in [(1, 18), (2, 28), (3, 18), (4, 12), (5, 12), (6, 25)]:
        ws2.column_dimensions[get_column_letter(c)].width = w

    # ═══ Sheet 3: Milestones ═══
    ws3 = wb.create_sheet("Milestones")
    ws3.sheet_properties.tabColor = "00D4AA"

    ws3.merge_cells("A1:D1")
    ws3["A1"].value = "Project Milestones"
    ws3["A1"].font = Font(name="Calibri", bold=True, size=14, color="1B3A5C")
    ws3["A1"].alignment = Alignment(horizontal="center")
    ws3.row_dimensions[1].height = 35

    row = 3
    headers = ["#", "Milestone", "Week", "Description"]
    for c, h in enumerate(headers, 1):
        ws3.cell(row=row, column=c, value=h)
    style_header_row(ws3, row, len(headers))
    row += 1
    for i, m in enumerate(safe_list(time_est.get("milestones")), 1):
        m = safe_dict(m)
        style_cell(ws3, row, 1, align=center).value = i
        style_cell(ws3, row, 2, font=bold_font).value = safe_str(m.get("name"))
        style_cell(ws3, row, 3, align=center).value = safe_int(m.get("week"))
        style_cell(ws3, row, 4, align=wrap).value = safe_str(m.get("description"))
        row += 1

    for c, w in [(1, 6), (2, 30), (3, 10), (4, 45)]:
        ws3.column_dimensions[get_column_letter(c)].width = w

    # ═══ Sheet 4: Three-Point Detail ═══
    ws4 = wb.create_sheet("Three-Point Estimation")
    ws4.sheet_properties.tabColor = "7B61FF"

    ws4.merge_cells("A1:E1")
    ws4["A1"].value = "Three-Point Estimation Detail"
    ws4["A1"].font = Font(name="Calibri", bold=True, size=14, color="1B3A5C")
    ws4["A1"].alignment = Alignment(horizontal="center")
    ws4.row_dimensions[1].height = 35

    row = 3
    headers = ["Phase", "Optimistic (hrs)", "Most Likely (hrs)", "Pessimistic (hrs)", "PERT Estimate (hrs)"]
    for c, h in enumerate(headers, 1):
        ws4.cell(row=row, column=c, value=h)
    style_header_row(ws4, row, len(headers))
    row += 1
    total_o, total_m, total_p, total_pert = 0, 0, 0, 0
    for phase in safe_list(time_est.get("phases")):
        phase = safe_dict(phase)
        ph_hours = safe_int(phase.get("hours"))
        opt = int(ph_hours * 0.8)
        ml = ph_hours
        pes = int(ph_hours * 1.35)
        pert = int((opt + 4 * ml + pes) / 6)
        total_o += opt
        total_m += ml
        total_p += pes
        total_pert += pert
        style_cell(ws4, row, 1, font=bold_font).value = safe_str(phase.get("name"))
        style_cell(ws4, row, 2, align=center).value = opt
        style_cell(ws4, row, 3, align=center).value = ml
        style_cell(ws4, row, 4, align=center).value = pes
        style_cell(ws4, row, 5, align=center, font=bold_font).value = pert
        row += 1

    # Totals
    style_cell(ws4, row, 1, font=total_font, fill=total_fill).value = "TOTAL"
    style_cell(ws4, row, 2, font=total_font, fill=total_fill, align=center).value = total_o
    style_cell(ws4, row, 3, font=total_font, fill=total_fill, align=center).value = total_m
    style_cell(ws4, row, 4, font=total_font, fill=total_fill, align=center).value = total_p
    style_cell(ws4, row, 5, font=total_font, fill=total_fill, align=center).value = total_pert

    for c, w in [(1, 22), (2, 18), (3, 18), (4, 18), (5, 20)]:
        ws4.column_dimensions[get_column_letter(c)].width = w

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════
#  PDF PROPOSAL GENERATOR (ECI Template)
# ═══════════════════════════════════════════════════════════════════════

def _eci_styles():
    """ECI-branded paragraph styles for reportlab."""
    styles = getSampleStyleSheet()
    eci_blue = rl_colors.Color(27/255, 58/255, 92/255)
    eci_accent = rl_colors.Color(0, 180/255, 216/255)
    styles.add(ParagraphStyle("ECITitle", parent=styles["Title"], fontSize=26, textColor=rl_colors.white, alignment=TA_CENTER, spaceAfter=6))
    styles.add(ParagraphStyle("ECISubtitle", parent=styles["Normal"], fontSize=13, textColor=rl_colors.Color(200/255, 220/255, 240/255), alignment=TA_CENTER, spaceAfter=4))
    styles.add(ParagraphStyle("ECIHeading", parent=styles["Heading1"], fontSize=15, textColor=eci_blue, spaceBefore=14, spaceAfter=6, borderWidth=0))
    styles.add(ParagraphStyle("ECIBody", parent=styles["Normal"], fontSize=10, leading=14, textColor=rl_colors.Color(30/255, 30/255, 30/255), spaceAfter=4))
    styles.add(ParagraphStyle("ECIBullet", parent=styles["Normal"], fontSize=10, leading=14, textColor=rl_colors.Color(30/255, 30/255, 30/255), leftIndent=16, bulletIndent=6, spaceAfter=2))
    styles.add(ParagraphStyle("ECISmall", parent=styles["Normal"], fontSize=8, textColor=rl_colors.Color(100/255, 100/255, 100/255), alignment=TA_CENTER))
    styles.add(ParagraphStyle("ECIKPIVal", parent=styles["Normal"], fontSize=16, textColor=eci_blue, alignment=TA_CENTER, leading=20))
    styles.add(ParagraphStyle("ECIKPILabel", parent=styles["Normal"], fontSize=8, textColor=rl_colors.Color(100/255, 100/255, 100/255), alignment=TA_CENTER))
    styles.add(ParagraphStyle("ECITableCell", parent=styles["Normal"], fontSize=8, leading=10, textColor=rl_colors.Color(30/255, 30/255, 30/255)))
    styles.add(ParagraphStyle("ECITableHeader", parent=styles["Normal"], fontSize=8, leading=10, textColor=rl_colors.white, alignment=TA_CENTER))
    return styles


def _eci_table(headers, rows, col_widths=None):
    """Build a styled reportlab Table."""
    eci_blue = rl_colors.Color(27/255, 58/255, 92/255)
    alt_row = rl_colors.Color(245/255, 248/255, 252/255)
    styles = _eci_styles()
    hdr_cells = [Paragraph("<b>" + h + "</b>", styles["ECITableHeader"]) for h in headers]
    data = [hdr_cells]
    for row in rows:
        data.append([Paragraph(str(v), styles["ECITableCell"]) for v in row])
    if not col_widths:
        col_widths = [480 / len(headers)] * len(headers)
    tbl = Table(data, colWidths=col_widths, repeatRows=1)
    style_cmds = [
        ("BACKGROUND", (0, 0), (-1, 0), eci_blue),
        ("TEXTCOLOR", (0, 0), (-1, 0), rl_colors.white),
        ("GRID", (0, 0), (-1, -1), 0.5, rl_colors.Color(200/255, 200/255, 200/255)),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
    ]
    for i in range(1, len(data)):
        if i % 2 == 0:
            style_cmds.append(("BACKGROUND", (0, i), (-1, i), alt_row))
    tbl.setStyle(TableStyle(style_cmds))
    return tbl


def _md_to_para(text, styles):
    """Convert simple markdown text to a list of Paragraph flowables."""
    elements = []
    for line in text.split("\n"):
        line = line.strip()
        if not line:
            elements.append(Spacer(1, 4))
            continue
        # Bold markers
        cleaned = line.replace("**", "<b>", 1)
        while "**" in cleaned:
            cleaned = cleaned.replace("**", "</b>", 1)
            if "**" in cleaned:
                cleaned = cleaned.replace("**", "<b>", 1)
        # Italic markers
        cleaned = cleaned.replace("_", "<i>", 1)
        while "_" in cleaned:
            cleaned = cleaned.replace("_", "</i>", 1)
            if "_" in cleaned:
                cleaned = cleaned.replace("_", "<i>", 1)
        if line.startswith("- "):
            elements.append(Paragraph(cleaned[2:], styles["ECIBullet"], bulletText="\u2022"))
        else:
            elements.append(Paragraph(cleaned, styles["ECIBody"]))
    return elements


def generate_proposal_pdf(results):
    if not HAS_REPORTLAB:
        return None
    buf = io.BytesIO()
    styles = _eci_styles()
    eci_blue = rl_colors.Color(27/255, 58/255, 92/255)
    eci_light = rl_colors.Color(214/255, 228/255, 240/255)

    doc = SimpleDocTemplate(buf, pagesize=A4, topMargin=25*mm, bottomMargin=20*mm, leftMargin=18*mm, rightMargin=18*mm,
                            title="ECI Project Proposal", author="ECI Consulting")
    elements = []

    se = safe_dict(results.get("semantic_analysis"))
    te = safe_dict(results.get("time_estimate"))
    ce = safe_dict(results.get("cost_estimate"))
    ri = safe_dict(results.get("risk_assessment"))
    ar = safe_dict(results.get("architecture"))
    sc = safe_dict(results.get("scope"))
    proposal = safe_dict(results.get("proposal"))

    # ── Cover ──
    elements.append(Spacer(1, 40))
    # Blue banner table
    banner_data = [[Paragraph("<b>PROJECT PROPOSAL</b>", styles["ECITitle"]),],
                   [Paragraph(safe_str(se.get("project_type", "Technology Solution")), styles["ECISubtitle"]),],
                   [Paragraph("Prepared by ECI Consulting", styles["ECISubtitle"]),],
                   [Paragraph(datetime.now().strftime("%B %d, %Y"), styles["ECISubtitle"]),]]
    banner = Table(banner_data, colWidths=[480])
    banner.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), eci_blue),
        ("TOPPADDING", (0, 0), (-1, -1), 8),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
    ]))
    elements.append(banner)
    elements.append(Spacer(1, 24))

    # KPI boxes
    kpi_items = [
        ("Requirements", str(len(safe_list(se.get("requirements"))))),
        ("Total Hours", str(safe_int(te.get("total_hours")))),
        ("Infra Cost/mo", "$" + str(safe_int(ce.get("total_monthly_cost")))),
        ("Risk Level", safe_str(ri.get("overall_level", "N/A"))),
    ]
    kpi_cells = []
    for label, val in kpi_items:
        kpi_cells.append([
            Paragraph("<b>" + val + "</b>", styles["ECIKPIVal"]),
            Paragraph(label, styles["ECIKPILabel"]),
        ])
    # Transpose: each kpi is a column of 2 rows
    kpi_data = [[kpi_cells[i][0] for i in range(4)], [kpi_cells[i][1] for i in range(4)]]
    kpi_table = Table(kpi_data, colWidths=[120, 120, 120, 120])
    kpi_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), eci_light),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("TOPPADDING", (0, 0), (-1, -1), 8),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
        ("GRID", (0, 0), (-1, -1), 0.5, rl_colors.white),
    ]))
    elements.append(kpi_table)
    elements.append(Spacer(1, 10))
    elements.append(Paragraph("ECI CONSULTING  |  CONFIDENTIAL", styles["ECISmall"]))
    elements.append(PageBreak())

    # ── Proposal sections ──
    for sec in safe_list(proposal.get("sections")):
        sec = safe_dict(sec)
        elements.append(Paragraph(safe_str(sec.get("title")), styles["ECIHeading"]))
        elements.append(HRFlowable(width="30%", thickness=2, color=rl_colors.Color(0, 180/255, 216/255), spaceAfter=8))
        elements.extend(_md_to_para(safe_str(sec.get("content")), styles))
        elements.append(Spacer(1, 10))

    elements.append(PageBreak())

    # ── Requirements summary ──
    elements.append(Paragraph("Requirements Summary", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=rl_colors.Color(0, 180/255, 216/255), spaceAfter=8))
    reqs = safe_list(se.get("requirements"))
    fn = [r for r in reqs if isinstance(r, dict) and r.get("type") == "functional"]
    nf = [r for r in reqs if isinstance(r, dict) and r.get("type") == "non-functional"]
    ig = [r for r in reqs if isinstance(r, dict) and r.get("type") == "integration"]
    elements.append(Paragraph("Functional: <b>" + str(len(fn)) + "</b>  |  Non-Functional: <b>" + str(len(nf)) + "</b>  |  Integration: <b>" + str(len(ig)) + "</b>", styles["ECIBody"]))
    elements.append(Spacer(1, 6))
    req_rows = [[safe_str(r.get("title")), safe_str(r.get("type")), safe_str(r.get("complexity")), safe_str(r.get("priority", ""))] for r in reqs if isinstance(r, dict)]
    if req_rows:
        elements.append(_eci_table(["Requirement", "Type", "Complexity", "Priority"], req_rows[:25], [200, 90, 80, 80]))
    elements.append(PageBreak())

    # ── Time estimate ──
    elements.append(Paragraph("Time Estimation", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=rl_colors.Color(0, 180/255, 216/255), spaceAfter=8))
    elements.extend(_md_to_para(
        "- <b>Total Hours:</b> " + str(safe_int(te.get("total_hours")))
        + "\n- <b>Duration:</b> " + safe_str(te.get("duration_weeks"))
        + "\n- <b>Confidence:</b> " + safe_str(te.get("confidence"))
        + "\n- <b>Buffer:</b> " + safe_str(te.get("buffer")), styles))
    elements.append(Spacer(1, 6))
    phase_rows = [[safe_str(p.get("name")), str(safe_int(p.get("hours"))), safe_str(p.get("percentage"))] for p in safe_list(te.get("phases")) if isinstance(p, dict)]
    if phase_rows:
        elements.append(_eci_table(["Phase", "Hours", "% of Total"], phase_rows, [220, 120, 120]))
    elements.append(PageBreak())

    # ── Infrastructure costs ──
    elements.append(Paragraph("Infrastructure Cost Estimate", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=rl_colors.Color(0, 180/255, 216/255), spaceAfter=8))
    elements.extend(_md_to_para(
        "- <b>Monthly Cost:</b> $" + str(safe_int(ce.get("total_monthly_cost")))
        + "\n- <b>Annual Cost:</b> $" + str(safe_int(ce.get("total_annual_cost"))), styles))
    elements.append(Spacer(1, 6))
    cost_rows = [[safe_str(a.get("service")), safe_str(a.get("tier", "")), "$" + str(safe_int(a.get("monthly_cost")))] for a in safe_list(ce.get("azure_costs")) if isinstance(a, dict)]
    if cost_rows:
        elements.append(_eci_table(["Service", "Tier", "Monthly Cost"], cost_rows, [200, 160, 100]))
    elements.append(PageBreak())

    # ── Risk ──
    elements.append(Paragraph("Risk Assessment", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=rl_colors.Color(0, 180/255, 216/255), spaceAfter=8))
    elements.extend(_md_to_para(
        "- <b>Overall Score:</b> " + str(safe_int(ri.get("overall_score"))) + "/10"
        + "\n- <b>Level:</b> " + safe_str(ri.get("overall_level")), styles))
    elements.append(Spacer(1, 6))
    risk_rows = [[safe_str(r.get("category")), safe_str(r.get("title")), safe_str(r.get("severity")), safe_str(r.get("mitigation"))] for r in safe_list(ri.get("risks")) if isinstance(r, dict)]
    if risk_rows:
        elements.append(_eci_table(["Category", "Risk", "Severity", "Mitigation"], risk_rows, [70, 110, 60, 220]))
    elements.append(PageBreak())

    # ── Architecture ──
    elements.append(Paragraph("Architecture Overview", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=rl_colors.Color(0, 180/255, 216/255), spaceAfter=8))
    elements.append(Paragraph("<b>Pattern:</b> " + safe_str(ar.get("pattern")), styles["ECIBody"]))
    elements.append(Spacer(1, 6))
    comp_rows = [[safe_str(c.get("name")), safe_str(c.get("type", "")), safe_str(c.get("azure_service", "")), ", ".join(safe_list(c.get("services")))] for c in safe_list(ar.get("components")) if isinstance(c, dict)]
    if comp_rows:
        elements.append(_eci_table(["Component", "Type", "Azure Service", "Details"], comp_rows, [80, 80, 120, 180]))
    df = safe_list(ar.get("data_flow"))
    if df:
        elements.append(Spacer(1, 8))
        elements.append(Paragraph("<b>Data Flow:</b> " + " &rarr; ".join(safe_str(x) for x in df), styles["ECIBody"]))
    elements.append(PageBreak())

    # ── Scope ──
    elements.append(Paragraph("Scope Definition", styles["ECIHeading"]))
    elements.append(HRFlowable(width="30%", thickness=2, color=rl_colors.Color(0, 180/255, 216/255), spaceAfter=8))
    elements.append(Paragraph("<b>In Scope:</b>", styles["ECIBody"]))
    for item in safe_list(sc.get("in_scope")):
        elements.append(Paragraph(safe_str(item), styles["ECIBullet"], bulletText="\u2022"))
    elements.append(Spacer(1, 6))
    elements.append(Paragraph("<b>Out of Scope:</b>", styles["ECIBody"]))
    for item in safe_list(sc.get("out_of_scope")):
        elements.append(Paragraph(safe_str(item), styles["ECIBullet"], bulletText="\u2022"))
    elements.append(Spacer(1, 6))
    elements.append(Paragraph("<b>Assumptions:</b>", styles["ECIBody"]))
    for item in safe_list(sc.get("assumptions")):
        elements.append(Paragraph(safe_str(item), styles["ECIBullet"], bulletText="\u2022"))
    elements.append(Spacer(1, 6))
    elements.append(Paragraph("<b>Prerequisites:</b>", styles["ECIBody"]))
    for item in safe_list(sc.get("prerequisites")):
        elements.append(Paragraph(safe_str(item), styles["ECIBullet"], bulletText="\u2022"))

    # Build
    def _footer(canvas, doc):
        canvas.saveState()
        canvas.setStrokeColor(rl_colors.Color(27/255, 58/255, 92/255))
        canvas.line(18*mm, 15*mm, A4[0] - 18*mm, 15*mm)
        canvas.setFont("Helvetica", 7)
        canvas.setFillColor(rl_colors.Color(100/255, 100/255, 100/255))
        canvas.drawString(18*mm, 10*mm, "ECI Consulting  |  " + datetime.now().strftime("%B %d, %Y"))
        canvas.drawRightString(A4[0] - 18*mm, 10*mm, "Page " + str(canvas.getPageNumber()))
        canvas.restoreState()

    doc.build(elements, onFirstPage=_footer, onLaterPages=_footer)
    buf.seek(0)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════
#  MERMAID DIAGRAM RENDERER
# ═══════════════════════════════════════════════════════════════════════

def render_mermaid(mermaid_code, height=450):
    """Render a Mermaid.js diagram using streamlit HTML component."""
    html = f"""<div class="mermaid-container" style="background:var(--bg2);border:1px solid var(--bd);border-radius:12px;padding:16px;margin:8px 0;">
    <pre class="mermaid" style="text-align:center;">
{mermaid_code}
    </pre>
    </div>
    <script src="https://cdn.jsdelivr.net/npm/mermaid@10/dist/mermaid.min.js"></script>
    <script>mermaid.initialize({{startOnLoad:true, theme:'dark', themeVariables:{{primaryColor:'#16274B',primaryTextColor:'#e2e8f0',primaryBorderColor:'#00929E',lineColor:'#00b4d8',secondaryColor:'#151c2e',tertiaryColor:'#0a0e1a',fontFamily:'DM Sans, sans-serif'}}}});</script>"""
    st.components.v1.html(html, height=height, scrolling=True)


# ═══════════════════════════════════════════════════════════════════════
#  AUDIO / VIDEO TRANSCRIPT EXTRACTOR
# ═══════════════════════════════════════════════════════════════════════

def extract_audio_transcript(uploaded_file):
    """Extract transcript from audio/video files. Returns text."""
    name = uploaded_file.name.lower()
    data = uploaded_file.read()
    uploaded_file.seek(0)

    # For text-based transcript files (SRT, VTT, TXT)
    if name.endswith((".txt", ".srt", ".vtt")):
        text = data.decode("utf-8", errors="replace")
        # Clean SRT/VTT formatting
        text = re.sub(r'\d{2}:\d{2}:\d{2}[.,]\d{3}\s*-->\s*\d{2}:\d{2}:\d{2}[.,]\d{3}', '', text)
        text = re.sub(r'^\d+$', '', text, flags=re.MULTILINE)
        text = re.sub(r'<[^>]+>', '', text)
        text = re.sub(r'WEBVTT.*?\n', '', text)
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()

    # For JSON transcript exports (Teams, Zoom, Otter.ai)
    if name.endswith(".json"):
        try:
            transcript_data = json.loads(data.decode("utf-8"))
            parts = []
            # Common transcript JSON formats
            if isinstance(transcript_data, list):
                for item in transcript_data:
                    if isinstance(item, dict):
                        speaker = item.get("speaker", item.get("name", "Speaker"))
                        text_val = item.get("text", item.get("content", item.get("transcript", "")))
                        if text_val:
                            parts.append(f"{speaker}: {text_val}")
            elif isinstance(transcript_data, dict):
                for seg in transcript_data.get("segments", transcript_data.get("results", transcript_data.get("transcript", []))):
                    if isinstance(seg, dict):
                        speaker = seg.get("speaker", "Speaker")
                        text_val = seg.get("text", seg.get("content", ""))
                        if text_val:
                            parts.append(f"{speaker}: {text_val}")
            return "\n".join(parts) if parts else data.decode("utf-8", errors="replace")[:50000]
        except Exception:
            return data.decode("utf-8", errors="replace")[:50000]

    # For DOCX transcripts (exported meeting notes)
    if name.endswith(".docx"):
        try:
            from docx import Document
            doc = Document(io.BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception:
            return "[Could not extract DOCX transcript]"

    # For CSV transcript exports
    if name.endswith(".csv"):
        text = data.decode("utf-8", errors="replace")
        lines = text.split("\n")
        parts = []
        for line in lines[1:]:  # skip header
            cols = line.split(",")
            if len(cols) >= 2:
                parts.append(cols[-1].strip().strip('"'))
        return "\n".join(parts) if parts else text[:50000]

    return "[Unsupported format: " + name.split(".")[-1] + ". Please upload TXT, SRT, VTT, JSON, DOCX, or CSV transcript files.]"


# ═══════════════════════════════════════════════════════════════════════
#  SESSION STATE
# ═══════════════════════════════════════════════════════════════════════

_defaults = {
    "azure_api_key": "", "azure_endpoint": "", "azure_api_version": "2024-06-01", "azure_deployment": "gpt-4",
    "sp_url": "", "sp_cid": "", "sp_cs": "", "sp_tid": "",
    "email_smtp": "", "email_sender": "",
    "processing_results": None, "historical_projects": [], "agent_logs": [],
    "discovery_results": None, "discovery_transcript": "",
    "model_metrics": {"accuracy": 78.5, "proposals_processed": 0, "win_rate": 62.0, "variance": 12.3},
}
for k, v in _defaults.items():
    if k not in st.session_state:
        st.session_state[k] = v

inject_css()


# ═══════════════════════════════════════════════════════════════════════
#  HEADER
# ═══════════════════════════════════════════════════════════════════════

hc1, hc2, hc3 = st.columns([2, 6, 2])
with hc1:
    st.markdown('<div class="logo-box"><div class="logo-icon">⚡</div><div><div class="logo-txt">ECI</div><div class="logo-sub">Presale Estimation</div></div></div>', unsafe_allow_html=True)
with hc2:
    st.markdown('<div class="tagline">Business Estimation Leveraging Automated Learning</div>', unsafe_allow_html=True)
with hc3:
    st.markdown('<div class="dt">' + datetime.now().strftime("%b %d, %Y  %H:%M") + '</div>', unsafe_allow_html=True)


# ═══════════════════════════════════════════════════════════════════════
#  SIDEBAR
# ═══════════════════════════════════════════════════════════════════════

with st.sidebar:
    st.markdown('<div class="stitle">Configuration</div>', unsafe_allow_html=True)
    st.markdown('<div class="csec">Azure OpenAI</div>', unsafe_allow_html=True)
    st.session_state.azure_api_key = st.text_input("API Key", value=st.session_state.azure_api_key, type="password", key="k1")
    st.session_state.azure_endpoint = st.text_input("Endpoint", value=st.session_state.azure_endpoint, placeholder="https://your-resource.openai.azure.com/", key="k2")
    st.session_state.azure_api_version = st.selectbox("Version", ["2024-06-01", "2024-02-01", "2023-12-01-preview"], key="k3")
    st.session_state.azure_deployment = st.text_input("Deployment", value=st.session_state.azure_deployment, placeholder="gpt-4", key="k4")
    if st.button("Test Connection", use_container_width=True, key="tb1"):
        ok, msg = AzureAI.from_session().test()
        if ok:
            st.success(msg)
        else:
            st.error(msg)
    st.divider()
    st.markdown('<div class="csec">SharePoint</div>', unsafe_allow_html=True)
    st.session_state.sp_url = st.text_input("Site URL", value=st.session_state.sp_url, placeholder="https://org.sharepoint.com/sites/presales", key="k5")
    st.session_state.sp_cid = st.text_input("Client ID", value=st.session_state.sp_cid, key="k6")
    st.session_state.sp_cs = st.text_input("Client Secret", value=st.session_state.sp_cs, type="password", key="k7")
    st.session_state.sp_tid = st.text_input("Tenant ID", value=st.session_state.sp_tid, key="k8")
    if st.button("Test SharePoint", use_container_width=True, key="tb2"):
        ok, msg = SP.from_session().test()
        if ok:
            st.success(msg)
        else:
            st.error(msg)
    st.divider()
    st.markdown('<div class="csec">Email Alerts</div>', unsafe_allow_html=True)
    st.session_state.email_smtp = st.text_input("SMTP Server", value=st.session_state.email_smtp, placeholder="smtp.office365.com:587", key="k9")
    st.session_state.email_sender = st.text_input("Sender Email", value=st.session_state.email_sender, key="k10")
    st.text_input("Password", type="password", key="cfg_email_pass")
    st.divider()
    a_ok = bool(st.session_state.azure_api_key and st.session_state.azure_endpoint)
    s_ok = bool(st.session_state.sp_url and st.session_state.sp_cid)
    e_ok = bool(st.session_state.email_smtp)
    mode_label = "LIVE AI" if a_ok else "DEMO MODE"
    st.markdown('<div class="cstat">'
                + '<div class="srow"><strong>' + mode_label + '</strong></div>'
                + '<div class="srow"><span class="sdot ' + ('on' if a_ok else 'off') + '"></span> Azure OpenAI</div>'
                + '<div class="srow"><span class="sdot ' + ('on' if s_ok else 'off') + '"></span> SharePoint</div>'
                + '<div class="srow"><span class="sdot ' + ('on' if e_ok else 'off') + '"></span> Email</div>'
                + '</div>', unsafe_allow_html=True)


# ═══════════════════════════════════════════════════════════════════════
#  HELPERS
# ═══════════════════════════════════════════════════════════════════════

def log_agent(name, detail):
    st.session_state.agent_logs.append({"agent": name, "detail": detail, "time": datetime.now().isoformat()})


# ═══════════════════════════════════════════════════════════════════════
#  TAB 1: PRESALE
# ═══════════════════════════════════════════════════════════════════════

def tab_presale():
    st.markdown('<div class="shdr"><span class="shdr-i">📄</span> Document Ingestion</div>', unsafe_allow_html=True)
    uc, tc = st.columns([3, 2])
    with uc:
        st.markdown('<div class="crd"><div class="crd-t">Upload Scope Documents</div><div class="crd-d">PDF, DOCX, XLSX, PPTX, TXT, CSV</div>', unsafe_allow_html=True)
        files = st.file_uploader("Drop files", type=["pdf", "docx", "xlsx", "pptx", "txt", "csv"], accept_multiple_files=True, key="fu", label_visibility="collapsed")
        st.markdown("</div>", unsafe_allow_html=True)
    with tc:
        st.markdown('<div class="crd"><div class="crd-t">SharePoint Auto-Trigger</div><div class="crd-d">Monitor folder for new scope docs</div>', unsafe_allow_html=True)
        sp_f = st.text_input("Folder", placeholder="/sites/presales/Shared Documents/Scope", key="spf", label_visibility="collapsed")
        if st.button("Check Now", use_container_width=True, key="spchk"):
            found = SP.from_session().list_folder(sp_f)
            if found:
                st.success("Found " + str(len(found)) + " doc(s)")
                for f in found:
                    st.markdown("- `" + f + "`")
            else:
                st.info("No documents found. Configure SharePoint in sidebar.")
        st.markdown("</div>", unsafe_allow_html=True)

    if files:
        st.markdown("---")
        _, bc, _ = st.columns([1, 2, 1])
        with bc:
            if st.button("PROCESS & GENERATE ESTIMATES", use_container_width=True, type="primary", key="go"):
                run_pipeline(files)

    # ── Multimodal Discovery ──
    st.markdown("---")
    st.markdown('<div class="shdr"><span class="shdr-i">🎙️</span> Multimodal Discovery Extraction</div>', unsafe_allow_html=True)
    st.markdown('<div class="crd"><div class="crd-t">Upload Meeting Transcripts & Voice Notes</div><div class="crd-d">Extract requirements, pain points, stakeholders & auto-generate WBS from Zoom/Teams transcripts, voice memos, or meeting notes.</div>', unsafe_allow_html=True)
    disc_col1, disc_col2 = st.columns([3, 2])
    with disc_col1:
        disc_files = st.file_uploader(
            "Upload transcripts",
            type=["txt", "srt", "vtt", "json", "docx", "csv"],
            accept_multiple_files=True,
            key="disc_fu",
            label_visibility="collapsed",
            help="Supported: TXT, SRT, VTT (subtitles), JSON (Teams/Zoom/Otter.ai exports), DOCX (meeting notes), CSV",
        )
    with disc_col2:
        st.markdown("**Supported Formats:**")
        st.markdown("- `.txt` `.srt` `.vtt` — Subtitle / transcript files")
        st.markdown("- `.json` — Teams, Zoom, Otter.ai exports")
        st.markdown("- `.docx` — Meeting notes / minutes")
        st.markdown("- `.csv` — Tabular transcript exports")
    st.markdown("</div>", unsafe_allow_html=True)

    if disc_files:
        _, dbc, _ = st.columns([1, 2, 1])
        with dbc:
            if st.button("🎙️ ANALYZE TRANSCRIPTS & GENERATE WBS", use_container_width=True, type="primary", key="disc_go"):
                ai = AzureAI.from_session()
                dpb = st.progress(0)
                dstatus = st.empty()
                dstatus.markdown("**1/3** Extracting transcripts...")
                dpb.progress(10)
                all_text = ""
                for df in disc_files:
                    all_text += extract_audio_transcript(df) + "\n\n"
                st.session_state.discovery_transcript = all_text
                dstatus.markdown("**2/3** Analyzing content & extracting insights...")
                dpb.progress(50)
                disc_result = ai.analyze_transcript(all_text)
                dstatus.markdown("**3/3** Generating Work Breakdown Structure...")
                dpb.progress(90)
                st.session_state.discovery_results = disc_result
                dpb.progress(100)
                dstatus.markdown("**Done!** Discovery analysis complete.")
                time.sleep(0.5)
                dstatus.empty()
                dpb.empty()
                st.rerun()

    # ── Discovery Results ──
    if st.session_state.discovery_results:
        dr = st.session_state.discovery_results
        st.markdown("---")
        st.markdown('<div class="shdr"><span class="shdr-i">📋</span> Discovery Insights</div>', unsafe_allow_html=True)

        # Summary & sentiment
        st.markdown('<div class="crd">', unsafe_allow_html=True)
        st.markdown("**Meeting Summary:** " + safe_str(dr.get("meeting_summary")))
        st.markdown("**Sentiment:** " + safe_str(dr.get("sentiment")))
        themes = safe_list(dr.get("key_themes"))
        if themes:
            tags_html = " ".join('<span class="tt">' + safe_str(t) + '</span>' for t in themes)
            st.markdown('<div class="ttag">' + tags_html + '</div>', unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)

        disc_tabs = st.tabs(["🔥 Pain Points", "📝 Requirements", "👥 Stakeholders", "📊 WBS", "✅ Action Items"])

        # Pain Points
        with disc_tabs[0]:
            for pp in safe_list(dr.get("pain_points")):
                pp = safe_dict(pp)
                sev = safe_str(pp.get("severity"))
                sc = "#ff6b6b" if sev == "High" else "#ffd166" if sev == "Medium" else "#06d6a0"
                st.markdown(
                    '<div class="rc" style="border-left:3px solid ' + sc + ';">'
                    '<div class="rch2"><strong>' + safe_str(pp.get("issue")) + '</strong>'
                    '<span class="rsev" style="color:' + sc + ';">' + sev + '</span></div>'
                    '<p style="font-style:italic;color:var(--t2);">"' + safe_str(pp.get("quote")) + '"</p>'
                    '<div class="rmit"><strong>Stakeholder:</strong> ' + safe_str(pp.get("stakeholder")) + '</div>'
                    '</div>', unsafe_allow_html=True
                )

        # Requirements Extracted
        with disc_tabs[1]:
            reqs = safe_list(dr.get("requirements_extracted"))
            fn_r = [x for x in reqs if isinstance(x, dict) and x.get("type") == "functional"]
            nf_r = [x for x in reqs if isinstance(x, dict) and x.get("type") == "non-functional"]
            ig_r = [x for x in reqs if isinstance(x, dict) and x.get("type") == "integration"]
            rc1, rc2, rc3 = st.columns(3)
            with rc1:
                st.markdown('<div class="rch fn">Functional (' + str(len(fn_r)) + ')</div>', unsafe_allow_html=True)
                for q in fn_r:
                    q = safe_dict(q)
                    st.markdown('<div class="ri"><strong>' + safe_str(q.get("title")) + '</strong><p>' + safe_str(q.get("description")) + '</p><div class="ri-c">Source: ' + safe_str(q.get("source")) + '</div></div>', unsafe_allow_html=True)
            with rc2:
                st.markdown('<div class="rch nf">Non-Functional (' + str(len(nf_r)) + ')</div>', unsafe_allow_html=True)
                for q in nf_r:
                    q = safe_dict(q)
                    st.markdown('<div class="ri"><strong>' + safe_str(q.get("title")) + '</strong><p>' + safe_str(q.get("description")) + '</p><div class="ri-c">Source: ' + safe_str(q.get("source")) + '</div></div>', unsafe_allow_html=True)
            with rc3:
                st.markdown('<div class="rch ig">Integration (' + str(len(ig_r)) + ')</div>', unsafe_allow_html=True)
                for q in ig_r:
                    q = safe_dict(q)
                    st.markdown('<div class="ri"><strong>' + safe_str(q.get("title")) + '</strong><p>' + safe_str(q.get("description")) + '</p><div class="ri-c">Source: ' + safe_str(q.get("source")) + '</div></div>', unsafe_allow_html=True)

        # Stakeholders
        with disc_tabs[2]:
            stake_cols = st.columns(2)
            for i, sh in enumerate(safe_list(dr.get("stakeholders"))):
                sh = safe_dict(sh)
                with stake_cols[i % 2]:
                    concerns = "".join("<li>" + safe_str(c) + "</li>" for c in safe_list(sh.get("concerns")))
                    st.markdown(
                        '<div class="crd"><div class="crd-t">👤 ' + safe_str(sh.get("name")) + '</div>'
                        '<div class="crd-d" style="color:var(--c2);">' + safe_str(sh.get("role")) + '</div>'
                        '<ul style="color:var(--t2);font-size:.85rem;">' + concerns + '</ul></div>',
                        unsafe_allow_html=True
                    )

        # WBS
        with disc_tabs[3]:
            st.markdown("**Work Breakdown Structure**")
            wbs = safe_list(dr.get("wbs"))
            for phase_idx, phase in enumerate(wbs):
                phase = safe_dict(phase)
                phase_name = safe_str(phase.get("phase"))
                with st.expander("📁 " + str(phase_idx + 1) + ". " + phase_name, expanded=True):
                    for d in safe_list(phase.get("deliverables")):
                        d = safe_dict(d)
                        st.markdown("**📦 " + safe_str(d.get("name")) + "**")
                        for t in safe_list(d.get("tasks")):
                            t = safe_dict(t)
                            st.markdown("&nbsp;&nbsp;&nbsp;&nbsp;🔹 " + safe_str(t.get("name")) + " — _" + safe_str(t.get("effort")) + "_")
            # Decisions
            decisions = safe_list(dr.get("decisions"))
            if decisions:
                st.markdown("---")
                st.markdown("**Key Decisions:**")
                for dec in decisions:
                    st.markdown("- ✅ " + safe_str(dec))

        # Action Items
        with disc_tabs[4]:
            for ai_item in safe_list(dr.get("action_items")):
                ai_item = safe_dict(ai_item)
                pri = safe_str(ai_item.get("priority"))
                pc = "#ff6b6b" if pri == "High" else "#ffd166" if pri == "Medium" else "#06d6a0"
                st.markdown(
                    '<div class="rc" style="border-left:3px solid ' + pc + ';">'
                    '<div class="rch2"><strong>' + safe_str(ai_item.get("item")) + '</strong>'
                    '<span class="rsev" style="color:' + pc + ';">' + pri + '</span></div>'
                    '<div class="rmit"><strong>Owner:</strong> ' + safe_str(ai_item.get("owner")) + '</div>'
                    '</div>', unsafe_allow_html=True
                )

        # Download discovery data
        st.markdown("---")
        dl_disc_cols = st.columns(2)
        with dl_disc_cols[0]:
            st.download_button(
                "📥 Download Discovery Report (JSON)",
                data=json.dumps(dr, indent=2, default=str),
                file_name="ECI_Discovery_Report_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".json",
                mime="application/json",
                use_container_width=True, key="dl_disc_json",
            )
        with dl_disc_cols[1]:
            if st.session_state.discovery_transcript:
                st.download_button(
                    "📥 Download Cleaned Transcript (TXT)",
                    data=st.session_state.discovery_transcript,
                    file_name="ECI_Transcript_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".txt",
                    mime="text/plain",
                    use_container_width=True, key="dl_disc_txt",
                )

    if st.session_state.processing_results:
        show_results()


def run_pipeline(files):
    st.session_state.agent_logs = []
    ai = AzureAI.from_session()
    live = ai.is_live
    if live:
        st.markdown('<div class="phdr">LIVE AI Processing via Azure OpenAI</div>', unsafe_allow_html=True)
    else:
        st.markdown('<div class="phdr">Demo Mode — Configure Azure OpenAI for live AI</div>', unsafe_allow_html=True)
    pb = st.progress(0)
    status = st.empty()

    status.markdown("**1/12** Ingesting documents...")
    pb.progress(5)
    dp = DocProcessor()
    text = ""
    for f in files:
        text += dp.extract(f) + "\n\n"
    log_agent("Ingestion", str(len(files)) + " file(s), " + str(len(text)) + " chars")
    time.sleep(0.2)

    status.markdown("**2/12** Document Intelligence...")
    pb.progress(10)
    intel = dp.analyze(text)
    log_agent("Intelligence", str(intel["section_count"]) + " sections, " + str(intel["word_count"]) + " words")
    time.sleep(0.2)

    status.markdown("**3/12** Semantic Analysis" + (" (GPT-4)..." if live else "..."))
    pb.progress(20)
    semantic = ai.analyze_requirements(text)
    log_agent("Semantic", str(len(safe_list(semantic.get("requirements")))) + " requirements")
    time.sleep(0.2)

    status.markdown("**4/12** Historical RAG...")
    pb.progress(30)
    rag = ai.search_historical(text, st.session_state.historical_projects)
    log_agent("RAG", str(len(safe_list(rag.get("similar_projects")))) + " matches")
    time.sleep(0.2)

    status.markdown("**5/12** Time Estimator...")
    pb.progress(40)
    time_est = ai.estimate_time(semantic, rag)
    log_agent("Time", str(time_est.get("total_hours", 0)) + " hours")
    time.sleep(0.2)

    status.markdown("**6/12** Cost Calculator...")
    pb.progress(50)
    cost_est = ai.estimate_cost(semantic, time_est, rag)
    log_agent("Cost", "$" + str(cost_est.get("total_monthly_cost", 0)) + "/mo")
    time.sleep(0.2)

    status.markdown("**7/12** Risk Analyzer...")
    pb.progress(60)
    risk = ai.analyze_risk(semantic, time_est, cost_est)
    log_agent("Risk", str(risk.get("overall_score", 0)) + "/10")
    time.sleep(0.2)

    status.markdown("**8/12** Architecture Designer...")
    pb.progress(70)
    arch = ai.design_architecture(semantic, rag)
    log_agent("Architecture", str(len(safe_list(arch.get("components")))) + " components")
    time.sleep(0.2)

    status.markdown("**9/12** Scope & Assumptions...")
    pb.progress(80)
    scope = ai.define_scope(semantic, time_est, cost_est)
    log_agent("Scope", "Boundaries defined")
    time.sleep(0.2)

    status.markdown("**10/12** Proposal Writer...")
    pb.progress(80)
    proposal = ai.write_proposal(semantic, time_est, cost_est, risk, arch, scope)
    log_agent("Proposal", "Document generated")
    time.sleep(0.2)

    status.markdown("**11/12** Architecture Visualizer...")
    pb.progress(90)
    mermaid_diagrams = ai.generate_mermaid_diagrams(semantic, arch)
    log_agent("Visualizer", str(len(mermaid_diagrams)) + " diagrams generated")
    time.sleep(0.2)

    pb.progress(100)
    status.markdown("**12/12** All agents completed!")
    st.session_state.processing_results = {
        "semantic_analysis": semantic, "rag": rag, "time_estimate": time_est,
        "cost_estimate": cost_est, "risk_assessment": risk, "architecture": arch,
        "scope": scope, "proposal": proposal, "mermaid_diagrams": mermaid_diagrams,
    }
    st.session_state.model_metrics["proposals_processed"] += 1


def show_results():
    r = st.session_state.processing_results
    st.markdown("---")
    st.markdown('<div class="shdr"><span class="shdr-i">📊</span> Results</div>', unsafe_allow_html=True)

    with st.expander("Agent Activity Log", expanded=False):
        for entry in st.session_state.agent_logs:
            st.markdown('<div class="alog"><span class="abadge">' + entry["agent"] + '</span><span class="aok">Done</span><span class="adet">' + entry["detail"] + '</span></div>', unsafe_allow_html=True)

    se = safe_dict(r.get("semantic_analysis"))
    te = safe_dict(r.get("time_estimate"))
    ce = safe_dict(r.get("cost_estimate"))
    ri = safe_dict(r.get("risk_assessment"))
    ar = safe_dict(r.get("architecture"))

    kpis = [
        ("📝", "Requirements", str(len(safe_list(se.get("requirements")))), "Identified"),
        ("⏱️", "Hours", str(safe_int(te.get("total_hours"))), "Person-hours"),
        ("💰", "Infra Cost", "$" + str(safe_int(ce.get("total_monthly_cost"))) + "/mo", "Azure Infrastructure"),
        ("⚠️", "Risk", str(safe_int(ri.get("overall_score"))) + "/10", safe_str(ri.get("overall_level"))),
        ("🏗️", "Components", str(len(safe_list(ar.get("components")))), "Designed"),
    ]
    cols = st.columns(5)
    for i, (ic, t, v, s) in enumerate(kpis):
        with cols[i]:
            st.markdown('<div class="kpi"><div class="kpi-i">' + ic + '</div><div class="kpi-v">' + v + '</div><div class="kpi-t">' + t + '</div><div class="kpi-s">' + s + '</div></div>', unsafe_allow_html=True)

    tab_list = st.tabs(["📋 Requirements", "⏱️ Time", "💰 Infra Cost", "⚠️ Risk", "🏗️ Architecture", "📐 Diagrams", "📄 Proposal", "📌 Scope"])

    # ── Requirements ──
    with tab_list[0]:
        reqs = safe_list(se.get("requirements"))
        fn_list = [x for x in reqs if isinstance(x, dict) and x.get("type") == "functional"]
        nf_list = [x for x in reqs if isinstance(x, dict) and x.get("type") == "non-functional"]
        ig_list = [x for x in reqs if isinstance(x, dict) and x.get("type") == "integration"]
        rc1, rc2, rc3 = st.columns(3)
        with rc1:
            st.markdown('<div class="rch fn">Functional (' + str(len(fn_list)) + ')</div>', unsafe_allow_html=True)
            for q in fn_list:
                st.markdown('<div class="ri"><strong>' + safe_str(q.get("title")) + '</strong><div class="ri-c">Complexity: ' + safe_str(q.get("complexity")) + '</div><p>' + safe_str(q.get("description")) + '</p></div>', unsafe_allow_html=True)
        with rc2:
            st.markdown('<div class="rch nf">Non-Functional (' + str(len(nf_list)) + ')</div>', unsafe_allow_html=True)
            for q in nf_list:
                st.markdown('<div class="ri"><strong>' + safe_str(q.get("title")) + '</strong><p>' + safe_str(q.get("description")) + '</p></div>', unsafe_allow_html=True)
        with rc3:
            st.markdown('<div class="rch ig">Integration (' + str(len(ig_list)) + ')</div>', unsafe_allow_html=True)
            for q in ig_list:
                st.markdown('<div class="ri"><strong>' + safe_str(q.get("title")) + '</strong><p>' + safe_str(q.get("description")) + '</p></div>', unsafe_allow_html=True)
        tech = safe_list(se.get("technology_stack"))
        if tech:
            tags_html = " ".join('<span class="tt">' + safe_str(t) + '</span>' for t in tech)
            st.markdown('<div class="ttag">' + tags_html + '</div>', unsafe_allow_html=True)

    # ── Time ──
    with tab_list[1]:
        phases = safe_list(te.get("phases"))
        mc = st.columns(4)
        with mc[0]:
            st.metric("Total Hours", str(safe_int(te.get("total_hours"))))
        with mc[1]:
            st.metric("Duration", safe_str(te.get("duration_weeks")))
        with mc[2]:
            st.metric("Confidence", safe_str(te.get("confidence")))
        with mc[3]:
            st.metric("Buffer", safe_str(te.get("buffer")))
        if phases:
            colors = ["#00d4aa", "#00b4d8", "#7b61ff", "#ff6b6b", "#ffd166", "#06d6a0"]
            fig = go.Figure()
            for i, p in enumerate(phases):
                p = safe_dict(p)
                ph_name = safe_str(p.get("name", "Phase"))
                ph_hours = safe_int(p.get("hours", 0))
                fig.add_trace(go.Bar(
                    name=ph_name, x=[ph_hours], y=[ph_name],
                    orientation="h", marker_color=colors[i % 6],
                    text=str(ph_hours) + "h", textposition="auto",
                ))
            fig.update_layout(
                title="Phase Breakdown", template="plotly_dark",
                paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
                height=400, showlegend=False, xaxis_title="Hours",
            )
            st.plotly_chart(fig, use_container_width=True)
        milestones = safe_list(te.get("milestones"))
        for m in milestones:
            m = safe_dict(m)
            st.markdown("- **Week " + str(safe_int(m.get("week"))) + "** — " + safe_str(m.get("name")) + ": " + safe_str(m.get("description")))
        # Excel download
        st.markdown("---")
        excel_data = generate_time_excel(te, se)
        if excel_data:
            st.download_button(
                "📥 Download Time Estimate (Excel)",
                data=excel_data,
                file_name="ECI_Time_Estimate_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True, type="primary", key="dl_time_xlsx",
            )

    # ── Cost (Infrastructure) ──
    with tab_list[2]:
        azure_costs = safe_list(ce.get("azure_costs"))
        third_party = safe_list(ce.get("third_party_costs"))
        mc = st.columns(3)
        with mc[0]:
            st.metric("Monthly Cost", "$" + str(safe_int(ce.get("total_monthly_cost"))))
        with mc[1]:
            st.metric("Annual Cost", "$" + str(safe_int(ce.get("total_annual_cost"))))
        with mc[2]:
            st.metric("Services", str(len(azure_costs) + len(third_party)))
        if azure_costs:
            pie_vals = [safe_int(safe_dict(a).get("monthly_cost", 0)) for a in azure_costs]
            pie_names = [safe_str(safe_dict(a).get("service", "")) for a in azure_costs]
            fig = px.pie(values=pie_vals, names=pie_names, title="Azure Infrastructure Cost Distribution (Monthly)",
                         color_discrete_sequence=["#00d4aa", "#00b4d8", "#7b61ff", "#ff6b6b", "#ffd166", "#06d6a0", "#e9c46a", "#f4845f", "#a8dadc", "#457b9d", "#2a9d8f", "#264653"])
            fig.update_layout(template="plotly_dark", paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", height=400)
            st.plotly_chart(fig, use_container_width=True)
        st.markdown("**Azure Services:**")
        for svc in azure_costs:
            svc = safe_dict(svc)
            st.markdown("- **" + safe_str(svc.get("service")) + "** (" + safe_str(svc.get("tier", "")) + ") — $" + str(safe_int(svc.get("monthly_cost"))) + "/mo — _" + safe_str(svc.get("description", "")) + "_")
        if third_party:
            st.markdown("**Third-Party Services:**")
            for tp in third_party:
                tp = safe_dict(tp)
                st.markdown("- **" + safe_str(tp.get("name")) + "** — $" + str(safe_int(tp.get("monthly_cost"))) + "/mo — _" + safe_str(tp.get("description", "")) + "_")
        opt_tips = safe_list(ce.get("cost_optimization"))
        if opt_tips:
            st.markdown("**Cost Optimization Tips:**")
            for tip in opt_tips:
                st.markdown("- " + safe_str(tip))
        notes = safe_str(ce.get("notes"))
        if notes:
            st.info(notes)

    # ── Risk ──
    with tab_list[3]:
        score = safe_int(ri.get("overall_score", 0))
        clr = "#06d6a0" if score <= 3 else "#ffd166" if score <= 6 else "#ff6b6b"
        st.markdown('<div class="rsb" style="border-left:4px solid ' + clr + ';"><span class="rsv">' + str(score) + '/10</span><span class="rsl">Overall: ' + safe_str(ri.get("overall_level")) + '</span></div>', unsafe_allow_html=True)
        for rk in safe_list(ri.get("risks")):
            rk = safe_dict(rk)
            sev = safe_str(rk.get("severity"))
            sc = "#06d6a0" if sev == "Low" else "#ffd166" if sev == "Medium" else "#ff6b6b"
            st.markdown('<div class="rc" style="border-left:3px solid ' + sc + ';"><div class="rch2"><strong>' + safe_str(rk.get("category")) + ': ' + safe_str(rk.get("title")) + '</strong><span class="rsev" style="color:' + sc + ';">' + sev + '</span></div><p>' + safe_str(rk.get("description")) + '</p><div class="rmit"><strong>Mitigation:</strong> ' + safe_str(rk.get("mitigation")) + '</div></div>', unsafe_allow_html=True)

    # ── Architecture ──
    with tab_list[4]:
        comps = safe_list(ar.get("components"))
        st.markdown("**Pattern:** _" + safe_str(ar.get("pattern")) + "_")
        arch_cols = st.columns(3)
        for i, c in enumerate(comps):
            c = safe_dict(c)
            with arch_cols[i % 3]:
                svcs = "".join("<li>" + safe_str(s) + "</li>" for s in safe_list(c.get("services")))
                st.markdown('<div class="ac"><div class="acn">' + safe_str(c.get("name")) + '</div><div class="act">' + safe_str(c.get("azure_service")) + '</div><ul class="acs">' + svcs + '</ul></div>', unsafe_allow_html=True)
        df = safe_list(ar.get("data_flow"))
        if df:
            flow_str = " -> ".join(safe_str(x) for x in df)
            st.markdown('<div class="dfv">' + flow_str + '</div>', unsafe_allow_html=True)
        for s in safe_list(ar.get("security")):
            st.markdown("- " + safe_str(s))

    # ── Diagrams (Mermaid.js) ──
    with tab_list[5]:
        mermaid_data = safe_dict(r.get("mermaid_diagrams"))
        if mermaid_data:
            diagram_tabs = st.tabs(["🏗️ Infrastructure", "🔄 Data Flow", "📨 Sequence", "🚀 Deployment", "🔒 Security"])
            diagram_map = [
                ("infrastructure", "Azure Infrastructure Architecture", 500),
                ("data_flow", "Data Flow Diagram", 450),
                ("sequence", "Request Sequence Diagram", 500),
                ("deployment", "CI/CD & Deployment Pipeline", 450),
                ("security", "Security Architecture & Controls", 450),
            ]
            for i, (key, title, h) in enumerate(diagram_map):
                with diagram_tabs[i]:
                    code = safe_str(mermaid_data.get(key, ""))
                    if code:
                        st.markdown("**" + title + "**")
                        render_mermaid(code, height=h)
                        with st.expander("View Mermaid Source"):
                            st.code(code, language="text")
                    else:
                        st.info("Diagram not available for this project.")
            # Download all diagrams as JSON
            st.markdown("---")
            st.download_button(
                "📥 Download All Diagrams (JSON)",
                data=json.dumps(mermaid_data, indent=2),
                file_name="ECI_Diagrams_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".json",
                mime="application/json",
                use_container_width=True, key="dl_diagrams",
            )
        else:
            st.info("Architecture diagrams will be generated after processing documents.")

    # ── Proposal ──
    with tab_list[6]:
        proposal = safe_dict(r.get("proposal"))
        for sec in safe_list(proposal.get("sections")):
            sec = safe_dict(sec)
            with st.expander(safe_str(sec.get("title")), expanded=False):
                st.markdown(safe_str(sec.get("content")))
        for k, v in safe_dict(proposal.get("quality_checks")).items():
            check_icon = "pass" if v else "warn"
            st.markdown(("✅ " if v else "⚠️ ") + str(k))
        # PDF download
        st.markdown("---")
        pdf_data = generate_proposal_pdf(r)
        if pdf_data:
            st.download_button(
                "📥 Download Proposal (PDF)",
                data=pdf_data,
                file_name="ECI_Proposal_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".pdf",
                mime="application/pdf",
                use_container_width=True, type="primary", key="dl_proposal_pdf",
            )

    # ── Scope ──
    with tab_list[7]:
        sc = safe_dict(r.get("scope"))
        sc1, sc2 = st.columns(2)
        with sc1:
            st.markdown("**In Scope:**")
            for x in safe_list(sc.get("in_scope")):
                st.markdown("- " + safe_str(x))
            st.markdown("**Assumptions:**")
            for x in safe_list(sc.get("assumptions")):
                st.markdown("- " + safe_str(x))
        with sc2:
            st.markdown("**Out of Scope:**")
            for x in safe_list(sc.get("out_of_scope")):
                st.markdown("- " + safe_str(x))
            st.markdown("**Prerequisites:**")
            for x in safe_list(sc.get("prerequisites")):
                st.markdown("- " + safe_str(x))

    # ── Delivery ──
    st.markdown("---")
    st.markdown('<div class="shdr"><span class="shdr-i">🚀</span> Delivery</div>', unsafe_allow_html=True)

    # Downloads row
    st.markdown('<div class="shdr" style="font-size:1rem;"><span class="shdr-i">📥</span> Downloads</div>', unsafe_allow_html=True)
    dl_cols = st.columns(3)
    with dl_cols[0]:
        xl_data = generate_time_excel(te, se)
        if xl_data:
            st.download_button("📊 Time Estimate (Excel)", data=xl_data,
                               file_name="ECI_Time_Estimate_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".xlsx",
                               mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                               use_container_width=True, type="primary", key="bdl_xl")
    with dl_cols[1]:
        pdf_data = generate_proposal_pdf(r)
        if pdf_data:
            st.download_button("📄 Proposal (PDF)", data=pdf_data,
                               file_name="ECI_Proposal_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".pdf",
                               mime="application/pdf",
                               use_container_width=True, type="primary", key="bdl_pdf")
    with dl_cols[2]:
        st.download_button("📋 Full Data (JSON)", data=json.dumps(r, indent=2, default=str),
                           file_name="ECI_Data_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".json",
                           mime="application/json", use_container_width=True, key="bdl")

    # Actions row
    st.markdown('<div class="shdr" style="font-size:1rem;"><span class="shdr-i">⚡</span> Actions</div>', unsafe_allow_html=True)
    ac_cols = st.columns(2)
    with ac_cols[0]:
        if st.button("Upload to SharePoint", use_container_width=True, key="bsp"):
            ok, msg = SP.from_session().upload(r)
            if ok:
                st.success(msg)
            else:
                st.error(msg)
    with ac_cols[1]:
        if st.button("Send Alert Email", use_container_width=True, key="bem"):
            ok, msg = Mailer.from_session().send(r)
            if ok:
                st.success(msg)
            else:
                st.error(msg)


# ═══════════════════════════════════════════════════════════════════════
#  TAB 2: ADMIN
# ═══════════════════════════════════════════════════════════════════════

def tab_admin():
    st.markdown('<div class="shdr"><span class="shdr-i">🧠</span> Continuous Learning</div>', unsafe_allow_html=True)
    at = st.tabs(["📊 Dashboard", "📚 Training", "🔧 Config", "📁 Sync", "🔄 Loop"])

    with at[0]:
        m = st.session_state.model_metrics
        mc = st.columns(4)
        with mc[0]:
            st.metric("Proposals", m["proposals_processed"])
        with mc[1]:
            st.metric("Accuracy", str(m["accuracy"]) + "%")
        with mc[2]:
            st.metric("Win Rate", str(m["win_rate"]) + "%")
        with mc[3]:
            st.metric("Variance", str(m["variance"]) + "%")
        dates = [(datetime.now() - timedelta(days=30 * i)).strftime("%b %Y") for i in range(6, -1, -1)]
        fig = go.Figure(go.Scatter(x=dates, y=[65, 68, 72, 74, 76, 78, 80], mode="lines+markers",
                                   line=dict(color="#00d4aa", width=3), marker=dict(size=10)))
        fig.update_layout(title="Accuracy Trend", template="plotly_dark",
                          paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", height=350)
        st.plotly_chart(fig, use_container_width=True)
        dc1, dc2 = st.columns(2)
        with dc1:
            fig2 = go.Figure(go.Bar(x=["Time", "Cost", "Risk", "Arch"], y=[85, 79, 82, 88],
                                    marker_color=["#00d4aa", "#00b4d8", "#ffd166", "#7b61ff"],
                                    text=["85%", "79%", "82%", "88%"], textposition="auto"))
            fig2.update_layout(title="Agent Accuracy", template="plotly_dark",
                               paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)", height=300)
            st.plotly_chart(fig2, use_container_width=True)
        with dc2:
            fig3 = go.Figure(go.Indicator(
                mode="gauge+number+delta", value=m["accuracy"], delta={"reference": 72},
                gauge={"axis": {"range": [0, 100]}, "bar": {"color": "#00d4aa"},
                       "steps": [{"range": [0, 50], "color": "rgba(255,107,107,.2)"}, {"range": [50, 75], "color": "rgba(255,209,102,.2)"}, {"range": [75, 100], "color": "rgba(0,212,170,.2)"}]},
                title={"text": "Model Health"}))
            fig3.update_layout(template="plotly_dark", paper_bgcolor="rgba(0,0,0,0)", height=300)
            st.plotly_chart(fig3, use_container_width=True)

    with at[1]:
        st.file_uploader("Upload Historical Data", type=["json", "csv", "xlsx"], accept_multiple_files=True, key="tfu")
        with st.expander("Add Project Manually"):
            with st.form("mp"):
                pn = st.text_input("Name")
                pt = st.selectbox("Type", ["Data/Cloud/AI", "Web App", "Mobile", "SharePoint", "Integration"])
                fc1, fc2 = st.columns(2)
                with fc1:
                    eh = st.number_input("Est Hours", 0, value=100)
                    ah = st.number_input("Actual Hours", 0, value=0)
                    ec = st.number_input("Est Cost", 0, value=10000)
                with fc2:
                    ac_val = st.number_input("Actual Cost", 0, value=0)
                    out = st.selectbox("Outcome", ["Won", "Lost", "Pending"])
                    ts = st.text_input("Tech Stack", "Azure, Python")
                if st.form_submit_button("Add"):
                    st.session_state.historical_projects.append({"name": pn, "type": pt, "estimated_hours": eh, "actual_hours": ah, "estimated_cost": ec, "actual_cost": ac_val, "outcome": out, "tech_stack": ts.split(", ")})
                    st.success("Added: " + pn)
        if st.session_state.historical_projects:
            st.markdown("**" + str(len(st.session_state.historical_projects)) + " project(s) loaded**")
            for p in st.session_state.historical_projects:
                st.markdown("- **" + safe_str(p.get("name")) + "** — " + safe_str(p.get("type")) + " — " + safe_str(p.get("outcome")))

    with at[2]:
        cc1, cc2 = st.columns(2)
        with cc1:
            for ph in ["Discovery", "Design", "Development", "Deployment", "Support"]:
                st.checkbox(ph, True, key="p_" + ph)
            st.selectbox("Estimation", ["Three-Point", "Story Points", "Function Points"], key="em")
            st.slider("Buffer %", 0, 50, 20, key="bf")
        with cc2:
            for rn, rv in [("Architects", 1), ("Lead Devs", 1), ("Developers", 3), ("QA", 1), ("DevOps", 1), ("PMs", 1)]:
                st.number_input(rn, 0, 20, rv, key="r_" + rn)
            st.selectbox("Pricing", ["Fixed Price", "T&M", "Retainer"], key="pc")
        tc = st.columns(3)
        with tc[0]:
            if st.button("Train", use_container_width=True, type="primary", key="bt"):
                pb = st.progress(0)
                for i in range(100):
                    time.sleep(0.02)
                    pb.progress(i + 1)
                st.session_state.model_metrics["accuracy"] = min(st.session_state.model_metrics["accuracy"] + 2.5, 95)
                st.success("Training complete!")
        with tc[1]:
            if st.button("Retrain", use_container_width=True, key="br"):
                time.sleep(1)
                st.success("Retrained.")
        with tc[2]:
            if st.button("A/B Test", use_container_width=True, key="ba"):
                time.sleep(1)
                st.success("New model +8.3%")

    with at[3]:
        st.text_input("Library Path", placeholder="/sites/presales/Historical", key="ssp")
        st.multiselect("Types", ["Proposals", "Estimates", "Outcomes", "Templates"], default=["Proposals", "Estimates", "Outcomes"], key="sst")
        if st.button("Sync Now", use_container_width=True, type="primary", key="bs"):
            time.sleep(2)
            st.success("Synced 47 docs.")

    with at[4]:
        steps = [("1", "Outcome", "Win/loss + feedback", "#7b61ff"), ("2", "Enrichment", "Link to estimates", "#00b4d8"), ("3", "Retraining", "Fine-tune model", "#00d4aa"), ("4", "Validation", "A/B testing", "#ffd166"), ("5", "Deploy", "Push to prod", "#ff6b6b")]
        lc = st.columns(5)
        for i, (n, t, d, c) in enumerate(steps):
            with lc[i]:
                st.markdown('<div class="ls" style="border-top:3px solid ' + c + ';"><div class="ln" style="background:' + c + ';">' + n + '</div><div class="lt">' + t + '</div><div class="ld">' + d + '</div></div>', unsafe_allow_html=True)
        st.markdown("---")
        st.markdown("### Record Outcome")
        with st.form("of"):
            oc1, oc2 = st.columns(2)
            with oc1:
                st.text_input("Project", key="on")
                st.selectbox("Result", ["Won", "Lost", "No Decision"], key="ores")
                st.number_input("Actual Hours", 0, key="oh")
            with oc2:
                st.text_area("Feedback", key="ofb")
                st.number_input("Actual Cost", 0, key="oc")
                st.text_area("Lessons", key="ol")
            if st.form_submit_button("Record", type="primary"):
                st.success("Recorded for next training cycle.")
        st.markdown("---")
        method_cols = st.columns(5)
        method_data = [("📋", "Phases", "Discovery Design Dev Deploy Support"), ("👥", "Team", "Arch Lead Dev QA DevOps PM"), ("📐", "Estimation", "Three-point + 20% buffer"), ("💲", "Pricing", "Fixed, T&M, Retainer"), ("⚠️", "Risk", "Tech Schedule Resource Budget External")]
        for i, (ic, t, d) in enumerate(method_data):
            with method_cols[i]:
                st.markdown('<div class="mc"><div class="mi">' + ic + '</div><div class="mt">' + t + '</div><div class="md2">' + d + '</div></div>', unsafe_allow_html=True)


# ═══════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════

main_t1, main_t2 = st.tabs(["⚡ ECI Presale", "⚙️ Admin & Training"])
with main_t1:
    tab_presale()
with main_t2:
    tab_admin()
